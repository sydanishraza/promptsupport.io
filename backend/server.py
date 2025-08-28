#!/usr/bin/env python3
"""
PromptSupport Enhanced Content Engine Backend
FastAPI server with AI integrations for content processing
"""

import os
import uuid
import asyncio
import time
import hashlib
from datetime import datetime, timedelta
from typing import List, Optional, Dict, Any
import json
import io
import base64
from bson import ObjectId

# ObjectId serialization helper
def objectid_to_str(obj: Any) -> Any:
    """Convert ObjectId to string for JSON serialization"""
    if isinstance(obj, ObjectId):
        return str(obj)
    elif isinstance(obj, dict):
        return {key: objectid_to_str(value) for key, value in obj.items()}
    elif isinstance(obj, list):
        return [objectid_to_str(item) for item in obj]
    else:
        return obj

from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Depends, Header, Request
from fastapi.responses import JSONResponse, FileResponse, StreamingResponse, Response
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
import motor.motor_asyncio
from pymongo import MongoClient
import aiofiles
from dotenv import load_dotenv
import requests
import httpx
import re
import markdown
from bs4 import BeautifulSoup
# KE-PR3: Media intelligence now imported from engine.media

# New Engine Package Imports (KE-PR1: scaffolding)
try:
    import sys
    import os
    # Add parent directory to Python path for engine package
    parent_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    if parent_dir not in sys.path:
        sys.path.insert(0, parent_dir)
    
    from engine.models import RawBundle, QAReport, MediaAsset
    from engine.logging_util import stage_log, logger
    from config.settings import settings
    
    # KE-PR2: Import linking modules
    from engine.linking.anchors import stable_slug, anchor_id, assign_heading_ids, validate_heading_ladder
    from engine.linking.toc import build_toc, build_minitoc, anchors_resolve
    from engine.linking.bookmarks import extract_headings_registry, generate_doc_uid, generate_doc_slug, backfill_registry, get_registry
    from engine.linking.links import build_href, get_default_route_map, build_link
    
    # KE-PR3: Import media and assets modules
    from engine.media import media_intelligence
    from engine.stores.assets import save_bytes, save_file, read_file, get_asset_path, hash_bytes
    
    # KE-PR4: Import V2 engine classes
    from engine.v2.analyzer import V2MultiDimensionalAnalyzer, get_v2_analyzer
    from engine.v2.outline import V2GlobalOutlinePlanner, V2PerArticleOutlinePlanner
    from engine.v2.prewrite import V2PrewriteSystem
    from engine.v2.style import V2StyleProcessor
    from engine.v2.related import V2RelatedLinksSystem
    from engine.v2.gaps import V2GapFillingSystem
    from engine.v2.evidence import V2EvidenceTaggingSystem
    from engine.v2.code_norm import V2CodeNormalizationSystem
    from engine.v2.generator import V2ArticleGenerator
    from engine.v2.validate import V2ValidationSystem
    from engine.v2.crossqa import V2CrossArticleQASystem
    from engine.v2.adapt import V2AdaptiveAdjustmentSystem
    from engine.v2.publish import V2PublishingSystem
    from engine.v2.versioning import V2VersioningSystem
    from engine.v2.review import V2ReviewSystem
    from engine.v2.extractor import V2ContentExtractor
    from engine.v2.media import V2MediaManager
    
    # KE-PR5: Import V2 Pipeline Orchestrator
    from engine.v2.pipeline import Pipeline, get_pipeline
    
    # KE-PR6: Import centralized LLM client
    from engine.llm.client import get_llm_client
    
    print("✅ Engine package modules loaded successfully")
    print("✅ KE-PR2: Linking modules loaded successfully")
    print("✅ KE-PR3: Media and assets modules loaded successfully")
    print("✅ KE-PR4: V2 engine classes loaded successfully")
    print("✅ KE-PR5: V2 pipeline orchestrator loaded successfully")
    print("✅ KE-PR6: Centralized LLM client loaded successfully")
except ImportError as e:
    print(f"⚠️ Engine package import failed: {e}")
    # Fallback - create dummy objects to prevent crashes
    class RawBundle: pass
    class QAReport: pass  
    class MediaAsset: pass
    def stage_log(name): return lambda f: f
    logger = None
    settings = None
    
    # KE-PR2: Fallback linking functions
    def stable_slug(text, max_len=60): return text.lower().replace(' ', '-')
    def anchor_id(text, prefix=None): return stable_slug(text)
    def assign_heading_ids(html): return html
    def validate_heading_ladder(html): return True
    def build_toc(headings): return headings
    def build_minitoc(html): return html
    def anchors_resolve(html): return True
    def extract_headings_registry(html): return []
    def generate_doc_uid(): return "fallback-uid"
    def generate_doc_slug(title): return title.replace(' ', '-').lower()
    async def backfill_registry(limit=None): return {"status": "fallback"}
    async def get_registry(doc_uid): return {}
    def build_href(doc, anchor, route_map): return f"#{anchor}"
    def get_default_route_map(env): return {}
    # KE-PR3: Fallback media and assets functions  
    import hashlib
    
    class MediaIntelligenceService:
        async def analyze_media_comprehensive(self, *args, **kwargs): return {}
        def create_enhanced_media_html(self, *args, **kwargs): return ""
        def generate_contextual_placement(self, *args, **kwargs): return {}
    
    media_intelligence = MediaIntelligenceService()
    
    def save_bytes(data, filename, upload_dir="static/uploads"): 
        h = hashlib.md5(data).hexdigest()[:8]
        return h, f"{h}_{filename}"
    def save_file(temp_path, upload_dir="static/uploads"): return "fallback", "fallback.tmp"
    def read_file(path): return b""
    def get_asset_path(filename, upload_dir="static/uploads"): return f"{upload_dir}/{filename}"
    def hash_bytes(data): return hashlib.md5(data).hexdigest()[:8]
    
    # KE-PR4: Fallback V2 engine classes
    class V2MultiDimensionalAnalyzer: 
        async def analyze_normalized_document(self, *args, **kwargs): return {}
    class V2GlobalOutlinePlanner: 
        async def create_global_outline(self, *args, **kwargs): return {}
    class V2PerArticleOutlinePlanner: 
        async def create_article_outline(self, *args, **kwargs): return {}
    class V2PrewriteSystem: 
        async def extract_prewrite_data(self, *args, **kwargs): return {}
    class V2StyleProcessor: 
        def process_style(self, *args, **kwargs): return {}
    class V2RelatedLinksSystem: 
        async def generate_related_links(self, *args, **kwargs): return []
    class V2GapFillingSystem: 
        async def fill_gaps(self, *args, **kwargs): return ""
    class V2EvidenceTaggingSystem: 
        def tag_evidence(self, *args, **kwargs): return ""
    class V2CodeNormalizationSystem: 
        def normalize_code_blocks(self, *args, **kwargs): return ""
    class V2ArticleGenerator: 
        async def generate_article(self, *args, **kwargs): return {}
    class V2ValidationSystem: 
        def validate_content(self, *args, **kwargs): return True
    class V2CrossArticleQASystem: 
        async def perform_cross_article_qa(self, *args, **kwargs): return {}
    class V2AdaptiveAdjustmentSystem: 
        async def adjust_article_balance(self, *args, **kwargs): return {}
    class V2PublishingSystem: 
        async def publish_v2_content(self, *args, **kwargs): return {}
    class V2VersioningSystem: 
        async def create_version(self, *args, **kwargs): return {}
    class V2ReviewSystem: 
        def create_review_request(self, *args, **kwargs): return {}
        async def enqueue_for_review(self, *args, **kwargs): return {"review_id": "fallback", "review_status": "queued"}
    class V2ContentExtractor: 
        def extract_content(self, *args, **kwargs): return {}
    class V2MediaManager: 
        def process_media(self, *args, **kwargs): return {}
    
    v2_analyzer = V2MultiDimensionalAnalyzer()
    
    # KE-PR5: Fallback pipeline
    class Pipeline:
        async def run(self, job_id, content, metadata): return [], None, f"fallback_{job_id}"
    def get_pipeline(llm_client=None): return Pipeline()
    
    # KE-PR6: Fallback LLM client
    class LLMClient:
        def __init__(self, provider=None, **kwargs):
            self.provider = provider or "fallback"
        async def complete(self, *args, **kwargs): return "Fallback LLM response"
        async def moderate(self, *args, **kwargs): return {"ok": True}
        async def analyze_content(self, *args, **kwargs): return {"analysis": "fallback", "success": False}
    def get_llm_client(provider=None, **kwargs): return LLMClient(provider, **kwargs)

# HTML preprocessing pipeline imports
import mammoth
import pypandoc
# from pdfminer.six import extract_text as pdf_extract_text
from lxml import etree
from lxml.html import fromstring as html_fromstring, tostring as html_tostring
import shutil
import subprocess

# Load environment variables
load_dotenv()

# Helper function to convert ObjectId to string for JSON serialization
def objectid_to_str(data):
    """Convert ObjectId to string for JSON serialization"""
    if isinstance(data, ObjectId):
        return str(data)
    elif isinstance(data, dict):
        return {key: objectid_to_str(value) for key, value in data.items()}
    elif isinstance(data, list):
        return [objectid_to_str(item) for item in data]
    return data

# Configuration
MONGO_URL = os.getenv("MONGO_URL", "mongodb://localhost:27017/")
DATABASE_NAME = os.getenv("DATABASE_NAME", "promptsupport_db")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
QDRANT_API_KEY = os.getenv("QDRANT_API_KEY")
ASSEMBLYAI_API_KEY = os.getenv("ASSEMBLYAI_API_KEY")
QDRANT_HOST = os.getenv("QDRANT_HOST", "localhost")
QDRANT_PORT = int(os.getenv("QDRANT_PORT", "6333"))

# Initialize FastAPI app
app = FastAPI(
    title="PromptSupport Enhanced Content Engine",
    description="AI-native content processing and management system with KE-PR8 API Router",
    version="2.0.0"
)

# KE-PR9: Import MongoDB Repository Layer
try:
    import sys
    import os
    
    # Add engine path for repository imports
    engine_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'engine')
    if engine_path not in sys.path:
        sys.path.insert(0, engine_path)
    
    from stores.mongo import (
        RepositoryFactory, 
        upsert_content, 
        fetch_article_by_slug, 
        fetch_article_by_uid,
        update_article_headings, 
        update_article_xrefs,
        test_mongo_roundtrip
    )
    print("✅ KE-PR9: MongoDB repository layer imported successfully")
    mongo_repo_available = True
except ImportError as e:
    print(f"⚠️ KE-PR9: MongoDB repository layer import failed: {e}")
    mongo_repo_available = False
try:
    # Import and include the organized API router
    sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    from api.router import router as api_router
    
    app.include_router(api_router)
    print("✅ KE-PR8: API Router included successfully")
    print(f"✅ KE-PR8: {len(api_router.routes)} routes organized by domain")
    
except Exception as e:
    print(f"❌ KE-PR8: Failed to include API router - {e}")
    # Fallback: keep essential routes in server.py if router fails
    @app.get("/api/health")
    def fallback_health():
        return {"status": "fallback", "error": "API router failed to load"}

# Mount static files for serving uploaded images under /api/static route
import os
static_dir = os.path.join(os.path.dirname(__file__), "static")
app.mount("/api/static", StaticFiles(directory=static_dir), name="static")

# CORS middleware - Enhanced configuration for preview environments
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "*",  # Allow all origins for development
        "https://mongo-repo-refactor.preview.emergentagent.com",
        "https://mongo-repo-refactor.preview.emergentagent.com", 
        "http://localhost:3000",
        "http://127.0.0.1:3000"
    ],
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS", "HEAD", "PATCH"],
    allow_headers=["*"],  # Allow all headers
    expose_headers=["*"]  # Expose all headers
)

# === DOCUMENT PROCESSING HELPER FUNCTIONS ===

def is_table_of_contents(text: str, position: int = 0) -> bool:
    """Detect Table of Contents sections"""
    toc_indicators = [
        'table of contents',
        'contents',
        'index',
        '...........',  # Dotted lines common in TOCs
    ]
    
    text_lower = text.lower()
    
    # Check for TOC title
    if any(indicator in text_lower for indicator in toc_indicators[:3]):
        return True
        
    # Check for dotted lines or page numbers (TOC formatting)
    if ('.' * 5 in text) or (re.search(r'\d+$', text.strip())):
        return True
        
    # Check if this looks like a TOC entry (short line with numbers)
    if len(text) < 100 and re.search(r'^\w.*\d+$', text.strip()):
        return True
        
    return False

def is_header_footer_or_page_number(text: str) -> bool:
    """Detect headers, footers, and page numbers"""
    # Page numbers (just numbers or "Page X")
    if re.match(r'^(page\s+)?\d+$', text.lower().strip()):
        return True
        
    # Very short text that might be headers/footers
    if len(text.strip()) < 10 and not any(c.isalpha() for c in text):
        return True
        
    # Common header/footer patterns
    header_footer_patterns = [
        r'^\d+$',  # Just numbers
        r'^page \d+',  # Page numbers
        r'confidential',
        r'proprietary',
        r'draft',
        r'version \d+',
    ]
    
    return any(re.search(pattern, text.lower()) for pattern in header_footer_patterns)

def is_legal_disclaimer(text: str) -> bool:
    """Detect legal disclaimers and copyright notices"""
    legal_patterns = [
        r'copyright',
        r'©',
        r'all rights reserved',
        r'proprietary and confidential',
        r'trademark',
        r'disclaimer',
        r'terms of use',
        r'privacy policy',
        r'legal notice'
    ]
    
    text_lower = text.lower()
    return any(re.search(pattern, text_lower) for pattern in legal_patterns)

def is_cover_page(paragraphs_sample: list) -> bool:
    """Detect if the first few paragraphs constitute a cover page"""
    if not paragraphs_sample:
        return False
        
    first_page_text = '\n'.join([p.text for p in paragraphs_sample[:5] if hasattr(p, 'text')])
    
    cover_indicators = [
        len(first_page_text.split('\n')) <= 5 and any(word in first_page_text.upper() for word in ['GUIDE', 'MANUAL', 'DOCUMENT', 'REPORT']),
        first_page_text.count('\n') <= 3 and len(first_page_text) < 200,
        any(pattern in first_page_text.upper() for pattern in ['COPYRIGHT', '©', 'ALL RIGHTS RESERVED'])
    ]
    
    return any(cover_indicators)

def detect_heading_level_from_text(text: str) -> int:
    """Detect heading level from text patterns"""
    text = text.strip()
    
    # Check for numbered headings
    if re.match(r'^\d+\.?\s+', text):
        return 2
    if re.match(r'^\d+\.\d+\.?\s+', text):
        return 3
    if re.match(r'^\d+\.\d+\.\d+\.?\s+', text):
        return 4
        
    # Check for text patterns that suggest headings
    if text.isupper() and len(text) < 100:
        return 1
    if text.endswith(':') and len(text) < 80:
        return 3
        
    return None

async def regenerate_articles_with_enhanced_context(extracted_content: dict, contextual_images: list, template_data: dict, training_session: dict) -> list:
    """Phase 2: Enhanced article regeneration with contextual analysis"""
    try:
        print(f"🪄 Phase 2: Starting enhanced article regeneration")
        
        # Analyze content structure for semantic splitting
        articles_data = []
        current_article = {
            "content_blocks": [],
            "images": [],
            "headings": [],
            "title": ""
        }
        
        h1_count = 0
        
        for block in extracted_content['structure']:
            block_type = block['type']
            content = block['content']
            
            # Major heading changes trigger new articles
            if block_type == 'h1':
                # Save current article if it has content
                if current_article['content_blocks']:
                    articles_data.append(current_article)
                
                # Start new article
                h1_count += 1
                current_article = {
                    "content_blocks": [block],
                    "images": [],
                    "headings": [content],
                    "title": content or f"Article {h1_count}"
                }
            elif block_type == 'h2' and len(current_article['content_blocks']) > 0:
                # H2 might trigger new article if current is getting long
                current_word_count = sum(len(b['content'].split()) for b in current_article['content_blocks'])
                if current_word_count > 8000:  # PERFORMANCE BOOST: Increased from 3000 to 8000 words per article for fewer LLM calls
                    # Save current article and start new one
                    articles_data.append(current_article)
                    current_article = {
                        "content_blocks": [block],
                        "images": [],
                        "headings": [content],
                        "title": content or f"Article {len(articles_data) + 1}"
                    }
                else:
                    current_article['content_blocks'].append(block)
                    current_article['headings'].append(content)
            else:
                current_article['content_blocks'].append(block)
        
        # Add final article
        if current_article['content_blocks']:
            articles_data.append(current_article)
        
        # Distribute images across articles based on context
        images_per_article = len(contextual_images) // max(1, len(articles_data))
        for i, article_data in enumerate(articles_data):
            start_idx = i * images_per_article
            end_idx = start_idx + images_per_article if i < len(articles_data) - 1 else len(contextual_images)
            article_data['images'] = contextual_images[start_idx:end_idx]
        
        # Generate final articles
        final_articles = []
        
        for i, article_data in enumerate(articles_data):
            # Create structured content
            html_content = await generate_enhanced_html_content(article_data, template_data)
            markdown_content = await generate_enhanced_markdown_content(article_data, template_data)
            
            # TICKET 3: Generate universal bookmark and link data
            # Generate document identifiers
            import time
            import random
            import string
            import re
            import unicodedata
            
            # Generate doc_uid (ULID-style)
            timestamp = int(time.time() * 1000)  # milliseconds
            random_part = ''.join(random.choices(string.ascii_uppercase + string.digits, k=16))
            timestamp_b36 = format(timestamp, 'x').upper()[-8:]  # Last 8 chars of hex timestamp
            doc_uid = f"01JZ{timestamp_b36}{random_part[:8]}"
            
            # Generate doc_slug from title
            article_title = f"Article {i+1} From {training_session['filename']}" if not article_data['title'] else article_data['title']
            norm = unicodedata.normalize("NFKD", article_title).encode("ascii", "ignore").decode("ascii")
            doc_slug = re.sub(r"\s+", "-", norm.lower())
            doc_slug = re.sub(r"[^a-z0-9-]", "", doc_slug)
            doc_slug = re.sub(r"-{2,}", "-", doc_slug).strip("-")[:50]  # Limit length
            
            # Extract headings registry from HTML content
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            headings_registry = []
            for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                if heading.get_text(strip=True):
                    # Generate anchor ID
                    heading_text = heading.get_text(strip=True)
                    anchor_id = re.sub(r'[^a-zA-Z0-9\s-]', '', heading_text).strip()
                    anchor_id = re.sub(r'\s+', '-', anchor_id).lower()[:50]
                    
                    headings_registry.append({
                        'level': int(heading.name[1]),
                        'text': heading_text,
                        'anchor': anchor_id,
                        'id': f"{doc_uid}#{anchor_id}"
                    })
            
            # Create media array with placement info
            media_array = []
            for img in article_data['images']:
                media_array.append({
                    "url": img['url'],
                    "alt": img['alt_text'],
                    "caption": img.get('caption', ''),
                    "placement": img.get('placement', 'inline'),
                    "filename": img['filename']
                })
            
            article = {
                "id": str(uuid.uuid4()),
                "title": article_title,
                "html": html_content,
                "markdown": markdown_content,
                "content": html_content,  # For backward compatibility
                # TICKET 3: Universal bookmark and link data
                "doc_uid": doc_uid,
                "doc_slug": doc_slug,
                "headings_registry": headings_registry,
                "xrefs": article_data.get('xrefs', []),
                "related_links": article_data.get('related_links', []),
                "media": media_array,
                "tags": ["extracted", "generated", "enhanced"],
                "status": "training",
                "template_id": training_session['template_id'],
                "session_id": training_session['session_id'],
                "word_count": len(html_content.split()),
                "image_count": len(media_array),
                "format": "html",
                "created_at": datetime.utcnow().isoformat(),
                "ai_processed": True,
                "ai_model": "gpt-4o-mini (with claude + local llm fallback)",
                "training_mode": True,
                "metadata": {
                    "article_number": i + 1,
                    "source_filename": training_session['filename'],
                    "template_applied": training_session['template_id'],
                    "phase": "enhanced_extraction"
                }
            }
            
            final_articles.append(article)
            print(f"📄 Generated enhanced article {i+1}: {len(html_content)} chars, {len(media_array)} images")
        
        print(f"✅ Phase 2 Complete: Generated {len(final_articles)} enhanced articles")
        return final_articles
        
    except Exception as e:
        print(f"❌ Enhanced regeneration error: {e}")
        import traceback
        traceback.print_exc()
        return []

async def generate_enhanced_html_content(article_data: dict, template_data: dict) -> str:
    """Generate HTML content for an article"""
    html_parts = []
    
    # Skip adding title as H1 - let frontend handle article title display
    # Article title should be displayed by frontend components, not embedded in content
    
    # Process content blocks
    image_index = 0
    
    for block in article_data['content_blocks']:
        block_type = block['type']
        content = block['content']
        
        if block_type.startswith('h'):
            level = int(block_type[1])
            # Convert any H1 to H2 to avoid duplicate H1s in content
            if level == 1:
                html_parts.append(f"<h2>{content}</h2>")
            else:
                html_parts.append(f"<{block_type}>{content}</{block_type}>")
        elif block_type == 'paragraph':
            html_parts.append(f"<p>{content}</p>")
            
            # Insert image after some paragraphs
            if image_index < len(article_data['images']) and len(html_parts) % 3 == 0:
                img = article_data['images'][image_index]
                img_html = f'<figure class="embedded-image"><img src="{img["url"]}" alt="{img["alt_text"]}" style="max-width: 100%; height: auto; margin: 1rem 0;"><figcaption>{img.get("caption", f"Figure {image_index + 1}")}</figcaption></figure>'
                html_parts.append(img_html)
                image_index += 1
    
    # Add remaining images at the end
    while image_index < len(article_data['images']):
        img = article_data['images'][image_index]
        img_html = f'<figure class="embedded-image"><img src="{img["url"]}" alt="{img["alt_text"]}" style="max-width: 100%; height: auto; margin: 1rem 0;"><figcaption>{img.get("caption", f"Figure {image_index + 1}")}</figcaption></figure>'
        html_parts.append(img_html)
        image_index += 1
    
    return '\n\n'.join(html_parts)

async def generate_enhanced_markdown_content(article_data: dict, template_data: dict) -> str:
    """Generate Markdown content for an article"""
    md_parts = []
    
    # Add title
    if article_data['title']:
        md_parts.append(f"# {article_data['title']}")
    
    # Process content blocks
    image_index = 0
    
    for block in article_data['content_blocks']:
        block_type = block['type']
        content = block['content']
        
        if block_type == 'h1':
            md_parts.append(f"# {content}")
        elif block_type == 'h2':
            md_parts.append(f"## {content}")
        elif block_type == 'h3':
            md_parts.append(f"### {content}")
        elif block_type == 'h4':
            md_parts.append(f"#### {content}")
        elif block_type == 'paragraph':
            # Convert HTML formatting to Markdown
            md_content = content.replace('<strong>', '**').replace('</strong>', '**')
            md_content = md_content.replace('<em>', '*').replace('</em>', '*')
            md_content = md_content.replace('<u>', '_').replace('</u>', '_')
            md_parts.append(md_content)
            
            # Insert image after some paragraphs
            if image_index < len(article_data['images']) and len(md_parts) % 3 == 0:
                img = article_data['images'][image_index]
                md_parts.append(f"![{img['alt_text']}]({img['url']})")
                if img.get('caption'):
                    md_parts.append(f"*{img['caption']}*")
                image_index += 1
    
    # Add remaining images at the end
    while image_index < len(article_data['images']):
        img = article_data['images'][image_index]
        md_parts.append(f"![{img['alt_text']}]({img['url']})")
        if img.get('caption'):
            md_parts.append(f"*{img['caption']}*")
        image_index += 1
    
    return '\n\n'.join(md_parts)

async def generate_comprehensive_outline(content: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
    """Step 2: Generate a comprehensive topic-specific outline after content analysis"""
    try:
        print(f"📋 STEP 2: GENERATING COMPREHENSIVE OUTLINE from {len(content)} characters")
        
        system_message = """You are a technical documentation expert. Analyze the provided content and create a comprehensive, topic-specific outline.

CRITICAL INSTRUCTIONS:
1. Read and analyze the ENTIRE content thoroughly
2. Identify ALL major topics, subtopics, and important details
3. Create a logical, hierarchical structure that covers 100% of the content
4. Each outline item should be substantial enough for a detailed article
5. Ensure NO content is missed or summarized away
6. Group related topics together logically

OUTPUT FORMAT - Return valid JSON:
{
  "document_analysis": {
    "title": "Main document title/subject",
    "type": "user guide|technical manual|API documentation|tutorial|etc",
    "complexity": "basic|intermediate|advanced",
    "estimated_reading_time": 15
  },
  "comprehensive_outline": [
    {
      "topic_title": "Introduction and Overview",
      "topic_type": "overview",
      "content_focus": "What this covers and why it matters",
      "key_points": ["point1", "point2", "point3"],
      "estimated_length": "medium",
      "priority": "high"
    },
    {
      "topic_title": "Getting Started Guide", 
      "topic_type": "how-to",
      "content_focus": "Step-by-step setup and configuration",
      "key_points": ["installation", "setup", "first steps"],
      "estimated_length": "long",
      "priority": "high"
    }
  ],
  "outline_summary": {
    "total_topics": 8,
    "overview_topics": 1,
    "how_to_topics": 3,
    "reference_topics": 2,
    "concept_topics": 2
  }
}

Create as many topics as needed for comprehensive coverage. Do not limit yourself."""

        # Generate outline using LLM
        outline_response = await call_llm_with_fallback(
            system_message=system_message,
            user_message=f"Analyze this content and create a comprehensive topic-specific outline:\n\n{content[:25000]}"
        )
        
        if outline_response:
            try:
                # Clean the response - remove code block markers if present
                cleaned_response = outline_response.strip()
                if cleaned_response.startswith('```json'):
                    cleaned_response = cleaned_response[7:]  # Remove ```json
                if cleaned_response.endswith('```'):
                    cleaned_response = cleaned_response[:-3]  # Remove ```
                cleaned_response = cleaned_response.strip()
                
                outline_data = json.loads(cleaned_response)
                total_topics = len(outline_data.get('comprehensive_outline', []))
                print(f"✅ COMPREHENSIVE OUTLINE GENERATED: {total_topics} topics planned")
                return outline_data
            except json.JSONDecodeError as e:
                print(f"⚠️ Outline JSON parsing error: {e}")
                print(f"Raw response: {outline_response[:500]}...")
                
        print("⚠️ Outline generation failed")
        return None
        
    except Exception as e:
        print(f"❌ Error generating comprehensive outline: {e}")
        return None

async def create_articles_from_outline(content: str, outline: Dict[str, Any], metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Create articles one by one based on the comprehensive outline"""
    try:
        articles = []
        outline_items = outline.get('comprehensive_outline', [])
        
        print(f"🏗️ CREATING ARTICLES FROM OUTLINE: {len(outline_items)} articles planned")
        
        for i, item in enumerate(outline_items):
            try:
                print(f"📄 Creating article {i+1}/{len(outline_items)}: {item.get('article_title', 'Untitled')}")
                
                # Create targeted system message for this specific article
                system_message = f"""You are a technical writer creating a comprehensive article based on a specific section of a larger document.

ARTICLE REQUIREMENTS:
- Title: {item.get('article_title', 'Untitled')}
- Type: {item.get('article_type', 'informational')}
- Content Focus: {item.get('content_summary', 'General content')}
- Key Topics: {', '.join(item.get('key_topics', []))}

CRITICAL INSTRUCTIONS:
1. Extract ALL relevant content from the source document for this article
2. Create a comprehensive, standalone article that covers the topic thoroughly
3. Include ALL technical details, procedures, and specifications
4. Use proper HTML formatting with headings, lists, and emphasis
5. Ensure the article is complete and self-contained
6. Do NOT reference other articles or create placeholder links

CRITICAL OUTPUT FORMAT:
- Return ONLY article content HTML (headings, paragraphs, lists, etc.)
- Do NOT include document structure tags: <!DOCTYPE>, <html>, <head>, <body>
- Do NOT wrap content in ```html code blocks
- Start directly with content (e.g., <h2>Introduction</h2><p>Content...</p>)
- Use semantic HTML: <h2>, <h3>, <p>, <ul>, <ol>, <li>, <strong>, <em>, <blockquote>"""

                # Extract relevant content for this specific article
                article_content = await call_llm_with_fallback(
                    system_message=system_message,
                    user_message=f"Create a comprehensive article from this source document:\n\n{content}"
                )
                
                if article_content:
                    # CRITICAL FIX: Clean HTML content to remove document structure
                    cleaned_content = clean_article_html_content(article_content)
                    
                    # CONTENT VALIDATION: Ensure we have substantial content
                    content_text = re.sub(r'<[^>]+>', '', cleaned_content).strip()
                    if len(content_text) < 50:
                        print(f"⚠️ Article content too short ({len(content_text)} chars), regenerating...")
                        # Try to regenerate with enhanced prompt
                        enhanced_system_message = f"""You are a technical writer creating a comprehensive article based on a specific section of a larger document.

ARTICLE REQUIREMENTS:
- Title: {item.get('article_title', 'Untitled')}
- Type: {item.get('article_type', 'informational')}
- Content Focus: {item.get('content_summary', 'General content')}
- Key Topics: {', '.join(item.get('key_topics', []))}

CRITICAL INSTRUCTIONS:
1. Extract ALL relevant content from the source document for this article
2. Create a comprehensive, standalone article that covers the topic thoroughly
3. Include ALL technical details, procedures, and specifications
4. Use proper HTML formatting with headings, lists, and emphasis
5. Ensure the article is complete and self-contained
6. MINIMUM 200 words of actual content

CRITICAL OUTPUT FORMAT:
- Return ONLY article content HTML (headings, paragraphs, lists, etc.)
- Do NOT include document structure tags: <!DOCTYPE>, <html>, <head>, <body>
- Do NOT wrap content in ```html code blocks
- Start directly with content (e.g., <h2>Introduction</h2><p>Content...</p>)
- Use semantic HTML: <h2>, <h3>, <p>, <ul>, <ol>, <li>, <strong>, <em>, <blockquote>
- ENSURE substantial content - at least 200 words"""

                        article_content = await call_llm_with_fallback(
                            system_message=enhanced_system_message,
                            user_message=f"Create a comprehensive article from this source document:\n\n{content}"
                        )
                        if article_content:
                            cleaned_content = clean_article_html_content(article_content)
                            content_text = re.sub(r'<[^>]+>', '', cleaned_content).strip()
                    
                    # Final validation before saving
                    if len(content_text) >= 50:
                        # Create article object
                        article = {
                            "id": str(uuid.uuid4()),
                            "title": item.get('article_title', f'Article {i+1}'),
                            "content": cleaned_content,
                            "status": "published",
                            "article_type": item.get('article_type', 'informational'),
                            "source_document": metadata.get("original_filename", "Unknown"),
                            "estimated_length": item.get('estimated_length', 'medium'),
                            "key_topics": item.get('key_topics', []),
                            "created_at": datetime.utcnow(),
                            "metadata": {
                                "outline_based": True,
                                "article_number": i + 1,
                                "total_articles": len(outline_items),
                                "content_validated": True,
                                "content_length": len(content_text),
                                **metadata
                            }
                        }
                        
                        # KE-PR9: Add TICKET-3 fields
                        article = ensure_ticket3_fields(article)
                        
                        # KE-PR9: Use repository pattern for consistent data access
                        if mongo_repo_available:
                            content_repo = RepositoryFactory.get_content_library()
                            await content_repo.insert_article(article)
                        else:
                            # KE-PR9.3: Fallback to repository pattern
                    content_repo = RepositoryFactory.get_content_library()
                    target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
                
                if not target_doc:
                    broken_links.append({
                        "type": "xref",
                        "target_doc_uid": target_doc_uid,
                        "anchor_id": anchor_id,
                        "label": label,
                        "reason": "target_document_not_found"
                    })
                    continue
                
                # Check if anchor exists in target document headings
                target_headings = target_doc.get("headings", [])
                anchor_exists = any(h.get("id") == anchor_id for h in target_headings)
                
                if not anchor_exists:
                    broken_links.append({
                        "type": "xref", 
                        "target_doc_uid": target_doc_uid,
                        "anchor_id": anchor_id,
                        "label": label,
                        "reason": "anchor_not_found_in_target",
                        "available_anchors": [h.get("id") for h in target_headings[:5]]  # First 5 for debugging
                    })
            
            # Validate related_links (similar process)
            for related in related_links:
                target_doc_uid = related.get("doc_uid")
                anchor_id = related.get("anchor_id", "")
                
                # KE-PR9: Find target document using repository pattern
                if mongo_repo_available:
                    content_repo = RepositoryFactory.get_content_library()
                    target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
                else:
                    # KE-PR9.3: Fallback to repository pattern
                    content_repo = RepositoryFactory.get_content_library()
                    target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
                
                if not target_doc:
                    broken_links.append({
                        "type": "related_link",
                        "target_doc_uid": target_doc_uid,
                        "anchor_id": anchor_id,
                        "reason": "target_document_not_found"
                    })
                    continue
                
                if anchor_id:  # Only validate anchor if specified
                    target_headings = target_doc.get("headings", [])
                    anchor_exists = any(h.get("id") == anchor_id for h in target_headings)
                    
                    if not anchor_exists:
                        broken_links.append({
                            "type": "related_link",
                            "target_doc_uid": target_doc_uid, 
                            "anchor_id": anchor_id,
                            "reason": "anchor_not_found_in_target"
                        })
            
            resolution_rate = ((total_links - len(broken_links)) / total_links * 100) if total_links > 0 else 100
            
            print(f"🔍 TICKET 3: Link validation complete - {resolution_rate:.1f}% resolved ({len(broken_links)} broken)")
            
            return {
                "total_links": total_links,
                "broken_links": broken_links,
                "resolution_rate": resolution_rate,
                "links_resolve": len(broken_links) == 0
            }
            
        except Exception as e:
            print(f"❌ TICKET 3: Error validating cross-document links - {e}")
            return {
                "total_links": 0,
                "broken_links": [],
                "resolution_rate": 0,
                "links_resolve": False,
                "error": str(e)
            }
    
    def _apply_bookmark_registry(self, content: str, article_title: str) -> dict:
        """TICKET 3: Apply bookmark registry extraction for universal links"""
        try:
            print(f"📖 TICKET 3: Starting bookmark registry for '{article_title[:50]}...'")
            
            # Extract headings from content using V2ValidationSystem method
            v2_validator = V2ValidationSystem()
            headings = v2_validator.extract_headings_registry(content)
            
            # Generate document identifiers
            doc_uid = self.generate_doc_uid()
            doc_slug = self.generate_doc_slug(article_title)
            
            print(f"📖 TICKET 3: Bookmark registry complete - {len(headings)} headings, doc_uid: {doc_uid}")
            
            return {
                'headings': headings,
                'doc_uid': doc_uid,
                'doc_slug': doc_slug,
                'bookmark_count': len(headings),
                'changes_applied': [f"Extracted {len(headings)} bookmarks", f"Generated doc_uid: {doc_uid}", f"Generated doc_slug: {doc_slug}"]
            }
            
        except Exception as e:
            print(f"❌ TICKET 3: Error in bookmark registry - {e}")
            return {
                'headings': [],
                'doc_uid': None,
                'doc_slug': None,
                'bookmark_count': 0,
                'changes_applied': [f"Bookmark registry error: {str(e)}"],
                'error': str(e)
            }
    
    async def backfill_bookmark_registry(self, limit: int = None) -> dict:
        """TICKET 3: Backfill existing v2 articles with bookmark registry data"""
        try:
            from motor.motor_asyncio import AsyncIOMotorClient
            import os
            
            mongo_url = os.environ.get('MONGO_URL')
            client = AsyncIOMotorClient(mongo_url)
            db = client.promptsupport
            
            print(f"🔄 TICKET 3: Starting bookmark registry backfill for existing v2 articles")
            
            # KE-PR9: Use repository pattern to find articles needing backfill
            try:
                if mongo_repo_available:
                    content_repo = RepositoryFactory.get_content_library()
                    # Get all V2 articles and filter those needing backfill
                    all_articles = await content_repo.find_by_engine("v2", limit=limit)
                    
                    # Filter articles that need backfilling
                    articles = []
                    for article in all_articles:
                        needs_backfill = (
                            not article.get("doc_uid") or 
                            not article.get("headings") or 
                            (isinstance(article.get("headings"), list) and len(article.get("headings")) == 0)
                        )
                        if needs_backfill:
                            articles.append(article)
                else:
                    # Fallback to direct database query
                    query = {
                        "$or": [
                            {"metadata.engine": "v2", "doc_uid": {"$exists": False}},
                            {"metadata.engine": "v2", "headings": {"$exists": False}},
                            {"metadata.engine": "v2", "headings": []}  # Empty headings array
                        ]
                    }
                    
                    if limit:
                        articles_cursor = db.content_library.find(query).limit(limit)
                    else:
                        articles_cursor = db.content_library.find(query)
                    
                    articles = await articles_cursor.to_list(None)
                    
            except Exception as repo_error:
                print(f"⚠️ KE-PR9: Backfill query fallback to direct DB: {repo_error}")
                # Final fallback to direct database query
                query = {
                    "$or": [
                        {"metadata.engine": "v2", "doc_uid": {"$exists": False}},
                        {"metadata.engine": "v2", "headings": {"$exists": False}},
                        {"metadata.engine": "v2", "headings": []}  # Empty headings array
                    ]
                }
                
                if limit:
                    articles_cursor = db.content_library.find(query).limit(limit)
                else:
                    articles_cursor = db.content_library.find(query)
                
                articles = await articles_cursor.to_list(None)
                
            total_articles = len(articles)
            
            if total_articles == 0:
                print("✅ TICKET 3: No articles need backfilling")
                return {"articles_processed": 0, "success": True}
            
            processed_count = 0
            error_count = 0
            
            for article in articles:
                try:
                    article_id = article.get('_id')
                    title = article.get('title', 'Untitled')
                    content = article.get('content', '') or article.get('html', '')
                    
                    if not content:
                        print(f"⚠️ TICKET 3: Skipping article '{title}' - no content found")
                        continue
                    
                    # Generate doc_uid and doc_slug if missing
                    doc_uid = article.get('doc_uid')
                    if not doc_uid:
                        doc_uid = self.generate_doc_uid()
                    
                    doc_slug = article.get('doc_slug')  
                    if not doc_slug:
                        doc_slug = self.generate_doc_slug(title)
                    
                    # Apply Ticket 2 stable anchors if needed (ensure IDs exist)
                    processed_content = self.assign_heading_ids(content)
                    
                    # Extract headings registry
                    headings = self.extract_headings_registry(processed_content)
                    
                    # Update article in database
                    update_data = {
                        "doc_uid": doc_uid,
                        "doc_slug": doc_slug,
                        "headings": headings,
                        "content": processed_content,  # Updated content with IDs
                        "html": processed_content     # Sync html field
                    }
                    
                    # Initialize xrefs and related_links if missing
                    if "xrefs" not in article:
                        update_data["xrefs"] = []
                    if "related_links" not in article:
                        update_data["related_links"] = []
                    
                    await db.content_library.update_one(
                        {"_id": article_id},
                        {"$set": update_data}
                    )
                    
                    processed_count += 1
                    print(f"📖 TICKET 3: Backfilled article '{title[:50]}...' - doc_uid: {doc_uid}, {len(headings)} headings")
                    
                except Exception as article_error:
                    error_count += 1
                    print(f"❌ TICKET 3: Error backfilling article '{article.get('title', 'Unknown')}' - {article_error}")
            
            success_rate = (processed_count / total_articles * 100) if total_articles > 0 else 100
            
            print(f"✅ TICKET 3: Backfill complete - {processed_count}/{total_articles} articles processed ({success_rate:.1f}% success)")
            
            return {
                "articles_found": total_articles,
                "articles_processed": processed_count,
                "articles_failed": error_count,
                "success_rate": success_rate,
                "success": error_count == 0
            }
            
        except Exception as e:
            print(f"❌ TICKET 3: Error in backfill process - {e}")
            return {
                "articles_found": 0,
                "articles_processed": 0,
                "articles_failed": 0,
                "success_rate": 0,
                "success": False,
                "error": str(e)
            }

# Global V2 Style Processor instance
v2_style_processor = V2StyleProcessor()

# ========================================
# V2 ENGINE: RELATED LINKS SYSTEM
# ========================================

class V2RelatedLinksSystem:
    """V2 Engine: Enhanced related links system with content library indexing and similarity matching"""
    
    def __init__(self):
        self.content_index = {}  # Cache for content library index
        self.index_last_updated = None
        
    async def generate_related_links(self, article: dict, source_content: str, 
                                   source_blocks: list, run_id: str) -> dict:
        """Generate comprehensive related links for an article"""
        try:
            article_title = article.get('title', 'Untitled')
            print(f"🔗 V2 RELATED LINKS: Generating related links for '{article_title}' - engine=v2")
            
            # Step 1: Build/update content library index
            await self._update_content_index()
            
            # Step 2: Find related internal articles
            internal_links = await self._find_internal_related_articles(article)
            
            # Step 3: Extract external links from source content and blocks
            external_links = await self._extract_source_external_links(source_content, source_blocks)
            
            # Step 4: Merge and format final related links
            final_related_links = self._merge_and_format_links(internal_links, external_links)
            
            related_links_count = len(final_related_links)
            internal_count = len([l for l in final_related_links if l.get('type') == 'internal'])
            external_count = len([l for l in final_related_links if l.get('type') == 'external'])
            
            print(f"✅ V2 RELATED LINKS: Generated {related_links_count} links ({internal_count} internal, {external_count} external) for '{article_title}' - engine=v2")
            
            return {
                "related_links_id": f"related_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "article_title": article_title,
                "related_links_status": "success",
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2",
                
                # Related links data
                "related_links": final_related_links,
                "internal_links_count": internal_count,
                "external_links_count": external_count,
                "total_links_count": related_links_count,
                
                # Metadata
                "content_library_articles_indexed": len(self.content_index),
                "similarity_method": "keyword_and_semantic"
            }
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error generating related links - {e} - engine=v2")
            return {
                "related_links_id": f"related_error_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "related_links_status": "error",
                "error": str(e),
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
    
    async def _update_content_index(self):
        """Build/update lightweight vector/keyword index over content_library"""
        try:
            # Check if index needs updating (update every 5 minutes)
            now = datetime.utcnow()
            if (self.index_last_updated and 
                (now - self.index_last_updated).total_seconds() < 300):
                return  # Index is still fresh
            
            print(f"🔍 V2 RELATED LINKS: Updating content library index - engine=v2")
            
            new_index = {}
            
            # Get all articles from content_library
            articles_cursor = db.content_library.find({"engine": "v2"})
            articles_count = 0
            
            async for article in articles_cursor:
                try:
                    article_id = str(article.get('_id', ''))
                    title = article.get('title', '').strip()
                    content = article.get('content', '') or article.get('html', '')
                    
                    if not title or not content:
                        continue
                    
                    # Extract H2 headings for indexing
                    h2_headings = self._extract_headings(content, level=2)
                    
                    # Create summary from first paragraph or first 200 chars
                    summary = self._extract_summary(content)
                    
                    # Build keyword index
                    keywords = self._extract_keywords(title, summary, h2_headings)
                    
                    new_index[article_id] = {
                        "title": title,
                        "summary": summary,
                        "h2_headings": h2_headings,
                        "keywords": keywords,
                        "created_at": article.get('created_at'),
                        "url": f"/content-library/article/{article_id}"
                    }
                    
                    articles_count += 1
                    
                except Exception as article_error:
                    print(f"❌ V2 RELATED LINKS: Error indexing article {article.get('_id')} - {article_error}")
                    continue
            
            self.content_index = new_index
            self.index_last_updated = now
            
            print(f"✅ V2 RELATED LINKS: Content index updated - {articles_count} articles indexed - engine=v2")
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error updating content index - {e} - engine=v2")
    
    def _extract_headings(self, content: str, level: int = 2) -> list:
        """Extract headings of specified level from HTML/markdown content"""
        import re
        from bs4 import BeautifulSoup
        
        headings = []
        
        try:
            # Try HTML parsing first
            soup = BeautifulSoup(content, 'html.parser')
            heading_tags = soup.find_all(f'h{level}')
            
            for tag in heading_tags:
                heading_text = tag.get_text().strip()
                if heading_text:
                    headings.append(heading_text)
            
            # Also try markdown parsing
            markdown_pattern = rf'^{"#" * level}\s+(.+)$'
            markdown_headings = re.findall(markdown_pattern, content, re.MULTILINE)
            headings.extend(markdown_headings)
            
            # Remove duplicates while preserving order
            seen = set()
            unique_headings = []
            for heading in headings:
                if heading.lower() not in seen:
                    seen.add(heading.lower())
                    unique_headings.append(heading)
            
            return unique_headings[:10]  # Limit to first 10 headings
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error extracting headings - {e}")
            return []
    
    def _extract_summary(self, content: str) -> str:
        """Extract summary from article content (first paragraph or first 200 chars)"""
        try:
            from bs4 import BeautifulSoup
            
            # Try to get first paragraph
            soup = BeautifulSoup(content, 'html.parser')
            
            # Look for first paragraph
            first_para = soup.find('p')
            if first_para:
                summary = first_para.get_text().strip()
                if len(summary) > 50:
                    return summary[:200] + "..." if len(summary) > 200 else summary
            
            # Fallback: clean text and take first 200 chars
            clean_text = soup.get_text().strip()
            if clean_text:
                return clean_text[:200] + "..." if len(clean_text) > 200 else clean_text
            
            return ""
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error extracting summary - {e}")
            return ""
    
    def _extract_keywords(self, title: str, summary: str, headings: list) -> set:
        """Extract keywords from title, summary, and headings for matching"""
        import re
        
        keywords = set()
        
        # Common stop words to filter out
        stop_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 
            'by', 'from', 'as', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had',
            'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might', 'can',
            'this', 'that', 'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they'
        }
        
        # Extract from title (higher weight)
        title_words = re.findall(r'\b[a-zA-Z]{3,}\b', title.lower())
        keywords.update(word for word in title_words if word not in stop_words)
        
        # Extract from summary
        summary_words = re.findall(r'\b[a-zA-Z]{3,}\b', summary.lower())
        keywords.update(word for word in summary_words[:20] if word not in stop_words)  # Limit summary words
        
        # Extract from headings
        for heading in headings:
            heading_words = re.findall(r'\b[a-zA-Z]{3,}\b', heading.lower())
            keywords.update(word for word in heading_words if word not in stop_words)
        
        return keywords
    
    async def _find_internal_related_articles(self, article: dict) -> list:
        """Find top 5 related internal articles using keyword similarity"""
        try:
            article_title = article.get('title', '')
            article_content = article.get('content', '') or article.get('html', '')
            
            if not article_title or not self.content_index:
                return []
            
            # Extract keywords from current article
            article_h2s = self._extract_headings(article_content, level=2)
            article_summary = self._extract_summary(article_content)
            article_keywords = self._extract_keywords(article_title, article_summary, article_h2s)
            
            # Calculate similarity scores
            similarity_scores = []
            
            for indexed_id, indexed_article in self.content_index.items():
                # Skip if it's the same article (by title comparison)
                if indexed_article['title'].lower().strip() == article_title.lower().strip():
                    continue
                
                indexed_keywords = indexed_article['keywords']
                
                # Calculate keyword overlap similarity
                if not article_keywords or not indexed_keywords:
                    similarity = 0
                else:
                    common_keywords = article_keywords.intersection(indexed_keywords)
                    similarity = len(common_keywords) / max(len(article_keywords), len(indexed_keywords))
                
                # Boost score for title word matches
                title_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', article_title.lower()))
                indexed_title_words = set(re.findall(r'\b[a-zA-Z]{3,}\b', indexed_article['title'].lower()))
                title_overlap = len(title_words.intersection(indexed_title_words))
                
                final_score = similarity + (title_overlap * 0.3)  # Boost for title matches
                
                if final_score > 0.1:  # Only include articles with some relevance
                    similarity_scores.append({
                        "id": indexed_id,
                        "title": indexed_article['title'],
                        "url": indexed_article['url'],
                        "summary": indexed_article['summary'],
                        "similarity_score": final_score,
                        "common_keywords": len(article_keywords.intersection(indexed_keywords)) if article_keywords and indexed_keywords else 0
                    })
            
            # Sort by similarity score and take top 5
            similarity_scores.sort(key=lambda x: x['similarity_score'], reverse=True)
            top_related = similarity_scores[:5]
            
            # Filter same-topic duplicates (very similar titles)
            filtered_related = []
            used_title_words = set()
            
            for related in top_related:
                related_title_words = set(re.findall(r'\b[a-zA-Z]{4,}\b', related['title'].lower()))
                
                # Check if this is too similar to already included articles
                is_duplicate = False
                for used_words in used_title_words:
                    overlap = len(related_title_words.intersection(used_words))
                    if overlap >= 2:  # 2+ shared significant words = duplicate topic
                        is_duplicate = True
                        break
                
                if not is_duplicate:
                    filtered_related.append({
                        "title": related['title'],
                        "url": related['url'],
                        "type": "internal",
                        "description": related['summary'][:100] + "..." if len(related['summary']) > 100 else related['summary'],
                        "similarity_score": related['similarity_score']
                    })
                    used_title_words.add(related_title_words)
                
                if len(filtered_related) >= 5:
                    break
            
            return filtered_related
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error finding internal related articles - {e}")
            return []
    
    async def _extract_source_external_links(self, source_content: str, source_blocks: list) -> list:
        """Extract external links that are actually present in source content/blocks"""
        try:
            from bs4 import BeautifulSoup
            import re
            
            external_links = []
            found_urls = set()
            
            # Extract from source content
            if source_content:
                external_links.extend(self._extract_external_links_from_text(source_content, found_urls))
            
            # Extract from source blocks
            for block in source_blocks[:50]:  # Limit to first 50 blocks for performance
                block_content = block.get('content', '') or block.get('text', '')
                if block_content:
                    external_links.extend(self._extract_external_links_from_text(block_content, found_urls))
            
            # Limit to reasonable number and prioritize quality
            unique_links = []
            for link in external_links:
                if link['url'] not in [existing['url'] for existing in unique_links]:
                    unique_links.append(link)
            
            # Sort by title length (longer titles usually more descriptive)
            unique_links.sort(key=lambda x: len(x.get('title', '')), reverse=True)
            
            return unique_links[:10]  # Limit to 10 external links
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error extracting source external links - {e}")
            return []
    
    def _extract_external_links_from_text(self, text: str, found_urls: set) -> list:
        """Extract external links from a piece of text"""
        try:
            from bs4 import BeautifulSoup
            import re
            
            links = []
            
            # Try HTML parsing first
            try:
                soup = BeautifulSoup(text, 'html.parser')
                html_links = soup.find_all('a', href=True)
                
                for link in html_links:
                    href = link.get('href', '').strip()
                    link_text = link.get_text().strip()
                    
                    if self._is_valid_external_url(href) and href not in found_urls:
                        found_urls.add(href)
                        links.append({
                            "title": link_text or self._extract_domain_name(href),
                            "url": href,
                            "type": "external",
                            "description": f"External link: {self._extract_domain_name(href)}"
                        })
            except:
                pass
            
            # Also extract plain URLs from text
            url_pattern = r'https?://[^\s<>"\'()\\]+[^\s<>"\'()\\.,;:]'
            plain_urls = re.findall(url_pattern, text)
            
            for url in plain_urls:
                if self._is_valid_external_url(url) and url not in found_urls:
                    found_urls.add(url)
                    links.append({
                        "title": self._extract_domain_name(url),
                        "url": url,
                        "type": "external", 
                        "description": f"External resource: {self._extract_domain_name(url)}"
                    })
            
            return links[:5]  # Limit per text block
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error extracting links from text - {e}")
            return []
    
    def _is_valid_external_url(self, url: str) -> bool:
        """Check if URL is a valid external URL"""
        if not url:
            return False
        
        # Must be http/https
        if not url.startswith(('http://', 'https://')):
            return False
        
        # Skip internal/relative links  
        if url.startswith('#') or url.startswith('/'):
            return False
        
        # Skip common non-content URLs
        skip_domains = {'example.com', 'test.com', 'localhost', 'placeholder'}
        domain = self._extract_domain_name(url).lower()
        
        if any(skip in domain for skip in skip_domains):
            return False
        
        # Must have valid domain structure
        if '.' not in domain or len(domain) < 4:
            return False
        
        return True
    
    def _extract_domain_name(self, url: str) -> str:
        """Extract domain name from URL"""
        try:
            from urllib.parse import urlparse
            parsed = urlparse(url)
            domain = parsed.netloc or parsed.path.split('/')[0]
            return domain.replace('www.', '')
        except:
            return url[:50]  # Fallback to first 50 chars
    
    def _merge_and_format_links(self, internal_links: list, external_links: list) -> list:
        """Merge internal and external links and format for final output"""
        try:
            merged_links = []
            
            # Add internal links (3-6 as specified)
            for link in internal_links[:6]:
                merged_links.append({
                    "title": link['title'],
                    "url": link['url'],
                    "type": "internal",
                    "description": link.get('description', ''),
                    "priority": "high"  # Internal links get high priority
                })
            
            # Add external links (any valid ones found in source)
            for link in external_links:
                merged_links.append({
                    "title": link['title'],
                    "url": link['url'],
                    "type": "external",
                    "description": link.get('description', ''),
                    "priority": "medium"  # External links get medium priority
                })
            
            # Ensure we have 3-6 total links by adjusting if needed
            if len(merged_links) < 3 and len(internal_links) > 0:
                # If we have fewer than 3, try to add more internal links
                additional_internal = internal_links[len(internal_links[:6]):6]
                for link in additional_internal:
                    if len(merged_links) >= 6:
                        break
                    merged_links.append({
                        "title": link['title'],
                        "url": link['url'],
                        "type": "internal",
                        "description": link.get('description', ''),
                        "priority": "medium"
                    })
            
            return merged_links
            
        except Exception as e:
            print(f"❌ V2 RELATED LINKS: Error merging links - {e}")
            return []

# Global V2 Related Links System instance
v2_related_links_system = V2RelatedLinksSystem()

# ========================================
# V2 ENGINE: GAP FILLING SYSTEM
# ========================================

class V2GapFillingSystem:
    """V2 Engine: Intelligent gap filling system to replace [MISSING] placeholders with in-corpus retrieval"""
    
    def __init__(self):
        self.gap_patterns = [
            r'\[MISSING\]',
            r'\[PLACEHOLDER\]',
            r'\[TBD\]',
            r'\[TODO\]',
            r'\[FILL\]'
        ]
        
    async def fill_content_gaps(self, articles: list, source_content: str, 
                               source_blocks: list, run_id: str, 
                               enrich_mode: str = "internal") -> dict:
        """Fill gaps in articles using in-corpus retrieval and pattern synthesis"""
        try:
            print(f"🔍 V2 GAP FILLING: Starting gap filling process - mode: {enrich_mode} - engine=v2")
            
            gap_filling_results = []
            total_gaps_found = 0
            total_gaps_filled = 0
            
            for i, article in enumerate(articles):
                try:
                    article_title = article.get('title', f'Article {i+1}')
                    article_content = article.get('content', '') or article.get('html', '')
                    
                    if not article_content:
                        continue
                    
                    # Step 1: Detect gaps in the article
                    gaps_detected = self._detect_gaps(article_content, article_title)
                    
                    if not gaps_detected:
                        # No gaps found
                        gap_filling_results.append({
                            "article_index": i,
                            "article_title": article_title,
                            "gap_filling_status": "no_gaps",
                            "gaps_found": 0,
                            "gaps_filled": 0,
                            "patches_applied": []
                        })
                        continue
                    
                    total_gaps_found += len(gaps_detected)
                    
                    # Step 2: Retrieve relevant content for gap filling
                    retrieval_results = await self._retrieve_gap_context(
                        gaps_detected, source_content, source_blocks, enrich_mode
                    )
                    
                    # Step 3: Generate patches using LLM
                    patches = await self._generate_gap_patches(
                        gaps_detected, retrieval_results, enrich_mode
                    )
                    
                    # Step 4: Apply patches to article content
                    patched_content, applied_patches = self._apply_patches(
                        article_content, patches
                    )
                    
                    # Update article with patched content
                    article['content'] = patched_content
                    article['html'] = patched_content
                    
                    # Track gap filling metadata
                    gaps_filled_count = len(applied_patches)
                    total_gaps_filled += gaps_filled_count
                    
                    gap_filling_result = {
                        "article_index": i,
                        "article_title": article_title,
                        "gap_filling_status": "success" if gaps_filled_count > 0 else "no_patches",
                        "gaps_found": len(gaps_detected),
                        "gaps_filled": gaps_filled_count,
                        "patches_applied": applied_patches,
                        "retrieval_sources": len(retrieval_results),
                        "enrich_mode": enrich_mode
                    }
                    
                    gap_filling_results.append(gap_filling_result)
                    
                    print(f"✅ V2 GAP FILLING: Processed '{article_title[:50]}...' - {gaps_filled_count}/{len(gaps_detected)} gaps filled - engine=v2")
                    
                except Exception as article_error:
                    print(f"❌ V2 GAP FILLING: Error processing article {i+1} - {article_error} - engine=v2")
                    gap_filling_results.append({
                        "article_index": i,
                        "article_title": article.get('title', f'Article {i+1}'),
                        "gap_filling_status": "error",
                        "error": str(article_error),
                        "gaps_found": 0,
                        "gaps_filled": 0
                    })
            
            # Calculate success metrics
            successful_articles = len([r for r in gap_filling_results if r.get('gap_filling_status') == 'success'])
            gap_fill_rate = (total_gaps_filled / total_gaps_found * 100) if total_gaps_found > 0 else 100
            
            return {
                "gap_filling_id": f"gaps_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "gap_filling_status": "success",
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2",
                
                # Gap filling metrics
                "articles_processed": len(articles),
                "articles_with_gaps": len([r for r in gap_filling_results if r.get('gaps_found', 0) > 0]),
                "successful_gap_filling": successful_articles,
                "total_gaps_found": total_gaps_found,
                "total_gaps_filled": total_gaps_filled,
                "gap_fill_rate": gap_fill_rate,
                "enrich_mode": enrich_mode,
                
                # Detailed results
                "gap_filling_results": gap_filling_results
            }
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error in gap filling process - {e} - engine=v2")
            return {
                "gap_filling_id": f"gaps_error_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "gap_filling_status": "error",
                "error": str(e),
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
    
    def _detect_gaps(self, content: str, article_title: str) -> list:
        """Detect gap placeholders in article content"""
        try:
            import re
            gaps_found = []
            
            for pattern in self.gap_patterns:
                matches = list(re.finditer(pattern, content, re.IGNORECASE))
                
                for match in matches:
                    # Get surrounding context for better gap understanding
                    start_pos = max(0, match.start() - 100)
                    end_pos = min(len(content), match.end() + 100)
                    context = content[start_pos:end_pos]
                    
                    # Try to infer what type of content is missing
                    gap_type = self._infer_gap_type(context, match.group())
                    
                    gaps_found.append({
                        "pattern": match.group(),
                        "position": match.start(),
                        "context": context,
                        "gap_type": gap_type,
                        "section": self._extract_section_name(content, match.start())
                    })
            
            return gaps_found
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error detecting gaps - {e}")
            return []
    
    def _infer_gap_type(self, context: str, gap_placeholder: str) -> str:
        """Infer the type of content that should fill the gap based on context"""
        context_lower = context.lower()
        
        # API-related gaps
        if any(keyword in context_lower for keyword in ['api', 'endpoint', 'request', 'response']):
            return "api_detail"
        
        # Code-related gaps
        if any(keyword in context_lower for keyword in ['code', 'function', 'method', 'class']):
            return "code_example"
        
        # Configuration gaps
        if any(keyword in context_lower for keyword in ['config', 'setting', 'parameter', 'option']):
            return "configuration"
        
        # Authentication gaps
        if any(keyword in context_lower for keyword in ['auth', 'token', 'key', 'credential']):
            return "authentication"
        
        # Step or procedure gaps
        if any(keyword in context_lower for keyword in ['step', 'follow', 'process', 'procedure']):
            return "procedure_step"
        
        # Default to generic content
        return "generic_content"
    
    def _extract_section_name(self, content: str, position: int) -> str:
        """Extract the section name where the gap occurs"""
        try:
            import re
            
            # Look for the nearest heading before the gap position
            content_before = content[:position]
            
            # Find H2/H3 headings
            heading_pattern = r'<h[23][^>]*>(.*?)</h[23]>|^#{1,3}\s+(.+)$'
            headings = list(re.finditer(heading_pattern, content_before, re.MULTILINE | re.IGNORECASE))
            
            if headings:
                last_heading = headings[-1]
                heading_text = last_heading.group(1) or last_heading.group(2)
                return heading_text.strip() if heading_text else "Unknown Section"
            
            return "Introduction"
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error extracting section name - {e}")
            return "Unknown Section"
    
    async def _retrieve_gap_context(self, gaps: list, source_content: str, 
                                   source_blocks: list, enrich_mode: str) -> list:
        """Retrieve relevant content to fill gaps from available sources"""
        try:
            retrieval_results = []
            
            for gap in gaps:
                gap_type = gap.get('gap_type', 'generic_content')
                section = gap.get('section', 'Unknown Section')
                context = gap.get('context', '')
                
                # Extract keywords from gap context for retrieval
                keywords = self._extract_gap_keywords(context, gap_type)
                
                # Search in source content and blocks
                relevant_blocks = self._search_source_blocks(
                    keywords, source_blocks, gap_type
                )
                
                # If internal mode, also search content library
                if enrich_mode == "internal":
                    library_results = await self._search_content_library(keywords, gap_type)
                    relevant_blocks.extend(library_results)
                
                retrieval_results.append({
                    "gap_context": context,
                    "gap_type": gap_type,
                    "section": section,
                    "keywords": keywords,
                    "relevant_blocks": relevant_blocks[:10],  # Limit to top 10 results
                    "source_count": len(relevant_blocks)
                })
            
            return retrieval_results
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error in gap context retrieval - {e}")
            return []
    
    def _extract_gap_keywords(self, context: str, gap_type: str) -> list:
        """Extract relevant keywords from gap context for retrieval"""
        import re
        
        # Remove common words
        stop_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with',
            'by', 'from', 'as', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had'
        }
        
        # Extract meaningful words
        words = re.findall(r'\b[a-zA-Z]{3,}\b', context.lower())
        keywords = [word for word in words if word not in stop_words]
        
        # Add gap-type specific keywords
        type_keywords = {
            'api_detail': ['api', 'endpoint', 'request', 'response', 'method'],
            'code_example': ['code', 'function', 'example', 'syntax'],
            'configuration': ['config', 'setting', 'parameter', 'value'],
            'authentication': ['auth', 'token', 'key', 'credential', 'login'],
            'procedure_step': ['step', 'process', 'procedure', 'follow']
        }
        
        if gap_type in type_keywords:
            keywords.extend(type_keywords[gap_type])
        
        return list(set(keywords))[:10]  # Return unique keywords, limit to 10
    
    def _search_source_blocks(self, keywords: list, source_blocks: list, gap_type: str) -> list:
        """Search source blocks for content relevant to filling the gap"""
        try:
            relevant_blocks = []
            
            for i, block in enumerate(source_blocks[:100]):  # Limit search scope
                block_content = block.get('content', '') or block.get('text', '')
                
                if not block_content:
                    continue
                
                # Calculate relevance score
                relevance_score = 0
                content_lower = block_content.lower()
                
                for keyword in keywords:
                    if keyword.lower() in content_lower:
                        relevance_score += 1
                
                # Bonus for gap type specific content
                if gap_type == 'api_detail' and any(term in content_lower for term in ['api', 'endpoint', 'http']):
                    relevance_score += 2
                elif gap_type == 'code_example' and any(term in content_lower for term in ['```', 'code', 'function']):
                    relevance_score += 2
                elif gap_type == 'authentication' and any(term in content_lower for term in ['token', 'auth', 'key']):
                    relevance_score += 2
                
                if relevance_score > 0:
                    relevant_blocks.append({
                        "block_id": f"b{i}",
                        "content": block_content[:500],  # Limit block size
                        "relevance_score": relevance_score,
                        "block_type": block.get('block_type', 'text')
                    })
            
            # Sort by relevance score
            relevant_blocks.sort(key=lambda x: x['relevance_score'], reverse=True)
            return relevant_blocks
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error searching source blocks - {e}")
            return []
    
    async def _search_content_library(self, keywords: list, gap_type: str) -> list:
        """Search content library for relevant gap-filling content"""
        try:
            library_results = []
            
            # Search content library for articles with matching keywords
            articles_cursor = db.content_library.find({"engine": "v2"}).limit(20)
            
            async for article in articles_cursor:
                try:
                    content = article.get('content', '') or article.get('html', '')
                    title = article.get('title', '')
                    
                    if not content:
                        continue
                    
                    # Calculate keyword relevance
                    content_lower = content.lower()
                    relevance_score = 0
                    
                    for keyword in keywords:
                        if keyword.lower() in content_lower:
                            relevance_score += 1
                        if keyword.lower() in title.lower():
                            relevance_score += 2  # Title matches get higher score
                    
                    if relevance_score > 0:
                        # Extract relevant snippet
                        snippet = self._extract_relevant_snippet(content, keywords)
                        
                        library_results.append({
                            "block_id": f"lib_{str(article.get('_id', ''))}",
                            "content": snippet,
                            "relevance_score": relevance_score,
                            "block_type": "library_article",
                            "source_title": title
                        })
                
                except Exception as article_error:
                    continue
            
            # Sort by relevance
            library_results.sort(key=lambda x: x['relevance_score'], reverse=True)
            return library_results[:5]  # Return top 5 library results
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error searching content library - {e}")
            return []
    
    def _extract_relevant_snippet(self, content: str, keywords: list) -> str:
        """Extract a relevant snippet from content based on keywords"""
        try:
            # Find the best paragraph that contains multiple keywords
            paragraphs = content.split('\n')
            
            best_paragraph = ""
            best_score = 0
            
            for paragraph in paragraphs:
                if len(paragraph.strip()) < 20:
                    continue
                
                score = 0
                paragraph_lower = paragraph.lower()
                
                for keyword in keywords:
                    if keyword.lower() in paragraph_lower:
                        score += 1
                
                if score > best_score:
                    best_score = score
                    best_paragraph = paragraph
            
            if best_paragraph:
                return best_paragraph[:400] + "..." if len(best_paragraph) > 400 else best_paragraph
            
            # Fallback to first 400 characters
            return content[:400] + "..." if len(content) > 400 else content
            
        except Exception as e:
            return content[:400] + "..." if len(content) > 400 else content
    
    async def _generate_gap_patches(self, gaps: list, retrieval_results: list, 
                                   enrich_mode: str) -> list:
        """Generate patches for gaps using LLM with retrieved context"""
        try:
            patches = []
            
            for i, gap in enumerate(gaps):
                try:
                    retrieval_result = retrieval_results[i] if i < len(retrieval_results) else {}
                    relevant_blocks = retrieval_result.get('relevant_blocks', [])
                    
                    if not relevant_blocks and enrich_mode == "internal":
                        # No relevant content found for internal mode - skip
                        continue
                    
                    # Create LLM prompt for gap patching
                    patch_result = await self._create_gap_patch(
                        gap, relevant_blocks, enrich_mode
                    )
                    
                    if patch_result:
                        patches.append(patch_result)
                
                except Exception as gap_error:
                    print(f"❌ V2 GAP FILLING: Error generating patch for gap - {gap_error}")
                    continue
            
            return patches
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error generating gap patches - {e}")
            return []
    
    async def _create_gap_patch(self, gap: dict, relevant_blocks: list, enrich_mode: str) -> dict:
        """Create a single gap patch using LLM"""
        try:
            gap_context = gap.get('context', '')
            gap_type = gap.get('gap_type', 'generic_content')
            section = gap.get('section', 'Unknown Section')
            
            # Create evidence text from relevant blocks
            evidence_text = ""
            support_block_ids = []
            
            if relevant_blocks:
                for block in relevant_blocks[:5]:  # Use top 5 blocks
                    evidence_text += f"Block {block['block_id']}: {block['content']}\n\n"
                    support_block_ids.append(block['block_id'])
            
            # Create LLM prompt based on enrich mode
            if enrich_mode == "external" and not evidence_text.strip():
                # External mode with no evidence - use standard patterns
                system_message = """You are a gap-filler. Create concise, generic content for missing information using standard API patterns and best practices.

RULES FOR EXTERNAL MODE:
- Use only generic, industry-standard patterns
- No vendor-specific facts or details
- Mark as "(Assumed Standard Practice)" 
- Keep to 1-2 sentences maximum
- Set confidence to "low" when using generic patterns

Output format:
{"text": "content", "confidence": "low", "support_block_ids": [], "reasoning": "explanation"}"""
                
                user_message = f"""Gap Context: {gap_context}
Gap Type: {gap_type}
Section: {section}

Create generic filler content for this gap using standard industry practices. Mark as assumed standard practice."""

            else:
                # Internal mode with evidence
                system_message = """You are a gap-filler. Propose concise content to replace gaps strictly from the provided evidence.

RULES FOR INTERNAL MODE:
- If evidence supports a specific value, output 1-2 sentence patch with high confidence
- If evidence is generic patterns, output template sentence with low confidence
- Only use information from provided evidence blocks
- Reference support_block_ids for evidence used

Output format:
{"text": "content", "confidence": "high|low", "support_block_ids": ["b1","b2"], "reasoning": "explanation"}"""

                user_message = f"""Gap Context: {gap_context}
Gap Type: {gap_type}
Section: {section}

Available Evidence:
{evidence_text}

Create a patch for this gap using only the provided evidence. If evidence is insufficient, return null."""
            
            # Call LLM for gap patch
            llm_response = await call_llm_with_fallback(system_message, user_message)
            
            if not llm_response:
                return None
            
            # Parse LLM response
            try:
                import json
                patch_data = json.loads(llm_response)
                
                return {
                    "location": f"Section: {section}",
                    "text": patch_data.get('text', ''),
                    "support_block_ids": patch_data.get('support_block_ids', support_block_ids),
                    "confidence": patch_data.get('confidence', 'medium'),
                    "gap_pattern": gap.get('pattern', '[MISSING]'),
                    "gap_position": gap.get('position', 0),
                    "reasoning": patch_data.get('reasoning', ''),
                    "enrich_mode": enrich_mode
                }
                
            except json.JSONDecodeError:
                # Fallback parsing
                return {
                    "location": f"Section: {section}",
                    "text": llm_response.strip(),
                    "support_block_ids": support_block_ids,
                    "confidence": "medium",
                    "gap_pattern": gap.get('pattern', '[MISSING]'),
                    "gap_position": gap.get('position', 0),
                    "reasoning": "LLM response parsing fallback",
                    "enrich_mode": enrich_mode
                }
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error creating gap patch - {e}")
            return None
    
    def _apply_patches(self, content: str, patches: list) -> tuple:
        """Apply patches to content, replacing gaps with generated content"""
        try:
            patched_content = content
            applied_patches = []
            
            # Sort patches by position in reverse order to maintain positions
            sorted_patches = sorted(patches, key=lambda x: x.get('gap_position', 0), reverse=True)
            
            for patch in sorted_patches:
                gap_pattern = patch.get('gap_pattern', '[MISSING]')
                replacement_text = patch.get('text', '')
                confidence = patch.get('confidence', 'medium')
                
                if not replacement_text:
                    continue
                
                # Add confidence indicator if low confidence
                if confidence == 'low':
                    replacement_text += " (Assumed Standard Practice)"
                
                # Replace the gap pattern with the patch
                if gap_pattern in patched_content:
                    patched_content = patched_content.replace(gap_pattern, replacement_text, 1)
                    applied_patches.append({
                        "location": patch.get('location', ''),
                        "original": gap_pattern,
                        "replacement": replacement_text,
                        "confidence": confidence,
                        "support_block_ids": patch.get('support_block_ids', []),
                        "reasoning": patch.get('reasoning', '')
                    })
            
            return patched_content, applied_patches
            
        except Exception as e:
            print(f"❌ V2 GAP FILLING: Error applying patches - {e}")
            return content, []

# Global V2 Gap Filling System instance
v2_gap_filling_system = V2GapFillingSystem()

# ========================================
# V2 ENGINE: EVIDENCE TAGGING SYSTEM
# ========================================

class V2EvidenceTaggingSystem:
    """V2 Engine: Evidence tagging system to enforce fidelity by mapping paragraphs to source blocks"""
    
    def __init__(self):
        self.faq_indicators = ['faq', 'frequently asked questions', 'q:', 'question:', 'a:', 'answer:']
        
    async def tag_content_with_evidence(self, articles: list, source_blocks: list, 
                                       prewrite_data: dict, run_id: str) -> dict:
        """Tag paragraphs in articles with evidence block IDs to enforce fidelity"""
        try:
            print(f"🏷️ V2 EVIDENCE TAGGING: Starting evidence tagging process - engine=v2")
            
            evidence_tagging_results = []
            total_paragraphs = 0
            total_tagged_paragraphs = 0
            
            for i, article in enumerate(articles):
                try:
                    article_title = article.get('title', f'Article {i+1}')
                    article_content = article.get('content', '') or article.get('html', '')
                    
                    if not article_content:
                        continue
                    
                    # Step 1: Parse paragraphs from article content
                    paragraphs = self._parse_paragraphs(article_content)
                    
                    # Step 2: Get prewrite facts for this article
                    article_prewrite = self._get_article_prewrite_data(article, prewrite_data)
                    
                    # Step 3: Tag paragraphs with evidence
                    tagged_content, tagging_stats = self._tag_paragraphs_with_evidence(
                        article_content, paragraphs, source_blocks, article_prewrite
                    )
                    
                    # Update article with tagged content
                    article['content'] = tagged_content
                    article['html'] = tagged_content
                    
                    # Track statistics
                    article_total_paragraphs = tagging_stats['total_paragraphs']
                    article_tagged_paragraphs = tagging_stats['tagged_paragraphs']
                    total_paragraphs += article_total_paragraphs
                    total_tagged_paragraphs += article_tagged_paragraphs
                    
                    tagging_rate = (article_tagged_paragraphs / article_total_paragraphs * 100) if article_total_paragraphs > 0 else 0
                    
                    evidence_tagging_result = {
                        "article_index": i,
                        "article_title": article_title,
                        "evidence_tagging_status": "success",
                        "total_paragraphs": article_total_paragraphs,
                        "tagged_paragraphs": article_tagged_paragraphs,
                        "untagged_paragraphs": article_total_paragraphs - article_tagged_paragraphs,
                        "tagging_rate": tagging_rate,
                        "evidence_mapping": tagging_stats['evidence_mapping'],
                        "faq_paragraphs_skipped": tagging_stats['faq_paragraphs_skipped']
                    }
                    
                    evidence_tagging_results.append(evidence_tagging_result)
                    
                    print(f"✅ V2 EVIDENCE TAGGING: Tagged '{article_title[:50]}...' - {article_tagged_paragraphs}/{article_total_paragraphs} paragraphs ({tagging_rate:.1f}%) - engine=v2")
                    
                except Exception as article_error:
                    print(f"❌ V2 EVIDENCE TAGGING: Error processing article {i+1} - {article_error} - engine=v2")
                    evidence_tagging_results.append({
                        "article_index": i,
                        "article_title": article.get('title', f'Article {i+1}'),
                        "evidence_tagging_status": "error",
                        "error": str(article_error),
                        "total_paragraphs": 0,
                        "tagged_paragraphs": 0
                    })
            
            # Calculate overall statistics
            overall_tagging_rate = (total_tagged_paragraphs / total_paragraphs * 100) if total_paragraphs > 0 else 0
            successful_articles = len([r for r in evidence_tagging_results if r.get('evidence_tagging_status') == 'success'])
            
            return {
                "evidence_tagging_id": f"evidence_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "evidence_tagging_status": "success",
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2",
                
                # Evidence tagging metrics
                "articles_processed": len(articles),
                "successful_articles": successful_articles,
                "total_paragraphs": total_paragraphs,
                "tagged_paragraphs": total_tagged_paragraphs,
                "untagged_paragraphs": total_paragraphs - total_tagged_paragraphs,
                "overall_tagging_rate": overall_tagging_rate,
                "target_achieved": overall_tagging_rate >= 95.0,
                
                # Detailed results
                "evidence_tagging_results": evidence_tagging_results,
                "source_blocks_used": len(source_blocks)
            }
            
        except Exception as e:
            print(f"❌ V2 EVIDENCE TAGGING: Error in evidence tagging process - {e} - engine=v2")
            return {
                "evidence_tagging_id": f"evidence_error_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "evidence_tagging_status": "error",
                "error": str(e),
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
    
    def _parse_paragraphs(self, content: str) -> list:
        """Parse paragraphs from HTML/markdown content"""
        try:
            from bs4 import BeautifulSoup
            import re
            
            paragraphs = []
            
            # Try HTML parsing first
            try:
                soup = BeautifulSoup(content, 'html.parser')
                
                # Find all paragraph elements
                p_tags = soup.find_all('p')
                
                for i, p_tag in enumerate(p_tags):
                    paragraph_text = p_tag.get_text().strip()
                    
                    if len(paragraph_text) < 20:  # Skip very short paragraphs
                        continue
                    
                    # Check if this is a FAQ paragraph (skip tagging)
                    is_faq = self._is_faq_paragraph(paragraph_text, str(p_tag))
                    
                    paragraphs.append({
                        "index": i,
                        "text": paragraph_text,
                        "html": str(p_tag),
                        "is_faq": is_faq,
                        "start_pos": content.find(str(p_tag)),
                        "end_pos": content.find(str(p_tag)) + len(str(p_tag))
                    })
                
            except Exception as html_error:
                # Fallback to simple text parsing
                text_paragraphs = content.split('\n\n')
                
                for i, paragraph in enumerate(text_paragraphs):
                    paragraph_text = paragraph.strip()
                    
                    if len(paragraph_text) < 20:
                        continue
                    
                    is_faq = self._is_faq_paragraph(paragraph_text, paragraph)
                    
                    paragraphs.append({
                        "index": i,
                        "text": paragraph_text,
                        "html": paragraph,
                        "is_faq": is_faq,
                        "start_pos": content.find(paragraph),
                        "end_pos": content.find(paragraph) + len(paragraph)
                    })
            
            return paragraphs
            
        except Exception as e:
            print(f"❌ V2 EVIDENCE TAGGING: Error parsing paragraphs - {e}")
            return []
    
    def _is_faq_paragraph(self, paragraph_text: str, paragraph_html: str) -> bool:
        """Check if a paragraph is part of FAQ section (should skip evidence tagging)"""
        text_lower = paragraph_text.lower()
        html_lower = paragraph_html.lower()
        
        # Check for FAQ indicators
        for indicator in self.faq_indicators:
            if indicator in text_lower or indicator in html_lower:
                return True
        
        # Check for FAQ patterns (Q: / A:)
        if re.match(r'^(q:|question:|a:|answer:)', text_lower.strip()):
            return True
        
        # Check if in FAQ section
        if 'faq' in html_lower or 'question' in html_lower:
            return True
        
        return False
    
    def _get_article_prewrite_data(self, article: dict, prewrite_data: dict) -> dict:
        """Get prewrite facts and evidence for a specific article"""
        try:
            article_title = article.get('title', '')
            
            # Try to find matching prewrite data by title
            prewrite_results = prewrite_data.get('prewrite_results', [])
            
            for prewrite_result in prewrite_results:
                if prewrite_result.get('article_title', '').strip().lower() == article_title.strip().lower():
                    return prewrite_result.get('prewrite_data', {})
            
            # Fallback: return first prewrite data if no match
            if prewrite_results:
                return prewrite_results[0].get('prewrite_data', {})
            
            return {}
            
        except Exception as e:
            print(f"❌ V2 EVIDENCE TAGGING: Error getting prewrite data - {e}")
            return {}
    
    def _tag_paragraphs_with_evidence(self, content: str, paragraphs: list, 
                                     source_blocks: list, prewrite_data: dict) -> tuple:
        """Tag paragraphs with evidence block IDs"""
        try:
            tagged_content = content
            tagged_count = 0
            faq_skipped = 0
            evidence_mapping = []
            
            # Sort paragraphs by position in reverse order to maintain positions during replacement
            sorted_paragraphs = sorted(paragraphs, key=lambda x: x.get('start_pos', 0), reverse=True)
            
            for paragraph in sorted_paragraphs:
                try:
                    # Skip FAQ paragraphs
                    if paragraph.get('is_faq', False):
                        faq_skipped += 1
                        continue
                    
                    paragraph_text = paragraph.get('text', '')
                    paragraph_html = paragraph.get('html', '')
                    
                    if not paragraph_text:
                        continue
                    
                    # Find relevant evidence blocks for this paragraph
                    evidence_blocks = self._find_evidence_for_paragraph(
                        paragraph_text, source_blocks, prewrite_data
                    )
                    
                    if evidence_blocks:
                        # Create evidence attribute
                        block_ids = [block['block_id'] for block in evidence_blocks]
                        evidence_attr = f'data-evidence="{json.dumps(block_ids)}"'
                        
                        # Add evidence attribute to paragraph
                        if paragraph_html.startswith('<p'):
                            # HTML paragraph - add attribute to <p> tag
                            if '>' in paragraph_html:
                                tag_end = paragraph_html.find('>')
                                tagged_paragraph = (paragraph_html[:tag_end] + ' ' + 
                                                  evidence_attr + paragraph_html[tag_end:])
                            else:
                                tagged_paragraph = paragraph_html
                        else:
                            # Plain text paragraph - add as HTML comment
                            tagged_paragraph = f'<!-- {evidence_attr} -->\n{paragraph_html}'
                        
                        # Replace in content
                        tagged_content = tagged_content.replace(paragraph_html, tagged_paragraph)
                        tagged_count += 1
                        
                        # Track evidence mapping
                        evidence_mapping.append({
                            "paragraph_index": paragraph.get('index'),
                            "paragraph_preview": paragraph_text[:100],
                            "evidence_blocks": block_ids,
                            "evidence_count": len(block_ids),
                            "confidence_score": self._calculate_evidence_confidence(evidence_blocks)
                        })
                
                except Exception as paragraph_error:
                    print(f"❌ V2 EVIDENCE TAGGING: Error tagging paragraph - {paragraph_error}")
                    continue
            
            return tagged_content, {
                "total_paragraphs": len(paragraphs),
                "tagged_paragraphs": tagged_count,
                "faq_paragraphs_skipped": faq_skipped,
                "evidence_mapping": evidence_mapping
            }
            
        except Exception as e:
            print(f"❌ V2 EVIDENCE TAGGING: Error in paragraph tagging - {e}")
            return content, {
                "total_paragraphs": len(paragraphs),
                "tagged_paragraphs": 0,
                "faq_paragraphs_skipped": 0,
                "evidence_mapping": []
            }
    
    def _find_evidence_for_paragraph(self, paragraph_text: str, source_blocks: list, 
                                    prewrite_data: dict) -> list:
        """Find relevant evidence blocks for a paragraph using prewrite facts and block matching"""
        try:
            import re
            evidence_blocks = []
            
            # Extract keywords from paragraph
            paragraph_keywords = self._extract_paragraph_keywords(paragraph_text)
            
            # Try to match with prewrite facts first
            prewrite_facts = prewrite_data.get('facts', [])
            
            for fact in prewrite_facts:
                fact_text = fact.get('text', '')
                fact_blocks = fact.get('source_blocks', [])
                
                if not fact_text or not fact_blocks:
                    continue
                
                # Calculate relevance between paragraph and fact
                relevance_score = self._calculate_text_relevance(paragraph_text, fact_text)
                
                if relevance_score > 0.3:  # Threshold for relevance
                    # Add source blocks from matching fact
                    for block_id in fact_blocks:
                        evidence_blocks.append({
                            "block_id": block_id,
                            "relevance_score": relevance_score,
                            "source": "prewrite_fact",
                            "fact_text": fact_text[:100]
                        })
            
            # If no prewrite matches, try direct block matching
            if not evidence_blocks:
                for i, block in enumerate(source_blocks[:50]):  # Limit for performance
                    block_content = block.get('content', '') or block.get('text', '')
                    
                    if not block_content:
                        continue
                    
                    # Calculate relevance between paragraph and block
                    relevance_score = self._calculate_text_relevance(paragraph_text, block_content)
                    
                    if relevance_score > 0.2:  # Lower threshold for direct block matching
                        evidence_blocks.append({
                            "block_id": f"b{i}",
                            "relevance_score": relevance_score,
                            "source": "direct_block",
                            "block_preview": block_content[:100]
                        })
            
            # Sort by relevance score and return top matches
            evidence_blocks.sort(key=lambda x: x['relevance_score'], reverse=True)
            return evidence_blocks[:3]  # Return top 3 evidence blocks
            
        except Exception as e:
            print(f"❌ V2 EVIDENCE TAGGING: Error finding evidence - {e}")
            return []
    
    def _extract_paragraph_keywords(self, paragraph_text: str) -> list:
        """Extract keywords from paragraph text for matching"""
        import re
        
        # Remove common stop words
        stop_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with',
            'by', 'from', 'as', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 'has', 'had',
            'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might', 'can',
            'this', 'that', 'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they'
        }
        
        # Extract meaningful words
        words = re.findall(r'\b[a-zA-Z]{3,}\b', paragraph_text.lower())
        keywords = [word for word in words if word not in stop_words]
        
        return keywords[:10]  # Return top 10 keywords
    
    def _calculate_text_relevance(self, text1: str, text2: str) -> float:
        """Calculate relevance score between two pieces of text"""
        try:
            keywords1 = set(self._extract_paragraph_keywords(text1))
            keywords2 = set(self._extract_paragraph_keywords(text2))
            
            if not keywords1 or not keywords2:
                return 0.0
            
            # Calculate Jaccard similarity
            intersection = len(keywords1.intersection(keywords2))
            union = len(keywords1.union(keywords2))
            
            return intersection / union if union > 0 else 0.0
            
        except Exception as e:
            return 0.0
    
    def _calculate_evidence_confidence(self, evidence_blocks: list) -> float:
        """Calculate confidence score for evidence blocks"""
        if not evidence_blocks:
            return 0.0
        
        # Base confidence on number of blocks and relevance scores
        avg_relevance = sum(block.get('relevance_score', 0) for block in evidence_blocks) / len(evidence_blocks)
        block_count_factor = min(len(evidence_blocks) / 3.0, 1.0)  # Max factor at 3 blocks
        
        return avg_relevance * block_count_factor

# Global V2 Evidence Tagging System instance
v2_evidence_tagging_system = V2EvidenceTaggingSystem()

# ========================================
# V2 ENGINE: CODE NORMALIZATION SYSTEM
# ========================================

class V2CodeNormalizationSystem:
    """V2 Engine: Code normalization and beautification system for Prism-ready code blocks"""
    
    def __init__(self):
        self.language_mappings = {
            # Shell/Curl mappings
            'bash': 'language-bash',
            'sh': 'language-bash', 
            'shell': 'language-bash',
            'curl': 'language-bash',
            
            # HTTP mappings
            'http': 'language-http',
            'https': 'language-http',
            'request': 'language-http',
            'response': 'language-http',
            
            # Data format mappings
            'json': 'language-json',
            'yaml': 'language-yaml',
            'yml': 'language-yaml',
            'xml': 'language-xml',
            
            # Query language mappings
            'graphql': 'language-graphql',
            'sql': 'language-sql',
            'mysql': 'language-sql',
            'postgresql': 'language-sql',
            
            # Programming language mappings
            'javascript': 'language-javascript',
            'js': 'language-javascript',
            'typescript': 'language-typescript',
            'ts': 'language-typescript',
            'python': 'language-python',
            'py': 'language-python',
            
            # Other languages
            'css': 'language-css',
            'html': 'language-html',
            'markdown': 'language-markdown',
            'md': 'language-markdown'
        }
        
    async def normalize_code_blocks(self, articles: list, source_blocks: list, 
                                   prewrite_data: dict, run_id: str) -> dict:
        """Normalize and beautify code blocks for Prism rendering"""
        try:
            print(f"🎨 V2 CODE NORMALIZATION: Starting code block normalization - engine=v2")
            
            code_normalization_results = []
            total_code_blocks = 0
            total_normalized_blocks = 0
            
            for i, article in enumerate(articles):
                try:
                    article_title = article.get('title', f'Article {i+1}')
                    article_content = article.get('content', '') or article.get('html', '')
                    
                    if not article_content:
                        continue
                    
                    # Step 0: Sanitize legacy code blocks to minimal markup
                    sanitized_content = self.sanitize_legacy_code_blocks(article_content)
                    
                    # Step 1: Detect and extract code blocks
                    code_blocks = self._extract_code_blocks(sanitized_content)
                    
                    # Step 2: Normalize and beautify each code block
                    normalized_content, normalization_stats = self._normalize_article_code_blocks(
                        sanitized_content, code_blocks, source_blocks, prewrite_data
                    )
                    
                    # Update article with normalized content
                    article['content'] = normalized_content
                    article['html'] = normalized_content
                    
                    # Track statistics
                    article_code_blocks = normalization_stats['total_code_blocks']
                    article_normalized_blocks = normalization_stats['normalized_blocks']
                    total_code_blocks += article_code_blocks
                    total_normalized_blocks += article_normalized_blocks
                    
                    normalization_rate = (article_normalized_blocks / article_code_blocks * 100) if article_code_blocks > 0 else 100
                    
                    code_normalization_result = {
                        "article_index": i,
                        "article_title": article_title,
                        "code_normalization_status": "success",
                        "total_code_blocks": article_code_blocks,
                        "normalized_blocks": article_normalized_blocks,
                        "normalization_rate": normalization_rate,
                        "language_distribution": normalization_stats['language_distribution'],
                        "beautification_applied": normalization_stats['beautification_applied']
                    }
                    
                    code_normalization_results.append(code_normalization_result)
                    
                    print(f"✅ V2 CODE NORMALIZATION: Processed '{article_title[:50]}...' - {article_normalized_blocks}/{article_code_blocks} blocks normalized ({normalization_rate:.1f}%) - engine=v2")
                    
                except Exception as article_error:
                    print(f"❌ V2 CODE NORMALIZATION: Error processing article {i+1} - {article_error} - engine=v2")
                    code_normalization_results.append({
                        "article_index": i,
                        "article_title": article.get('title', f'Article {i+1}'),
                        "code_normalization_status": "error",
                        "error": str(article_error),
                        "total_code_blocks": 0,
                        "normalized_blocks": 0
                    })
            
            # Calculate overall statistics
            overall_normalization_rate = (total_normalized_blocks / total_code_blocks * 100) if total_code_blocks > 0 else 100
            successful_articles = len([r for r in code_normalization_results if r.get('code_normalization_status') == 'success'])
            
            return {
                "code_normalization_id": f"code_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "code_normalization_status": "success",
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2",
                
                # Code normalization metrics
                "articles_processed": len(articles),
                "successful_articles": successful_articles,
                "total_code_blocks": total_code_blocks,
                "normalized_blocks": total_normalized_blocks,
                "overall_normalization_rate": overall_normalization_rate,
                
                # Detailed results
                "code_normalization_results": code_normalization_results,
                "source_blocks_used": len(source_blocks)
            }
            
        except Exception as e:
            print(f"❌ V2 CODE NORMALIZATION: Error in code normalization process - {e} - engine=v2")
            return {
                "code_normalization_id": f"code_error_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "code_normalization_status": "error",
                "error": str(e),
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
    
    def _extract_code_blocks(self, content: str) -> list:
        """Extract code blocks from HTML/markdown content"""
        try:
            from bs4 import BeautifulSoup
            import re
            
            code_blocks = []
            
            # Try HTML parsing first
            try:
                soup = BeautifulSoup(content, 'html.parser')
                
                # Find existing code blocks (pre > code)
                pre_tags = soup.find_all('pre')
                for i, pre_tag in enumerate(pre_tags):
                    code_tag = pre_tag.find('code')
                    if code_tag:
                        code_text = code_tag.get_text()
                        
                        # Extract language from class
                        code_classes = code_tag.get('class', [])
                        detected_lang = self._extract_language_from_classes(code_classes)
                        
                        code_blocks.append({
                            "index": i,
                            "type": "html_pre_code",
                            "content": code_text,
                            "detected_language": detected_lang,
                            "original_html": str(pre_tag),
                            "start_pos": content.find(str(pre_tag)),
                            "end_pos": content.find(str(pre_tag)) + len(str(pre_tag))
                        })
                
                # Find standalone code tags
                standalone_code_tags = soup.find_all('code')
                for i, code_tag in enumerate(standalone_code_tags):
                    # Skip if already found in pre tag
                    if code_tag.find_parent('pre'):
                        continue
                    
                    code_text = code_tag.get_text()
                    if len(code_text) > 50:  # Only process longer code snippets
                        code_classes = code_tag.get('class', [])
                        detected_lang = self._extract_language_from_classes(code_classes)
                        
                        code_blocks.append({
                            "index": len(pre_tags) + i,
                            "type": "html_code",
                            "content": code_text,
                            "detected_language": detected_lang,
                            "original_html": str(code_tag),
                            "start_pos": content.find(str(code_tag)),
                            "end_pos": content.find(str(code_tag)) + len(str(code_tag))
                        })
                
            except Exception as html_error:
                print(f"❌ V2 CODE NORMALIZATION: HTML parsing error - {html_error}")
            
            # Also find markdown code blocks
            markdown_code_blocks = re.findall(r'```(\w*)\n(.*?)\n```', content, re.DOTALL)
            for i, (lang, code_content) in enumerate(markdown_code_blocks):
                code_blocks.append({
                    "index": len(code_blocks),
                    "type": "markdown",
                    "content": code_content.strip(),
                    "detected_language": lang.lower() if lang else None,
                    "original_html": f"```{lang}\n{code_content}\n```",
                    "start_pos": content.find(f"```{lang}\n{code_content}\n```"),
                    "end_pos": content.find(f"```{lang}\n{code_content}\n```") + len(f"```{lang}\n{code_content}\n```")
                })
            
            return code_blocks
            
        except Exception as e:
            print(f"❌ V2 CODE NORMALIZATION: Error extracting code blocks - {e}")
            return []
    
    def _extract_language_from_classes(self, classes: list) -> str:
        """Extract language from CSS classes"""
        for class_name in classes:
            if class_name.startswith('language-'):
                return class_name.replace('language-', '')
            elif class_name.startswith('lang-'):
                return class_name.replace('lang-', '')
        return None
    
    def _normalize_article_code_blocks(self, content: str, code_blocks: list, 
                                      source_blocks: list, prewrite_data: dict) -> tuple:
        """Normalize and beautify code blocks in article content"""
        try:
            normalized_content = content
            normalized_count = 0
            language_distribution = {}
            beautification_applied = []
            
            # Sort code blocks by position in reverse order to maintain positions during replacement
            sorted_code_blocks = sorted(code_blocks, key=lambda x: x.get('start_pos', 0), reverse=True)
            
            for code_block in sorted_code_blocks:
                try:
                    # Step 1: Detect language
                    language = self._detect_code_language(code_block)
                    
                    # Step 2: Beautify code content
                    beautified_code = self._beautify_code(code_block['content'], language)
                    
                    # Step 3: Generate evidence mapping
                    evidence_blocks = self._find_code_evidence(code_block, source_blocks, prewrite_data)
                    
                    # Step 4: Create Prism-ready HTML
                    prism_html = self._create_prism_html(
                        beautified_code, language, evidence_blocks, 
                        filename=None, caption=None
                    )
                    
                    # Step 5: Replace in content
                    original_html = code_block['original_html']
                    if original_html in normalized_content:
                        normalized_content = normalized_content.replace(original_html, prism_html)
                        normalized_count += 1
                        
                        # Track language distribution
                        language_distribution[language] = language_distribution.get(language, 0) + 1
                        beautification_applied.append(language)
                
                except Exception as block_error:
                    print(f"❌ V2 CODE NORMALIZATION: Error normalizing code block - {block_error}")
                    continue
            
            return normalized_content, {
                "total_code_blocks": len(code_blocks),
                "normalized_blocks": normalized_count,
                "language_distribution": language_distribution,
                "beautification_applied": beautification_applied
            }
            
        except Exception as e:
            print(f"❌ V2 CODE NORMALIZATION: Error in article code normalization - {e}")
            return content, {
                "total_code_blocks": len(code_blocks),
                "normalized_blocks": 0,
                "language_distribution": {},
                "beautification_applied": []
            }
    
    def _detect_code_language(self, code_block: dict) -> str:
        """Detect programming language from code block"""
        # Use detected language if available
        detected_lang = code_block.get('detected_language')
        if detected_lang:
            return detected_lang.lower()
        
        # Fallback to content sniffing
        code_content = code_block['content'].strip()
        
        # JSON detection
        if code_content.startswith(('{', '[')):
            try:
                import json
                json.loads(code_content)
                return 'json'
            except:
                pass
        
        # YAML detection
        if ':' in code_content and not code_content.startswith('<'):
            return 'yaml'
        
        # XML detection
        if code_content.startswith('<') and '>' in code_content:
            return 'xml'
        
        # GraphQL detection
        if any(keyword in code_content for keyword in ['query', 'mutation', 'subscription', 'fragment']):
            return 'graphql'
        
        # SQL detection
        if any(keyword in code_content.upper() for keyword in ['SELECT', 'INSERT', 'UPDATE', 'DELETE', 'CREATE']):
            return 'sql'
        
        # Curl detection
        if code_content.startswith('curl'):
            return 'bash'
        
        # Python detection
        if any(keyword in code_content for keyword in ['def ', 'import ', 'from ', 'print(']):
            return 'python'
        
        # JavaScript detection
        if any(keyword in code_content for keyword in ['function', 'const ', 'let ', 'var ', '=>']):
            return 'javascript'
        
        # Default to generic
        return 'text'
    
    def _beautify_code(self, code_content: str, language: str) -> str:
        """Beautify code content based on language"""
        try:
            if language == 'json':
                return self._beautify_json(code_content)
            elif language in ['yaml', 'yml']:
                return self._beautify_yaml(code_content)
            elif language == 'xml':
                return self._beautify_xml(code_content)
            elif language == 'sql':
                return self._beautify_sql(code_content)
            elif language == 'bash' and 'curl' in code_content:
                return self._beautify_curl(code_content)
            else:
                return self._beautify_generic(code_content)
                
        except Exception as e:
            print(f"❌ V2 CODE NORMALIZATION: Error beautifying {language} code - {e}")
            return self._beautify_generic(code_content)
    
    def _beautify_json(self, code_content: str) -> str:
        """Beautify JSON content"""
        try:
            import json
            parsed = json.loads(code_content)
            return json.dumps(parsed, indent=2, ensure_ascii=False)
        except:
            return code_content.strip()
    
    def _beautify_yaml(self, code_content: str) -> str:
        """Beautify YAML content"""
        try:
            import yaml
            parsed = yaml.safe_load(code_content)
            return yaml.dump(parsed, indent=2, default_flow_style=False, allow_unicode=True)
        except:
            return code_content.strip()
    
    def _beautify_xml(self, code_content: str) -> str:
        """Beautify XML content"""
        try:
            import xml.dom.minidom
            dom = xml.dom.minidom.parseString(code_content)
            return dom.toprettyxml(indent="  ")[23:]  # Skip XML declaration
        except:
            return code_content.strip()
    
    def _beautify_sql(self, code_content: str) -> str:
        """Beautify SQL content"""
        try:
            import sqlparse
            return sqlparse.format(code_content, reindent=True, keyword_case='upper', indent_width=2)
        except:
            return code_content.strip()
    
    def _beautify_curl(self, code_content: str) -> str:
        """Beautify curl commands with proper line breaks"""
        try:
            lines = code_content.strip().split('\n')
            beautified_lines = []
            
            for line in lines:
                line = line.strip()
                if line.startswith('curl'):
                    beautified_lines.append(line + ' \\')
                elif line.startswith('-'):
                    beautified_lines.append('  ' + line + ' \\')
                else:
                    beautified_lines.append('  ' + line)
            
            # Remove trailing backslash from last line
            if beautified_lines and beautified_lines[-1].endswith(' \\'):
                beautified_lines[-1] = beautified_lines[-1][:-2]
            
            return '\n'.join(beautified_lines)
        except:
            return code_content.strip()
    
    def _beautify_generic(self, code_content: str) -> str:
        """Generic code beautification"""
        # Normalize whitespace and indentation
        lines = code_content.split('\n')
        normalized_lines = []
        
        for line in lines:
            # Convert tabs to 2 spaces
            line = line.expandtabs(2)
            # Remove trailing whitespace
            line = line.rstrip()
            normalized_lines.append(line)
        
        # Ensure final newline
        result = '\n'.join(normalized_lines).strip() + '\n'
        return result
    
    def _find_code_evidence(self, code_block: dict, source_blocks: list, prewrite_data: dict) -> list:
        """Find evidence blocks for code examples"""
        try:
            code_content = code_block['content']
            evidence_blocks = []
            
            # Simple keyword matching for code evidence
            code_keywords = self._extract_code_keywords(code_content)
            
            for i, block in enumerate(source_blocks[:20]):  # Limit for performance
                block_content = block.get('content', '') or block.get('text', '')
                
                if not block_content:
                    continue
                
                # Check if this block contains code or technical content
                if any(keyword in block_content.lower() for keyword in code_keywords):
                    evidence_blocks.append(f"b{i}")
                
                if len(evidence_blocks) >= 2:  # Limit to 2 evidence blocks per code
                    break
            
            return evidence_blocks
            
        except Exception as e:
            print(f"❌ V2 CODE NORMALIZATION: Error finding code evidence - {e}")
            return []
    
    def _extract_code_keywords(self, code_content: str) -> list:
        """Extract keywords from code content for evidence matching"""
        import re
        
        # Extract meaningful identifiers and function names
        keywords = re.findall(r'\b[a-zA-Z_][a-zA-Z0-9_]{2,}\b', code_content)
        
        # Filter out common words and keep significant terms
        significant_keywords = []
        for keyword in keywords:
            if len(keyword) > 3 and keyword.lower() not in ['function', 'return', 'const', 'true', 'false']:
                significant_keywords.append(keyword.lower())
        
        return list(set(significant_keywords))[:5]  # Return top 5 unique keywords
    
    def _create_prism_html(self, code_content: str, language: str, evidence_blocks: list, 
                          filename: str = None, caption: str = None) -> str:
        """Create minimal Prism-ready HTML markup for code blocks"""
        try:
            import html
            
            # HTML escape the code content
            escaped_code = html.escape(code_content)
            
            # Map language to Prism class
            prism_class = self.language_mappings.get(language, f'language-{language}')
            
            # Create evidence comment if evidence blocks exist
            evidence_comment = ""
            if evidence_blocks:
                evidence_comment = f'<!-- evidence: {json.dumps(evidence_blocks)} -->\n'
            
            # Create language display name for data-lang attribute
            display_lang = language.upper()
            
            # Build data attributes
            data_attrs = []
            data_attrs.append(f'data-lang="{display_lang}"')
            data_attrs.append('data-start="1"')
            if filename:
                data_attrs.append(f'data-filename="{filename}"')
            
            # Build minimal HTML structure - just pre and code
            prism_html = f'''{evidence_comment}<pre class="line-numbers" {' '.join(data_attrs)}>
<code class="{prism_class}">{escaped_code}</code>
</pre>'''
            
            # Add caption as separate paragraph if provided
            if caption:
                prism_html += f'\n<p class="code-caption"><em>{html.escape(caption)}</em></p>'
            
            return prism_html
            
        except Exception as e:
            print(f"❌ V2 CODE NORMALIZATION: Error creating Prism HTML - {e}")
            # Fallback to simple code block
            return f'<pre class="line-numbers" data-start="1"><code class="language-{language}">{html.escape(code_content)}</code></pre>'
    
    def sanitize_legacy_code_blocks(self, content: str) -> str:
        """
        Migration sanitizer to rewrite legacy code blocks to minimal markup.
        Removes figure wrappers, toolbars, and nested pre elements.
        """
        try:
            from bs4 import BeautifulSoup
            import re
            
            print(f"🧹 V2 CODE NORMALIZATION: Sanitizing legacy code blocks")
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Find all figure.code-block elements (legacy format)
            figures = soup.find_all('figure', class_='code-block')
            
            for figure in figures:
                try:
                    # Extract data attributes from figure
                    data_lang = figure.get('data-lang', 'TEXT')
                    data_filename = figure.get('data-filename')
                    
                    # Find the inner pre and code elements
                    pre_element = figure.find('pre')
                    if not pre_element:
                        continue
                        
                    code_element = pre_element.find('code')
                    if not code_element:
                        continue
                    
                    # Extract code content
                    code_content = code_element.get_text()
                    
                    # Extract language class
                    code_classes = code_element.get('class', [])
                    language_class = 'language-text'
                    for cls in code_classes:
                        if cls.startswith('language-'):
                            language_class = cls
                            break
                    
                    # Build data attributes for new pre element
                    data_attrs = []
                    data_attrs.append(f'data-lang="{data_lang}"')
                    data_attrs.append('data-start="1"')
                    if data_filename:
                        data_attrs.append(f'data-filename="{data_filename}"')
                    
                    # Create new minimal markup
                    new_pre = soup.new_tag('pre', **{
                        'class': 'line-numbers',
                        'data-lang': data_lang,
                        'data-start': '1'
                    })
                    if data_filename:
                        new_pre['data-filename'] = data_filename
                    
                    new_code = soup.new_tag('code', **{'class': language_class})
                    new_code.string = code_content
                    new_pre.append(new_code)
                    
                    # Handle caption if exists
                    figcaption = figure.find('figcaption', class_='code-caption')
                    if figcaption:
                        caption_p = soup.new_tag('p', **{'class': 'code-caption'})
                        caption_em = soup.new_tag('em')
                        caption_em.string = figcaption.get_text()
                        caption_p.append(caption_em)
                        figure.insert_after(caption_p)
                    
                    # Replace figure with minimal pre element
                    figure.replace_with(new_pre)
                    
                except Exception as figure_error:
                    print(f"❌ V2 CODE NORMALIZATION: Error sanitizing figure - {figure_error}")
                    continue
            
            # Find any nested pre elements (outer pre wrapping figure)
            nested_pres = soup.find_all('pre')
            for pre in nested_pres:
                # Check if this pre contains a figure or another pre
                inner_figure = pre.find('figure', class_='code-block')
                inner_pre = pre.find('pre')
                
                if inner_figure or inner_pre:
                    # This is a wrapper pre, extract the inner content
                    if inner_figure:
                        # Extract from figure as above
                        continue
                    elif inner_pre and inner_pre != pre:
                        # Replace wrapper with inner pre
                        pre.replace_with(inner_pre)
            
            sanitized_content = str(soup)
            
            # Clean up any extra whitespace and formatting
            sanitized_content = re.sub(r'\n\s*\n\s*\n', '\n\n', sanitized_content)
            
            print(f"✅ V2 CODE NORMALIZATION: Legacy code block sanitization complete")
            return sanitized_content
            
        except Exception as e:
            print(f"❌ V2 CODE NORMALIZATION: Error in legacy sanitization - {e}")
            return content  # Return original content if sanitization fails

# Global V2 Code Normalization System instance
v2_code_normalization_system = V2CodeNormalizationSystem()

# ========================================
# V2 ENGINE: ARTICLE GENERATOR SYSTEM
# ========================================

class V2ArticleGenerator:
    """V2 Engine: Final article generation with strict format and audience-aware styling"""
    
    def __init__(self):
        self.required_structure = [
            "h1_title",
            "intro_paragraph", 
            "mini_toc",
            "main_body",
            "faqs",
            "related_links"
        ]
        
        self.audience_styles = {
            "developer": {
                "tone": "technical and precise",
                "focus": "implementation details, code examples, technical specifications",
                "language": "technical terminology, specific APIs, development concepts"
            },
            "business": {
                "tone": "strategic and outcome-focused",
                "focus": "business value, ROI, strategic implications, competitive advantages",
                "language": "business terminology, metrics, strategic concepts"
            },
            "admin": {
                "tone": "procedural and authoritative",
                "focus": "configuration steps, system management, best practices",
                "language": "administrative terminology, system concepts, operational procedures"
            },
            "end_user": {
                "tone": "friendly and accessible",
                "focus": "practical usage, step-by-step guidance, user benefits",
                "language": "plain language, minimal jargon, user-friendly explanations"
            }
        }
    
    async def generate_final_articles(self, normalized_doc, per_article_outlines: list, analysis: dict, run_id: str) -> dict:
        """V2 Engine: Generate final articles with strict format and audience-aware styling"""
        try:
            print(f"📝 V2 ARTICLE GEN: Generating final articles with strict format - engine=v2")
            
            generated_articles = []
            audience = analysis.get('audience', 'end_user')
            
            for article_outline_data in per_article_outlines:
                article_id = article_outline_data.get('article_id', 'unknown')
                outline = article_outline_data.get('outline', {})
                
                if not outline:
                    continue
                
                print(f"📝 V2 ARTICLE GEN: Generating article '{outline.get('title', 'Untitled')}' for {audience} audience - engine=v2")
                
                # Generate final article with strict format
                article_result = await self._generate_single_article(
                    normalized_doc, 
                    article_id, 
                    outline, 
                    analysis, 
                    audience
                )
                
                if article_result:
                    generated_articles.append({
                        "article_id": article_id,
                        "article_data": article_result
                    })
            
            # Store generated articles
            stored_articles = await self._store_generated_articles(generated_articles, run_id, normalized_doc.doc_id)
            
            print(f"✅ V2 ARTICLE GEN: Generated {len(generated_articles)} final articles with strict format - engine=v2")
            return stored_articles
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error generating final articles - {e} - engine=v2")
            return {"generated_articles": [], "run_id": run_id, "doc_id": normalized_doc.doc_id}
    
    async def _generate_single_article(self, normalized_doc, article_id: str, outline: dict, analysis: dict, audience: str) -> dict:
        """Generate a single final article with strict format"""
        try:
            # Get all blocks referenced in the outline
            article_blocks = self._extract_blocks_from_outline(normalized_doc, outline)
            
            if not article_blocks:
                print(f"⚠️ V2 ARTICLE GEN: No blocks found for article {article_id} - engine=v2")
                return None
            
            # Create comprehensive article input for LLM
            article_input = self._create_article_generation_input(outline, article_blocks, analysis, audience)
            
            # Generate article using LLM
            article_result = await self._perform_llm_article_generation(article_input, audience)
            
            if article_result:
                # Validate and enhance the generated article
                validated_article = await self._validate_and_enhance_article(
                    article_result, outline, article_blocks, audience
                )
                
                # Convert HTML to Markdown
                markdown_content = await self._convert_html_to_markdown(validated_article.get('html', ''))
                validated_article['markdown'] = markdown_content
                
                return validated_article
            else:
                # Fallback to rule-based article generation
                print(f"🔄 V2 ARTICLE GEN: LLM failed, using rule-based fallback for article {article_id} - engine=v2")
                return await self._rule_based_article_generation(outline, article_blocks, audience)
                
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error generating single article {article_id} - {e} - engine=v2")
            return None
    
    def _extract_blocks_from_outline(self, normalized_doc, outline: dict) -> list:
        """Extract all blocks referenced in the article outline"""
        try:
            article_blocks = []
            sections = outline.get('sections', [])
            
            for section in sections:
                subsections = section.get('subsections', [])
                for subsection in subsections:
                    block_ids = subsection.get('block_ids', [])
                    
                    for block_id in block_ids:
                        try:
                            block_index = int(block_id.split('_')[1]) - 1
                            if 0 <= block_index < len(normalized_doc.blocks):
                                article_blocks.append({
                                    "block_id": block_id,
                                    "block": normalized_doc.blocks[block_index],
                                    "section": section.get('heading', ''),
                                    "subsection": subsection.get('heading', '')
                                })
                        except (IndexError, ValueError):
                            print(f"⚠️ V2 ARTICLE GEN: Invalid block_id {block_id} - engine=v2")
                            continue
            
            return article_blocks
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error extracting blocks from outline - {e}")
            return []
    
    def _create_article_generation_input(self, outline: dict, article_blocks: list, analysis: dict, audience: str) -> str:
        """Create comprehensive input for LLM article generation"""
        try:
            input_parts = []
            
            # Article metadata
            input_parts.append(f"ARTICLE_TITLE: {outline.get('title', 'Untitled Article')}")
            input_parts.append(f"TARGET_AUDIENCE: {audience}")
            input_parts.append(f"CONTENT_TYPE: {analysis.get('content_type', 'conceptual')}")
            input_parts.append(f"COMPLEXITY: {analysis.get('complexity', 'intermediate')}")
            
            # Audience styling guidance
            style_guide = self.audience_styles.get(audience, self.audience_styles['end_user'])
            input_parts.append(f"TONE: {style_guide['tone']}")
            input_parts.append(f"FOCUS: {style_guide['focus']}")
            input_parts.append(f"LANGUAGE_STYLE: {style_guide['language']}")
            
            # Article structure outline
            input_parts.append("\nARTICLE_OUTLINE:")
            sections = outline.get('sections', [])
            for i, section in enumerate(sections, 1):
                input_parts.append(f"{i}. {section.get('heading', f'Section {i}')}")
                subsections = section.get('subsections', [])
                for j, subsection in enumerate(subsections, 1):
                    input_parts.append(f"   {i}.{j} {subsection.get('heading', f'Subsection {j}')}")
                    block_ids = subsection.get('block_ids', [])
                    if block_ids:
                        input_parts.append(f"       Blocks: {', '.join(block_ids)}")
            
            # Source blocks with detailed content
            input_parts.append("\nSOURCE_BLOCKS (all must be used):")
            for item in article_blocks:
                block_id = item['block_id']
                block = item['block']
                section = item.get('section', '')
                subsection = item.get('subsection', '')
                
                block_info = f"ID:{block_id} | SECTION:{section} | SUBSECTION:{subsection} | TYPE:{block.block_type}"
                
                if hasattr(block, 'level') and block.level:
                    block_info += f" | LEVEL:{block.level}"
                
                if hasattr(block, 'language') and block.language:
                    block_info += f" | LANG:{block.language}"
                
                # Full content (not truncated for final generation)
                block_info += f"\nCONTENT: {block.content}"
                input_parts.append(block_info)
                input_parts.append("---")
            
            # FAQs from outline
            faqs = outline.get('faq_suggestions', [])
            if faqs:
                input_parts.append("\nFAQ_SUGGESTIONS:")
                for i, faq in enumerate(faqs, 1):
                    input_parts.append(f"Q{i}: {faq.get('q', '')}")
                    input_parts.append(f"A{i}: {faq.get('a', '')}")
            
            # Related links from outline
            related_links = outline.get('related_link_suggestions', [])
            if related_links:
                input_parts.append("\nRELATED_LINKS:")
                for i, link in enumerate(related_links, 1):
                    input_parts.append(f"{i}. {link.get('label', 'Link')} - {link.get('url', '')}")
            
            return "\n".join(input_parts)
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error creating article input - {e}")
            return f"ARTICLE_TITLE: {outline.get('title', 'Error')}\nERROR: Could not create article input"
    
    async def _perform_llm_article_generation(self, article_input: str, audience: str) -> dict:
        """Perform LLM-based article generation using specified format"""
        try:
            system_message = f"""You are a professional technical writer. Generate a full article based on the outline and source blocks.

Create articles with EXACT structure and audience-appropriate styling for {audience} readers.

CRITICAL REQUIREMENTS:
1. Follow the EXACT article structure (NO H1 in content, Intro, Mini-TOC, Main Body, FAQs, Related Links)
2. Cover ALL assigned block_ids with 100% coverage - every source block must be reflected in content
3. Style and tone must match the target audience ({audience})
4. Insert [MISSING] where source information is insufficient
5. Do NOT embed media - only reference media IDs if needed
6. Create working Mini-TOC with clickable anchor links (#section-anchors)
7. Use ordered lists (OL) for procedural/sequential content
8. Consolidate related code blocks instead of fragmenting them

EXACT ARTICLE STRUCTURE:
1. NO H1 TITLE (title handled by frontend - start with intro paragraph)
2. Intro Paragraph (overview, context, what reader will learn)
3. Mini-TOC as simple bullet list: <ul><li>Section Name</li></ul> (links will be added automatically)
4. Main Body (H2/H3 sections WITHOUT id attributes - IDs will be added automatically)
5. FAQs (Q&A format addressing common questions)
6. Related Links (bulleted list of internal and external references)

AUDIENCE STYLING FOR {audience.upper()}:
- Tone: {self.audience_styles.get(audience, {}).get('tone', 'professional')}
- Focus: {self.audience_styles.get(audience, {}).get('focus', 'practical guidance')}
- Language: {self.audience_styles.get(audience, {}).get('language', 'appropriate terminology')}

CONTENT REQUIREMENTS:
- All source blocks must be incorporated into appropriate sections
- Code blocks should be properly formatted with syntax highlighting
- Tables should be HTML tables with proper structure
- Lists should use appropriate HTML list formatting
- Headings should have anchor IDs for Mini-TOC linking"""

            user_message = f"""Generate a complete article using ALL source blocks provided.

{article_input}

REQUIREMENTS:
- Follow the exact article structure above
- Cover **all** assigned block_ids (100% coverage)
- Style matches the detected audience ({audience})
- Insert [MISSING] if info is absent
- Do **not** embed media; only reference media IDs
- Return JSON only: {{"html":"...","summary":"..."}}

MANDATORY HTML STRUCTURE (follow EXACTLY):

EXAMPLE of CORRECT HTML format:
{{
  "html": "<p>This tutorial demonstrates how to build a basic Google Map using its JavaScript API. You will learn how to create an HTML page, add a map with a custom marker, and authenticate the map using an API key.</p>
<ul>
  <li>Using Google Map Javascript API</li>
  <li>Create an HTML Page</li>
  <li>Add a Map with a Custom Marker</li>
  <li>Authenticate the Map</li>
  <li>Result</li>
</ul>
<h2>Using Google Map Javascript API</h2>
<p>Section content here...</p>
<h2>Create an HTML Page</h2>
<h3>Steps to Create the HTML Page</h3>
<ol>
  <li>Use any text editor of your choice and add a basic HTML structure</li>
  <li>Add the following meta tag inside the head element</li>
  <li>Add a title for the HTML page inside the head element</li>
</ol>
<pre class='line-numbers' data-lang='HTML' data-start='1'>
<code class='language-html'>&lt;!DOCTYPE html&gt;
&lt;html&gt;
&lt;head&gt;
  &lt;meta charset='UTF8'&gt;
  &lt;title&gt;Google Maps JavaScript API Tutorial&lt;/title&gt;
&lt;/head&gt;
&lt;body&gt;
  &lt;div id='my_map' style='height:900px; width:100%'&gt;&lt;/div&gt;
&lt;/body&gt;
&lt;/html&gt;</code>
</pre>",
  "summary": "Brief summary of what this article covers and its key takeaways"
}}

CRITICAL REQUIREMENTS:
1. NEVER use <h1> tags - START with introduction paragraph <p>
2. Mini-TOC MUST be simple bullet list - links will be added automatically
3. Use <ol> for procedural steps (create, add, configure, install, etc.)
4. Use <ul> only for non-procedural lists (and Mini-TOC)
5. Consolidate code into single <pre><code> blocks
6. Section headings use plain <h2>, <h3> tags without id attributes"""

            print(f"🤖 V2 ARTICLE GEN: Sending article generation request to LLM - {audience} audience - engine=v2")
            
            # Use existing LLM system
            ai_response = await call_llm_with_fallback(system_message, user_message)
            
            if ai_response:
                # Parse JSON response
                import json
                import re
                
                # Clean response and extract JSON
                cleaned_response = re.sub(r'[-\x1f\x7f-\x9f]', '', ai_response)
                
                # Try to extract JSON from response
                json_match = re.search(r'\{.*\}', cleaned_response, re.DOTALL)
                if json_match:
                    json_str = json_match.group(0)
                    article_data = json.loads(json_str)
                    
                    # Validate required fields
                    if 'html' in article_data and 'summary' in article_data:
                        html_length = len(article_data.get('html', ''))
                        print(f"🎯 V2 ARTICLE GEN: LLM article generation successful - {html_length} chars HTML - engine=v2")
                        return article_data
                    else:
                        print(f"⚠️ V2 ARTICLE GEN: Missing required fields in LLM response - engine=v2")
                        return None
                else:
                    print(f"⚠️ V2 ARTICLE GEN: No JSON found in LLM article response - engine=v2")
                    return None
            else:
                print(f"❌ V2 ARTICLE GEN: No response from LLM for article generation - engine=v2")
                return None
                
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error in LLM article generation - {e} - engine=v2")
            return None
    
    async def _validate_and_enhance_article(self, article_result: dict, outline: dict, article_blocks: list, audience: str) -> dict:
        """Validate and enhance LLM-generated article"""
        try:
            enhanced_article = article_result.copy()
            html_content = enhanced_article.get('html', '')
            
            # Validate structure elements
            required_elements = ['<h1>', '<ul>', '<h2', '<h3']  # Basic structure validation
            missing_elements = []
            
            for element in required_elements:
                if element not in html_content:
                    missing_elements.append(element)
            
            if missing_elements:
                print(f"⚠️ V2 ARTICLE GEN: Missing structural elements: {missing_elements} - engine=v2")
            
            # Ensure Mini-TOC has anchor links
            if '<ul>' in html_content and 'href="#' not in html_content:
                print(f"⚠️ V2 ARTICLE GEN: Mini-TOC missing anchor links - engine=v2")
            
            # Validate block coverage (basic check)
            block_ids = [item['block_id'] for item in article_blocks]
            coverage_check = {
                'total_blocks': len(article_blocks),
                'html_length': len(html_content),
                'has_faqs': 'FAQ' in html_content or 'Q:' in html_content,
                'has_toc': '<ul>' in html_content and '<li>' in html_content,
                'has_sections': html_content.count('<h2') >= 1
            }
            
            # Add validation metadata
            enhanced_article['validation_metadata'] = {
                **coverage_check,
                'audience': audience,
                'block_ids_assigned': block_ids,
                'missing_elements': missing_elements,
                'validation_method': 'llm_enhanced',
                'engine': 'v2'
            }
            
            # Ensure no media embedding
            import re
            img_tags = re.findall(r'<img[^>]*>', html_content)
            if img_tags:
                print(f"⚠️ V2 ARTICLE GEN: Found {len(img_tags)} img tags, removing embedded media - engine=v2")
                # Remove img tags
                html_content = re.sub(r'<img[^>]*>', '[MEDIA_REFERENCE]', html_content)
                enhanced_article['html'] = html_content
            
            return enhanced_article
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error validating article - {e} - engine=v2")
            return article_result
    
    async def _convert_html_to_markdown(self, html_content: str) -> str:
        """Convert HTML content to Markdown format"""
        try:
            import re
            
            markdown_content = html_content
            
            # Convert headings
            markdown_content = re.sub(r'<h1[^>]*>(.*?)</h1>', r'# \1', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'<h2[^>]*id="([^"]*)"[^>]*>(.*?)</h2>', r'## \2 {#\1}', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'<h2[^>]*>(.*?)</h2>', r'## \1', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'<h3[^>]*>(.*?)</h3>', r'### \1', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'<h4[^>]*>(.*?)</h4>', r'#### \1', markdown_content, flags=re.IGNORECASE)
            
            # Convert paragraphs
            markdown_content = re.sub(r'<p[^>]*>(.*?)</p>', r'\1\n', markdown_content, flags=re.IGNORECASE | re.DOTALL)
            
            # Convert lists
            markdown_content = re.sub(r'<ul[^>]*>', '', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'</ul>', '\n', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'<ol[^>]*>', '', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'</ol>', '\n', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'<li[^>]*>(.*?)</li>', r'- \1', markdown_content, flags=re.IGNORECASE | re.DOTALL)
            
            # Convert links
            markdown_content = re.sub(r'<a[^>]*href="([^"]*)"[^>]*>(.*?)</a>', r'[\2](\1)', markdown_content, flags=re.IGNORECASE)
            
            # Convert code blocks
            markdown_content = re.sub(r'<pre[^>]*><code[^>]*>(.*?)</code></pre>', r'```\n\1\n```', markdown_content, flags=re.IGNORECASE | re.DOTALL)
            markdown_content = re.sub(r'<code[^>]*>(.*?)</code>', r'`\1`', markdown_content, flags=re.IGNORECASE)
            
            # Convert emphasis
            markdown_content = re.sub(r'<strong[^>]*>(.*?)</strong>', r'**\1**', markdown_content, flags=re.IGNORECASE)
            markdown_content = re.sub(r'<em[^>]*>(.*?)</em>', r'*\1*', markdown_content, flags=re.IGNORECASE)
            
            # Clean up extra whitespace
            markdown_content = re.sub(r'\n\n\n+', '\n\n', markdown_content)
            markdown_content = markdown_content.strip()
            
            return markdown_content
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error converting HTML to Markdown - {e}")
            return html_content  # Return original HTML if conversion fails
    
    async def _rule_based_article_generation(self, outline: dict, article_blocks: list, audience: str) -> dict:
        """Fallback rule-based article generation"""
        try:
            print(f"🔧 V2 ARTICLE GEN: Creating rule-based article for {audience} audience - engine=v2")
            
            title = outline.get('title', 'Article Title')
            sections = outline.get('sections', [])
            faqs = outline.get('faq_suggestions', [])
            related_links = outline.get('related_link_suggestions', [])
            
            html_parts = []
            
            # 1. H1 Title - REMOVED (frontend will handle title as H1)
            # html_parts.append(f'<h1>{title}</h1>')
            
            # 2. Intro paragraph
            audience_intro = {
                'developer': f'This technical guide provides comprehensive implementation details for {title.lower()}.',
                'business': f'This strategic overview covers the business implications and value of {title.lower()}.',
                'admin': f'This administrative guide provides step-by-step procedures for {title.lower()}.',
                'end_user': f'This user-friendly guide helps you understand and use {title.lower()}.'
            }
            intro = audience_intro.get(audience, f'This comprehensive guide covers {title.lower()}.')
            html_parts.append(f'<p>{intro}</p>')
            
            # 3. Mini-TOC with clickable anchor links
            if sections:
                html_parts.append('<ul>')
                for i, section in enumerate(sections, 1):
                    section_anchor = f"section{i}"  # Match existing ID format
                    section_heading = section.get('heading', f'Section {i}')
                    html_parts.append(f'<li><a href="#{section_anchor}" class="toc-link">{section_heading}</a></li>')
                html_parts.append('</ul>')
            
            # 4. Main Body
            for i, section in enumerate(sections, 1):
                section_anchor = f"section{i}"  # Match TOC anchor format
                section_heading = section.get('heading', f'Section {i}')
                html_parts.append(f'<h2 id="{section_anchor}">{section_heading}</h2>')
                
                # Add subsections
                subsections = section.get('subsections', [])
                for j, subsection in enumerate(subsections, 1):
                    subsection_heading = subsection.get('heading', f'Subsection {j}')
                    html_parts.append(f'<h3>{subsection_heading}</h3>')
                    
                    # Add content from blocks
                    block_ids = subsection.get('block_ids', [])
                    for item in article_blocks:
                        if item['block_id'] in block_ids:
                            block = item['block']
                            if block.block_type == 'paragraph':
                                html_parts.append(f'<p>{block.content}</p>')
                            elif block.block_type == 'code':
                                html_parts.append(f'<pre><code>{block.content}</code></pre>')
                            elif block.block_type == 'list':
                                # Detect if this should be an ordered list
                                list_content = block.content
                                is_procedural = any(word in list_content.lower() for word in [
                                    'step', 'create', 'add', 'use', 'open', 'save', 'click',
                                    'go to', 'select', 'copy', 'replace', 'locate', 'generate',
                                    'first', 'second', 'third', 'next', 'then', 'finally'
                                ])
                                
                                list_tag = 'ol' if is_procedural else 'ul'
                                html_parts.append(f'<{list_tag}>')
                                for line in list_content.split('\n'):
                                    if line.strip():
                                        html_parts.append(f'<li>{line.strip()}</li>')
                                html_parts.append(f'</{list_tag}>')
                            else:
                                html_parts.append(f'<p>{block.content}</p>')
            
            # 5. FAQs
            if faqs:
                html_parts.append('<h2 id="faqs">Frequently Asked Questions</h2>')
                for faq in faqs:
                    question = faq.get('q', 'Question')
                    answer = faq.get('a', 'Answer')
                    html_parts.append(f'<h3>Q: {question}</h3>')
                    html_parts.append(f'<p>A: {answer}</p>')
            
            # 6. Related Links
            if related_links:
                html_parts.append('<h2 id="related-links">Related Links</h2>')
                html_parts.append('<ul>')
                for link in related_links:
                    label = link.get('label', 'Link')
                    url = link.get('url', '#')
                    html_parts.append(f'<li><a href="{url}">{label}</a></li>')
                html_parts.append('</ul>')
            
            html_content = '\n'.join(html_parts)
            
            # Apply code block consolidation and cleanup
            html_content = await self._consolidate_code_blocks(html_content)
            
            # Convert to Markdown
            markdown_content = await self._convert_html_to_markdown(html_content)
            
            article_data = {
                'html': html_content,
                'summary': f'A comprehensive {audience}-focused guide covering {title.lower()} with detailed sections, FAQs, and related resources.',
                'markdown': markdown_content,
                'validation_metadata': {
                    'total_blocks': len(article_blocks),
                    'html_length': len(html_content),
                    'has_faqs': len(faqs) > 0,
                    'has_toc': len(sections) > 0,
                    'has_sections': len(sections) > 0,
                    'analysis_method': 'rule_based_fallback',
                    'engine': 'v2'
                }
            }
            
            print(f"🎯 V2 ARTICLE GEN: Rule-based article complete - {len(html_content)} chars HTML - engine=v2")
            return article_data
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error in rule-based generation - {e} - engine=v2")
            return {
                'html': f'<h2>{outline.get("title", "Error")}</h2><p>Error generating article content.</p>',
                'summary': 'Error occurred during article generation.',
                'markdown': f'# {outline.get("title", "Error")}\n\nError generating article content.'
            }
    
    async def _store_generated_articles(self, generated_articles: list, run_id: str, doc_id: str) -> dict:
        """Store generated articles with processing run"""
        try:
            # Create comprehensive articles record
            articles_record = {
                "articles_id": str(uuid.uuid4()),
                "run_id": run_id,
                "doc_id": doc_id,
                "generated_articles": generated_articles,
                "total_articles": len(generated_articles),
                "created_at": datetime.utcnow().isoformat(),
                "engine": "v2",
                "version": "2.0"
            }
            
            # Store in generated articles collection
            await db.v2_generated_articles.insert_one(articles_record)
            
            print(f"📊 V2 ARTICLE GEN: Generated articles stored with run {run_id} - engine=v2")
            return articles_record
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error storing generated articles - {e} - engine=v2")
            return {"generated_articles": generated_articles, "run_id": run_id, "doc_id": doc_id}
    
    async def get_generated_articles_for_run(self, run_id: str) -> dict:
        """Retrieve stored generated articles for a processing run"""
        try:
            articles_record = await db.v2_generated_articles.find_one({"run_id": run_id})
            if articles_record:
                return articles_record
            else:
                print(f"⚠️ V2 ARTICLE GEN: No generated articles found for run {run_id} - engine=v2")
                return None
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error retrieving generated articles - {e} - engine=v2")
            return None
    
    def _extract_title_from_html(self, html_content: str, fallback_title: str = 'Generated Article') -> str:
        """Extract title from HTML content, with fallback to provided title"""
        import re
        
        try:
            # Try to extract from h1 tag first
            h1_match = re.search(r'<h1[^>]*>(.*?)</h1>', html_content, re.IGNORECASE | re.DOTALL)
            if h1_match:
                title_text = re.sub(r'<[^>]+>', '', h1_match.group(1)).strip()
                if 5 < len(title_text) < 120:
                    return title_text
            
            # Try to extract from h2 tag as fallback
            h2_match = re.search(r'<h2[^>]*>(.*?)</h2>', html_content, re.IGNORECASE | re.DOTALL)
            if h2_match:
                title_text = re.sub(r'<[^>]+>', '', h2_match.group(1)).strip()
                if 5 < len(title_text) < 120:
                    return title_text
            
            # Try to extract from title tag
            title_match = re.search(r'<title[^>]*>(.*?)</title>', html_content, re.IGNORECASE | re.DOTALL)
            if title_match:
                title_text = re.sub(r'<[^>]+>', '', title_match.group(1)).strip()
                if 5 < len(title_text) < 120:
                    return title_text
            
            # If no suitable title found, return fallback
            return fallback_title
            
        except Exception as e:
            print(f"⚠️ V2 ARTICLE GEN: Error extracting title from HTML - {e}")
            return fallback_title
    
    async def _consolidate_code_blocks(self, html_content: str) -> str:
        """Consolidate fragmented code blocks in generated HTML"""
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Find consecutive pre elements and consolidate them
            pre_elements = soup.find_all('pre')
            consolidated_count = 0
            
            i = 0
            while i < len(pre_elements):
                current_pre = pre_elements[i]
                consecutive_pres = [current_pre]
                
                # Look for consecutive pre elements
                j = i + 1
                while j < len(pre_elements):
                    next_pre = pre_elements[j]
                    # Check if elements are consecutive in the DOM
                    if self._are_consecutive_elements(current_pre, next_pre):
                        consecutive_pres.append(next_pre)
                        current_pre = next_pre
                        j += 1
                    else:
                        break
                
                # If we have multiple consecutive pres, consolidate them
                if len(consecutive_pres) > 1:
                    # Combine all code content
                    combined_lines = []
                    for pre in consecutive_pres:
                        code_elem = pre.find('code')
                        if code_elem:
                            combined_lines.append(code_elem.get_text())
                    
                    # Create consolidated pre element with Prism classes
                    new_pre = soup.new_tag('pre', **{
                        'class': 'line-numbers',
                        'data-lang': 'HTML',
                        'data-start': '1'
                    })
                    
                    new_code = soup.new_tag('code', **{'class': 'language-html'})
                    new_code.string = '\n'.join(combined_lines)
                    new_pre.append(new_code)
                    
                    # Replace first pre with consolidated version
                    consecutive_pres[0].replace_with(new_pre)
                    
                    # Remove the other pres
                    for pre in consecutive_pres[1:]:
                        pre.decompose()
                    
                    consolidated_count += len(consecutive_pres) - 1
                
                i = j if j > i else i + 1
            
            if consolidated_count > 0:
                print(f"🔧 V2 ARTICLE GEN: Consolidated {consolidated_count} code blocks")
            
            return str(soup)
            
        except Exception as e:
            print(f"❌ V2 ARTICLE GEN: Error consolidating code blocks - {e}")
            return html_content
    
    def _are_consecutive_elements(self, elem1, elem2) -> bool:
        """Check if two elements are consecutive in the DOM"""
        try:
            # Simple check: if they have the same parent and elem2 follows elem1
            if elem1.parent == elem2.parent:
                siblings = list(elem1.parent.children)
                try:
                    idx1 = siblings.index(elem1)
                    idx2 = siblings.index(elem2)
                    # Allow some gap for whitespace/text nodes
                    return 0 < (idx2 - idx1) <= 3
                except ValueError:
                    return False
            return False
        except:
            return False

# Global V2 Article Generator instance
v2_article_generator = V2ArticleGenerator()

# ========================================
# V2 ENGINE: VALIDATION SYSTEM
# ========================================

class V2ValidationSystem:
    """V2 Engine: Comprehensive validation system for fidelity, coverage, placeholders, and style"""
    
    def __init__(self):
        self.required_sections = [
            "h1_title",
            "intro_paragraph",
            "mini_toc", 
            "main_body",
            "faqs",
            "related_links"
        ]
        
        self.quality_thresholds = {
            "coverage_percent": 100,
            "fidelity_score": 0.9,
            "redundancy_score": 0.3,  # Lower is better
            "max_placeholders": 0
        }
        
        # TICKET 1 FIX: Style validation rules including H1 prohibition
        self.style_validation_rules = {
            "no_h1_in_body": True,
            "require_mini_toc": True,
            "require_structured_headings": True
        }
    
    # KE-PR2: Use extracted linking modules instead of inline methods
    def stable_slug(self, text: str, max_len: int = 60) -> str:
        """TICKET 2: Generate deterministic, URL-safe slugs - delegated to engine.linking.anchors"""
        return stable_slug(text, max_len)
    
    def assign_heading_ids(self, html: str) -> str:
        """TICKET 2: Assign deterministic IDs to headings - delegated to engine.linking.anchors"""
        return assign_heading_ids(html)
    
    def validate_heading_ladder(self, html: str) -> bool:
        """TICKET 2: Validate proper heading hierarchy - delegated to engine.linking.anchors"""
        return validate_heading_ladder(html)
    
    def build_minitoc(self, html: str) -> str:
        """TICKET 2: Build Mini-TOC with clickable links - delegated to engine.linking.toc"""
        return build_minitoc(html)
    
    def anchors_resolve(self, html: str) -> bool:
        """TICKET 2: Validate TOC links resolve - delegated to engine.linking.toc"""
        return anchors_resolve(html)
    
    def extract_headings_registry(self, html: str) -> list:
        """TICKET 3: Extract headings for bookmark registry - delegated to engine.linking.bookmarks"""
        return extract_headings_registry(html)
    
    # KE-PR2: Use extracted linking modules 
    def generate_doc_uid(self) -> str:
        """TICKET 3: Generate document UID - delegated to engine.linking.bookmarks"""
        return generate_doc_uid()
    
    def generate_doc_slug(self, title: str) -> str:
        """TICKET 3: Generate document slug - delegated to engine.linking.bookmarks"""
        return generate_doc_slug(title)
    
    def build_href(self, target_doc: dict, anchor_id: str, route_map: dict) -> str:
        """TICKET 3: Build environment-aware href - delegated to engine.linking.links"""
        return build_href(target_doc, anchor_id, route_map)
    
    def get_default_route_map(self, environment: str = "content_library") -> dict:
        """TICKET 3: Get default route map - delegated to engine.linking.links"""
        return get_default_route_map(environment)
    
    async def backfill_bookmark_registry(self, limit: int = None) -> dict:
        """TICKET 3: Backfill existing v2 articles with bookmark registry data"""
        try:
            from motor.motor_asyncio import AsyncIOMotorClient
            import os
            
            mongo_url = os.environ.get('MONGO_URL')
            client = AsyncIOMotorClient(mongo_url)
            db = client.promptsupport
            
            print(f"🔄 TICKET 3: Starting bookmark registry backfill for existing v2 articles")
            
            # KE-PR9: Use repository pattern to find articles needing backfill
            try:
                if mongo_repo_available:
                    content_repo = RepositoryFactory.get_content_library()
                    # Get all V2 articles and filter those needing backfill
                    all_articles = await content_repo.find_by_engine("v2", limit=limit)
                    
                    # Filter articles that need backfilling
                    articles = []
                    for article in all_articles:
                        needs_backfill = (
                            not article.get("doc_uid") or 
                            not article.get("headings") or 
                            (isinstance(article.get("headings"), list) and len(article.get("headings")) == 0)
                        )
                        if needs_backfill:
                            articles.append(article)
                else:
                    # Fallback to direct database query
                    query = {
                        "$or": [
                            {"metadata.engine": "v2", "doc_uid": {"$exists": False}},
                            {"metadata.engine": "v2", "headings": {"$exists": False}},
                            {"metadata.engine": "v2", "headings": []}  # Empty headings array
                        ]
                    }
                    
                    if limit:
                        articles_cursor = db.content_library.find(query).limit(limit)
                    else:
                        articles_cursor = db.content_library.find(query)
                    
                    articles = await articles_cursor.to_list(None)
                    
            except Exception as repo_error:
                print(f"⚠️ KE-PR9: Backfill query fallback to direct DB: {repo_error}")
                # Final fallback to direct database query
                query = {
                    "$or": [
                        {"metadata.engine": "v2", "doc_uid": {"$exists": False}},
                        {"metadata.engine": "v2", "headings": {"$exists": False}},
                        {"metadata.engine": "v2", "headings": []}  # Empty headings array
                    ]
                }
                
                if limit:
                    articles_cursor = db.content_library.find(query).limit(limit)
                else:
                    articles_cursor = db.content_library.find(query)
                
                articles = await articles_cursor.to_list(None)
                
            total_articles = len(articles)
            
            if total_articles == 0:
                print("✅ TICKET 3: No articles need backfilling")
                return {"articles_processed": 0, "success": True}
            
            processed_count = 0
            error_count = 0
            
            for article in articles:
                try:
                    article_id = article.get('_id')
                    title = article.get('title', 'Untitled')
                    content = article.get('content', '') or article.get('html', '')
                    
                    if not content:
                        print(f"⚠️ TICKET 3: Skipping article '{title}' - no content found")
                        continue
                    
                    # Generate doc_uid and doc_slug if missing
                    doc_uid = article.get('doc_uid')
                    if not doc_uid:
                        doc_uid = self.generate_doc_uid()
                    
                    doc_slug = article.get('doc_slug')  
                    if not doc_slug:
                        doc_slug = self.generate_doc_slug(title)
                    
                    # Apply Ticket 2 stable anchors if needed (ensure IDs exist)
                    processed_content = self.assign_heading_ids(content)
                    
                    # Extract headings registry
                    headings = self.extract_headings_registry(processed_content)
                    
                    # Update article in database
                    update_data = {
                        "doc_uid": doc_uid,
                        "doc_slug": doc_slug,
                        "headings": headings,
                        "content": processed_content,  # Updated content with IDs
                        "html": processed_content     # Sync html field
                    }
                    
                    # Initialize xrefs and related_links if missing
                    if "xrefs" not in article:
                        update_data["xrefs"] = []
                    if "related_links" not in article:
                        update_data["related_links"] = []
                    
                    await db.content_library.update_one(
                        {"_id": article_id},
                        {"$set": update_data}
                    )
                    
                    processed_count += 1
                    print(f"📖 TICKET 3: Backfilled article '{title[:50]}...' - doc_uid: {doc_uid}, {len(headings)} headings")
                    
                except Exception as article_error:
                    error_count += 1
                    print(f"❌ TICKET 3: Error backfilling article '{article.get('title', 'Unknown')}' - {article_error}")
            
            success_rate = (processed_count / total_articles * 100) if total_articles > 0 else 100
            
            print(f"✅ TICKET 3: Backfill complete - {processed_count}/{total_articles} articles processed ({success_rate:.1f}% success)")
            
            return {
                "articles_found": total_articles,
                "articles_processed": processed_count,
                "articles_failed": error_count,
                "success_rate": success_rate,
                "success": error_count == 0
            }
            
        except Exception as e:
            print(f"❌ TICKET 3: Error in backfill process - {e}")
            return {
                "articles_found": 0,
                "articles_processed": 0,
                "articles_failed": 0,
                "success_rate": 0,
                "success": False,
                "error": str(e)
            }
    
    async def validate_cross_document_links(self, doc_uid: str, xrefs: list, related_links: list) -> dict:
        """TICKET 3: Validate that cross-document links resolve properly"""
        try:
            from motor.motor_asyncio import AsyncIOMotorClient
            import os
            
            # Get database connection
            mongo_url = os.environ.get('MONGO_URL')
            client = AsyncIOMotorClient(mongo_url)
            db = client.promptsupport
            
            broken_links = []
            total_links = len(xrefs) + len(related_links)
            
            print(f"🔍 TICKET 3: Validating {total_links} cross-document links for doc {doc_uid}")
            
            # Validate xrefs
            for xref in xrefs:
                target_doc_uid = xref.get("doc_uid")
                anchor_id = xref.get("anchor_id")
                label = xref.get("label", "Unknown")
                
                # KE-PR9: Find target document using repository pattern
                if mongo_repo_available:
                    content_repo = RepositoryFactory.get_content_library()
                    target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
                else:
                    # KE-PR9.3: Fallback to repository pattern
                    content_repo = RepositoryFactory.get_content_library()
                    target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
                
                if not target_doc:
                    broken_links.append({
                        "type": "xref",
                        "target_doc_uid": target_doc_uid,
                        "anchor_id": anchor_id,
                        "label": label,
                        "reason": "target_document_not_found"
                    })
                    continue
                
                # Check if anchor exists in target document headings
                target_headings = target_doc.get("headings", [])
                anchor_exists = any(h.get("id") == anchor_id for h in target_headings)
                
                if not anchor_exists:
                    broken_links.append({
                        "type": "xref", 
                        "target_doc_uid": target_doc_uid,
                        "anchor_id": anchor_id,
                        "label": label,
                        "reason": "anchor_not_found_in_target",
                        "available_anchors": [h.get("id") for h in target_headings[:5]]  # First 5 for debugging
                    })
            
            # Validate related_links (similar process)
            for related in related_links:
                target_doc_uid = related.get("doc_uid")
                anchor_id = related.get("anchor_id", "")
                
                # KE-PR9: Find target document using repository pattern
                if mongo_repo_available:
                    content_repo = RepositoryFactory.get_content_library()
                    target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
                else:
                    # KE-PR9.3: Fallback to repository pattern
                    content_repo = RepositoryFactory.get_content_library()
                    target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
                
                if not target_doc:
                    broken_links.append({
                        "type": "related_link",
                        "target_doc_uid": target_doc_uid,
                        "anchor_id": anchor_id,
                        "reason": "target_document_not_found"
                    })
                    continue
                
                if anchor_id:  # Only validate anchor if specified
                    target_headings = target_doc.get("headings", [])
                    anchor_exists = any(h.get("id") == anchor_id for h in target_headings)
                    
                    if not anchor_exists:
                        broken_links.append({
                            "type": "related_link",
                            "target_doc_uid": target_doc_uid, 
                            "anchor_id": anchor_id,
                            "reason": "anchor_not_found_in_target"
                        })
            
            resolution_rate = ((total_links - len(broken_links)) / total_links * 100) if total_links > 0 else 100
            
            print(f"🔍 TICKET 3: Link validation complete - {resolution_rate:.1f}% resolved ({len(broken_links)} broken)")
            
            return {
                "total_links": total_links,
                "broken_links": broken_links,
                "resolution_rate": resolution_rate,
                "links_resolve": len(broken_links) == 0
            }
            
        except Exception as e:
            print(f"❌ TICKET 3: Error validating cross-document links - {e}")
            return {
                "total_links": 0,
                "broken_links": [],
                "resolution_rate": 0,
                "links_resolve": False,
                "error": str(e)
            }
    
    def _apply_bookmark_registry(self, content: str, article_title: str) -> dict:
        """TICKET 3: Apply bookmark registry extraction for universal links"""
        try:
            print(f"📖 TICKET 3: Starting bookmark registry for '{article_title[:50]}...'")
            
            # Extract headings from content
            headings = self.extract_headings_registry(content)
            
            # Generate document identifiers
            doc_uid = self.generate_doc_uid()
            doc_slug = self.generate_doc_slug(article_title)
            
            print(f"📖 TICKET 3: Bookmark registry complete - {len(headings)} headings, doc_uid: {doc_uid}")
            
            return {
                'headings': headings,
                'doc_uid': doc_uid,
                'doc_slug': doc_slug,
                'bookmark_count': len(headings),
                'changes_applied': [f"Extracted {len(headings)} bookmarks", f"Generated doc_uid: {doc_uid}", f"Generated doc_slug: {doc_slug}"]
            }
            
        except Exception as e:
            print(f"❌ TICKET 3: Error in bookmark registry - {e}")
            return {
                'headings': [],
                'doc_uid': None,
                'doc_slug': None,
                'bookmark_count': 0,
                'changes_applied': [f"Bookmark registry error: {str(e)}"],
                'error': str(e)
            }

    
    async def validate_generated_articles(self, normalized_doc, generated_articles_result: dict, analysis: dict, run_id: str) -> dict:
        """V2 Engine: Comprehensive validation of generated articles"""
        try:
            print(f"🔍 V2 VALIDATION: Starting comprehensive validation - run {run_id} - engine=v2")
            
            generated_articles = generated_articles_result.get('generated_articles', [])
            if not generated_articles:
                print(f"⚠️ V2 VALIDATION: No articles to validate - run {run_id} - engine=v2")
                return self._create_validation_result("no_articles", run_id, {})
            
            # Step 1: Fidelity and Coverage Validation
            fidelity_coverage_result = await self._validate_fidelity_and_coverage(
                normalized_doc, generated_articles, run_id
            )
            
            # Step 2: Placeholder Detection
            placeholder_result = await self._detect_placeholders(generated_articles, run_id)
            
            # Step 3: Style Guard Validation
            style_result = await self._validate_style_guard(generated_articles, run_id)
            
            # Step 4: Evidence Tagging Validation
            evidence_result = await self._validate_evidence_tagging(generated_articles, normalized_doc, run_id)
            
            # Step 5: Metrics Calculation
            metrics_result = await self._calculate_validation_metrics(
                normalized_doc, generated_articles, analysis, run_id
            )
            
            # Step 6: Overall Validation Decision
            validation_result = self._consolidate_validation_results(
                fidelity_coverage_result,
                placeholder_result,
                style_result,
                evidence_result,
                metrics_result,
                run_id
            )
            
            # Log validation outcome
            status = validation_result.get('validation_status', 'unknown')
            if status == 'passed':
                print(f"✅ V2 VALIDATION: All validation checks passed - run {run_id} - engine=v2")
            else:
                print(f"⚠️ V2 VALIDATION: Validation failed with status '{status}' - run {run_id} - engine=v2")
            
            return validation_result
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error during validation - {e} - run {run_id} - engine=v2")
            return self._create_validation_result("error", run_id, {"error": str(e)})
    
    async def _validate_fidelity_and_coverage(self, normalized_doc, generated_articles: list, run_id: str) -> dict:
        """V2 Engine: Validate fidelity and coverage using LLM"""
        try:
            print(f"🔍 V2 VALIDATION: Fidelity & Coverage validation - run {run_id} - engine=v2")
            
            # Prepare source blocks for validation
            source_blocks = []
            for i, block in enumerate(normalized_doc.blocks):
                source_blocks.append({
                    "block_id": f"block_{i+1}",
                    "block_type": block.block_type,
                    "content": block.content[:500],  # Truncate for LLM input
                    "full_length": len(block.content)
                })
            
            # Prepare generated articles with mapped block_ids
            article_summaries = []
            for generated_article in generated_articles:
                article_data = generated_article.get('article_data', {})
                validation_metadata = article_data.get('validation_metadata', {})
                
                article_summaries.append({
                    "article_id": generated_article.get('article_id', 'unknown'),
                    "title": article_data.get('title', 'Unknown'),
                    "html_content": article_data.get('html', '')[:1000],  # Truncate for LLM
                    "mapped_block_ids": validation_metadata.get('block_ids_assigned', []),
                    "content_length": len(article_data.get('html', ''))
                })
            
            # Create LLM prompt for fidelity and coverage validation
            system_message = """You are a validation agent. Check fidelity and coverage vs. source.

Your task is to:
1. Compute fidelity_score (0-1): How well do the generated articles stick to the source content without hallucinating?
2. Compute coverage_percent (0-100): What percentage of source blocks are covered in the generated articles?
3. Identify hallucinated_claims: Content in articles that doesn't appear in the source blocks
4. Identify uncovered_blocks: Source blocks that weren't used in any generated article

Fidelity scoring:
- 1.0: Perfect fidelity, no hallucinated content
- 0.9: Minor additions that don't contradict source
- 0.7: Some unsupported claims but generally accurate
- 0.5: Significant hallucinations or contradictions
- 0.0: Mostly hallucinated content

Coverage scoring:
- 100%: All source blocks are represented in generated articles
- 80%: Most important blocks covered, some minor omissions
- 60%: Key blocks covered but significant gaps
- 40%: Major omissions in coverage
- 0%: No source blocks covered

Return ONLY JSON in the exact format specified."""

            user_message = f"""Validate these generated articles against their source blocks.

SOURCE BLOCKS:
{json.dumps(source_blocks, indent=2)}

GENERATED ARTICLES:
{json.dumps(article_summaries, indent=2)}

Compute fidelity_score (0-1) and coverage_percent (0-100).
List hallucinated_claims and uncovered_blocks.

Return ONLY JSON in this exact format:
{{
  "fidelity_score": 0.95,
  "coverage_percent": 100,
  "hallucinated_claims": [],
  "uncovered_blocks": []
}}"""

            # Call LLM for validation
            print(f"🤖 V2 VALIDATION: Sending fidelity & coverage request to LLM - run {run_id} - engine=v2")
            ai_response = await call_llm_with_fallback(system_message, user_message)
            
            if ai_response:
                # Parse JSON response
                import re
                json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
                if json_match:
                    validation_data = json.loads(json_match.group(0))
                    
                    # Validate required fields
                    required_fields = ['fidelity_score', 'coverage_percent', 'hallucinated_claims', 'uncovered_blocks']
                    if all(field in validation_data for field in required_fields):
                        print(f"✅ V2 VALIDATION: Fidelity={validation_data['fidelity_score']}, Coverage={validation_data['coverage_percent']}% - run {run_id} - engine=v2")
                        return validation_data
                    else:
                        print(f"⚠️ V2 VALIDATION: Missing fields in LLM response - run {run_id} - engine=v2")
                        return self._create_fallback_fidelity_coverage(normalized_doc, generated_articles)
                else:
                    print(f"⚠️ V2 VALIDATION: No JSON found in LLM response - run {run_id} - engine=v2")
                    return self._create_fallback_fidelity_coverage(normalized_doc, generated_articles)
            else:
                print(f"❌ V2 VALIDATION: No LLM response for fidelity validation - run {run_id} - engine=v2")
                return self._create_fallback_fidelity_coverage(normalized_doc, generated_articles)
                
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error in fidelity & coverage validation - {e} - run {run_id} - engine=v2")
            return self._create_fallback_fidelity_coverage(normalized_doc, generated_articles)
    
    def _create_fallback_fidelity_coverage(self, normalized_doc, generated_articles: list) -> dict:
        """Create fallback fidelity and coverage scores"""
        try:
            # Basic coverage calculation
            total_blocks = len(normalized_doc.blocks)
            covered_blocks = set()
            
            for generated_article in generated_articles:
                article_data = generated_article.get('article_data', {})
                validation_metadata = article_data.get('validation_metadata', {})
                mapped_blocks = validation_metadata.get('block_ids_assigned', [])
                
                for block_id in mapped_blocks:
                    covered_blocks.add(block_id)
            
            coverage_percent = (len(covered_blocks) / total_blocks * 100) if total_blocks > 0 else 0
            
            # Conservative fidelity score for fallback
            fidelity_score = 0.85 if len(generated_articles) > 0 else 0.0
            
            # Identify uncovered blocks
            all_block_ids = set(f"block_{i+1}" for i in range(total_blocks))
            uncovered_blocks = list(all_block_ids - covered_blocks)
            
            return {
                "fidelity_score": fidelity_score,
                "coverage_percent": coverage_percent,
                "hallucinated_claims": [],  # Can't detect without LLM
                "uncovered_blocks": uncovered_blocks
            }
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error in fallback fidelity calculation - {e}")
            return {
                "fidelity_score": 0.0,
                "coverage_percent": 0.0,
                "hallucinated_claims": [],
                "uncovered_blocks": []
            }
    
    async def _detect_placeholders(self, generated_articles: list, run_id: str) -> dict:
        """V2 Engine: Detect placeholders and incomplete content using LLM"""
        try:
            print(f"🔍 V2 VALIDATION: Placeholder detection - run {run_id} - engine=v2")
            
            # Prepare HTML content for placeholder detection
            articles_html = []
            for generated_article in generated_articles:
                article_data = generated_article.get('article_data', {})
                html_content = article_data.get('html', '')
                
                articles_html.append({
                    "article_id": generated_article.get('article_id', 'unknown'),
                    "title": article_data.get('title', 'Unknown'),
                    "html_content": html_content
                })
            
            # Create LLM prompt for placeholder detection
            system_message = """You are a style checker. Detect placeholders and incomplete content.

Your task is to find any incomplete or placeholder content including:
- [MISSING] markers
- TODO items
- Lorem ipsum text
- Placeholder text like "Insert content here", "Add details", etc.
- Empty sections or incomplete sentences
- Generic placeholder content

For each placeholder found, provide:
- article_id: Which article contains it
- location: Which section or area (e.g., "Section: Setup", "FAQ section", "Introduction")
- text: The exact placeholder text found

Return ONLY JSON in the exact format specified."""

            user_message = f"""Check these generated articles for placeholders and incomplete content.

GENERATED ARTICLES HTML:
{json.dumps(articles_html, indent=2)}

Report any [MISSING], TODO, lorem ipsum, or other placeholder content with article_id and location.

Return ONLY JSON in this exact format:
{{
  "placeholders": [
    {{
      "article_id": "a1",
      "location": "Section: Setup",
      "text": "[MISSING]"
    }}
  ]
}}"""

            # Call LLM for placeholder detection
            print(f"🤖 V2 VALIDATION: Sending placeholder detection request to LLM - run {run_id} - engine=v2")
            ai_response = await call_llm_with_fallback(system_message, user_message)
            
            if ai_response:
                # Parse JSON response
                import re
                json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
                if json_match:
                    placeholder_data = json.loads(json_match.group(0))
                    
                    if 'placeholders' in placeholder_data:
                        placeholder_count = len(placeholder_data['placeholders'])
                        print(f"🔍 V2 VALIDATION: Found {placeholder_count} placeholders - run {run_id} - engine=v2")
                        return placeholder_data
                    else:
                        print(f"⚠️ V2 VALIDATION: Missing placeholders field in LLM response - run {run_id} - engine=v2")
                        return self._create_fallback_placeholder_detection(generated_articles)
                else:
                    print(f"⚠️ V2 VALIDATION: No JSON found in placeholder response - run {run_id} - engine=v2")
                    return self._create_fallback_placeholder_detection(generated_articles)
            else:
                print(f"❌ V2 VALIDATION: No LLM response for placeholder detection - run {run_id} - engine=v2")
                return self._create_fallback_placeholder_detection(generated_articles)
                
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error in placeholder detection - {e} - run {run_id} - engine=v2")
            return self._create_fallback_placeholder_detection(generated_articles)
    
    def _create_fallback_placeholder_detection(self, generated_articles: list) -> dict:
        """Create fallback placeholder detection using regex"""
        try:
            import re
            
            placeholder_patterns = [
                r'\[MISSING\]',
                r'TODO',
                r'lorem ipsum',
                r'Lorem Ipsum',
                r'placeholder',
                r'INSERT_\w+',
                r'ADD_\w+',
                r'\[.*PLACEHOLDER.*\]'
            ]
            
            placeholders = []
            
            for generated_article in generated_articles:
                article_id = generated_article.get('article_id', 'unknown')
                article_data = generated_article.get('article_data', {})
                html_content = article_data.get('html', '')
                
                for pattern in placeholder_patterns:
                    matches = re.finditer(pattern, html_content, re.IGNORECASE)
                    for match in matches:
                        placeholders.append({
                            "article_id": article_id,
                            "location": "Content scan",
                            "text": match.group(0)
                        })
            
            print(f"🔍 V2 VALIDATION: Fallback placeholder detection found {len(placeholders)} placeholders")
            return {"placeholders": placeholders}
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error in fallback placeholder detection - {e}")
            return {"placeholders": []}
    
    async def _validate_style_guard(self, generated_articles: list, run_id: str) -> dict:
        """V2 Engine: Validate article structure and style compliance"""
        try:
            print(f"🔍 V2 VALIDATION: Style guard validation - run {run_id} - engine=v2")
            
            style_results = []
            
            for generated_article in generated_articles:
                article_id = generated_article.get('article_id', 'unknown')
                article_data = generated_article.get('article_data', {})
                html_content = article_data.get('html', '')
                
                # TICKET 1 & 2 FIX: Check for structural elements, H1 prohibition, and anchor validation
                structure_check = {
                    "no_h1_in_body": self.validate_no_h1_in_body(html_content),  # HARD FAIL if H1 found
                    "intro_paragraph": bool(re.search(r'<p[^>]*>.*?</p>', html_content, re.IGNORECASE | re.DOTALL)),
                    "mini_toc": bool(re.search(r'<ul[^>]*class="mini-toc"', html_content, re.IGNORECASE)),
                    "main_body": bool(re.search(r'<h2[^>]*>.*?</h2>', html_content, re.IGNORECASE | re.DOTALL)),
                    "faqs": bool(re.search(r'FAQ|Q:|Question', html_content, re.IGNORECASE)),
                    "related_links": bool(re.search(r'<ul[^>]*>.*?<li[^>]*>.*?<a[^>]*href=', html_content, re.IGNORECASE | re.DOTALL))
                }
                
                # TICKET 2 FIX: Add anchor validation
                anchor_validation = {
                    "heading_ladder": self.validate_heading_ladder_structure(html_content),
                    "anchors_resolve": self.validate_anchor_resolution(html_content)
                }
                
                # Combine structural and anchor validation
                structure_check.update(anchor_validation)
                
                # TICKET 1 FIX: Hard fail if H1 found in body content
                has_h1_violation = not structure_check["no_h1_in_body"]
                if has_h1_violation:
                    print(f"❌ V2 VALIDATION: HARD FAIL - H1 tag found in article {article_id} body content - PROHIBITED")
                    # This will cause compliance score to be 0 and fail validation
                
                # Calculate compliance score
                total_elements = len(structure_check)
                passed_elements = sum(structure_check.values())
                compliance_score = passed_elements / total_elements if total_elements > 0 else 0
                
                # Identify missing elements
                missing_elements = [element for element, present in structure_check.items() if not present]
                
                style_results.append({
                    "article_id": article_id,
                    "structure_check": structure_check,
                    "compliance_score": compliance_score,
                    "missing_elements": missing_elements,
                    "content_length": len(html_content)
                })
            
            # Overall style compliance
            overall_compliance = sum(result['compliance_score'] for result in style_results) / len(style_results) if style_results else 0
            
            print(f"🎯 V2 VALIDATION: Style guard compliance: {overall_compliance:.2f} - run {run_id} - engine=v2")
            
            return {
                "overall_compliance": overall_compliance,
                "article_results": style_results,
                "validation_method": "programmatic_style_guard"
            }
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error in style guard validation - {e} - run {run_id} - engine=v2")
            return {
                "overall_compliance": 0.0,
                "article_results": [],
                "validation_method": "error"
            }
    
    def validate_no_h1_in_body(self, html: str) -> bool:
        """TICKET 1 FIX: Hard gate validation - no H1 tags allowed in body content"""
        import re
        h1_matches = re.findall(r'<h1\b[^>]*>', html, re.IGNORECASE)
        return len(h1_matches) == 0
    
    def validate_heading_ladder_structure(self, html: str) -> bool:
        """TICKET 2: Validate proper heading hierarchy (H2->H3->H4)"""
        from bs4 import BeautifulSoup
        
        soup = BeautifulSoup(html, 'html.parser')
        levels = []
        
        for tag in soup.find_all(["h2", "h3", "h4"]):
            level = int(tag.name[1])
            levels.append(level)
            
            # Check for proper progression
            if len(levels) > 1:
                prev_level = levels[-2]
                # H3 should not appear without H2, and levels shouldn't skip
                if (level == 3 and 2 not in levels) or (level - prev_level > 1):
                    return False
        
        return True
    
    def validate_anchor_resolution(self, html: str) -> bool:
        """TICKET 2: Validate that all TOC links resolve to actual heading IDs"""
        from bs4 import BeautifulSoup
        
        soup = BeautifulSoup(html, 'html.parser')
        
        # Get all existing IDs in the document
        existing_ids = {tag.get("id") for tag in soup.find_all(attrs={"id": True}) if tag.get("id")}
        
        # Check all Mini-TOC links
        broken_links = []
        for link in soup.select(".mini-toc a[href^='#']"):
            target_id = link.get("href", "")[1:]  # Remove the #
            if target_id not in existing_ids:
                broken_links.append(target_id)
        
        return len(broken_links) == 0
    
    async def _calculate_validation_metrics(self, normalized_doc, generated_articles: list, analysis: dict, run_id: str) -> dict:
        """V2 Engine: Calculate validation metrics"""
        try:
            print(f"🔍 V2 VALIDATION: Calculating validation metrics - run {run_id} - engine=v2")
            
            # Redundancy Score Calculation
            redundancy_score = await self._calculate_redundancy_score(generated_articles)
            
            # Granularity Alignment Score
            granularity_alignment = await self._calculate_granularity_alignment(generated_articles, analysis)
            
            # Complexity Alignment Score
            complexity_alignment = await self._calculate_complexity_alignment(normalized_doc, generated_articles, analysis)
            
            metrics = {
                "redundancy_score": redundancy_score,
                "granularity_alignment_score": granularity_alignment,
                "complexity_alignment_score": complexity_alignment,
                "total_articles": len(generated_articles),
                "total_source_blocks": len(normalized_doc.blocks),
                "average_article_length": self._calculate_average_article_length(generated_articles)
            }
            
            print(f"📊 V2 VALIDATION: Metrics calculated - Redundancy: {redundancy_score:.2f}, Granularity: {granularity_alignment:.2f}, Complexity: {complexity_alignment:.2f} - run {run_id} - engine=v2")
            
            return metrics
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error calculating validation metrics - {e} - run {run_id} - engine=v2")
            return {
                "redundancy_score": 0.5,
                "granularity_alignment_score": 0.5,
                "complexity_alignment_score": 0.5,
                "total_articles": len(generated_articles),
                "total_source_blocks": len(normalized_doc.blocks) if normalized_doc else 0,
                "average_article_length": 0
            }
    
    async def _calculate_redundancy_score(self, generated_articles: list) -> float:
        """Calculate content redundancy between articles"""
        try:
            if len(generated_articles) <= 1:
                return 0.0  # No redundancy with single article
            
            # Simple overlap detection based on content similarity
            article_contents = []
            for generated_article in generated_articles:
                article_data = generated_article.get('article_data', {})
                html_content = article_data.get('html', '')
                # Remove HTML tags for text comparison
                import re
                text_content = re.sub(r'<[^>]+>', '', html_content).lower()
                article_contents.append(set(text_content.split()))
            
            # Calculate pairwise similarity
            total_comparisons = 0
            total_overlap = 0
            
            for i in range(len(article_contents)):
                for j in range(i + 1, len(article_contents)):
                    content_a = article_contents[i]
                    content_b = article_contents[j]
                    
                    if len(content_a) > 0 and len(content_b) > 0:
                        overlap = len(content_a & content_b)
                        union = len(content_a | content_b)
                        similarity = overlap / union if union > 0 else 0
                        
                        total_overlap += similarity
                        total_comparisons += 1
            
            # Average redundancy score (lower is better)
            redundancy_score = total_overlap / total_comparisons if total_comparisons > 0 else 0.0
            return min(redundancy_score, 1.0)
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error calculating redundancy score - {e}")
            return 0.3  # Default moderate redundancy
    
    async def _calculate_granularity_alignment(self, generated_articles: list, analysis: dict) -> float:
        """Calculate how well article granularity matches analysis expectations"""
        try:
            expected_granularity = analysis.get('granularity', 'shallow')
            article_count = len(generated_articles)
            
            # Granularity expectations
            granularity_expectations = {
                'shallow': {'min_articles': 1, 'max_articles': 3},
                'moderate': {'min_articles': 2, 'max_articles': 8},
                'deep': {'min_articles': 5, 'max_articles': 20}
            }
            
            expectations = granularity_expectations.get(expected_granularity, {'min_articles': 1, 'max_articles': 10})
            
            # Calculate alignment score
            if expectations['min_articles'] <= article_count <= expectations['max_articles']:
                alignment_score = 1.0
            elif article_count < expectations['min_articles']:
                # Too few articles
                alignment_score = article_count / expectations['min_articles']
            else:
                # Too many articles
                excess = article_count - expectations['max_articles']
                penalty = min(0.5, excess * 0.1)
                alignment_score = max(0.1, 1.0 - penalty)
            
            return alignment_score
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error calculating granularity alignment - {e}")
            return 0.7  # Default moderate alignment
    
    async def _calculate_complexity_alignment(self, normalized_doc, generated_articles: list, analysis: dict) -> float:
        """Calculate how well article complexity matches source complexity"""
        try:
            expected_complexity = analysis.get('complexity', 'intermediate')
            
            # Calculate source complexity indicators
            total_blocks = len(normalized_doc.blocks)
            code_blocks = sum(1 for block in normalized_doc.blocks if block.block_type == 'code')
            
            source_complexity_score = 0
            if code_blocks > 0:
                source_complexity_score += min(0.3, code_blocks / total_blocks)
            if total_blocks > 20:
                source_complexity_score += 0.2
            if len(normalized_doc.media) > 0:
                source_complexity_score += 0.1
            
            # Calculate generated complexity
            total_content_length = sum(
                len(gen_article.get('article_data', {}).get('html', ''))
                for gen_article in generated_articles
            )
            
            generated_complexity_score = 0
            if total_content_length > 5000:
                generated_complexity_score += 0.3
            if len(generated_articles) > 3:
                generated_complexity_score += 0.2
            
            # Alignment calculation
            complexity_difference = abs(source_complexity_score - generated_complexity_score)
            alignment_score = max(0.1, 1.0 - complexity_difference * 2)
            
            return alignment_score
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error calculating complexity alignment - {e}")
            return 0.6  # Default moderate complexity alignment
    
    def _calculate_average_article_length(self, generated_articles: list) -> int:
        """Calculate average article content length"""
        try:
            if not generated_articles:
                return 0
            
            total_length = sum(
                len(gen_article.get('article_data', {}).get('html', ''))
                for gen_article in generated_articles
            )
            
            return total_length // len(generated_articles)
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error calculating average article length - {e}")
            return 0
    
    def _consolidate_validation_results(self, fidelity_coverage: dict, placeholder: dict, style: dict, evidence: dict, metrics: dict, run_id: str) -> dict:
        """Consolidate all validation results into final validation decision"""
        try:
            # Extract key metrics
            fidelity_score = fidelity_coverage.get('fidelity_score', 0.0)
            coverage_percent = fidelity_coverage.get('coverage_percent', 0.0)
            placeholder_count = len(placeholder.get('placeholders', []))
            style_compliance = style.get('overall_compliance', 0.0)
            evidence_tagging_rate = evidence.get('overall_tagging_rate', 0.0)
            
            # Apply quality thresholds
            fidelity_passed = fidelity_score >= self.quality_thresholds['fidelity_score']
            coverage_passed = coverage_percent >= self.quality_thresholds['coverage_percent']
            placeholder_passed = placeholder_count <= self.quality_thresholds['max_placeholders']
            style_passed = style_compliance >= 0.8  # 80% structural compliance required
            evidence_passed = evidence.get('validation_passed', True)  # ≥95% evidence tagging required
            
            # Overall validation status
            all_passed = all([fidelity_passed, coverage_passed, placeholder_passed, style_passed, evidence_passed])
            
            if all_passed:
                validation_status = "passed"
                status_message = "All validation checks passed"
            else:
                validation_status = "partial"
                failed_checks = []
                if not fidelity_passed:
                    failed_checks.append(f"fidelity_score ({fidelity_score:.2f} < {self.quality_thresholds['fidelity_score']})")
                if not coverage_passed:
                    failed_checks.append(f"coverage_percent ({coverage_percent:.1f}% < {self.quality_thresholds['coverage_percent']}%)")
                if not placeholder_passed:
                    failed_checks.append(f"placeholder_count ({placeholder_count} > {self.quality_thresholds['max_placeholders']})")
                if not style_passed:
                    failed_checks.append(f"style_compliance ({style_compliance:.2f} < 0.8)")
                if not evidence_passed:
                    failed_checks.append(f"evidence_tagging ({evidence_tagging_rate:.1f}% < 95%)")
                
                status_message = f"Failed checks: {', '.join(failed_checks)}"
            
            # Create comprehensive validation result
            validation_result = {
                "validation_id": f"validation_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "validation_status": validation_status,
                "status_message": status_message,
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2",
                
                # Detailed validation results
                "fidelity_coverage": fidelity_coverage,
                "placeholder_detection": placeholder,
                "style_guard": style,
                "evidence_tagging": evidence,
                "validation_metrics": metrics,
                
                # Summary scores
                "summary_scores": {
                    "fidelity_score": fidelity_score,
                    "coverage_percent": coverage_percent,  
                    "placeholder_count": placeholder_count,
                    "style_compliance": style_compliance,
                    "evidence_tagging_rate": evidence_tagging_rate,
                    "redundancy_score": metrics.get('redundancy_score', 0.0),
                    "granularity_alignment": metrics.get('granularity_alignment_score', 0.0),
                    "complexity_alignment": metrics.get('complexity_alignment_score', 0.0)
                },
                
                # Threshold compliance
                "threshold_compliance": {
                    "fidelity_passed": fidelity_passed,
                    "coverage_passed": coverage_passed,
                    "placeholder_passed": placeholder_passed,
                    "style_passed": style_passed,
                    "evidence_passed": evidence_passed
                },
                
                # Actionable diagnostics for failed runs
                "diagnostics": self._generate_actionable_diagnostics(
                    fidelity_coverage, placeholder, style, evidence, metrics, validation_status
                )
            }
            
            return validation_result
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error consolidating validation results - {e} - run {run_id} - engine=v2")
            return self._create_validation_result("error", run_id, {"consolidation_error": str(e)})
    
    async def _validate_evidence_tagging(self, generated_articles: list, normalized_doc, run_id: str) -> dict:
        """Validate that paragraphs have proper evidence tagging for fidelity enforcement"""
        try:
            print(f"🏷️ V2 VALIDATION: Validating evidence tagging - run {run_id} - engine=v2")
            
            evidence_results = []
            total_paragraphs = 0
            total_tagged_paragraphs = 0
            total_untagged_paragraphs = 0
            
            for generated_article in generated_articles:
                article_id = generated_article.get('article_id', 'unknown')
                article_data = generated_article.get('article_data', {})
                html_content = article_data.get('html', '')
                
                if not html_content:
                    continue
                
                # Parse paragraphs from content
                paragraphs = self._parse_validation_paragraphs(html_content)
                
                # Check evidence tagging for each paragraph
                tagged_count = 0
                untagged_paragraphs = []
                
                for paragraph in paragraphs:
                    if paragraph.get('is_faq', False):
                        continue  # Skip FAQ paragraphs
                    
                    total_paragraphs += 1
                    
                    # Check for evidence attributes
                    has_evidence = self._check_paragraph_evidence(paragraph)
                    
                    if has_evidence:
                        tagged_count += 1
                        total_tagged_paragraphs += 1
                    else:
                        untagged_paragraphs.append({
                            "text": paragraph.get('text', '')[:100],
                            "position": paragraph.get('position', 0)
                        })
                        total_untagged_paragraphs += 1
                
                # Calculate tagging rate for this article
                article_paragraphs = len([p for p in paragraphs if not p.get('is_faq', False)])
                article_tagging_rate = (tagged_count / article_paragraphs * 100) if article_paragraphs > 0 else 100
                
                evidence_results.append({
                    "article_id": article_id,
                    "total_paragraphs": article_paragraphs,
                    "tagged_paragraphs": tagged_count,
                    "untagged_paragraphs": len(untagged_paragraphs),
                    "tagging_rate": article_tagging_rate,
                    "untagged_examples": untagged_paragraphs[:3],  # First 3 examples
                    "validation_passed": article_tagging_rate >= 95.0
                })
            
            # Calculate overall evidence validation metrics
            overall_tagging_rate = (total_tagged_paragraphs / total_paragraphs * 100) if total_paragraphs > 0 else 100
            validation_passed = overall_tagging_rate >= 95.0
            
            # Check if source blocks are available for evidence
            source_blocks_available = len(normalized_doc.blocks) > 0
            
            evidence_validation_result = {
                "validation_type": "evidence_tagging",
                "validation_passed": validation_passed,
                "total_paragraphs": total_paragraphs,
                "tagged_paragraphs": total_tagged_paragraphs,
                "untagged_paragraphs": total_untagged_paragraphs,
                "overall_tagging_rate": overall_tagging_rate,
                "target_threshold": 95.0,
                "source_blocks_available": source_blocks_available,
                "source_blocks_count": len(normalized_doc.blocks),
                "article_results": evidence_results,
                "summary": f"Evidence tagging: {total_tagged_paragraphs}/{total_paragraphs} paragraphs tagged ({overall_tagging_rate:.1f}%)"
            }
            
            # Fail validation if evidence tagging is insufficient and source blocks exist
            if not validation_passed and source_blocks_available:
                evidence_validation_result["validation_message"] = f"Evidence tagging below 95% threshold ({overall_tagging_rate:.1f}%) - fidelity enforcement failed"
                print(f"❌ V2 VALIDATION: Evidence tagging failed - {overall_tagging_rate:.1f}% < 95% threshold - run {run_id} - engine=v2")
            else:
                evidence_validation_result["validation_message"] = f"Evidence tagging passed - {overall_tagging_rate:.1f}% paragraphs tagged"
                print(f"✅ V2 VALIDATION: Evidence tagging passed - {overall_tagging_rate:.1f}% paragraphs tagged - run {run_id} - engine=v2")
            
            return evidence_validation_result
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error validating evidence tagging - {e} - run {run_id} - engine=v2")
            return {
                "validation_type": "evidence_tagging",
                "validation_passed": False,
                "error": str(e),
                "summary": "Evidence tagging validation failed due to error"
            }
    
    def _parse_validation_paragraphs(self, html_content: str) -> list:
        """Parse paragraphs from HTML content for validation"""
        try:
            from bs4 import BeautifulSoup
            import re
            
            paragraphs = []
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Find all paragraph elements
            p_tags = soup.find_all('p')
            
            for i, p_tag in enumerate(p_tags):
                paragraph_text = p_tag.get_text().strip()
                
                if len(paragraph_text) < 20:  # Skip very short paragraphs
                    continue
                
                # Check if this is a FAQ paragraph
                is_faq = any(indicator in paragraph_text.lower() for indicator in 
                           ['faq', 'q:', 'question:', 'a:', 'answer:'])
                
                paragraphs.append({
                    "index": i,
                    "text": paragraph_text,
                    "html": str(p_tag),
                    "is_faq": is_faq,
                    "position": html_content.find(str(p_tag))
                })
            
            return paragraphs
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error parsing paragraphs for validation - {e}")
            return []
    
    def _check_paragraph_evidence(self, paragraph: dict) -> bool:
        """Check if a paragraph has proper evidence attribution"""
        try:
            paragraph_html = paragraph.get('html', '')
            
            # Check for data-evidence attribute
            if 'data-evidence=' in paragraph_html:
                return True
            
            # Check for HTML comment with evidence
            if '<!-- data-evidence=' in paragraph_html:
                return True
            
            # Check for any block ID references
            import re
            if re.search(r'data-evidence="\[.*?\]"', paragraph_html):
                return True
            
            return False
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error checking paragraph evidence - {e}")
            return False
    
    def _generate_actionable_diagnostics(self, fidelity_coverage: dict, placeholder: dict, style: dict, evidence: dict, metrics: dict, status: str) -> list:
        """Generate actionable diagnostics for validation failures"""
        diagnostics = []
        
        try:
            if status == "passed":
                diagnostics.append({
                    "type": "success",
                    "message": "All validation checks passed successfully",
                    "action": "No action required"
                })
                return diagnostics
            
            # Fidelity diagnostics
            fidelity_score = fidelity_coverage.get('fidelity_score', 0.0)
            if fidelity_score < self.quality_thresholds['fidelity_score']:
                hallucinated_claims = fidelity_coverage.get('hallucinated_claims', [])
                diagnostics.append({
                    "type": "fidelity_warning",
                    "message": f"Fidelity score ({fidelity_score:.2f}) below threshold ({self.quality_thresholds['fidelity_score']})",
                    "action": f"Review and remove {len(hallucinated_claims)} hallucinated claims",
                    "details": hallucinated_claims[:3]  # Show first 3 examples
                })
            
            # Coverage diagnostics
            coverage_percent = fidelity_coverage.get('coverage_percent', 0.0)
            if coverage_percent < self.quality_thresholds['coverage_percent']:
                uncovered_blocks = fidelity_coverage.get('uncovered_blocks', [])
                diagnostics.append({
                    "type": "coverage_warning",
                    "message": f"Coverage ({coverage_percent:.1f}%) below required 100%",
                    "action": f"Include content from {len(uncovered_blocks)} uncovered source blocks",
                    "details": uncovered_blocks[:5]  # Show first 5 examples
                })
            
            # Placeholder diagnostics
            placeholders = placeholder.get('placeholders', [])
            if len(placeholders) > self.quality_thresholds['max_placeholders']:
                diagnostics.append({
                    "type": "placeholder_warning",
                    "message": f"Found {len(placeholders)} placeholders (max allowed: {self.quality_thresholds['max_placeholders']})",
                    "action": "Replace placeholder content with actual information",
                    "details": [p.get('text', 'Unknown') for p in placeholders[:3]]  # Show first 3
                })
            
            # Style compliance diagnostics
            style_compliance = style.get('overall_compliance', 0.0)
            if style_compliance < 0.8:
                missing_elements = []
                for article_result in style.get('article_results', []):
                    missing_elements.extend(article_result.get('missing_elements', []))
                
                unique_missing = list(set(missing_elements))
                diagnostics.append({
                    "type": "style_warning",
                    "message": f"Style compliance ({style_compliance:.2f}) below 80%",
                    "action": f"Add missing structural elements: {', '.join(unique_missing)}",
                    "details": unique_missing
                })
            
            # Metrics diagnostics
            redundancy_score = metrics.get('redundancy_score', 0.0)
            if redundancy_score > self.quality_thresholds['redundancy_score']:
                diagnostics.append({
                    "type": "redundancy_warning",
                    "message": f"High content redundancy ({redundancy_score:.2f}) detected",
                    "action": "Review articles for duplicate content and consolidate where appropriate",
                    "details": []
                })
            
            # Evidence tagging diagnostics
            evidence_tagging_rate = evidence.get('overall_tagging_rate', 100.0)
            if evidence_tagging_rate < 95.0:
                untagged_count = evidence.get('untagged_paragraphs', 0)
                diagnostics.append({
                    "type": "evidence_warning",
                    "message": f"Evidence tagging below 95% threshold ({evidence_tagging_rate:.1f}%)",
                    "action": f"Add evidence block IDs to {untagged_count} untagged paragraphs for fidelity enforcement",
                    "details": [f"{untagged_count} paragraphs missing evidence attribution"]
                })
            
            return diagnostics
            
        except Exception as e:
            print(f"❌ V2 VALIDATION: Error generating diagnostics - {e}")
            return [{
                "type": "error",
                "message": "Error generating diagnostics",
                "action": "Review validation manually",
                "details": [str(e)]
            }]
    
    def _create_validation_result(self, status: str, run_id: str, additional_data: dict) -> dict:
        """Create a standard validation result structure"""
        return {
            "validation_id": f"validation_{run_id}_{int(datetime.utcnow().timestamp())}",
            "run_id": run_id,
            "validation_status": status,
            "timestamp": datetime.utcnow().isoformat(),
            "engine": "v2",
            **additional_data
        }

# Global V2 Validation System instance
v2_validation_system = V2ValidationSystem()

# ========================================
# V2 ENGINE: CROSS-ARTICLE QA SYSTEM
# ========================================

class V2CrossArticleQASystem:
    """V2 Engine: Cross-article quality assurance for coherence, deduplication, and consistency"""
    
    def __init__(self):
        self.duplicate_threshold = 0.8  # Similarity threshold for duplicate detection
        self.terminology_patterns = [
            # Common API terminology variations
            {"standard": "API key", "variations": ["Api key", "APIKey", "api key", "API-key", "api_key"]},
            {"standard": "API endpoint", "variations": ["Api endpoint", "APIEndpoint", "api endpoint", "API-endpoint", "api_endpoint"]},
            {"standard": "HTTP request", "variations": ["Http request", "HTTPRequest", "http request", "HTTP-request", "http_request"]},
            {"standard": "JSON response", "variations": ["Json response", "JSONResponse", "json response", "JSON-response", "json_response"]},
            {"standard": "OAuth token", "variations": ["Oauth token", "OAuthToken", "oauth token", "OAuth-token", "oauth_token"]},
            # Add more patterns as needed
        ]
    
    async def perform_cross_article_qa(self, generated_articles_result: dict, run_id: str) -> dict:
        """V2 Engine: Perform comprehensive cross-article quality assurance"""
        try:
            print(f"🔍 V2 CROSS-ARTICLE QA: Starting comprehensive QA analysis - run {run_id} - engine=v2")
            
            generated_articles = generated_articles_result.get('generated_articles', [])
            if len(generated_articles) < 2:
                print(f"⚠️ V2 CROSS-ARTICLE QA: Less than 2 articles, skipping cross-article analysis - run {run_id} - engine=v2")
                return self._create_qa_result("insufficient_articles", run_id, {"article_count": len(generated_articles)})
            
            # Prepare article set for analysis
            article_set = self._prepare_article_set(generated_articles)
            
            # Step 1: LLM-based cross-article analysis
            llm_qa_result = await self._perform_llm_cross_article_analysis(article_set, run_id)
            
            # Step 2: Programmatic validation and enhancement
            programmatic_qa_result = await self._perform_programmatic_qa_analysis(article_set, run_id)
            
            # Step 3: Consolidate findings
            consolidated_qa_result = self._consolidate_qa_findings(llm_qa_result, programmatic_qa_result, run_id)
            
            # Step 4: Perform consolidation pass
            consolidation_result = await self._perform_consolidation_pass(
                generated_articles, consolidated_qa_result, run_id
            )
            
            # Step 5: Create final QA result
            final_qa_result = {
                **consolidated_qa_result,
                "consolidation_result": consolidation_result,
                "qa_id": f"qa_{run_id}_{int(datetime.utcnow().timestamp())}",
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
            
            print(f"✅ V2 CROSS-ARTICLE QA: Analysis complete - Found {len(consolidated_qa_result.get('duplicates', []))} duplicates, {len(consolidated_qa_result.get('invalid_related_links', []))} invalid links - run {run_id} - engine=v2")
            return final_qa_result
            
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error in cross-article QA - {e} - run {run_id} - engine=v2")
            return self._create_qa_result("error", run_id, {"error": str(e)})
    
    def _prepare_article_set(self, generated_articles: list) -> dict:
        """Prepare article set for cross-article analysis"""
        try:
            article_set = {
                "articles": [],
                "total_count": len(generated_articles)
            }
            
            for generated_article in generated_articles:
                article_id = generated_article.get('article_id', 'unknown')
                article_data = generated_article.get('article_data', {})
                html_content = article_data.get('html', '')
                
                # Extract structured data from HTML
                structured_article = self._extract_structured_data(article_id, html_content)
                article_set["articles"].append(structured_article)
            
            return article_set
            
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error preparing article set - {e}")
            return {"articles": [], "total_count": 0}
    
    def _extract_structured_data(self, article_id: str, html_content: str) -> dict:
        """Extract structured data from HTML content"""
        try:
            import re
            from bs4 import BeautifulSoup
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract title
            title_tag = soup.find('h1')
            title = title_tag.get_text().strip() if title_tag else "Untitled"
            
            # Extract sections
            sections = []
            section_tags = soup.find_all(['h2', 'h3'])
            for tag in section_tags:
                sections.append({
                    "level": tag.name,
                    "heading": tag.get_text().strip(),
                    "id": tag.get('id', ''),
                    "content": self._get_section_content(tag)
                })
            
            # Extract FAQs
            faqs = []
            faq_patterns = [
                r'Q:\s*(.*?)\s*A:\s*(.*?)(?=Q:|$)',
                r'<h3[^>]*>Q:\s*(.*?)</h3>\s*<p[^>]*>A:\s*(.*?)</p>',
                r'Question:\s*(.*?)\s*Answer:\s*(.*?)(?=Question:|$)'
            ]
            
            for pattern in faq_patterns:
                matches = re.finditer(pattern, html_content, re.IGNORECASE | re.DOTALL)
                for match in matches:
                    faqs.append({
                        "question": match.group(1).strip(),
                        "answer": match.group(2).strip()
                    })
            
            # Extract related links
            related_links = []
            link_tags = soup.find_all('a', href=True)
            for link in link_tags:
                href = link.get('href', '')
                if href.startswith('#') or href.startswith('/') or href.startswith('http'):
                    related_links.append({
                        "label": link.get_text().strip(),
                        "url": href,
                        "is_internal": href.startswith('#') or href.startswith('/')
                    })
            
            return {
                "article_id": article_id,
                "title": title,
                "sections": sections,
                "faqs": faqs,
                "related_links": related_links,
                "content_length": len(html_content),
                "html_content": html_content[:1000]  # Truncated for analysis
            }
            
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error extracting structured data from {article_id} - {e}")
            return {
                "article_id": article_id,
                "title": "Error",
                "sections": [],
                "faqs": [],
                "related_links": [],
                "content_length": 0,
                "html_content": ""
            }
    
    def _get_section_content(self, heading_tag) -> str:
        """Get content under a heading tag"""
        try:
            content_parts = []
            next_element = heading_tag.next_sibling
            
            while next_element:
                if hasattr(next_element, 'name'):
                    if next_element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                        break
                    content_parts.append(str(next_element))
                next_element = next_element.next_sibling
            
            return ' '.join(content_parts)[:500]  # Truncate for analysis
            
        except Exception as e:
            return f"Error extracting content: {str(e)}"
    
    async def _perform_llm_cross_article_analysis(self, article_set: dict, run_id: str) -> dict:
        """V2 Engine: LLM-based cross-article analysis"""
        try:
            print(f"🤖 V2 CROSS-ARTICLE QA: Starting LLM cross-article analysis - run {run_id} - engine=v2")
            
            # Create system message
            system_message = """You are a documentation reviewer ensuring coherence across articles.

Your task is to analyze a set of articles and identify:
1. Duplicates: Repeated content (intros/sections) across articles
2. Invalid related links: Links that don't point to existing articles/sections
3. Duplicate FAQs: Same questions appearing across multiple articles
4. Terminology issues: Inconsistent usage of terms across articles

Be thorough but precise. Focus on actual duplicates and issues, not minor variations.

Return ONLY JSON in the exact format specified."""

            # Create user message with article set
            user_message = f"""Identify duplicates, invalid related links, duplicate FAQs, and terminology issues. Return JSON.

ARTICLE SET:
{json.dumps(article_set, indent=2)}

Analyze these articles and return ONLY JSON in this exact format:
{{
  "duplicates": [
    {{
      "article_id": "a1",
      "other_article_id": "a2", 
      "section": "Intro",
      "similarity_score": 0.95,
      "duplicate_type": "identical_content"
    }}
  ],
  "invalid_related_links": [
    {{
      "article_id": "a3",
      "label": "Nonexistent Link",
      "url": "/kb/missing",
      "issue": "target_not_found"
    }}
  ],
  "duplicate_faqs": [
    {{
      "question": "How to install X?",
      "article_ids": ["a1", "a2"],
      "identical_answer": true
    }}
  ],
  "terminology_issues": [
    {{
      "term": "API key",
      "inconsistent_usages": ["Api key", "APIKey", "api_key"],
      "suggested_standard": "API key",
      "article_ids": ["a1", "a2", "a3"]
    }}
  ]
}}"""

            # Call LLM for analysis
            print(f"🤖 V2 CROSS-ARTICLE QA: Sending cross-article analysis request to LLM - run {run_id} - engine=v2")
            ai_response = await call_llm_with_fallback(system_message, user_message)
            
            if ai_response:
                # Parse JSON response
                import re
                json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
                if json_match:
                    qa_data = json.loads(json_match.group(0))
                    
                    # Validate required fields
                    required_fields = ['duplicates', 'invalid_related_links', 'duplicate_faqs', 'terminology_issues']
                    if all(field in qa_data for field in required_fields):
                        duplicates_count = len(qa_data.get('duplicates', []))
                        invalid_links_count = len(qa_data.get('invalid_related_links', []))
                        duplicate_faqs_count = len(qa_data.get('duplicate_faqs', []))
                        terminology_issues_count = len(qa_data.get('terminology_issues', []))
                        
                        print(f"🔍 V2 CROSS-ARTICLE QA: LLM found {duplicates_count} duplicates, {invalid_links_count} invalid links, {duplicate_faqs_count} duplicate FAQs, {terminology_issues_count} terminology issues - run {run_id} - engine=v2")
                        return qa_data
                    else:
                        print(f"⚠️ V2 CROSS-ARTICLE QA: Missing fields in LLM response - run {run_id} - engine=v2")
                        return self._create_fallback_qa_analysis(article_set)
                else:
                    print(f"⚠️ V2 CROSS-ARTICLE QA: No JSON found in LLM response - run {run_id} - engine=v2")
                    return self._create_fallback_qa_analysis(article_set)
            else:
                print(f"❌ V2 CROSS-ARTICLE QA: No LLM response for cross-article analysis - run {run_id} - engine=v2")
                return self._create_fallback_qa_analysis(article_set)
                
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error in LLM cross-article analysis - {e} - run {run_id} - engine=v2")
            return self._create_fallback_qa_analysis(article_set)
    
    def _create_fallback_qa_analysis(self, article_set: dict) -> dict:
        """Create fallback QA analysis using programmatic methods"""
        try:
            articles = article_set.get('articles', [])
            
            # Basic duplicate detection
            duplicates = []
            for i, article_a in enumerate(articles):
                for j, article_b in enumerate(articles[i+1:], i+1):
                    # Check title similarity
                    title_a = article_a.get('title', '').lower()
                    title_b = article_b.get('title', '').lower()
                    
                    if title_a and title_b and self._calculate_similarity(title_a, title_b) > self.duplicate_threshold:
                        duplicates.append({
                            "article_id": article_a.get('article_id'),
                            "other_article_id": article_b.get('article_id'),
                            "section": "title",
                            "similarity_score": self._calculate_similarity(title_a, title_b),
                            "duplicate_type": "similar_title"
                        })
            
            # Basic FAQ duplicate detection
            duplicate_faqs = []
            faq_questions = {}
            
            for article in articles:
                article_id = article.get('article_id')
                faqs = article.get('faqs', [])
                
                for faq in faqs:
                    question = faq.get('question', '').lower().strip()
                    if question:
                        if question in faq_questions:
                            faq_questions[question].append(article_id)
                        else:
                            faq_questions[question] = [article_id]
            
            for question, article_ids in faq_questions.items():
                if len(article_ids) > 1:
                    duplicate_faqs.append({
                        "question": question,
                        "article_ids": article_ids,
                        "identical_answer": False  # Can't determine without deeper analysis
                    })
            
            # Basic terminology detection
            terminology_issues = []
            for pattern in self.terminology_patterns:
                standard = pattern["standard"]
                variations = pattern["variations"]
                found_variations = []
                found_in_articles = []
                
                for article in articles:
                    html_content = article.get('html_content', '').lower()
                    for variation in variations:
                        if variation.lower() in html_content:
                            if variation not in found_variations:
                                found_variations.append(variation)
                            if article.get('article_id') not in found_in_articles:
                                found_in_articles.append(article.get('article_id'))
                
                if len(found_variations) > 1:
                    terminology_issues.append({
                        "term": standard,
                        "inconsistent_usages": found_variations,
                        "suggested_standard": standard,
                        "article_ids": found_in_articles
                    })
            
            print(f"🔧 V2 CROSS-ARTICLE QA: Fallback analysis found {len(duplicates)} duplicates, {len(duplicate_faqs)} duplicate FAQs, {len(terminology_issues)} terminology issues")
            
            return {
                "duplicates": duplicates,
                "invalid_related_links": [],  # Requires URL resolution
                "duplicate_faqs": duplicate_faqs,
                "terminology_issues": terminology_issues,
                "analysis_method": "fallback_programmatic"
            }
            
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error in fallback QA analysis - {e}")
            return {
                "duplicates": [],
                "invalid_related_links": [],
                "duplicate_faqs": [],
                "terminology_issues": [],
                "analysis_method": "error"
            }
    
    def _calculate_similarity(self, text_a: str, text_b: str) -> float:
        """Calculate text similarity using simple word overlap"""
        try:
            words_a = set(text_a.lower().split())
            words_b = set(text_b.lower().split())
            
            if not words_a or not words_b:
                return 0.0
            
            intersection = len(words_a & words_b)
            union = len(words_a | words_b)
            
            return intersection / union if union > 0 else 0.0
            
        except Exception as e:
            return 0.0
    
    async def _perform_programmatic_qa_analysis(self, article_set: dict, run_id: str) -> dict:
        """V2 Engine: Programmatic QA analysis for validation"""
        try:
            print(f"🔍 V2 CROSS-ARTICLE QA: Performing programmatic QA validation - run {run_id} - engine=v2")
            
            articles = article_set.get('articles', [])
            
            # Validate related links
            invalid_related_links = []
            existing_article_ids = set(article.get('article_id', '') for article in articles)
            existing_sections = set()
            
            # Build section index
            for article in articles:
                article_id = article.get('article_id', '')
                sections = article.get('sections', [])
                for section in sections:
                    section_id = section.get('id', '')
                    if section_id:
                        existing_sections.add(f"#{section_id}")
                        existing_sections.add(f"{article_id}#{section_id}")
            
            # Check related links validity
            for article in articles:
                article_id = article.get('article_id', '')
                related_links = article.get('related_links', [])
                
                for link in related_links:
                    url = link.get('url', '')
                    label = link.get('label', '')
                    is_internal = link.get('is_internal', False)
                    
                    if is_internal:
                        # Check if internal link exists
                        if url.startswith('#'):
                            if url not in existing_sections:
                                invalid_related_links.append({
                                    "article_id": article_id,
                                    "label": label,
                                    "url": url,
                                    "issue": "section_anchor_not_found"
                                })
                        elif url.startswith('/'):
                            # Check if it's a reference to another article
                            if not any(article_ref in url for article_ref in existing_article_ids):
                                invalid_related_links.append({
                                    "article_id": article_id,
                                    "label": label,
                                    "url": url,
                                    "issue": "internal_article_not_found"
                                })
            
            # Additional consistency checks
            title_consistency = self._check_title_consistency(articles)
            section_consistency = self._check_section_consistency(articles)
            
            programmatic_result = {
                "invalid_related_links_validated": invalid_related_links,
                "title_consistency": title_consistency,
                "section_consistency": section_consistency,
                "analysis_method": "programmatic_validation"
            }
            
            print(f"🔍 V2 CROSS-ARTICLE QA: Programmatic validation found {len(invalid_related_links)} invalid links - run {run_id} - engine=v2")
            return programmatic_result
            
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error in programmatic QA analysis - {e} - run {run_id} - engine=v2")
            return {
                "invalid_related_links_validated": [],
                "title_consistency": {"consistent": True, "issues": []},
                "section_consistency": {"consistent": True, "issues": []},
                "analysis_method": "error"
            }
    
    def _check_title_consistency(self, articles: list) -> dict:
        """Check title formatting consistency"""
        try:
            title_patterns = []
            for article in articles:
                title = article.get('title', '')
                if title:
                    pattern = {
                        "has_numbers": any(char.isdigit() for char in title),
                        "has_colons": ':' in title,
                        "has_dashes": '-' in title,
                        "title_case": title.istitle(),
                        "length": len(title)
                    }
                    title_patterns.append(pattern)
            
            # Analyze consistency
            if not title_patterns:
                return {"consistent": True, "issues": []}
            
            consistency_issues = []
            title_case_count = sum(1 for p in title_patterns if p['title_case'])
            if 0 < title_case_count < len(title_patterns):
                consistency_issues.append("Inconsistent title case formatting")
            
            return {
                "consistent": len(consistency_issues) == 0,
                "issues": consistency_issues,
                "title_patterns": title_patterns
            }
            
        except Exception as e:
            return {
                "consistent": False,
                "issues": [f"Error checking title consistency: {str(e)}"]
            }
    
    def _check_section_consistency(self, articles: list) -> dict:
        """Check section heading consistency"""
        try:
            section_patterns = []
            all_headings = []
            
            for article in articles:
                sections = article.get('sections', [])
                for section in sections:
                    heading = section.get('heading', '')
                    if heading:
                        all_headings.append(heading)
                        pattern = {
                            "level": section.get('level', 'h2'),
                            "has_numbers": any(char.isdigit() for char in heading),
                            "title_case": heading.istitle(),
                            "length": len(heading)
                        }
                        section_patterns.append(pattern)
            
            # Check for consistency issues
            consistency_issues = []
            if len(all_headings) > 1:
                title_case_count = sum(1 for p in section_patterns if p['title_case'])
                if 0 < title_case_count < len(section_patterns):
                    consistency_issues.append("Inconsistent section heading case")
            
            return {
                "consistent": len(consistency_issues) == 0,
                "issues": consistency_issues,
                "total_sections": len(section_patterns)
            }
            
        except Exception as e:
            return {
                "consistent": False,
                "issues": [f"Error checking section consistency: {str(e)}"]
            }
    
    def _consolidate_qa_findings(self, llm_result: dict, programmatic_result: dict, run_id: str) -> dict:
        """Consolidate LLM and programmatic QA findings"""
        try:
            # Merge LLM and programmatic results
            consolidated = {
                "duplicates": llm_result.get('duplicates', []),
                "invalid_related_links": llm_result.get('invalid_related_links', []),
                "duplicate_faqs": llm_result.get('duplicate_faqs', []),
                "terminology_issues": llm_result.get('terminology_issues', []),
                "run_id": run_id,
                "analysis_methods": []
            }
            
            # Add programmatic validation for invalid links
            programmatic_invalid_links = programmatic_result.get('invalid_related_links_validated', [])
            for link in programmatic_invalid_links:
                # Check if not already found by LLM
                if not any(existing_link.get('url') == link.get('url') 
                         for existing_link in consolidated['invalid_related_links']):
                    consolidated['invalid_related_links'].append(link)
            
            # Add analysis methods used
            if llm_result.get('analysis_method') != 'error':
                consolidated['analysis_methods'].append('llm_analysis')
            if programmatic_result.get('analysis_method') != 'error':
                consolidated['analysis_methods'].append('programmatic_validation')
            
            # Add additional findings
            consolidated['title_consistency'] = programmatic_result.get('title_consistency', {})
            consolidated['section_consistency'] = programmatic_result.get('section_consistency', {})
            
            # Summary statistics
            consolidated['summary'] = {
                "total_duplicates": len(consolidated['duplicates']),
                "total_invalid_links": len(consolidated['invalid_related_links']),
                "total_duplicate_faqs": len(consolidated['duplicate_faqs']),
                "total_terminology_issues": len(consolidated['terminology_issues']),
                "issues_found": (
                    len(consolidated['duplicates']) + 
                    len(consolidated['invalid_related_links']) + 
                    len(consolidated['duplicate_faqs']) + 
                    len(consolidated['terminology_issues'])
                )
            }
            
            print(f"🔍 V2 CROSS-ARTICLE QA: Consolidated findings - {consolidated['summary']['issues_found']} total issues - run {run_id} - engine=v2")
            return consolidated
            
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error consolidating QA findings - {e} - run {run_id} - engine=v2")
            return {
                "duplicates": [],
                "invalid_related_links": [],
                "duplicate_faqs": [],
                "terminology_issues": [],
                "summary": {"issues_found": 0}
            }
    
    async def _perform_consolidation_pass(self, generated_articles: list, qa_findings: dict, run_id: str) -> dict:
        """V2 Engine: Perform consolidation pass to resolve QA issues"""
        try:
            print(f"🔧 V2 CROSS-ARTICLE QA: Starting consolidation pass - run {run_id} - engine=v2")
            
            consolidation_actions = []
            
            # Handle duplicates
            duplicates = qa_findings.get('duplicates', [])
            for duplicate in duplicates:
                action = await self._handle_duplicate_content(generated_articles, duplicate, run_id)
                consolidation_actions.append(action)
            
            # Handle invalid related links
            invalid_links = qa_findings.get('invalid_related_links', [])
            for invalid_link in invalid_links:
                action = await self._handle_invalid_related_link(generated_articles, invalid_link, run_id)
                consolidation_actions.append(action)
            
            # Handle duplicate FAQs
            duplicate_faqs = qa_findings.get('duplicate_faqs', [])
            for duplicate_faq in duplicate_faqs:
                action = await self._handle_duplicate_faq(generated_articles, duplicate_faq, run_id)
                consolidation_actions.append(action)
            
            # Handle terminology issues
            terminology_issues = qa_findings.get('terminology_issues', [])
            for terminology_issue in terminology_issues:
                action = await self._handle_terminology_issue(generated_articles, terminology_issue, run_id)
                consolidation_actions.append(action)
            
            consolidation_result = {
                "actions_taken": consolidation_actions,
                "total_actions": len(consolidation_actions),
                "successful_actions": len([a for a in consolidation_actions if a.get('status') == 'success']),
                "failed_actions": len([a for a in consolidation_actions if a.get('status') == 'failed']),
                "consolidation_method": "automated_pass"
            }
            
            print(f"🔧 V2 CROSS-ARTICLE QA: Consolidation complete - {consolidation_result['successful_actions']}/{consolidation_result['total_actions']} successful actions - run {run_id} - engine=v2")
            return consolidation_result
            
        except Exception as e:
            print(f"❌ V2 CROSS-ARTICLE QA: Error in consolidation pass - {e} - run {run_id} - engine=v2")
            return {
                "actions_taken": [],
                "total_actions": 0,
                "successful_actions": 0,
                "failed_actions": 0,
                "consolidation_method": "error"
            }
    
    async def _handle_duplicate_content(self, generated_articles: list, duplicate: dict, run_id: str) -> dict:
        """Handle duplicate content between articles"""
        try:
            article_id = duplicate.get('article_id')
            other_article_id = duplicate.get('other_article_id')
            section = duplicate.get('section', 'unknown')
            
            # For now, just record the action (full implementation would modify articles)
            action = {
                "type": "duplicate_content",
                "article_id": article_id,
                "other_article_id": other_article_id,
                "section": section,
                "action": "recorded_for_manual_review",
                "status": "success",
                "details": f"Duplicate {section} content identified between {article_id} and {other_article_id}"
            }
            
            print(f"📝 V2 CROSS-ARTICLE QA: Recorded duplicate content - {article_id} <-> {other_article_id} - run {run_id} - engine=v2")
            return action
            
        except Exception as e:
            return {
                "type": "duplicate_content",
                "action": "error",
                "status": "failed",
                "error": str(e)
            }
    
    async def _handle_invalid_related_link(self, generated_articles: list, invalid_link: dict, run_id: str) -> dict:
        """Handle invalid related links"""
        try:
            article_id = invalid_link.get('article_id')
            url = invalid_link.get('url')
            label = invalid_link.get('label')
            
            # Find the article and attempt to fix the link
            for generated_article in generated_articles:
                if generated_article.get('article_id') == article_id:
                    article_data = generated_article.get('article_data', {})
                    html_content = article_data.get('html', '')
                    
                    # Simple approach: remove the invalid link or mark it for review
                    # In a full implementation, this would intelligently fix or replace links
                    if url in html_content:
                        # Mark for manual review rather than auto-remove
                        action = {
                            "type": "invalid_related_link",
                            "article_id": article_id,
                            "url": url,
                            "label": label,
                            "action": "marked_for_manual_review",
                            "status": "success",
                            "details": f"Invalid link '{label}' -> '{url}' marked for review"
                        }
                        
                        print(f"🔗 V2 CROSS-ARTICLE QA: Marked invalid link for review - {article_id}: {url} - run {run_id} - engine=v2")
                        return action
            
            return {
                "type": "invalid_related_link",
                "article_id": article_id,
                "action": "article_not_found",
                "status": "failed"
            }
            
        except Exception as e:
            return {
                "type": "invalid_related_link",
                "action": "error", 
                "status": "failed",
                "error": str(e)
            }
    
    async def _handle_duplicate_faq(self, generated_articles: list, duplicate_faq: dict, run_id: str) -> dict:
        """Handle duplicate FAQ across articles"""
        try:
            question = duplicate_faq.get('question')
            article_ids = duplicate_faq.get('article_ids', [])
            
            # For now, record the duplication for potential consolidation
            action = {
                "type": "duplicate_faq",
                "question": question,
                "article_ids": article_ids,
                "action": "recorded_for_consolidation",
                "status": "success",
                "details": f"Duplicate FAQ '{question}' found in {len(article_ids)} articles",
                "consolidation_recommendation": "Consider creating a centralized FAQ section"
            }
            
            print(f"❓ V2 CROSS-ARTICLE QA: Recorded duplicate FAQ - '{question}' in {len(article_ids)} articles - run {run_id} - engine=v2")
            return action
            
        except Exception as e:
            return {
                "type": "duplicate_faq",
                "action": "error",
                "status": "failed", 
                "error": str(e)
            }
    
    async def _handle_terminology_issue(self, generated_articles: list, terminology_issue: dict, run_id: str) -> dict:
        """Handle terminology consistency issues"""
        try:
            term = terminology_issue.get('term')
            inconsistent_usages = terminology_issue.get('inconsistent_usages', [])
            suggested_standard = terminology_issue.get('suggested_standard')
            article_ids = terminology_issue.get('article_ids', [])
            
            # Record terminology standardization action
            action = {
                "type": "terminology_issue",
                "term": term,
                "inconsistent_usages": inconsistent_usages,
                "suggested_standard": suggested_standard,
                "article_ids": article_ids,
                "action": "recorded_for_standardization",
                "status": "success",
                "details": f"Terminology '{term}' has {len(inconsistent_usages)} variations across {len(article_ids)} articles",
                "standardization_recommendation": f"Standardize all variations to '{suggested_standard}'"
            }
            
            print(f"📝 V2 CROSS-ARTICLE QA: Recorded terminology issue - '{term}' needs standardization - run {run_id} - engine=v2")
            return action
            
        except Exception as e:
            return {
                "type": "terminology_issue",
                "action": "error",
                "status": "failed",
                "error": str(e)
            }
    
    def _create_qa_result(self, status: str, run_id: str, additional_data: dict) -> dict:
        """Create a standard QA result structure"""
        return {
            "qa_id": f"qa_{run_id}_{int(datetime.utcnow().timestamp())}",
            "run_id": run_id,
            "qa_status": status,
            "timestamp": datetime.utcnow().isoformat(),
            "engine": "v2",
            **additional_data
        }

# Global V2 Cross-Article QA System instance
v2_cross_article_qa_system = V2CrossArticleQASystem()

# ========================================
# V2 ENGINE: ADAPTIVE ADJUSTMENT SYSTEM
# ========================================

class V2AdaptiveAdjustmentSystem:
    """V2 Engine: Adaptive adjustment system for balancing article lengths and splits"""
    
    def __init__(self):
        self.word_count_thresholds = {
            "min_article_length": 300,  # Articles below this should be merged
            "max_section_length": 1200,  # Sections above this should be split
            "optimal_article_range": (500, 2000),  # Optimal article length range
            "optimal_section_range": (200, 800)   # Optimal section length range
        }
        
        self.granularity_expectations = {
            "shallow": {"min_articles": 1, "max_articles": 3, "target_length_per_article": 1500},
            "moderate": {"min_articles": 2, "max_articles": 8, "target_length_per_article": 1000},
            "deep": {"min_articles": 5, "max_articles": 20, "target_length_per_article": 800}
        }
    
    async def perform_adaptive_adjustment(self, generated_articles_result: dict, analysis: dict, run_id: str) -> dict:
        """V2 Engine: Perform adaptive adjustment for optimal article balance"""
        try:
            print(f"⚖️ V2 ADAPTIVE ADJUSTMENT: Starting length and split balancing - run {run_id} - engine=v2")
            
            generated_articles = generated_articles_result.get('generated_articles', [])
            if not generated_articles:
                print(f"⚠️ V2 ADAPTIVE ADJUSTMENT: No articles to adjust - run {run_id} - engine=v2")
                return self._create_adjustment_result("no_articles", run_id, {"article_count": 0})
            
            # Step 1: Analyze current word counts and structure
            word_count_analysis = await self._analyze_word_counts(generated_articles, run_id)
            
            # Step 2: LLM-based balancing analysis
            llm_adjustment_result = await self._perform_llm_balancing_analysis(
                word_count_analysis, analysis, run_id
            )
            
            # Step 3: Programmatic adjustment validation
            programmatic_adjustment_result = await self._perform_programmatic_adjustment_analysis(
                word_count_analysis, analysis, run_id
            )
            
            # Step 4: Consolidate adjustment recommendations
            consolidated_adjustments = self._consolidate_adjustment_recommendations(
                llm_adjustment_result, programmatic_adjustment_result, run_id
            )
            
            # Step 5: Apply adaptive adjustments
            adjustment_application_result = await self._apply_adaptive_adjustments(
                generated_articles, consolidated_adjustments, run_id
            )
            
            # Step 6: Create final adjustment result
            final_adjustment_result = {
                **consolidated_adjustments,
                "adjustment_application": adjustment_application_result,
                "adjustment_id": f"adjustment_{run_id}_{int(datetime.utcnow().timestamp())}",
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
            
            adjustments_applied = adjustment_application_result.get('total_adjustments', 0)
            granularity_status = consolidated_adjustments.get('granularity_check', 'unknown')
            
            print(f"✅ V2 ADAPTIVE ADJUSTMENT: Balancing complete - {adjustments_applied} adjustments applied, granularity: {granularity_status} - run {run_id} - engine=v2")
            return final_adjustment_result
            
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error in adaptive adjustment - {e} - run {run_id} - engine=v2")
            return self._create_adjustment_result("error", run_id, {"error": str(e)})
    
    async def _analyze_word_counts(self, generated_articles: list, run_id: str) -> dict:
        """Analyze word counts for articles and sections"""
        try:
            print(f"📊 V2 ADAPTIVE ADJUSTMENT: Analyzing word counts - run {run_id} - engine=v2")
            
            word_count_analysis = {
                "articles": [],
                "total_articles": len(generated_articles),
                "total_word_count": 0
            }
            
            for generated_article in generated_articles:
                article_id = generated_article.get('article_id', 'unknown')
                article_data = generated_article.get('article_data', {})
                html_content = article_data.get('html', '')
                
                # Extract text content and count words
                article_word_count = self._count_words_in_html(html_content)
                word_count_analysis['total_word_count'] += article_word_count
                
                # Analyze sections within the article
                sections_analysis = self._analyze_sections_word_count(html_content)
                
                article_analysis = {
                    "article_id": article_id,
                    "word_count": article_word_count,
                    "sections": sections_analysis,
                    "length_status": self._classify_article_length(article_word_count),
                    "needs_adjustment": self._needs_length_adjustment(article_word_count, sections_analysis)
                }
                
                word_count_analysis['articles'].append(article_analysis)
            
            # Calculate averages and statistics
            avg_word_count = word_count_analysis['total_word_count'] / len(generated_articles) if generated_articles else 0
            word_count_analysis['average_word_count'] = avg_word_count
            word_count_analysis['articles_needing_merge'] = len([a for a in word_count_analysis['articles'] if a['word_count'] < self.word_count_thresholds['min_article_length']])
            word_count_analysis['sections_needing_split'] = sum(len([s for s in a['sections'] if s['word_count'] > self.word_count_thresholds['max_section_length']]) for a in word_count_analysis['articles'])
            
            print(f"📊 V2 ADAPTIVE ADJUSTMENT: Analysis complete - Avg: {avg_word_count:.0f} words, {word_count_analysis['articles_needing_merge']} articles need merge, {word_count_analysis['sections_needing_split']} sections need split - run {run_id} - engine=v2")
            return word_count_analysis
            
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error analyzing word counts - {e} - run {run_id} - engine=v2")
            return {
                "articles": [],
                "total_articles": 0,
                "total_word_count": 0,
                "average_word_count": 0,
                "articles_needing_merge": 0,
                "sections_needing_split": 0
            }
    
    def _count_words_in_html(self, html_content: str) -> int:
        """Count words in HTML content by extracting text"""
        try:
            from bs4 import BeautifulSoup
            import re
            
            # Parse HTML and extract text
            soup = BeautifulSoup(html_content, 'html.parser')
            text_content = soup.get_text()
            
            # Count words (split on whitespace and filter empty strings)
            words = [word for word in re.split(r'\s+', text_content.strip()) if word]
            return len(words)
            
        except Exception as e:
            # Fallback: simple word count
            import re
            text_only = re.sub(r'<[^>]+>', '', html_content)
            words = [word for word in re.split(r'\s+', text_only.strip()) if word]
            return len(words)
    
    def _analyze_sections_word_count(self, html_content: str) -> list:
        """Analyze word counts for sections within an article"""
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(html_content, 'html.parser')
            sections = []
            
            # Find all section headings (h2, h3)
            section_headings = soup.find_all(['h2', 'h3'])
            
            for i, heading in enumerate(section_headings):
                section_heading = heading.get_text().strip()
                section_id = heading.get('id', f'section_{i+1}')
                
                # Get content under this heading
                section_content = self._get_section_content_for_word_count(heading)
                section_word_count = self._count_words_in_html(section_content)
                
                sections.append({
                    "section_id": section_id,
                    "heading": section_heading,
                    "word_count": section_word_count,
                    "length_status": self._classify_section_length(section_word_count),
                    "needs_split": section_word_count > self.word_count_thresholds['max_section_length']
                })
            
            return sections
            
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error analyzing sections - {e}")
            return []
    
    def _get_section_content_for_word_count(self, heading_tag) -> str:
        """Get content under a heading for word counting"""
        try:
            content_parts = []
            next_element = heading_tag.next_sibling
            
            while next_element:
                if hasattr(next_element, 'name'):
                    if next_element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                        break
                    content_parts.append(str(next_element))
                elif hasattr(next_element, 'string') and next_element.string:
                    content_parts.append(next_element.string)
                
                next_element = next_element.next_sibling
            
            return ''.join(content_parts)
            
        except Exception as e:
            return ""
    
    def _classify_article_length(self, word_count: int) -> str:
        """Classify article length status"""
        if word_count < self.word_count_thresholds['min_article_length']:
            return "too_short"
        elif word_count > self.word_count_thresholds['optimal_article_range'][1]:
            return "too_long"
        elif self.word_count_thresholds['optimal_article_range'][0] <= word_count <= self.word_count_thresholds['optimal_article_range'][1]:
            return "optimal"
        else:
            return "acceptable"
    
    def _classify_section_length(self, word_count: int) -> str:
        """Classify section length status"""
        if word_count > self.word_count_thresholds['max_section_length']:
            return "too_long"
        elif self.word_count_thresholds['optimal_section_range'][0] <= word_count <= self.word_count_thresholds['optimal_section_range'][1]:
            return "optimal"
        else:
            return "acceptable"
    
    def _needs_length_adjustment(self, article_word_count: int, sections_analysis: list) -> bool:
        """Determine if article needs length adjustment"""
        # Check if article is too short
        if article_word_count < self.word_count_thresholds['min_article_length']:
            return True
        
        # Check if any section is too long
        for section in sections_analysis:
            if section['word_count'] > self.word_count_thresholds['max_section_length']:
                return True
        
        return False
    
    async def _perform_llm_balancing_analysis(self, word_count_analysis: dict, analysis: dict, run_id: str) -> dict:
        """V2 Engine: LLM-based balancing analysis"""
        try:
            print(f"🤖 V2 ADAPTIVE ADJUSTMENT: Starting LLM balancing analysis - run {run_id} - engine=v2")
            
            granularity = analysis.get('granularity', 'moderate')
            
            # Create system message
            system_message = """You are a balancing agent. Suggest merges/splits to optimize readability.

Your task is to analyze article and section word counts and propose optimizations:
1. Merge suggestions: Articles under 300 words should be merged with neighbors
2. Split suggestions: Sections over 1200 words should be split into smaller sections
3. Granularity check: Ensure final article count aligns with granularity expectations

Consider readability and user experience when making recommendations.

Return ONLY JSON in the exact format specified."""

            # Create user message with word count data
            user_message = f"""Given word counts per article/section, propose merges/splits and confirm granularity.

WORD COUNT ANALYSIS:
{json.dumps(word_count_analysis, indent=2)}

CURRENT GRANULARITY: {granularity}

GRANULARITY EXPECTATIONS:
- shallow: 1-3 articles
- moderate: 2-8 articles  
- deep: 5-20 articles

Analyze the word counts and return ONLY JSON in this exact format:
{{
  "merge_suggestions": [
    {{
      "article_id": "a4",
      "merge_with": "a3",
      "reason": "Article too short (250 words)",
      "combined_word_count": 850
    }}
  ],
  "split_suggestions": [
    {{
      "article_id": "a2",
      "section": "Advanced Configuration",
      "reason": "Section too long (1450 words)",
      "suggested_split_points": ["Basic Settings", "Advanced Settings"]
    }}
  ],
  "granularity_check": "moderate",
  "granularity_alignment": "aligned",
  "optimization_priority": "readability"
}}"""

            # Call LLM for balancing analysis
            print(f"🤖 V2 ADAPTIVE ADJUSTMENT: Sending balancing analysis request to LLM - run {run_id} - engine=v2")
            ai_response = await call_llm_with_fallback(system_message, user_message)
            
            if ai_response:
                # Parse JSON response
                import re
                json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
                if json_match:
                    balancing_data = json.loads(json_match.group(0))
                    
                    # Validate required fields
                    required_fields = ['merge_suggestions', 'split_suggestions', 'granularity_check']
                    if all(field in balancing_data for field in required_fields):
                        merge_count = len(balancing_data.get('merge_suggestions', []))
                        split_count = len(balancing_data.get('split_suggestions', []))
                        granularity_check = balancing_data.get('granularity_check', 'unknown')
                        
                        print(f"🤖 V2 ADAPTIVE ADJUSTMENT: LLM suggests {merge_count} merges, {split_count} splits, granularity: {granularity_check} - run {run_id} - engine=v2")
                        return balancing_data
                    else:
                        print(f"⚠️ V2 ADAPTIVE ADJUSTMENT: Missing fields in LLM response - run {run_id} - engine=v2")
                        return self._create_fallback_balancing_analysis(word_count_analysis, analysis)
                else:
                    print(f"⚠️ V2 ADAPTIVE ADJUSTMENT: No JSON found in LLM response - run {run_id} - engine=v2")
                    return self._create_fallback_balancing_analysis(word_count_analysis, analysis)
            else:
                print(f"❌ V2 ADAPTIVE ADJUSTMENT: No LLM response for balancing analysis - run {run_id} - engine=v2")
                return self._create_fallback_balancing_analysis(word_count_analysis, analysis)
                
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error in LLM balancing analysis - {e} - run {run_id} - engine=v2")
            return self._create_fallback_balancing_analysis(word_count_analysis, analysis)
    
    def _create_fallback_balancing_analysis(self, word_count_analysis: dict, analysis: dict) -> dict:
        """Create fallback balancing analysis using programmatic rules"""
        try:
            articles = word_count_analysis.get('articles', [])
            granularity = analysis.get('granularity', 'moderate')
            
            # Generate merge suggestions for short articles
            merge_suggestions = []
            for i, article in enumerate(articles):
                if article['word_count'] < self.word_count_thresholds['min_article_length']:
                    # Find a neighbor to merge with (prefer next article, then previous)
                    merge_with = None
                    if i + 1 < len(articles):
                        merge_with = articles[i + 1]['article_id']
                    elif i > 0:
                        merge_with = articles[i - 1]['article_id']
                    
                    if merge_with:
                        merge_suggestions.append({
                            "article_id": article['article_id'],
                            "merge_with": merge_with,
                            "reason": f"Article too short ({article['word_count']} words)",
                            "combined_word_count": article['word_count'] + 400  # Estimated
                        })
            
            # Generate split suggestions for long sections
            split_suggestions = []
            for article in articles:
                for section in article.get('sections', []):
                    if section['word_count'] > self.word_count_thresholds['max_section_length']:
                        split_suggestions.append({
                            "article_id": article['article_id'],
                            "section": section['heading'],
                            "reason": f"Section too long ({section['word_count']} words)",
                            "suggested_split_points": ["Part 1", "Part 2"]
                        })
            
            # Check granularity alignment
            current_article_count = len(articles)
            expectations = self.granularity_expectations.get(granularity, {"min_articles": 1, "max_articles": 10})
            
            if expectations['min_articles'] <= current_article_count <= expectations['max_articles']:
                granularity_alignment = "aligned"
            else:
                granularity_alignment = "needs_adjustment"
            
            print(f"🔧 V2 ADAPTIVE ADJUSTMENT: Fallback analysis - {len(merge_suggestions)} merges, {len(split_suggestions)} splits")
            
            return {
                "merge_suggestions": merge_suggestions,
                "split_suggestions": split_suggestions,
                "granularity_check": granularity,
                "granularity_alignment": granularity_alignment,
                "optimization_priority": "readability",
                "analysis_method": "fallback_programmatic"
            }
            
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error in fallback balancing analysis - {e}")
            return {
                "merge_suggestions": [],
                "split_suggestions": [],
                "granularity_check": "unknown",
                "granularity_alignment": "unknown",
                "optimization_priority": "readability",
                "analysis_method": "error"
            }
    
    async def _perform_programmatic_adjustment_analysis(self, word_count_analysis: dict, analysis: dict, run_id: str) -> dict:
        """V2 Engine: Programmatic adjustment analysis for validation"""
        try:
            print(f"🔍 V2 ADAPTIVE ADJUSTMENT: Performing programmatic adjustment validation - run {run_id} - engine=v2")
            
            articles = word_count_analysis.get('articles', [])
            granularity = analysis.get('granularity', 'moderate')
            
            # Validate article count against granularity
            current_article_count = len(articles)
            expectations = self.granularity_expectations.get(granularity, {"min_articles": 1, "max_articles": 10})
            
            granularity_validation = {
                "current_count": current_article_count,
                "expected_range": f"{expectations['min_articles']}-{expectations['max_articles']}",
                "alignment": "aligned" if expectations['min_articles'] <= current_article_count <= expectations['max_articles'] else "out_of_range",
                "target_length_per_article": expectations.get('target_length_per_article', 1000)
            }
            
            # Calculate length distribution
            word_counts = [article['word_count'] for article in articles]
            length_distribution = {
                "min_length": min(word_counts) if word_counts else 0,
                "max_length": max(word_counts) if word_counts else 0,
                "average_length": sum(word_counts) / len(word_counts) if word_counts else 0,
                "median_length": sorted(word_counts)[len(word_counts)//2] if word_counts else 0,
                "total_length": sum(word_counts)
            }
            
            # Readability analysis
            readability_analysis = {
                "articles_too_short": len([a for a in articles if a['word_count'] < self.word_count_thresholds['min_article_length']]),
                "articles_too_long": len([a for a in articles if a['word_count'] > self.word_count_thresholds['optimal_article_range'][1]]),
                "articles_optimal": len([a for a in articles if self.word_count_thresholds['optimal_article_range'][0] <= a['word_count'] <= self.word_count_thresholds['optimal_article_range'][1]]),
                "sections_too_long": sum(len([s for s in article.get('sections', []) if s['word_count'] > self.word_count_thresholds['max_section_length']]) for article in articles),
                "readability_score": self._calculate_readability_score(articles)
            }
            
            programmatic_result = {
                "granularity_validation": granularity_validation,
                "length_distribution": length_distribution,
                "readability_analysis": readability_analysis,
                "adjustment_priority": self._determine_adjustment_priority(granularity_validation, readability_analysis),
                "analysis_method": "programmatic_validation"
            }
            
            print(f"🔍 V2 ADAPTIVE ADJUSTMENT: Programmatic validation complete - Readability score: {readability_analysis['readability_score']:.2f} - run {run_id} - engine=v2")
            return programmatic_result
            
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error in programmatic adjustment analysis - {e} - run {run_id} - engine=v2")
            return {
                "granularity_validation": {"alignment": "unknown"},
                "length_distribution": {},
                "readability_analysis": {"readability_score": 0.5},
                "adjustment_priority": "medium",
                "analysis_method": "error"
            }
    
    def _calculate_readability_score(self, articles: list) -> float:
        """Calculate readability score based on length distribution"""
        try:
            if not articles:
                return 0.0
            
            optimal_count = 0
            total_articles = len(articles)
            
            for article in articles:
                word_count = article['word_count']
                
                # Score based on article length
                if self.word_count_thresholds['optimal_article_range'][0] <= word_count <= self.word_count_thresholds['optimal_article_range'][1]:
                    optimal_count += 1
                
                # Penalty for sections that are too long
                long_sections = len([s for s in article.get('sections', []) if s['word_count'] > self.word_count_thresholds['max_section_length']])
                if long_sections > 0:
                    optimal_count -= 0.2 * long_sections  # Penalty for long sections
            
            # Calculate score (0.0 to 1.0)
            readability_score = max(0.0, min(1.0, optimal_count / total_articles))
            return readability_score
            
        except Exception as e:
            return 0.5  # Default moderate readability
    
    def _determine_adjustment_priority(self, granularity_validation: dict, readability_analysis: dict) -> str:
        """Determine adjustment priority based on analysis"""
        try:
            granularity_aligned = granularity_validation.get('alignment') == 'aligned'
            readability_score = readability_analysis.get('readability_score', 0.5)
            articles_too_short = readability_analysis.get('articles_too_short', 0)
            sections_too_long = readability_analysis.get('sections_too_long', 0)
            
            # High priority: Major readability issues or granularity misalignment
            if not granularity_aligned or readability_score < 0.3 or articles_too_short > 2 or sections_too_long > 3:
                return "high"
            
            # Medium priority: Moderate issues
            elif readability_score < 0.7 or articles_too_short > 0 or sections_too_long > 0:
                return "medium"
            
            # Low priority: Minor optimizations
            else:
                return "low"
                
        except Exception as e:
            return "medium"  # Default priority
    
    def _consolidate_adjustment_recommendations(self, llm_result: dict, programmatic_result: dict, run_id: str) -> dict:
        """Consolidate LLM and programmatic adjustment recommendations"""
        try:
            # Base consolidation on LLM results with programmatic validation
            consolidated = {
                "merge_suggestions": llm_result.get('merge_suggestions', []),
                "split_suggestions": llm_result.get('split_suggestions', []),
                "granularity_check": llm_result.get('granularity_check', 'moderate'),
                "run_id": run_id,
                "analysis_methods": []
            }
            
            # Add programmatic validation insights
            granularity_validation = programmatic_result.get('granularity_validation', {})
            readability_analysis = programmatic_result.get('readability_analysis', {})
            
            consolidated['granularity_alignment'] = granularity_validation.get('alignment', 'unknown')
            consolidated['readability_score'] = readability_analysis.get('readability_score', 0.5)
            consolidated['adjustment_priority'] = programmatic_result.get('adjustment_priority', 'medium')
            
            # Add analysis methods used
            if llm_result.get('analysis_method') != 'error':
                consolidated['analysis_methods'].append('llm_balancing_analysis')
            if programmatic_result.get('analysis_method') != 'error':
                consolidated['analysis_methods'].append('programmatic_validation')
            
            # Additional programmatic insights
            consolidated['length_distribution'] = programmatic_result.get('length_distribution', {})
            consolidated['readability_analysis'] = readability_analysis
            
            # Summary statistics
            consolidated['adjustment_summary'] = {
                "total_merge_suggestions": len(consolidated['merge_suggestions']),
                "total_split_suggestions": len(consolidated['split_suggestions']),
                "adjustment_priority": consolidated['adjustment_priority'],
                "readability_score": consolidated['readability_score'],
                "granularity_alignment": consolidated['granularity_alignment']
            }
            
            total_adjustments = len(consolidated['merge_suggestions']) + len(consolidated['split_suggestions'])
            print(f"⚖️ V2 ADAPTIVE ADJUSTMENT: Consolidated recommendations - {total_adjustments} total adjustments, priority: {consolidated['adjustment_priority']} - run {run_id} - engine=v2")
            return consolidated
            
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error consolidating adjustment recommendations - {e} - run {run_id} - engine=v2")
            return {
                "merge_suggestions": [],
                "split_suggestions": [],
                "granularity_check": "moderate",
                "adjustment_summary": {"total_adjustments": 0}
            }
    
    async def _apply_adaptive_adjustments(self, generated_articles: list, adjustment_recommendations: dict, run_id: str) -> dict:
        """V2 Engine: Apply adaptive adjustments to articles"""
        try:
            print(f"🔧 V2 ADAPTIVE ADJUSTMENT: Starting adjustment application - run {run_id} - engine=v2")
            
            adjustment_actions = []
            
            # Apply merge suggestions
            merge_suggestions = adjustment_recommendations.get('merge_suggestions', [])
            for merge_suggestion in merge_suggestions:
                action = await self._apply_merge_adjustment(generated_articles, merge_suggestion, run_id)
                adjustment_actions.append(action)
            
            # Apply split suggestions
            split_suggestions = adjustment_recommendations.get('split_suggestions', [])
            for split_suggestion in split_suggestions:
                action = await self._apply_split_adjustment(generated_articles, split_suggestion, run_id)
                adjustment_actions.append(action)
            
            adjustment_application_result = {
                "actions_applied": adjustment_actions,
                "total_adjustments": len(adjustment_actions),
                "successful_adjustments": len([a for a in adjustment_actions if a.get('status') == 'success']),
                "failed_adjustments": len([a for a in adjustment_actions if a.get('status') == 'failed']),
                "application_method": "automated_adjustment"
            }
            
            print(f"🔧 V2 ADAPTIVE ADJUSTMENT: Adjustment application complete - {adjustment_application_result['successful_adjustments']}/{adjustment_application_result['total_adjustments']} successful adjustments - run {run_id} - engine=v2")
            return adjustment_application_result
            
        except Exception as e:
            print(f"❌ V2 ADAPTIVE ADJUSTMENT: Error applying adaptive adjustments - {e} - run {run_id} - engine=v2")
            return {
                "actions_applied": [],
                "total_adjustments": 0,
                "successful_adjustments": 0,
                "failed_adjustments": 0,
                "application_method": "error"
            }
    
    async def _apply_merge_adjustment(self, generated_articles: list, merge_suggestion: dict, run_id: str) -> dict:
        """Apply merge adjustment for short articles"""
        try:
            article_id = merge_suggestion.get('article_id')
            merge_with = merge_suggestion.get('merge_with')
            
            # For now, record the merge action (full implementation would modify articles)
            action = {
                "type": "merge_adjustment",
                "article_id": article_id,
                "merge_with": merge_with,
                "action": "recorded_for_manual_review",
                "status": "success",
                "details": f"Merge suggestion recorded: {article_id} -> {merge_with}",
                "reason": merge_suggestion.get('reason', 'Article length optimization')
            }
            
            print(f"🔗 V2 ADAPTIVE ADJUSTMENT: Recorded merge adjustment - {article_id} -> {merge_with} - run {run_id} - engine=v2")
            return action
            
        except Exception as e:
            return {
                "type": "merge_adjustment",
                "action": "error",
                "status": "failed",
                "error": str(e)
            }
    
    async def _apply_split_adjustment(self, generated_articles: list, split_suggestion: dict, run_id: str) -> dict:
        """Apply split adjustment for long sections"""
        try:
            article_id = split_suggestion.get('article_id')
            section = split_suggestion.get('section')
            
            # For now, record the split action (full implementation would modify articles)
            action = {
                "type": "split_adjustment",
                "article_id": article_id,
                "section": section,
                "action": "recorded_for_manual_review",
                "status": "success",
                "details": f"Split suggestion recorded: {article_id} section '{section}'",
                "reason": split_suggestion.get('reason', 'Section length optimization'),
                "suggested_split_points": split_suggestion.get('suggested_split_points', [])
            }
            
            print(f"✂️ V2 ADAPTIVE ADJUSTMENT: Recorded split adjustment - {article_id}:'{section}' - run {run_id} - engine=v2")
            return action
            
        except Exception as e:
            return {
                "type": "split_adjustment",
                "action": "error",
                "status": "failed",
                "error": str(e)
            }
    
    def _create_adjustment_result(self, status: str, run_id: str, additional_data: dict) -> dict:
        """Create a standard adjustment result structure"""
        return {
            "adjustment_id": f"adjustment_{run_id}_{int(datetime.utcnow().timestamp())}",
            "run_id": run_id,
            "adjustment_status": status,
            "timestamp": datetime.utcnow().isoformat(),
            "engine": "v2",
            **additional_data
        }

# Global V2 Adaptive Adjustment System instance
v2_adaptive_adjustment_system = V2AdaptiveAdjustmentSystem()

# ========================================
# V2 ENGINE: PUBLISHING SYSTEM
# ========================================

class V2PublishingSystem:
    """V2 Engine: Publishing system for persisting finalized V2 content as single source of truth"""
    
    def __init__(self):
        self.required_fields = [
            'html', 'markdown', 'toc', 'faq', 'related_links', 
            'provenance_map', 'metrics', 'media_references'
        ]
        
        self.coverage_requirement = 100.0  # 100% coverage required for publishing
        
        self.v2_publishing_metadata = {
            "engine": "v2",
            "publishing_version": "1.0",
            "content_source": "v2_only",
            "quality_assured": True
        }
    
    async def publish_v2_content(self, articles: list, generated_articles_result: dict, 
                               validation_result: dict, qa_result: dict, 
                               adjustment_result: dict, run_id: str) -> dict:
        """V2 Engine: Publish finalized V2 content to content library"""
        try:
            print(f"📚 V2 PUBLISHING: Starting V2-only content publishing - run {run_id} - engine=v2")
            
            # KE-PR7: Check publish gate for P0 issues
            qa_report = validation_result.get('qa_report')
            if qa_report:
                # Import publish gate function
                try:
                    from app.engine.v2.validators import is_publishable
                    
                    is_pub, pub_message = is_publishable(qa_report)
                    if not is_pub:
                        print(f"🚫 V2 PUBLISHING: Publishing blocked - {pub_message} - run {run_id} - engine=v2")
                        return self._create_publishing_result("blocked", run_id, {
                            "block_reason": pub_message,
                            "qa_issues": len(qa_report.flags),
                            "p0_flags": len([f for f in qa_report.flags if f.severity == "P0"]),
                            "coverage_percent": qa_report.coverage_percent
                        })
                except ImportError as e:
                    print(f"⚠️ V2 PUBLISHING: Could not import publish gate - {e} - run {run_id} - engine=v2")
            
            if not articles:
                print(f"⚠️ V2 PUBLISHING: No articles to publish - run {run_id} - engine=v2")
                return self._create_publishing_result("no_articles", run_id, {"article_count": 0})
            
            # Step 1: Validate V2-only content
            v2_validation_result = await self._validate_v2_only_content(articles, validation_result, run_id)
            if not v2_validation_result['is_valid']:
                return self._create_publishing_result("validation_failed", run_id, v2_validation_result)
            
            # Step 2: Verify 100% coverage requirement
            coverage_verification_result = await self._verify_coverage_requirement(validation_result, run_id)
            if not coverage_verification_result['meets_requirement']:
                return self._create_publishing_result("coverage_insufficient", run_id, coverage_verification_result)
            
            # Step 3: Prepare comprehensive content library structure
            content_library_articles = await self._prepare_content_library_structure(
                articles, generated_articles_result, validation_result, qa_result, adjustment_result, run_id
            )
            
            # Step 4: Persist to content library
            publishing_results = await self._persist_to_content_library(content_library_articles, run_id)
            
            # Step 5: Create comprehensive publishing result
            final_publishing_result = {
                "publishing_id": f"publishing_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "publishing_status": "success",
                "published_articles": len(publishing_results.get('published_articles', [])),
                "coverage_achieved": coverage_verification_result.get('coverage_percent', 0),
                "v2_validation": v2_validation_result,
                "coverage_verification": coverage_verification_result,
                "publishing_results": publishing_results,
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
            
            published_count = final_publishing_result['published_articles']
            coverage = final_publishing_result['coverage_achieved']
            
            print(f"✅ V2 PUBLISHING: Publishing complete - {published_count} articles published with {coverage}% coverage - run {run_id} - engine=v2")
            return final_publishing_result
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error in V2 content publishing - {e} - run {run_id} - engine=v2")
            return self._create_publishing_result("error", run_id, {"error": str(e)})
    
    async def _validate_v2_only_content(self, articles: list, validation_result: dict, run_id: str) -> dict:
        """Validate that content is V2-only with no v1 contamination"""
        try:
            print(f"🔍 V2 PUBLISHING: Validating V2-only content - run {run_id} - engine=v2")
            
            validation_issues = []
            v2_articles_count = 0
            
            for article in articles:
                metadata = article.get('metadata', {})
                engine = metadata.get('engine', 'unknown')
                processing_version = metadata.get('processing_version', 'unknown')
                generated_by = metadata.get('generated_by', 'unknown')
                
                # Check for V2 engine processing
                if engine != 'v2':
                    validation_issues.append(f"Article {article.get('id', 'unknown')} not processed by V2 engine (engine: {engine})")
                
                # Check for V2 processing version
                if processing_version != '2.0':
                    validation_issues.append(f"Article {article.get('id', 'unknown')} not using V2 processing version (version: {processing_version})")
                    
                # Check for V2 article generator
                if 'v2_article_generator' not in generated_by:
                    validation_issues.append(f"Article {article.get('id', 'unknown')} not generated by V2ArticleGenerator (generated_by: {generated_by})")
                
                if engine == 'v2' and processing_version == '2.0':
                    v2_articles_count += 1
            
            # Check validation result origin
            validation_engine = validation_result.get('engine', 'unknown')
            if validation_engine != 'v2':
                validation_issues.append(f"Validation not performed by V2 engine (validation_engine: {validation_engine})")
            
            is_valid = len(validation_issues) == 0
            validation_success_rate = (v2_articles_count / len(articles)) * 100 if articles else 0
            
            v2_validation_result = {
                "is_valid": is_valid,
                "v2_articles_count": v2_articles_count,
                "total_articles": len(articles),
                "validation_success_rate": validation_success_rate,
                "validation_issues": validation_issues,
                "v2_only_compliance": validation_success_rate == 100.0
            }
            
            print(f"🔍 V2 PUBLISHING: V2 validation complete - {validation_success_rate:.1f}% V2 compliance, {len(validation_issues)} issues - run {run_id} - engine=v2")
            return v2_validation_result
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error validating V2-only content - {e} - run {run_id} - engine=v2")
            return {
                "is_valid": False,
                "validation_issues": [f"Validation error: {str(e)}"],
                "v2_only_compliance": False
            }
    
    async def _verify_coverage_requirement(self, validation_result: dict, run_id: str) -> dict:
        """Verify 100% coverage requirement for publishing"""
        try:
            print(f"📊 V2 PUBLISHING: Verifying 100% coverage requirement - run {run_id} - engine=v2")
            
            # Extract coverage from validation result
            fidelity_coverage = validation_result.get('fidelity_coverage', {})
            coverage_percent = fidelity_coverage.get('coverage_percent', 0.0)
            
            # Alternative: check summary scores
            if coverage_percent == 0.0:
                summary_scores = validation_result.get('summary_scores', {})
                coverage_percent = summary_scores.get('coverage_percent', 0.0)
            
            meets_requirement = coverage_percent >= self.coverage_requirement
            coverage_gap = self.coverage_requirement - coverage_percent if not meets_requirement else 0.0
            
            uncovered_blocks = fidelity_coverage.get('uncovered_blocks', [])
            
            coverage_verification = {
                "meets_requirement": meets_requirement,
                "coverage_percent": coverage_percent,
                "required_coverage": self.coverage_requirement,
                "coverage_gap": coverage_gap,
                "uncovered_blocks_count": len(uncovered_blocks),
                "uncovered_blocks": uncovered_blocks[:5],  # Show first 5 uncovered blocks
                "coverage_status": "sufficient" if meets_requirement else "insufficient"
            }
            
            status_message = f"Coverage: {coverage_percent}% ({'✅ SUFFICIENT' if meets_requirement else '❌ INSUFFICIENT'})"
            print(f"📊 V2 PUBLISHING: {status_message} - run {run_id} - engine=v2")
            
            return coverage_verification
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error verifying coverage requirement - {e} - run {run_id} - engine=v2")
            return {
                "meets_requirement": False,
                "coverage_percent": 0.0,
                "coverage_status": "error",
                "error": str(e)
            }
    
    async def _prepare_content_library_structure(self, articles: list, generated_articles_result: dict,
                                               validation_result: dict, qa_result: dict, 
                                               adjustment_result: dict, run_id: str) -> list:
        """Prepare comprehensive content library structure for V2 articles"""
        try:
            print(f"🏗️ V2 PUBLISHING: Preparing content library structure - run {run_id} - engine=v2")
            
            content_library_articles = []
            generated_articles = generated_articles_result.get('generated_articles', [])
            
            for i, article in enumerate(articles):
                # Find corresponding generated article for additional metadata
                generated_article = None
                if i < len(generated_articles):
                    generated_article = generated_articles[i]
                
                # Create comprehensive content library article
                content_library_article = await self._create_content_library_article(
                    article, generated_article, validation_result, qa_result, adjustment_result, run_id
                )
                
                content_library_articles.append(content_library_article)
            
            print(f"🏗️ V2 PUBLISHING: Content library structure prepared - {len(content_library_articles)} articles - run {run_id} - engine=v2")
            return content_library_articles
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error preparing content library structure - {e} - run {run_id} - engine=v2")
            return []
    
    async def _derive_markdown_from_html(self, final_html: str) -> str:
        """TICKET 1 FIX: Derive Markdown from HTML at publish time only"""
        try:
            import markdownify
            # Convert HTML to markdown preserving code fences, tables, and structure
            md = markdownify.markdownify(final_html, heading_style="ATX", bullets="-")
            return md.strip()
        except ImportError:
            # Fallback: Simple HTML to Markdown conversion
            import re
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(final_html, 'html.parser')
            
            # Convert headings
            for i in range(6, 0, -1):
                for heading in soup.find_all(f'h{i}'):
                    heading.replace_with(f"{'#' * i} {heading.get_text()}\n\n")
            
            # Convert paragraphs
            for p in soup.find_all('p'):
                p.replace_with(f"{p.get_text()}\n\n")
            
            # Convert lists
            for ul in soup.find_all('ul'):
                for li in ul.find_all('li'):
                    li.replace_with(f"- {li.get_text()}\n")
                ul.replace_with(f"{ul.get_text()}\n")
            
            for ol in soup.find_all('ol'):
                for i, li in enumerate(ol.find_all('li'), 1):
                    li.replace_with(f"{i}. {li.get_text()}\n")
                ol.replace_with(f"{ol.get_text()}\n")
            
            # Convert code blocks
            for pre in soup.find_all('pre'):
                code = pre.find('code')
                if code:
                    language = code.get('class')
                    if language and 'language-' in str(language):
                        lang = str(language).split('language-')[1].split()[0]
                        pre.replace_with(f"```{lang}\n{code.get_text()}\n```\n\n")
                    else:
                        pre.replace_with(f"```\n{code.get_text()}\n```\n\n")
                else:
                    pre.replace_with(f"```\n{pre.get_text()}\n```\n\n")
            
            return soup.get_text()
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error converting HTML to Markdown - {e}")
            return final_html  # Return HTML as fallback
    
    def _generate_fallback_doc_uid(self) -> str:
        """TICKET 3: Generate fallback doc_uid for existing articles"""
        import time
        import random
        import string
        
        timestamp = int(time.time() * 1000)
        random_part = ''.join(random.choices(string.ascii_uppercase + string.digits, k=16))
        timestamp_b36 = format(timestamp, 'x').upper()[-8:]
        return f"01JZ{timestamp_b36}{random_part[:8]}"
    
    def _generate_fallback_doc_slug(self, title: str) -> str:
        """TICKET 3: Generate fallback doc_slug from title"""
        import re, unicodedata
        
        norm = unicodedata.normalize("NFKD", title).encode("ascii", "ignore").decode("ascii")
        slug = re.sub(r"\s+", "-", norm.lower())
        slug = re.sub(r"[^a-z0-9-]", "", slug)
        slug = re.sub(r"-{2,}", "-", slug).strip("-")
        return slug[:80] if slug else "untitled-document"
    
    async def _create_content_library_article(self, article: dict, generated_article: dict,
                                            validation_result: dict, qa_result: dict,
                                            adjustment_result: dict, run_id: str) -> dict:
        """Create comprehensive content library article with all V2 metadata"""
        try:
            article_id = article.get('id', str(uuid.uuid4()))
            
            # Extract core content
            html_content = article.get('content', '')
            # TICKET 1 FIX: Generate markdown at publish time only
            markdown_content = await self._derive_markdown_from_html(html_content)
            
            # TICKET 3: Extract bookmark registry data
            headings_registry = article.get('headings_registry', [])
            doc_uid = article.get('doc_uid')
            doc_slug = article.get('doc_slug')
            
            # Generate doc_uid if missing (for existing articles)
            if not doc_uid:
                doc_uid = self._generate_fallback_doc_uid()
                print(f"📖 TICKET 3: Generated fallback doc_uid: {doc_uid}")
            
            # Generate doc_slug if missing
            if not doc_slug and article.get('title'):
                doc_slug = self._generate_fallback_doc_slug(article['title'])
                print(f"🏷️ TICKET 3: Generated fallback doc_slug: {doc_slug}")
            
            # Generate TOC with anchors
            toc = await self._generate_toc_with_anchors(html_content)
            
            # Extract FAQs from content  
            faq = await self._extract_faq_structure(html_content)
            
            # Extract related links
            related_links = await self._extract_related_links(html_content)
            
            # Create provenance map
            provenance_map = await self._create_provenance_map(article, generated_article)
            
            # Compile comprehensive metrics
            metrics = await self._compile_comprehensive_metrics(
                article, validation_result, qa_result, adjustment_result
            )
            
            # Extract media references (no embedding)
            media_references = await self._extract_media_references(html_content, article)
            
            # Create comprehensive content library article
            content_library_article = {
                "id": article_id,
                "title": article.get('title', 'Generated Article'),
                "summary": article.get('summary', ''),
                
                # Core content (required)
                "content": html_content,  # ADD MISSING CONTENT FIELD
                "html": html_content,
                "markdown": markdown_content,
                "toc": toc,
                "faq": faq,
                "related_links": related_links,
                "provenance_map": provenance_map,
                "metrics": metrics,
                "media_references": media_references,
                
                # Publishing metadata
                "status": "published",
                "published_at": datetime.utcnow().isoformat(),
                "created_at": article.get('created_at', datetime.utcnow().isoformat()),
                "updated_at": datetime.utcnow().isoformat(),
                
                # V2 Engine metadata
                "engine": "v2",
                "publishing_version": "1.0",
                "content_source": "v2_only",
                "quality_assured": True,
                
                # Processing pipeline metadata
                "processing_metadata": {
                    "run_id": run_id,
                    "processing_version": "2.0",
                    "generated_by": article.get('metadata', {}).get('generated_by', 'v2_article_generator'),
                    "validation_status": article.get('validation_status', 'unknown'),
                    "qa_status": article.get('qa_status', 'unknown'),
                    "adjustment_status": article.get('adjustment_status', 'unknown'),
                    "readability_score": article.get('readability_score', 0.5)
                },
                
                # Legacy compatibility (minimal)
                "tags": article.get('tags', []),
                "takeaways": article.get('takeaways', [])
            }
            
            # KE-PR9: Ensure TICKET-3 fields are present before publishing
            content_library_article = ensure_ticket3_fields(content_library_article)
            
            return content_library_article
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error creating content library article - {e}")
            return {
                "id": article.get('id', str(uuid.uuid4())),
                "title": article.get('title', 'Error Article'),
                "content": article.get('content', ''),  # ADD CONTENT FIELD FOR ERROR CASE TOO
                "html": article.get('content', ''),
                "error": str(e),
                "engine": "v2"
            }
    
    async def _generate_toc_with_anchors(self, html_content: str) -> list:
        """Generate table of contents with anchor links"""
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(html_content, 'html.parser')
            toc = []
            
            # Find all headings (h1, h2, h3, h4)
            headings = soup.find_all(['h1', 'h2', 'h3', 'h4'])
            
            for heading in headings:
                heading_text = heading.get_text().strip()
                heading_id = heading.get('id', '')
                
                # Generate ID if not present
                if not heading_id:
                    heading_id = heading_text.lower().replace(' ', '-').replace('?', '').replace('!', '')
                    heading['id'] = heading_id
                
                toc_entry = {
                    "level": int(heading.name[1]),  # h1 -> 1, h2 -> 2, etc.
                    "title": heading_text,
                    "anchor": f"#{heading_id}",
                    "id": heading_id
                }
                
                toc.append(toc_entry)
            
            return toc
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error generating TOC - {e}")
            return []
    
    async def _extract_faq_structure(self, html_content: str) -> list:
        """Extract FAQ structure from HTML content"""
        try:
            import re
            
            faqs = []
            
            # Pattern 1: Q: ... A: ...
            qa_pattern = r'Q:\s*(.*?)\s*A:\s*(.*?)(?=Q:|$)'
            matches = re.finditer(qa_pattern, html_content, re.IGNORECASE | re.DOTALL)
            
            for match in matches:
                question = match.group(1).strip()
                answer = match.group(2).strip()
                
                if question and answer:
                    faqs.append({
                        "question": question,
                        "answer": answer,
                        "id": f"faq_{len(faqs) + 1}"
                    })
            
            # Pattern 2: HTML structure (h3 with Q:, followed by p with A:)
            if not faqs:  # Only if no Q&A pattern found
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Look for FAQ sections
                faq_headings = soup.find_all(['h2', 'h3'], string=re.compile(r'FAQ|Questions?', re.I))
                
                for faq_heading in faq_headings:
                    # Find questions after this heading
                    current_element = faq_heading.next_sibling
                    
                    while current_element and current_element.name not in ['h1', 'h2']:
                        if hasattr(current_element, 'get_text'):
                            text = current_element.get_text().strip()
                            if text.startswith('Q:') or text.startswith('Question:'):
                                question = text.replace('Q:', '').replace('Question:', '').strip()
                                
                                # Look for answer in next element
                                answer_element = current_element.next_sibling
                                if answer_element and hasattr(answer_element, 'get_text'):
                                    answer_text = answer_element.get_text().strip()
                                    if answer_text.startswith('A:') or answer_text.startswith('Answer:'):
                                        answer = answer_text.replace('A:', '').replace('Answer:', '').strip()
                                        
                                        faqs.append({
                                            "question": question,
                                            "answer": answer,
                                            "id": f"faq_{len(faqs) + 1}"
                                        })
                        
                        current_element = current_element.next_sibling
            
            return faqs
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error extracting FAQ structure - {e}")
            return []
    
    async def _extract_related_links(self, html_content: str) -> list:
        """Extract related links from HTML content"""
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(html_content, 'html.parser')
            related_links = []
            
            # Find all links
            links = soup.find_all('a', href=True)
            
            for link in links:
                href = link.get('href', '')
                link_text = link.get_text().strip()
                
                if href and link_text:
                    link_entry = {
                        "title": link_text,
                        "url": href,
                        "type": "internal" if href.startswith('#') or href.startswith('/') else "external",
                        "description": link_text  # Use link text as description
                    }
                    
                    # Avoid duplicates
                    if not any(existing['url'] == href for existing in related_links):
                        related_links.append(link_entry)
            
            return related_links
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error extracting related links - {e}")
            return []
    
    async def _create_provenance_map(self, article: dict, generated_article: dict) -> dict:
        """Create provenance map linking article/sections to source blocks"""
        try:
            provenance_map = {
                "article_id": article.get('id', 'unknown'),
                "source_mapping": {},
                "coverage_summary": {}
            }
            
            if generated_article:
                article_data = generated_article.get('article_data', {})
                validation_metadata = article_data.get('validation_metadata', {})
                
                # Map assigned block IDs
                assigned_blocks = validation_metadata.get('block_ids_assigned', [])
                if assigned_blocks:
                    provenance_map["source_mapping"]["assigned_blocks"] = assigned_blocks
                    provenance_map["coverage_summary"]["total_blocks_assigned"] = len(assigned_blocks)
                
                # Add source document information
                article_metadata = article.get('metadata', {})
                normalized_doc_id = article_metadata.get('normalized_doc_id', 'unknown')
                
                provenance_map["source_mapping"]["normalized_doc_id"] = normalized_doc_id
                provenance_map["source_mapping"]["run_id"] = article_metadata.get('run_id', 'unknown')
                
                # Add section-level mapping if available
                # This would be enhanced in a full implementation to map specific sections to source blocks
                provenance_map["source_mapping"]["mapping_method"] = "v2_article_generator"
                provenance_map["source_mapping"]["extraction_method"] = article_metadata.get('extraction_method', 'v2_processing')
            
            return provenance_map
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error creating provenance map - {e}")
            return {
                "article_id": article.get('id', 'unknown'),
                "source_mapping": {},
                "error": str(e)
            }
    
    async def _compile_comprehensive_metrics(self, article: dict, validation_result: dict,
                                           qa_result: dict, adjustment_result: dict) -> dict:
        """Compile comprehensive metrics from all V2 pipeline steps"""
        try:
            metrics = {
                "compilation_timestamp": datetime.utcnow().isoformat(),
                "metrics_version": "v2.1"
            }
            
            # Validation metrics (Step 8)
            if validation_result:
                summary_scores = validation_result.get('summary_scores', {})
                metrics["fidelity_score"] = summary_scores.get('fidelity_score', 0.0)
                metrics["coverage_percent"] = summary_scores.get('coverage_percent', 0.0)
                metrics["style_compliance"] = summary_scores.get('style_compliance', 0.0)
                metrics["placeholder_count"] = summary_scores.get('placeholder_count', 0)
                
                validation_metrics = validation_result.get('validation_metrics', {})
                metrics["redundancy_score"] = validation_metrics.get('redundancy_score', 0.0)
                metrics["granularity_alignment_score"] = validation_metrics.get('granularity_alignment_score', 0.0)
                metrics["complexity_alignment_score"] = validation_metrics.get('complexity_alignment_score', 0.0)
            
            # QA metrics (Step 9)
            if qa_result:
                qa_summary = qa_result.get('summary', {})
                metrics["qa_issues_found"] = qa_summary.get('issues_found', 0)
                metrics["duplicates_found"] = qa_summary.get('total_duplicates', 0)
                metrics["invalid_links_found"] = qa_summary.get('total_invalid_links', 0)
                metrics["duplicate_faqs_found"] = qa_summary.get('total_duplicate_faqs', 0)
                metrics["terminology_issues_found"] = qa_summary.get('total_terminology_issues', 0)
            
            # Adjustment metrics (Step 10)
            if adjustment_result:
                adjustment_summary = adjustment_result.get('adjustment_summary', {})
                metrics["readability_score"] = adjustment_result.get('readability_score', 0.5)
                metrics["total_adjustments"] = adjustment_summary.get('total_adjustments', 0)
                metrics["merge_suggestions"] = adjustment_summary.get('total_merge_suggestions', 0)
                metrics["split_suggestions"] = adjustment_summary.get('total_split_suggestions', 0)
                metrics["granularity_alignment"] = adjustment_result.get('granularity_alignment', 'unknown')
            
            # Article-specific metrics
            article_metadata = article.get('metadata', {})
            metrics["validation_status"] = article.get('validation_status', 'unknown')
            metrics["qa_status"] = article.get('qa_status', 'unknown')
            metrics["adjustment_status"] = article.get('adjustment_status', 'unknown')
            
            # Content metrics
            content_length = len(article.get('content', ''))
            metrics["content_length_chars"] = content_length
            metrics["content_length_words"] = len(article.get('content', '').split()) if article.get('content') else 0
            
            # Overall quality score (composite metric)
            quality_components = [
                metrics.get('fidelity_score', 0) * 0.3,
                (metrics.get('coverage_percent', 0) / 100) * 0.3,
                metrics.get('style_compliance', 0) * 0.2,
                metrics.get('readability_score', 0) * 0.2
            ]
            metrics["overall_quality_score"] = sum(quality_components)
            
            return metrics
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error compiling comprehensive metrics - {e}")
            return {
                "compilation_error": str(e),
                "metrics_version": "error"
            }
    
    async def _extract_media_references(self, html_content: str, article: dict) -> list:
        """Extract media references (IDs/URLs + alt-text) without embedding"""
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(html_content, 'html.parser')
            media_references = []
            
            # Find all image references (but don't embed)
            images = soup.find_all('img')
            
            for img in images:
                src = img.get('src', '')
                alt_text = img.get('alt', '')
                
                if src:
                    media_ref = {
                        "type": "image",
                        "src": src,
                        "alt_text": alt_text,
                        "reference_only": True,  # Indicates no embedding
                        "media_id": src.split('/')[-1] if '/' in src else src  # Extract media ID from URL
                    }
                    
                    media_references.append(media_ref)
            
            # Find video references
            videos = soup.find_all('video')
            for video in videos:
                src = video.get('src', '')
                if src:
                    media_ref = {
                        "type": "video",
                        "src": src,
                        "reference_only": True,
                        "media_id": src.split('/')[-1] if '/' in src else src
                    }
                    media_references.append(media_ref)
            
            # Find audio references
            audios = soup.find_all('audio')
            for audio in audios:
                src = audio.get('src', '')
                if src:
                    media_ref = {
                        "type": "audio",
                        "src": src,
                        "reference_only": True,
                        "media_id": src.split('/')[-1] if '/' in src else src
                    }
                    media_references.append(media_ref)
            
            return media_references
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error extracting media references - {e}")
            return []
    
    async def _persist_to_content_library(self, content_library_articles: list, run_id: str) -> dict:
        """Persist V2 articles to content library collection"""
        try:
            print(f"💾 V2 PUBLISHING: Persisting to content library - {len(content_library_articles)} articles - run {run_id} - engine=v2")
            
            published_articles = []
            failed_articles = []
            
            for article in content_library_articles:
                try:
                    # Insert into content library using repository pattern
                    from engine.stores.mongo import RepositoryFactory
                    content_repo = RepositoryFactory.get_content_library()
                    result_id = await content_repo.insert_article(article)
                    
                    if result.inserted_id:
                        published_articles.append({
                            "article_id": article['id'],
                            "title": article['title'],
                            "inserted_id": str(result.inserted_id),
                            "status": "published"
                        })
                        print(f"📚 V2 PUBLISHING: Published article '{article['title']}' - ID: {article['id']} - engine=v2")
                    else:
                        failed_articles.append({
                            "article_id": article['id'],
                            "title": article['title'],
                            "error": "Insert operation returned no ID"
                        })
                        
                except Exception as article_error:
                    failed_articles.append({
                        "article_id": article.get('id', 'unknown'),
                        "title": article.get('title', 'unknown'),
                        "error": str(article_error)
                    })
                    print(f"❌ V2 PUBLISHING: Failed to publish article '{article.get('title', 'unknown')}' - {article_error}")
            
            persistence_result = {
                "published_articles": published_articles,
                "failed_articles": failed_articles,
                "total_published": len(published_articles),
                "total_failed": len(failed_articles),
                "success_rate": (len(published_articles) / len(content_library_articles)) * 100 if content_library_articles else 0
            }
            
            print(f"💾 V2 PUBLISHING: Persistence complete - {len(published_articles)}/{len(content_library_articles)} articles published ({persistence_result['success_rate']:.1f}% success) - run {run_id} - engine=v2")
            return persistence_result
            
        except Exception as e:
            print(f"❌ V2 PUBLISHING: Error persisting to content library - {e} - run {run_id} - engine=v2")
            return {
                "published_articles": [],
                "failed_articles": [],
                "total_published": 0,
                "total_failed": len(content_library_articles),
                "success_rate": 0.0,
                "error": str(e)
            }
    
    def _create_publishing_result(self, status: str, run_id: str, additional_data: dict) -> dict:
        """Create a standard publishing result structure"""
        return {
            "publishing_id": f"publishing_{run_id}_{int(datetime.utcnow().timestamp())}",
            "run_id": run_id,
            "publishing_status": status,
            "timestamp": datetime.utcnow().isoformat(),
            "engine": "v2",
            **additional_data
        }

# Global V2 Publishing System instance
v2_publishing_system = V2PublishingSystem()

# ========================================
# V2 ENGINE: VERSIONING & DIFF SYSTEM
# ========================================

class V2VersioningSystem:
    """V2 Engine: Versioning and diff system for reprocessing support and version comparison"""
    
    def __init__(self):
        self.version_metadata_fields = [
            'source_hash', 'version', 'supersedes', 'version_timestamp', 'change_summary'
        ]
        
        self.diff_comparison_fields = [
            'title', 'toc', 'sections', 'faq', 'related_links', 'content_changes'
        ]
    
    async def manage_versioning(self, content: str, content_type: str, articles: list, 
                              generated_articles_result: dict, publishing_result: dict, run_id: str) -> dict:
        """V2 Engine: Manage versioning for content and articles"""
        try:
            print(f"🔄 V2 VERSIONING: Starting versioning management - run {run_id} - engine=v2")
            
            # Step 1: Calculate source hash for change detection
            source_hash = self._calculate_source_hash(content, content_type)
            
            # Step 2: Check for existing versions
            existing_version_info = await self._find_existing_versions(source_hash, content_type, run_id)
            
            # Step 3: Determine version number and supersedes relationship
            version_metadata = await self._determine_version_metadata(
                source_hash, existing_version_info, run_id
            )
            
            # Step 4: Store version metadata with articles
            versioned_articles = await self._add_version_metadata_to_articles(
                articles, version_metadata, run_id
            )
            
            # Step 5: Create version record
            version_record = await self._create_version_record(
                version_metadata, generated_articles_result, publishing_result, run_id
            )
            
            # Step 6: Store version record
            await self._store_version_record(version_record, run_id)
            
            # Step 7: Generate diff if this is an update
            diff_result = None
            if version_metadata.get('supersedes'):
                diff_result = await self._generate_version_diff(
                    version_metadata['supersedes'], run_id, versioned_articles
                )
            
            versioning_result = {
                "versioning_id": f"versioning_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "versioning_status": "success",
                "version_metadata": version_metadata,
                "version_record": version_record,
                "diff_result": diff_result,
                "versioned_articles_count": len(versioned_articles),
                "timestamp": datetime.utcnow().isoformat(),
                "engine": "v2"
            }
            
            version_number = version_metadata.get('version', 1)
            is_update = version_metadata.get('supersedes') is not None
            
            update_text = "(update)" if is_update else "(new)"
            print(f"✅ V2 VERSIONING: Versioning complete - Version {version_number} {update_text} - run {run_id} - engine=v2")
            return versioning_result
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error in versioning management - {e} - run {run_id} - engine=v2")
            return self._create_versioning_result("error", run_id, {"error": str(e)})
    
    def _calculate_source_hash(self, content: str, content_type: str) -> str:
        """Calculate hash of source content for change detection"""
        try:
            import hashlib
            
            # Normalize content for consistent hashing
            normalized_content = content.strip().lower()
            
            # Include content type in hash to distinguish different types of same content
            hash_input = f"{content_type}:{normalized_content}"
            
            # Calculate SHA-256 hash
            source_hash = hashlib.sha256(hash_input.encode('utf-8')).hexdigest()
            
            hash_preview = source_hash[:16] + "..."
            print(f"📊 V2 VERSIONING: Source hash calculated - {hash_preview} (type: {content_type})")
            return source_hash
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error calculating source hash - {e}")
            # Fallback to timestamp-based hash
            import time
            fallback_hash = hashlib.sha256(f"{content_type}:{time.time()}".encode()).hexdigest()
            return fallback_hash
    
    async def _find_existing_versions(self, source_hash: str, content_type: str, run_id: str) -> dict:
        """Find existing versions of content based on source hash"""
        try:
            hash_preview = source_hash[:16] + "..."
            print(f"🔍 V2 VERSIONING: Searching for existing versions - hash: {hash_preview} - run {run_id} - engine=v2")
            
            # Search for existing version records with same source hash
            existing_versions = []
            async for version_record in db.v2_version_records.find({"source_hash": source_hash}).sort("version", -1):
                existing_versions.append(version_record)
            
            # Also search in content library for articles with same source hash
            content_library_versions = []
            async for article in db.content_library.find({"version_metadata.source_hash": source_hash}).sort("version_metadata.version", -1):
                if article.get('engine') == 'v2':  # Only V2 articles
                    content_library_versions.append(article)
            
            existing_version_info = {
                "has_existing_versions": len(existing_versions) > 0 or len(content_library_versions) > 0,
                "version_records": existing_versions,
                "content_library_versions": content_library_versions,
                "latest_version": existing_versions[0] if existing_versions else None,
                "total_versions": len(existing_versions)
            }
            
            if existing_version_info["has_existing_versions"]:
                latest_version = existing_version_info["latest_version"]
                if latest_version:
                    version_num = latest_version.get('version', 'unknown')
                    print(f"📋 V2 VERSIONING: Found {existing_version_info['total_versions']} existing versions - latest: v{version_num} - run {run_id} - engine=v2")
                else:
                    print(f"📋 V2 VERSIONING: Found existing content in library - run {run_id} - engine=v2")
            else:
                print(f"🆕 V2 VERSIONING: No existing versions found - this is a new content version - run {run_id} - engine=v2")
            
            return existing_version_info
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error finding existing versions - {e} - run {run_id} - engine=v2")
            return {
                "has_existing_versions": False,
                "version_records": [],
                "content_library_versions": [],
                "latest_version": None,
                "total_versions": 0
            }
    
    async def _determine_version_metadata(self, source_hash: str, existing_version_info: dict, run_id: str) -> dict:
        """Determine version number and supersedes relationship"""
        try:
            version_metadata = {
                "source_hash": source_hash,
                "version_timestamp": datetime.utcnow().isoformat(),
                "run_id": run_id
            }
            
            if existing_version_info["has_existing_versions"]:
                # This is an update to existing content
                latest_version = existing_version_info["latest_version"]
                
                if latest_version:
                    # Increment version number
                    previous_version = latest_version.get('version', 0)
                    version_metadata["version"] = previous_version + 1
                    version_metadata["supersedes"] = latest_version.get('run_id', 'unknown')
                    version_metadata["change_summary"] = "Content update detected"
                else:
                    # Found content in library but no version record
                    version_metadata["version"] = 2  # Assume this is version 2
                    version_metadata["supersedes"] = "content_library_existing"
                    version_metadata["change_summary"] = "Update to existing content library article"
            else:
                # This is new content
                version_metadata["version"] = 1
                version_metadata["supersedes"] = None
                version_metadata["change_summary"] = "New content version"
            
            supersedes_text = version_metadata.get('supersedes', 'none')
            if version_metadata.get('supersedes'):
                supersedes_info = f"(supersedes: {supersedes_text})"
            else:
                supersedes_info = "(new)"
            
            print(f"📊 V2 VERSIONING: Version metadata determined - v{version_metadata['version']} {supersedes_info} - run {run_id} - engine=v2")
            return version_metadata
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error determining version metadata - {e} - run {run_id} - engine=v2")
            return {
                "source_hash": source_hash,
                "version": 1,
                "supersedes": None,
                "version_timestamp": datetime.utcnow().isoformat(),
                "change_summary": "Version metadata error",
                "run_id": run_id
            }
    
    async def _add_version_metadata_to_articles(self, articles: list, version_metadata: dict, run_id: str) -> list:
        """Add version metadata to all articles"""
        try:
            print(f"🏷️ V2 VERSIONING: Adding version metadata to {len(articles)} articles - run {run_id} - engine=v2")
            
            versioned_articles = []
            
            for article in articles:
                # Create versioned article with metadata
                versioned_article = article.copy()
                
                # Add version metadata
                versioned_article['version_metadata'] = version_metadata.copy()
                
                # Add version-specific fields to main article metadata
                if 'metadata' not in versioned_article:
                    versioned_article['metadata'] = {}
                
                versioned_article['metadata']['version'] = version_metadata['version']
                versioned_article['metadata']['source_hash'] = version_metadata['source_hash']
                versioned_article['metadata']['version_timestamp'] = version_metadata['version_timestamp']
                
                if version_metadata.get('supersedes'):
                    versioned_article['metadata']['supersedes'] = version_metadata['supersedes']
                    versioned_article['metadata']['is_update'] = True
                else:
                    versioned_article['metadata']['is_update'] = False
                
                versioned_articles.append(versioned_article)
            
            print(f"🏷️ V2 VERSIONING: Version metadata added to all articles - v{version_metadata['version']} - run {run_id} - engine=v2")
            return versioned_articles
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error adding version metadata to articles - {e} - run {run_id} - engine=v2")
            return articles  # Return original articles on error
    
    async def _create_version_record(self, version_metadata: dict, generated_articles_result: dict, 
                                   publishing_result: dict, run_id: str) -> dict:
        """Create comprehensive version record for storage"""
        try:
            version_record = {
                "version_record_id": f"version_{run_id}_{int(datetime.utcnow().timestamp())}",
                "run_id": run_id,
                "engine": "v2",
                
                # Version metadata
                **version_metadata,
                
                # Processing metadata
                "processing_metadata": {
                    "articles_generated": len(generated_articles_result.get('generated_articles', [])),
                    "publishing_status": publishing_result.get('publishing_status', 'unknown'),
                    "published_articles": publishing_result.get('published_articles', 0),
                    "coverage_achieved": publishing_result.get('coverage_achieved', 0)
                },
                
                # Version chain metadata
                "version_chain": {
                    "is_initial_version": version_metadata.get('version', 1) == 1,
                    "is_update": version_metadata.get('supersedes') is not None,
                    "previous_run_id": version_metadata.get('supersedes'),
                    "version_number": version_metadata.get('version', 1)
                },
                
                # Storage metadata
                "created_at": datetime.utcnow().isoformat(),
                "record_type": "v2_version_record"
            }
            
            return version_record
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error creating version record - {e} - run {run_id} - engine=v2")
            return {
                "version_record_id": f"version_{run_id}_error",
                "run_id": run_id,
                "engine": "v2",
                "error": str(e)
            }
    
    async def _store_version_record(self, version_record: dict, run_id: str) -> bool:
        """Store version record in database"""
        try:
            result = await db.v2_version_records.insert_one(version_record)
            
            if result.inserted_id:
                print(f"💾 V2 VERSIONING: Version record stored - ID: {result.inserted_id} - run {run_id} - engine=v2")
                return True
            else:
                print(f"❌ V2 VERSIONING: Failed to store version record - run {run_id} - engine=v2")
                return False
                
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error storing version record - {e} - run {run_id} - engine=v2")
            return False
    
    async def _generate_version_diff(self, previous_run_id: str, current_run_id: str, current_articles: list) -> dict:
        """Generate diff between current and previous version"""
        try:
            print(f"🔍 V2 VERSIONING: Generating diff - current: {current_run_id} vs previous: {previous_run_id} - engine=v2")
            
            # Step 1: Find previous version articles
            previous_articles = []
            
            # Search in content library for previous version articles
            async for article in db.content_library.find({"metadata.run_id": previous_run_id, "engine": "v2"}):
                previous_articles.append(article)
            
            # Also search in v2_version_records for previous articles
            previous_version_record = await db.v2_version_records.find_one({"run_id": previous_run_id})
            
            if not previous_articles and not previous_version_record:
                print(f"⚠️ V2 VERSIONING: No previous version data found for diff - run {current_run_id} - engine=v2")
                return {
                    "diff_status": "no_previous_version",
                    "message": f"No previous version found for run_id: {previous_run_id}",
                    "current_run_id": current_run_id,
                    "previous_run_id": previous_run_id
                }
            
            # Step 2: Compare articles and generate diff
            diff_analysis = await self._compare_article_versions(previous_articles, current_articles, previous_run_id, current_run_id)
            
            # Step 3: Create diff result structure
            diff_result = {
                "diff_id": f"diff_{current_run_id}_{int(datetime.utcnow().timestamp())}",
                "current_run_id": current_run_id,
                "previous_run_id": previous_run_id,
                "diff_status": "success",
                "comparison_timestamp": datetime.utcnow().isoformat(),
                "engine": "v2",
                
                # Article count comparison
                "article_counts": {
                    "previous": len(previous_articles),
                    "current": len(current_articles),
                    "difference": len(current_articles) - len(previous_articles)
                },
                
                # Detailed diff analysis
                **diff_analysis
            }
            
            print(f"✅ V2 VERSIONING: Diff generated - {len(diff_analysis.get('title_changes', []))} title changes, {len(diff_analysis.get('content_changes', []))} content changes - run {current_run_id} - engine=v2")
            return diff_result
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error generating version diff - {e} - run {current_run_id} - engine=v2")
            return {
                "diff_status": "error",
                "error": str(e),
                "current_run_id": current_run_id,
                "previous_run_id": previous_run_id
            }
    
    async def _compare_article_versions(self, previous_articles: list, current_articles: list, 
                                      previous_run_id: str, current_run_id: str) -> dict:
        """Compare previous and current articles to identify differences"""
        try:
            comparison_result = {
                "title_changes": [],
                "toc_changes": [],
                "section_changes": [],
                "faq_changes": [],
                "related_links_changes": [],
                "content_changes": [],
                "new_articles": [],
                "removed_articles": [],
                "unchanged_articles": []
            }
            
            # Create title-based lookup for previous articles
            previous_lookup = {}
            for prev_article in previous_articles:
                title = prev_article.get('title', '').strip().lower()
                if title:
                    previous_lookup[title] = prev_article
            
            # Compare each current article with previous version
            for curr_article in current_articles:
                current_title = curr_article.get('title', '').strip()
                current_title_lower = current_title.lower()
                
                if current_title_lower in previous_lookup:
                    # Article exists in both versions - compare content
                    prev_article = previous_lookup[current_title_lower]
                    article_diff = await self._compare_individual_articles(prev_article, curr_article)
                    
                    if article_diff['has_changes']:
                        # Add changes to respective categories
                        if article_diff.get('title_changed'):
                            comparison_result["title_changes"].append(article_diff['title_diff'])
                        if article_diff.get('toc_changed'):
                            comparison_result["toc_changes"].append(article_diff['toc_diff'])
                        if article_diff.get('sections_changed'):
                            comparison_result["section_changes"].extend(article_diff['section_diffs'])
                        if article_diff.get('faq_changed'):
                            comparison_result["faq_changes"].append(article_diff['faq_diff'])
                        if article_diff.get('related_links_changed'):
                            comparison_result["related_links_changes"].append(article_diff['related_links_diff'])
                        if article_diff.get('content_changed'):
                            comparison_result["content_changes"].append(article_diff['content_diff'])
                    else:
                        comparison_result["unchanged_articles"].append({
                            "title": current_title,
                            "article_id": curr_article.get('id', 'unknown')
                        })
                    
                    # Remove from previous lookup to track removed articles
                    del previous_lookup[current_title_lower]
                else:
                    # New article in current version
                    comparison_result["new_articles"].append({
                        "title": current_title,
                        "article_id": curr_article.get('id', 'unknown'),
                        "content_preview": self._get_content_preview(curr_article.get('content', ''))
                    })
            
            # Any remaining articles in previous_lookup are removed articles
            for removed_title, removed_article in previous_lookup.items():
                comparison_result["removed_articles"].append({
                    "title": removed_article.get('title', ''),
                    "article_id": removed_article.get('id', 'unknown'),
                    "content_preview": self._get_content_preview(removed_article.get('content', ''))
                })
            
            return comparison_result
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error comparing article versions - {e} - engine=v2")
            return {"error": str(e)}
    
    async def _compare_individual_articles(self, prev_article: dict, curr_article: dict) -> dict:
        """Compare two individual articles and identify specific changes"""
        try:
            article_diff = {
                "has_changes": False,
                "article_title": curr_article.get('title', ''),
                "article_id": curr_article.get('id', 'unknown')
            }
            
            # Compare titles
            prev_title = prev_article.get('title', '').strip()
            curr_title = curr_article.get('title', '').strip()
            
            if prev_title != curr_title:
                article_diff["title_changed"] = True
                article_diff["title_diff"] = {
                    "previous": prev_title,
                    "current": curr_title,
                    "change_type": "title_modified"
                }
                article_diff["has_changes"] = True
            
            # Compare content for significant changes
            prev_content = prev_article.get('content', '')
            curr_content = curr_article.get('content', '')
            
            content_similarity = self._calculate_content_similarity(prev_content, curr_content)
            
            if content_similarity < 0.8:  # Significant content change threshold
                article_diff["content_changed"] = True
                article_diff["content_diff"] = {
                    "similarity_score": content_similarity,
                    "change_type": "significant_content_change",
                    "previous_preview": self._get_content_preview(prev_content),
                    "current_preview": self._get_content_preview(curr_content),
                    "word_count_change": self._count_words(curr_content) - self._count_words(prev_content)
                }
                article_diff["has_changes"] = True
            
            # Extract and compare table of contents
            prev_toc = self._extract_toc_from_content(prev_content)
            curr_toc = self._extract_toc_from_content(curr_content)
            
            if prev_toc != curr_toc:
                article_diff["toc_changed"] = True
                article_diff["toc_diff"] = {
                    "previous": prev_toc,
                    "current": curr_toc,
                    "change_type": "toc_structure_change"
                }
                article_diff["has_changes"] = True
            
            # Compare sections, FAQs, and related links would require more complex parsing
            # For now, we'll use content similarity as a proxy
            
            return article_diff
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error comparing individual articles - {e} - engine=v2")
            return {"has_changes": False, "error": str(e)}
    
    def _calculate_content_similarity(self, content1: str, content2: str) -> float:
        """Calculate similarity between two content strings"""
        try:
            # Simple word-based similarity calculation
            import re
            
            # Extract words from both contents
            words1 = set(re.findall(r'\w+', content1.lower()))
            words2 = set(re.findall(r'\w+', content2.lower()))
            
            if not words1 and not words2:
                return 1.0
            if not words1 or not words2:
                return 0.0
            
            # Calculate Jaccard similarity
            intersection = len(words1 & words2)
            union = len(words1 | words2)
            
            similarity = intersection / union if union > 0 else 0.0
            return similarity
            
        except Exception:
            return 0.5  # Fallback similarity score
    
    def _get_content_preview(self, content: str, max_length: int = 200) -> str:
        """Get a preview of content for diff display"""
        try:
            import re
            # Strip HTML tags for preview
            text_content = re.sub(r'<[^>]+>', '', content).strip()
            
            if len(text_content) <= max_length:
                return text_content
            
            return text_content[:max_length] + "..."
            
        except Exception:
            return "Content preview unavailable"
    
    def _extract_toc_from_content(self, content: str) -> list:
        """Extract table of contents (headings) from HTML content"""
        try:
            import re
            
            headings = []
            # Find all headings (h1, h2, h3, etc.)
            heading_pattern = r'<h([1-6])[^>]*>(.*?)</h[1-6]>'
            matches = re.findall(heading_pattern, content, re.IGNORECASE)
            
            for level, text in matches:
                # Clean heading text
                clean_text = re.sub(r'<[^>]+>', '', text).strip()
                headings.append({
                    "level": int(level),
                    "text": clean_text
                })
            
            return headings
            
        except Exception:
            return []
    
    def _count_words(self, content: str) -> int:
        """Count words in content"""
        try:
            import re
            # Strip HTML and count words
            text_content = re.sub(r'<[^>]+>', ' ', content)
            words = re.findall(r'\w+', text_content)
            return len(words)
        except Exception:
            return 0

    def _create_versioning_result(self, status: str, run_id: str, additional_data: dict) -> dict:
        """Create a standard versioning result structure"""
        return {
            "versioning_id": f"versioning_{run_id}_{int(datetime.utcnow().timestamp())}",
            "run_id": run_id,
            "versioning_status": status,
            "timestamp": datetime.utcnow().isoformat(),
            "engine": "v2",
            **additional_data
        }

    def _analyze_version_chains(self, versioning_results: list) -> dict:
        """Analyze version chains from versioning results"""
        try:
            # Group by source hash to find version chains
            source_chains = {}
            
            for result in versioning_results:
                source_hash = result.get('version_metadata', {}).get('source_hash')
                if source_hash:
                    if source_hash not in source_chains:
                        source_chains[source_hash] = []
                    source_chains[source_hash].append(result)
            
            # Analyze chains
            chain_analysis = {
                "total_source_content": len(source_chains),
                "content_with_multiple_versions": len([chain for chain in source_chains.values() if len(chain) > 1]),
                "longest_version_chain": max([len(chain) for chain in source_chains.values()]) if source_chains else 0,
                "average_versions_per_content": sum([len(chain) for chain in source_chains.values()]) / len(source_chains) if source_chains else 0
            }
            
            return chain_analysis
            
        except Exception as e:
            print(f"❌ V2 VERSIONING: Error analyzing version chains - {e}")
            return {"error": str(e)}
    
    async def create_version_from_articles(self, articles: list, run_id: str) -> dict:
        """Create version metadata from processed articles"""
        try:
            print(f"📦 KE-PR5: Creating version from {len(articles)} articles - run {run_id}")
            
            # Generate version ID
            version_id = f"v_{run_id}_{int(datetime.utcnow().timestamp())}_{uuid.uuid4().hex[:8]}"
            
            # Calculate content hash from all articles
            content_hash = self._calculate_articles_hash(articles)
            
            # Create version metadata
            version_metadata = {
                "version_id": version_id,
                "run_id": run_id,
                "content_hash": content_hash,
                "article_count": len(articles),
                "created_at": datetime.utcnow().isoformat(),
                "engine": "v2",
                "processing_version": "2.0"
            }
            
            # Add article summaries
            article_summaries = []
            for article in articles:
                summary = {
                    "article_id": article.get('id'),
                    "title": article.get('title'),
                    "content_length": len(article.get('content', '')),
                    "article_type": article.get('article_type', 'unknown'),
                    "status": article.get('status', 'draft')
                }
                article_summaries.append(summary)
            
            version_metadata["articles"] = article_summaries
            
            version_result = {
                "version_id": version_id,
                "version_metadata": version_metadata,
                "versioning_status": "success",
                "created_articles": len(articles),
                "timestamp": datetime.utcnow().isoformat()
            }
            
            print(f"✅ KE-PR5: Version created - {version_id} with {len(articles)} articles")
            return version_result
            
        except Exception as e:
            print(f"❌ KE-PR5: Error creating version from articles - {e}")
            return {
                "version_id": f"error_{run_id}",
                "versioning_status": "error",
                "error": str(e),
                "created_articles": len(articles) if articles else 0
            }
    
    def _calculate_articles_hash(self, articles: list) -> str:
        """Calculate hash from all article contents"""
        try:
            if not articles:
                return hashlib.sha256("empty".encode()).hexdigest()[:16]
            
            # Combine all article contents and titles for hashing
            combined_content = ""
            for article in articles:
                title = article.get('title', '')
                content = article.get('content', '')
                combined_content += f"{title}:{content}|"
            
            # Create hash
            content_hash = hashlib.sha256(combined_content.encode('utf-8')).hexdigest()[:16]
            return content_hash
            
        except Exception as e:
            print(f"⚠️ Error calculating articles hash: {e}")
            return f"hash_error_{int(datetime.utcnow().timestamp())}"

# Global V2 Versioning System instance
v2_versioning_system = V2VersioningSystem()

class V2ReviewSystem:
    """V2 Engine: Human-in-the-loop review and quality assurance system"""
    
    def __init__(self):
        self.review_statuses = ['pending_review', 'approved', 'rejected', 'published']
        self.rejection_reasons = [
            'quality_issues', 'incomplete_content', 'factual_errors', 
            'formatting_problems', 'missing_sections', 'redundancy_issues',
            'coverage_insufficient', 'fidelity_low', 'style_violations', 'other'
        ]
        
    async def get_runs_for_review(self, limit: int = 50, status_filter: str = None) -> dict:
        """Get list of processing runs available for review with quality badges"""
        try:
            print(f"📋 V2 REVIEW: Getting runs for review - limit: {limit} - engine=v2")
            
            # Get recent processing runs
            runs_query = {}
            if status_filter and status_filter in self.review_statuses:
                runs_query['review_status'] = status_filter
            
            # Get runs from various V2 collections
            runs_data = []
            
            # Get from validation results
            async for validation_result in db.v2_validation_results.find(runs_query).sort("timestamp", -1).limit(limit):
                # Convert ObjectId to string for serialization
                validation_result = objectid_to_str(validation_result)
                run_id = validation_result.get('run_id')
                if run_id:
                    run_data = await self._compile_run_data_for_review(run_id, validation_result)
                    if run_data:
                        runs_data.append(run_data)
            
            # Sort by timestamp
            runs_data.sort(key=lambda x: x.get('processing_timestamp', ''), reverse=True)
            
            # Compile summary statistics
            summary_stats = await self._compile_review_summary_stats(runs_data)
            
            review_response = {
                "review_system_status": "active",
                "engine": "v2",
                "review_data_generated_at": datetime.utcnow().isoformat(),
                
                # Summary statistics
                "summary": summary_stats,
                
                # Runs available for review
                "runs": runs_data[:limit]
            }
            
            print(f"✅ V2 REVIEW: Returning {len(runs_data)} runs for review - engine=v2")
            return review_response
            
        except Exception as e:
            print(f"❌ V2 REVIEW: Error getting runs for review - {e} - engine=v2")
            return {"error": str(e), "runs": []}
    
    async def _compile_run_data_for_review(self, run_id: str, validation_result: dict) -> dict:
        """Compile comprehensive data for a processing run for review"""
        try:
            # Get related results from other V2 collections
            qa_result = await db.v2_qa_results.find_one({"run_id": run_id})
            adjustment_result = await db.v2_adjustment_results.find_one({"run_id": run_id})
            publishing_result = await db.v2_publishing_results.find_one({"run_id": run_id})
            versioning_result = await db.v2_versioning_results.find_one({"run_id": run_id})
            
            # Convert ObjectIds to strings for serialization
            qa_result = objectid_to_str(qa_result) if qa_result else None
            adjustment_result = objectid_to_str(adjustment_result) if adjustment_result else None
            publishing_result = objectid_to_str(publishing_result) if publishing_result else None
            versioning_result = objectid_to_str(versioning_result) if versioning_result else None
            
            # Get articles from content library
            articles = []
            async for article in db.content_library.find({"metadata.run_id": run_id, "engine": "v2"}):
                # Convert ObjectId to string for serialization
                article = objectid_to_str(article)
                articles.append(article)
            
            # Calculate quality badges
            badges = self._calculate_quality_badges(validation_result, qa_result, adjustment_result)
            
            # Determine overall review status
            review_status = await self._determine_review_status(run_id, publishing_result)
            
            # Compile media references
            media_references = await self._compile_media_references(run_id)
            
            run_data = {
                "run_id": run_id,
                "review_status": review_status,
                "processing_timestamp": validation_result.get('timestamp', ''),
                
                # Quality badges for UI
                "badges": badges,
                
                # Article information
                "articles": {
                    "count": len(articles),
                    "titles": [article.get('title', 'Untitled') for article in articles],
                    "total_content_length": sum([len(article.get('content', '')) for article in articles]),
                    "articles_data": articles  # Full article data for preview
                },
                
                # Processing results summary
                "processing_results": {
                    "validation": self._summarize_validation_result(validation_result),
                    "qa": self._summarize_qa_result(qa_result),
                    "adjustment": self._summarize_adjustment_result(adjustment_result),
                    "publishing": self._summarize_publishing_result(publishing_result),
                    "versioning": self._summarize_versioning_result(versioning_result)
                },
                
                # Media library information
                "media": media_references,
                
                # Review metadata (if exists)
                "review_metadata": await self._get_review_metadata(run_id)
            }
            
            return run_data
            
        except Exception as e:
            print(f"❌ V2 REVIEW: Error compiling run data - {e} - engine=v2")
            return None
    
    def _calculate_quality_badges(self, validation_result: dict, qa_result: dict, adjustment_result: dict) -> dict:
        """Calculate quality badges for review UI display"""
        try:
            badges = {}
            
            # Coverage badge
            coverage = validation_result.get('summary_scores', {}).get('coverage_percent', 0)
            badges['coverage'] = {
                'value': f"{coverage}%",
                'status': 'excellent' if coverage >= 95 else 'good' if coverage >= 85 else 'warning',
                'tooltip': f'Content coverage: {coverage}% of source material'
            }
            
            # Fidelity badge  
            fidelity = validation_result.get('summary_scores', {}).get('fidelity_score', 0)
            badges['fidelity'] = {
                'value': f"{fidelity:.2f}",
                'status': 'excellent' if fidelity >= 0.9 else 'good' if fidelity >= 0.7 else 'warning',
                'tooltip': f'Content fidelity score: {fidelity:.2f} (accuracy to source)'
            }
            
            # Redundancy badge
            redundancy = validation_result.get('metrics', {}).get('redundancy_score', 0)
            badges['redundancy'] = {
                'value': f"{redundancy:.2f}",
                'status': 'excellent' if redundancy <= 0.2 else 'good' if redundancy <= 0.4 else 'warning',
                'tooltip': f'Content redundancy: {redundancy:.2f} (lower is better)'
            }
            
            # Granularity alignment badge
            granularity = validation_result.get('metrics', {}).get('granularity_alignment_score', 0)
            badges['granularity'] = {
                'value': f"{granularity:.2f}",
                'status': 'excellent' if granularity >= 0.8 else 'good' if granularity >= 0.6 else 'warning',
                'tooltip': f'Granularity alignment: {granularity:.2f}'
            }
            
            # Placeholders badge
            placeholders = validation_result.get('summary_scores', {}).get('placeholder_count', 0)
            badges['placeholders'] = {
                'value': str(placeholders),
                'status': 'excellent' if placeholders == 0 else 'warning',
                'tooltip': f'Placeholder content detected: {placeholders} instances'
            }
            
            # QA issues badge (if QA result available)
            if qa_result:
                issues = qa_result.get('summary', {}).get('issues_found', 0)
                badges['qa_issues'] = {
                    'value': str(issues),
                    'status': 'excellent' if issues == 0 else 'good' if issues <= 2 else 'warning',
                    'tooltip': f'Quality assurance issues: {issues} found'
                }
            
            # Readability badge (if adjustment result available)
            if adjustment_result:
                readability = adjustment_result.get('readability_score', 0.5)
                badges['readability'] = {
                    'value': f"{readability:.2f}",
                    'status': 'excellent' if readability >= 0.8 else 'good' if readability >= 0.6 else 'warning',
                    'tooltip': f'Content readability: {readability:.2f}'
                }
            
            return badges
            
        except Exception as e:
            print(f"❌ V2 REVIEW: Error calculating quality badges - {e}")
            return {}
    
    async def _determine_review_status(self, run_id: str, publishing_result: dict) -> str:
        """Determine the current review status for a processing run"""
        try:
            # Check if there's existing review metadata
            review_metadata = await db.v2_review_metadata.find_one({"run_id": run_id})
            
            if review_metadata:
                return review_metadata.get('review_status', 'pending_review')
            
            # Check publishing status to infer review status
            if publishing_result:
                publishing_status = publishing_result.get('publishing_status')
                if publishing_status == 'success':
                    return 'published'
                elif publishing_status in ['validation_failed', 'coverage_insufficient']:
                    return 'pending_review'
            
            return 'pending_review'
            
        except Exception as e:
            print(f"❌ V2 REVIEW: Error determining review status - {e}")
            return 'pending_review'
    
    async def _compile_media_references(self, run_id: str) -> dict:
        """Compile media references for review"""
        try:
            # In a full implementation, this would get media from the media library
            # For now, we'll return a placeholder structure
            return {
                "count": 0,
                "images": [],
                "videos": [],
                "documents": []
            }
            
        except Exception as e:
            print(f"❌ V2 REVIEW: Error compiling media references - {e}")
            return {"count": 0, "images": [], "videos": [], "documents": []}
    
    def _summarize_validation_result(self, validation_result: dict) -> dict:
        """Summarize validation result for review"""
        if not validation_result:
            return {"status": "not_available"}
        
        return {
            "status": validation_result.get('validation_status', 'unknown'),
            "coverage": validation_result.get('summary_scores', {}).get('coverage_percent', 0),
            "fidelity": validation_result.get('summary_scores', {}).get('fidelity_score', 0),
            "placeholders": validation_result.get('summary_scores', {}).get('placeholder_count', 0),
            "diagnostics_count": len(validation_result.get('diagnostics', []))
        }
    
    def _summarize_qa_result(self, qa_result: dict) -> dict:
        """Summarize QA result for review"""
        if not qa_result:
            return {"status": "not_available"}
        
        return {
            "status": qa_result.get('qa_status', 'unknown'),
            "issues_found": qa_result.get('summary', {}).get('issues_found', 0),
            "duplicates": qa_result.get('summary', {}).get('duplicates_found', 0),
            "invalid_links": qa_result.get('summary', {}).get('invalid_links_found', 0)
        }
    
    def _summarize_adjustment_result(self, adjustment_result: dict) -> dict:
        """Summarize adjustment result for review"""
        if not adjustment_result:
            return {"status": "not_available"}
        
        return {
            "status": adjustment_result.get('adjustment_status', 'unknown'),
            "readability_score": adjustment_result.get('readability_score', 0),
            "adjustments_made": adjustment_result.get('adjustment_summary', {}).get('total_adjustments', 0)
        }
    
    def _summarize_publishing_result(self, publishing_result: dict) -> dict:
        """Summarize publishing result for review"""
        if not publishing_result:
            return {"status": "not_available"}
        
        return {
            "status": publishing_result.get('publishing_status', 'unknown'),
            "published_articles": publishing_result.get('published_articles', 0),
            "coverage_achieved": publishing_result.get('coverage_achieved', 0)
        }
    
    def _summarize_versioning_result(self, versioning_result: dict) -> dict:
        """Summarize versioning result for review"""
        if not versioning_result:
            return {"status": "not_available"}
        
        return {
            "status": versioning_result.get('versioning_status', 'unknown'),
            "version_number": versioning_result.get('version_metadata', {}).get('version', 1),
            "is_update": versioning_result.get('version_metadata', {}).get('supersedes') is not None
        }
    
    async def _compile_review_summary_stats(self, runs_data: list) -> dict:
        """Compile summary statistics for review dashboard"""
        try:
            total_runs = len(runs_data)
            pending_review = len([r for r in runs_data if r.get('review_status') == 'pending_review'])
            approved = len([r for r in runs_data if r.get('review_status') == 'approved'])
            rejected = len([r for r in runs_data if r.get('review_status') == 'rejected'])
            published = len([r for r in runs_data if r.get('review_status') == 'published'])
            
            return {
                "total_runs": total_runs,
                "pending_review": pending_review,
                "approved": approved,
                "rejected": rejected,
                "published": published,
                "approval_rate": (approved + published) / total_runs * 100 if total_runs > 0 else 0
            }
            
        except Exception as e:
            print(f"❌ V2 REVIEW: Error compiling summary stats - {e}")
            return {"total_runs": 0, "pending_review": 0, "approved": 0, "rejected": 0, "published": 0, "approval_rate": 0}
    
    async def _get_review_metadata(self, run_id: str) -> dict:
        """Get existing review metadata for a run"""
        try:
            review_metadata = await db.v2_review_metadata.find_one({"run_id": run_id})
            return objectid_to_str(review_metadata) if review_metadata else {}
        except Exception as e:
            print(f"❌ V2 REVIEW: Error getting review metadata - {e}")
            return {}
    
    async def enqueue_for_review(self, run_id: str, articles: list, metadata: dict = None) -> dict:
        """Enqueue processing run for human review"""
        try:
            print(f"👥 V2 REVIEW: Enqueuing run for review - {run_id} - {len(articles)} articles")
            
            review_id = f"review_{run_id}_{int(time.time())}"
            
            # Create review entry
            review_entry = {
                "review_id": review_id,
                "run_id": run_id,
                "review_status": "pending_review",
                "articles_count": len(articles),
                "created_at": datetime.utcnow(),
                "metadata": metadata or {}
            }
            
            # Store in review queue
            await db.v2_review_queue.insert_one(review_entry)
            
            print(f"✅ V2 REVIEW: Run enqueued for review - {review_id}")
            
            return {
                "review_id": review_id,
                "review_status": "pending_review",
                "articles_count": len(articles)
            }
            
        except Exception as e:
            print(f"❌ V2 REVIEW: Error enqueuing for review - {e}")
            return {
                "review_id": f"error_{run_id}",
                "review_status": "error",
                "error": str(e)
            }

# Global V2 Review System instance
v2_review_system = V2ReviewSystem()

# Placeholder functions for Phase 6 features - to be implemented

async def create_high_quality_article_content(content: str, article_type: str, metadata: Dict[str, Any]) -> str:
    """Create high-quality article content with proper formatting and no duplication"""
    try:
        print(f"🎯 CREATING HIGH-QUALITY ARTICLE: {article_type}")
        
        doc_title = metadata.get('original_filename', 'Guide').replace('.docx', '').replace('.pdf', '').replace('_', ' ').replace('-', ' ')
        
        if "overview" in article_type.lower():
            # OVERVIEW ARTICLE - SOURCE CONTENT SUMMARY  
            system_message = f"""You are a content summarization specialist creating a high-level overview from source document content.

CORE PRINCIPLE: Create overview using ONLY information from the provided source content. Do NOT add generic sections or placeholder content.

TASK: Extract and summarize the key topics, sections, and information from the source to create a navigation-focused overview.

OVERVIEW APPROACH:
1. **Source Analysis**: Identify main topics, sections, and key concepts from source content
2. **Summarize Sections**: Create brief summaries of what each section covers based on source
3. **Extract Key Points**: Highlight important features, concepts, or steps mentioned in source
4. **Create Navigation**: Build roadmap based on actual source structure

HTML STRUCTURE:
- Wrap in `<div class="article-body">`
- Introductory paragraphs explaining what the source document covers
- `<h2>` sections based on source content structure
- Lists of key topics found in source material
- Brief descriptions of what each source section contains

CONTENT RESTRICTIONS:
- NO generic "Key Features" unless features are explicitly mentioned in source
- NO placeholder navigation tips unless guidance exists in source
- NO generic benefits lists unless benefits are stated in source
- NO template-style "What you'll learn" unless learning objectives exist in source

WYSIWYG ENHANCEMENTS:
- Add mini-TOC only if source has clear multiple sections
- Use proper heading hierarchy based on source structure
- Format source lists appropriately
- Add contextual notes only if source contains important warnings/tips

Focus on accurately representing what the source document contains, not what a generic overview should have."""

        else:
            # CONTEXTUAL CONTENT GENERATION - NO TEMPLATES
            system_message = f"""You are a content enhancement specialist that transforms raw document content into professional HTML articles for knowledge base display.

CORE PRINCIPLE: Use ONLY the provided source content. Do NOT add generic examples, placeholder FAQs, or template sections.

TASK: Transform the exact source content into well-structured HTML while preserving all original information, code examples, and specific details.

CONTENT PROCESSING APPROACH:
1. **Preserve Source Content**: Keep all original text, code examples, instructions, and specific details
2. **Enhance Structure**: Add proper HTML headings, lists, and semantic markup to organize existing content  
3. **Add Contextual Features**: Only add WYSIWYG features that enhance the actual source content
4. **No Generic Addition**: Never add placeholder FAQs, generic code examples, or template links

HTML STRUCTURE REQUIREMENTS:
- Wrap in `<div class="article-body">`
- Use `<h2 id="section-name">` for main sections from source content
- Use `<h3>` for subsections found in source
- Convert source lists to `<ul>` or `<ol>` with proper formatting
- Convert source code blocks to `<pre class="line-numbers"><code class="language-X">` format
- Use `<strong>` for emphasis on key terms found in source

WYSIWYG ENHANCEMENTS (only when source content warrants):
- Add mini-TOC if source has multiple clear sections
- Convert existing FAQ content to expandable format
- Add contextual notes only for complex technical content
- Use tables for source data that would benefit from tabular format

CRITICAL RULES:
- NEVER invent new FAQs, code examples, or related links
- NEVER add "Getting Started Guide" or "Best Practices" placeholder links  
- NEVER add generic "What are the benefits?" style questions
- NEVER add template JavaScript examples unless they exist in source
- Focus on making the SOURCE CONTENT look professional and well-structured

Your goal is to enhance presentation of existing content, not replace it with templates."""

        # Generate high-quality content
        response = await call_llm_with_fallback(
            system_message=system_message,
            user_message=f"""Create a high-quality {article_type} article for: {doc_title}

CRITICAL: Ensure NO text duplication, complete code examples, and proper WYSIWYG formatting.

Source content to base article on:
{content[:20000]}"""
        )
        
        if response:
            # Apply basic formatting and quality fixes only  
            clean_content = await apply_quality_fixes(response)
            
            # CONTENT VALIDATION: Ensure substantial content
            content_text = re.sub(r'<[^>]+>', '', clean_content).strip()
            if len(content_text) < 100:
                print(f"⚠️ Generated content too short ({len(content_text)} chars), regenerating...")
                
                # Enhanced regeneration with stricter requirements
                enhanced_system = f"""You are an expert technical writer. Create a comprehensive {article_type} article with MINIMUM 300 words.

STRICT REQUIREMENTS:
1. MINIMUM 300 words of actual content (not including HTML tags)
2. Multiple detailed sections with <h2> and <h3> headings
3. Include practical examples, code samples, or procedures from source
4. Add contextual information and explanations
5. Use semantic HTML formatting throughout
6. Create engaging, informative content users will find valuable

Source content to expand upon:
{content[:20000]}"""

                response = await call_llm_with_fallback(
                    system_message=enhanced_system,
                    user_message=f"Create substantial {article_type} content for: {doc_title}"
                )
                
                if response:
                    clean_content = await apply_quality_fixes(response)
                    content_text = re.sub(r'<[^>]+>', '', clean_content).strip()
            
            # Add WYSIWYG enhancements without template contamination
            if len(content_text) >= 100:
                enhanced_content = await add_wysiwyg_enhancements(clean_content, article_type)
                print(f"✅ High-quality content with WYSIWYG features: {len(content_text)} chars")
                return enhanced_content
            else:
                print(f"⚠️ Content still insufficient: {len(content_text)} chars")
                return clean_content
        else:
            return f"<h2>Error generating {article_type} content</h2><p>Please try again.</p>"
            
    except Exception as e:
        print(f"❌ Error creating high-quality article content: {e}")
        return f"<h2>Error</h2><p>Could not generate content: {e}</p>"

def clean_document_title(filename: str) -> str:
    """Clean source document titles from clutter and technical suffixes"""
    if not filename:
        return "Guide"
    
    # Remove file extensions
    title = filename.replace('.docx', '').replace('.pdf', '').replace('.doc', '').replace('.txt', '')
    
    # Replace underscores and hyphens with spaces
    title = title.replace('_', ' ').replace('-', ' ')
    
    # Remove common technical suffixes and prefixes
    suffixes_to_remove = [
        'v1', 'v2', 'v3', 'v4', 'v5', 'v6', 'v7', 'v8', 'v9',
        'version 1', 'version 2', 'version 3', 'version 4', 'version 5',
        'ver 1', 'ver 2', 'ver 3', 'ver 4', 'ver 5',
        'draft', 'final', 'final version', 'complete', 'full',
        'doc', 'document', 'documentation', 'guide book', 'manual doc',
        'user manual', 'admin guide', 'technical spec', 'specification',
        '2024', '2023', '2025', 'updated', 'latest', 'new',
        'FINAL', 'DRAFT', 'COMPLETE', 'FULL'
    ]
    
    # Convert to lowercase for comparison but preserve original case for title
    title_lower = title.lower().strip()
    
    for suffix in suffixes_to_remove:
        if title_lower.endswith(suffix.lower()):
            # Remove the suffix, preserving case
            title = title[:len(title)-len(suffix)].strip()
            title_lower = title.lower()
        if title_lower.startswith(suffix.lower()):
            # Remove prefix
            title = title[len(suffix):].strip()
            title_lower = title.lower()
    
    # Clean up multiple spaces and capitalize properly
    title = ' '.join(word.capitalize() for word in title.split() if word.strip())
    
    # If title becomes empty or too short, provide fallback
    if not title or len(title.strip()) < 3:
        return "Guide"
        
    return title.strip()

async def add_wysiwyg_enhancements(content: str, article_type: str = None) -> str:
    """Add WYSIWYG enhancements without template contamination"""
    try:
        print(f"🎨 ADDING WYSIWYG ENHANCEMENTS to {article_type or 'article'}")
        
        # Only enhance if content is substantial
        content_text = re.sub(r'<[^>]+>', '', content).strip()
        if len(content_text) < 100:
            print(f"🚫 Content too short for enhancements ({len(content_text)} chars)")
            return content
        
        enhanced = content
        
        # 1. Add article-body wrapper for WYSIWYG compatibility
        if '<div class="article-body">' not in enhanced:
            enhanced = f'<div class="article-body">\n{enhanced}\n</div>'
            print("✅ Added article-body wrapper")
        
        # 2. Enhance existing code blocks with proper classes
        def enhance_code_block(match):
            full_match = match.group(0)
            # Don't modify if already enhanced
            if 'line-numbers' in full_match:
                return full_match
            
            # Extract language and content
            language_match = re.search(r'class=["\']language-(\w+)["\']', full_match)
            language = language_match.group(1) if language_match else 'text'
            
            code_content_match = re.search(r'<code[^>]*>(.*?)</code>', full_match, re.DOTALL)
            code_content = code_content_match.group(1) if code_content_match else ''
            
            if code_content.strip():
                return f'<pre class="line-numbers"><code class="language-{language}">{code_content}</code></pre>'
            return full_match
        
        original_code_count = len(re.findall(r'<pre[^>]*><code', enhanced))
        if original_code_count > 0:
            enhanced = re.sub(r'<pre[^>]*><code[^>]*>.*?</code></pre>', enhance_code_block, enhanced, flags=re.DOTALL)
            print(f"✅ Enhanced {original_code_count} code blocks with line-numbers")
        
        # 3. Add heading IDs for navigation (only if not present)
        def add_heading_id(match):
            heading_tag = match.group(1)
            heading_text = match.group(3)
            existing_attrs = match.group(2) or ''
            
            # Skip if ID already exists
            if 'id=' in existing_attrs:
                return match.group(0)
            
            # Generate clean ID
            heading_id = heading_text.lower().strip()
            heading_id = re.sub(r'[^a-z0-9\s-]', '', heading_id)
            heading_id = re.sub(r'\s+', '-', heading_id)
            heading_id = heading_id[:50]  # Limit length
            
            return f'<{heading_tag}{existing_attrs} id="{heading_id}">{heading_text}</{heading_tag}>'
        
        heading_pattern = r'<(h[2-6])([^>]*)>([^<]*)</\1>'
        headings_added = len(re.findall(heading_pattern, enhanced))
        if headings_added > 0:
            enhanced = re.sub(heading_pattern, add_heading_id, enhanced)
            print(f"✅ Added navigation IDs to {headings_added} headings")
        
        # 4. Convert appropriate Q&A patterns to expandable sections (only if Q&A exists)
        qa_pattern = r'(Q:|Question:|FAQ:)\s*([^?]+\?)\s*(A:|Answer:)?\s*([^Q\n]+)'
        qa_matches = re.findall(qa_pattern, enhanced, re.IGNORECASE | re.MULTILINE)
        
        if qa_matches and len(qa_matches) > 1:  # Only if multiple Q&As
            print(f"📖 Converting {len(qa_matches)} Q&A items to expandable format")
            
            expandable_content = ""
            for match in qa_matches:
                question = match[1].strip()
                answer = match[3].strip()
                
                expandable_content += f'''<div class="expandable">
<div class="expandable-header"><span class="expandable-title">{question}</span></div>
<div class="expandable-content"><p>{answer}</p></div>
</div>
'''
            
            # Replace Q&A pattern with expandable format
            enhanced = re.sub(qa_pattern, '', enhanced, flags=re.IGNORECASE | re.MULTILINE)
            
            # Insert expandables after the first heading that mentions FAQ/Q&A
            if re.search(r'<h[2-6][^>]*[^>]*(?:faq|question)', enhanced, re.IGNORECASE):
                enhanced = re.sub(
                    r'(<h[2-6][^>]*[^>]*(?:faq|question)[^<]*</h[2-6]>)', 
                    rf'\1\n{expandable_content}', enhanced, 
                    flags=re.IGNORECASE, count=1
                )
                print(f"✅ Added {len(qa_matches)} expandable Q&A sections")
        
        # 5. Add contextual callouts only for technical content (sparingly)
        if ('api' in content_text.lower() or 'tutorial' in content_text.lower() or 'code' in content_text.lower()):
            if '<div class="note">' not in enhanced and '<blockquote class="tip">' not in enhanced:
                # Add ONE contextual note based on content type
                if 'api' in content_text.lower() and 'key' in content_text.lower():
                    note = '<div class="note">🔑 <strong>API Key:</strong> Ensure your API key is properly configured and has the necessary permissions.</div>\n\n'
                elif 'tutorial' in content_text.lower():
                    note = '<div class="note">📚 <strong>Tutorial:</strong> Follow each step carefully for successful implementation.</div>\n\n'
                else:
                    note = ''
                
                if note:
                    # Insert after first H2 section
                    if '<h2' in enhanced:
                        enhanced = re.sub(r'(<h2[^>]*>[^<]*</h2>)', rf'\1\n{note}', enhanced, count=1)
                        print("✅ Added contextual note")
        
        print(f"✅ WYSIWYG enhancements complete - enhanced content without template contamination")
        return enhanced
        
    except Exception as e:
        print(f"❌ Error adding WYSIWYG enhancements: {e}")
        return content

async def ensure_enhanced_features(content: str, article_type: str, doc_title: str) -> str:
    """Enhance existing content with WYSIWYG features - NO TEMPLATE INJECTION"""
    try:
        print(f"🎯 SELECTIVE ENHANCEMENT for {article_type} - enhancing existing content only")
        
        # Check content quality - skip if minimal content
        content_without_tags = re.sub(r'<[^>]+>', '', content).strip()
        
        if len(content_without_tags) < 100:
            print(f"🚫 Content too short ({len(content_without_tags)} chars) - returning as-is")
            return content
        
        # STEP 1: Article body wrapper (essential)
        if '<div class="article-body">' not in content:
            content = f'<div class="article-body">\n{content}\n</div>'
            print(f"✅ Added article-body wrapper")
        
        # STEP 2: Enhance code blocks only (no injection)
        def enhance_existing_code(match):
            code_element = match.group(0)
            if 'line-numbers' in code_element:
                return code_element
            
            # Extract language and content
            language_match = re.search(r'class=["\']language-(\w+)["\']', code_element)
            language = language_match.group(1) if language_match else 'text'
            
            code_content_match = re.search(r'<code[^>]*>(.*?)</code>', code_element, re.DOTALL)
            code_content = code_content_match.group(1) if code_content_match else ''
            
            return f'<pre class="line-numbers"><code class="language-{language}">{code_content}</code></pre>'
        
        # Only enhance existing code blocks
        original_code_count = len(re.findall(r'<pre[^>]*><code', content))
        content = re.sub(r'<pre[^>]*><code[^>]*>.*?</code></pre>', enhance_existing_code, content, flags=re.DOTALL)
        
        if original_code_count > 0:
            print(f"✅ Enhanced {original_code_count} existing code blocks with line-numbers class")
        
        # STEP 3: Convert existing FAQ patterns to expandable (no generic addition)
        faq_pattern = r'(Q:|Question:|FAQ:)\s*([^?]+\?)\s*(A:|Answer:)?\s*([^Q\n]+)'
        faq_matches = re.findall(faq_pattern, content, re.IGNORECASE | re.MULTILINE)
        
        if faq_matches and 'expandable' not in content:
            print(f"📖 Converting {len(faq_matches)} existing FAQ items to expandable format")
            
            expandable_content = ""
            for match in faq_matches:
                question = match[1].strip()
                answer = match[3].strip()
                
                expandable_content += f'''<div class="expandable">
<div class="expandable-header"><span class="expandable-title">{question}</span></div>
<div class="expandable-content"><p>{answer}</p></div>
</div>
'''
            
            # Replace FAQ pattern with expandable format
            content = re.sub(faq_pattern, '', content, flags=re.IGNORECASE | re.MULTILINE)
            
            if '<h2' in content and 'faq' in content.lower():
                content = re.sub(r'(<h2[^>]*[^>]*faq[^<]*</h2>)', rf'\1\n{expandable_content}', content, flags=re.IGNORECASE)
        
        # STEP 4: Add single contextual note only if highly technical content
        if ('<div class="note">' not in content and 
            ('api' in content.lower() and 'key' in content.lower()) or
            ('tutorial' in content.lower() and ('code' in content.lower() or 'example' in content.lower()))):
            
            # Add one highly contextual note
            if 'api' in content.lower() and 'key' in content.lower():
                note = '''<div class="note">🔑 <strong>API Key:</strong> Ensure your API key is properly configured and has the necessary permissions.</div>

'''
            elif 'tutorial' in content.lower():
                note = '''<div class="note">📚 <strong>Tutorial:</strong> Follow each step carefully for successful implementation.</div>

'''
            
            # Insert after first H2 section
            if '<h2' in content:
                content = re.sub(r'(<h2[^>]*>[^<]*</h2>)', rf'\1\n{note}', content, count=1)
                print(f"✅ Added 1 contextual note based on content type")
        
        # STEP 5: Add proper heading IDs for navigation
        content = re.sub(r'<h2([^>]*)>([^<]*)</h2>', 
                        lambda m: f'<h2{m.group(1)} id="{m.group(2).lower().replace(" ", "-").replace("#", "").replace(":", "").strip()}">{m.group(2)}</h2>', 
                        content)
        
        print(f"✅ Selective enhancement complete - enhanced existing content without template injection")
        return content
        
    except Exception as e:
        print(f"❌ Error in selective enhancement: {e}")
        return content

async def apply_quality_fixes(content: str) -> str:
    """Apply comprehensive quality fixes using proper BeautifulSoup text processing"""
    try:
        from bs4 import BeautifulSoup
        
        print(f"🔧 APPLYING PROPER QUALITY FIXES to content: {len(content)} chars")
        
        # Fix 0: COMPREHENSIVE HTML WRAPPER CLEANING - Remove ALL document structure
        
        # Remove markdown code block wrappers completely
        if '```html' in content:
            print(f"🚨 REMOVING HTML MARKDOWN WRAPPERS")
            content = re.sub(r'```html.*?```', '', content, flags=re.DOTALL | re.IGNORECASE)
        
        if '```' in content and ('<!DOCTYPE' in content or '<html' in content):
            print(f"🚨 REMOVING ANY MARKDOWN WRAPPERS WITH HTML")
            content = re.sub(r'```[^`]*<!DOCTYPE[^`]*```', '', content, flags=re.DOTALL | re.IGNORECASE)
            content = re.sub(r'```[^`]*<html[^`]*```', '', content, flags=re.DOTALL | re.IGNORECASE)
        
        # FIXED: Remove ONLY document wrapper elements, preserve content
        # Only remove full document wrappers, not content within them
        if '<!DOCTYPE' in content and '<html' in content and '</html>' in content:
            # Extract content from between body tags if present
            body_match = re.search(r'<body[^>]*>(.*?)</body>', content, re.DOTALL | re.IGNORECASE)
            if body_match:
                content = body_match.group(1)
                print(f"✅ Extracted content from <body> tags")
            else:
                # Just remove document structure tags
                content = re.sub(r'<!DOCTYPE[^>]*>', '', content, flags=re.IGNORECASE)
                content = re.sub(r'</?html[^>]*>', '', content, flags=re.IGNORECASE)
                content = re.sub(r'</?head[^>]*>', '', content, flags=re.IGNORECASE)
                content = re.sub(r'</?body[^>]*>', '', content, flags=re.IGNORECASE)
                content = re.sub(r'<head[^>]*>.*?</head>', '', content, flags=re.IGNORECASE | re.DOTALL)
        
        # Remove only metadata elements, not content
        content = re.sub(r'<meta[^>]*>', '', content, flags=re.IGNORECASE)
        content = re.sub(r'<title[^>]*>.*?</title>', '', content, flags=re.IGNORECASE | re.DOTALL)
        content = re.sub(r'<link[^>]*>', '', content, flags=re.IGNORECASE)
        
        # PRESERVE content styles and scripts that contain actual content/data
        # Only remove empty style/script blocks
        content = re.sub(r'<style[^>]*>\s*</style>', '', content, flags=re.IGNORECASE)
        content = re.sub(r'<script[^>]*>\s*</script>', '', content, flags=re.IGNORECASE)
        
        # Clean up wrapper artifacts
        content = re.sub(r'^[`\s]*', '', content.strip())
        content = re.sub(r'[`\s]*$', '', content.strip())
        
        print(f"✅ HTML wrapper cleaning completed")
        
        # Fix 1: PROPER TEXT DEDUPLICATION using BeautifulSoup
        
        soup = BeautifulSoup(content, 'html.parser')
        
        def remove_text_duplicates(text):
            """Remove duplicate sentences and phrases from text"""
            if not text or not text.strip():
                return text
            
            # Split by sentences (periods, exclamation marks, question marks)
            sentences = re.split(r'([.!?])', text)
            
            # Process sentence pairs (text + punctuation)
            cleaned_parts = []
            seen_sentences = set()
            
            i = 0
            while i < len(sentences):
                if i + 1 < len(sentences):
                    sentence = sentences[i].strip()
                    punctuation = sentences[i + 1] if i + 1 < len(sentences) else ''
                    
                    if sentence and sentence.lower() not in seen_sentences:
                        seen_sentences.add(sentence.lower())
                        cleaned_parts.append(sentence + punctuation)
                    
                    i += 2
                else:
                    # Handle last part without punctuation
                    sentence = sentences[i].strip()
                    if sentence and sentence.lower() not in seen_sentences:
                        cleaned_parts.append(sentence)
                    i += 1
            
            result = ''.join(cleaned_parts)
            
            # Additional cleanup for word-level duplications
            words = result.split()
            cleaned_words = []
            i = 0
            while i < len(words):
                if i + 1 < len(words) and words[i].lower() == words[i + 1].lower():
                    # Skip duplicate word
                    cleaned_words.append(words[i])
                    i += 2  # Skip the duplicate
                else:
                    cleaned_words.append(words[i])
                    i += 1
            
            final_result = ' '.join(cleaned_words)
            
            # Fix common duplication patterns
            final_result = re.sub(r'(\b\w+)\s+\1\b', r'\1', final_result, flags=re.IGNORECASE)
            
            return final_result
        
        # Apply deduplication to all text elements
        for element in soup.find_all(['p', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'div']):
            if element.string:  # Only process elements with direct text content
                original_text = element.string
                cleaned_text = remove_text_duplicates(original_text)
                if cleaned_text != original_text and cleaned_text.strip():
                    print(f"🔧 Fixed duplication: '{original_text[:50]}...' → '{cleaned_text[:50]}...'")
                    element.string = cleaned_text
            elif element.get_text() and not element.find_all():  # Text node without child elements
                original_text = element.get_text()
                cleaned_text = remove_text_duplicates(original_text)
                if cleaned_text != original_text and cleaned_text.strip():
                    print(f"🔧 Fixed duplication in {element.name}: '{original_text[:50]}...' → '{cleaned_text[:50]}...'")
                    element.clear()
                    element.append(cleaned_text)
        
        content = str(soup)
        
        # Fix 2: Clean up broken elements
        content = re.sub(r'<li>\s*</li>', '', content)
        content = re.sub(r'<p>\s*</p>', '', content)
        content = re.sub(r'<h[1-6]>\s*</h[1-6]>', '', content)
        
        # Fix 3: Fix broken UI references
        content = re.sub(r'click\s+\.\s*', 'click the button. ', content, flags=re.IGNORECASE)
        content = re.sub(r'navigate to\s+\.\s*', 'navigate to the section. ', content, flags=re.IGNORECASE)
        content = re.sub(r'Go to the\s+\.\s*', 'Go to the console. ', content, flags=re.IGNORECASE)
        
        # Fix 4: ENHANCED LIST FORMATTING with hierarchical support
        # Apply proper CSS classes for ordered lists with hierarchical numbering
        content = re.sub(r'<ol(?![^>]*class=)(?![^>]*style=)', '<ol class="doc-list doc-list-ordered"', content)
        
        # Handle nested ordered lists with different numbering styles
        # First level: 1, 2, 3... (default)
        # Second level: a, b, c... (lower-alpha) 
        # Third level: i, ii, iii... (lower-roman)
        nested_ol_count = 0
        def replace_nested_ol(match):
            nonlocal nested_ol_count
            nested_ol_count += 1
            if nested_ol_count % 3 == 1:
                return '<ol class="doc-list doc-list-nested doc-list-lower-alpha"'
            elif nested_ol_count % 3 == 2:
                return '<ol class="doc-list doc-list-nested doc-list-lower-roman"'
            else:
                return '<ol class="doc-list doc-list-nested"'
        
        # Apply nested styles to ordered lists within list items
        content = re.sub(r'<li[^>]*>([^<]*)<ol(?![^>]*class=)', 
                        lambda m: m.group(0).replace('<ol', f'<ol class="doc-list doc-list-nested"'), 
                        content)
        
        # Apply proper CSS classes for unordered lists with alternating styles
        content = re.sub(r'<ul(?![^>]*class=)', '<ul class="doc-list doc-list-unordered"', content)
        
        # Handle nested unordered lists with different bullet styles
        content = re.sub(r'<li[^>]*>([^<]*)<ul(?![^>]*class=)', 
                        lambda m: m.group(0).replace('<ul', '<ul class="doc-list doc-list-nested doc-list-circle"'), 
                        content)
        
        # Fix 5: ENHANCED HEADING HIERARCHY - Ensure no H1 tags in content (title field handles H1)
        # Convert ALL H1 tags to H2 since title field provides the H1
        content = re.sub(r'<h1([^>]*)>', r'<h2\1>', content)
        content = re.sub(r'</h1>', '</h2>', content)
        
        # Fix 5b: Remove duplicate title text that appears in content
        # Use a generic title for comparison
        doc_title = 'Guide'
        
        # Remove sentences that duplicate the title
        title_words = set(doc_title.lower().split())
        if len(title_words) > 1:  # Only if title has multiple words
            # Pattern to match opening paragraphs that repeat the title
            title_pattern = '|'.join(re.escape(word) for word in title_words if len(word) > 3)
            if title_pattern:
                # Remove paragraphs that are just title repetition
                content = re.sub(rf'<p[^>]*>\s*(?:(?:{title_pattern})\s*[-–—:]\s*)*(?:{title_pattern})?\s*</p>', 
                                '', content, flags=re.IGNORECASE)
        
        # Fix 6: Clean up excessive whitespace while preserving code blocks
        parts = re.split(r'(<pre[^>]*>.*?</pre>)', content, flags=re.DOTALL | re.IGNORECASE)
        cleaned_parts = []
        
        for i, part in enumerate(parts):
            if i % 2 == 0:  # Not a code block
                part = re.sub(r'\n\s*\n\s*\n', '\n\n', part)
                part = re.sub(r'[ \t]+', ' ', part)
            cleaned_parts.append(part)
        
        content = ''.join(cleaned_parts)
        
        print(f"✅ Quality fixes applied successfully: {len(content)} chars final")
        return content.strip()
        
    except Exception as e:
        print(f"❌ Error applying quality fixes: {e}")
        import traceback
        traceback.print_exc()
        
        # Simple fallback
        try:
            content = re.sub(r'<title[^>]*>.*?</title>', '', content, flags=re.IGNORECASE | re.DOTALL)
            content = re.sub(r'<!DOCTYPE[^>]*>', '', content, flags=re.IGNORECASE)
            content = re.sub(r'</?html[^>]*>', '', content, flags=re.IGNORECASE)
            content = re.sub(r'(\b\w+)\s+\1\b', r'\1', content, flags=re.IGNORECASE)
            return content.strip()
        except:
            return content

async def add_enhanced_cross_references(generated_articles: List[Dict[str, Any]], analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Placeholder for enhanced cross-reference addition"""
    # For now, use the existing cross-reference logic
    return await add_cross_references_and_related_links(generated_articles)

async def create_overview_article_with_sections(content: str, sections: List[Dict[str, Any]], metadata: Dict[str, Any], analysis: Dict[str, Any]) -> Dict[str, Any]:
    """Create overview article with sections - only if no Introduction section exists"""
    
    # Check if there's already an Introduction section to avoid duplication
    has_intro = any(
        section.get('title', '').lower() in ['introduction', 'overview', 'getting started', 'intro'] 
        for section in sections
    )
    
    if has_intro:
        print("🚫 SKIPPING overview creation - Introduction section already exists in sections")
        return None
    
    doc_title = clean_document_title(metadata.get('original_filename', 'Guide'))
    
    # Create mini-TOC for sections
    toc_items = []
    for i, section in enumerate(sections):
        section_title = section.get('title', f'Section {i+1}')
        section_id = section_title.lower().replace(' ', '-').replace('&', 'and')
        toc_items.append(f'<li><a href="#section-{section_id}">{section_title}</a></li>')
    
    mini_toc = f"""<div class="mini-toc">
<h3>📋 Guide Contents</h3>
<ul>
{chr(10).join(toc_items)}
</ul>
</div>"""
    
    overview_content = f"""{mini_toc}

<h2 id="guide-overview">📖 {doc_title} - Guide Overview</h2>

<div class="callout callout-info">
<div class="callout-title">ℹ️ About This Guide</div>
<div class="callout-content">This comprehensive guide covers all aspects of <strong>{doc_title}</strong> with detailed sections for easy navigation and reference.</div>
</div>

<p>The content is organized into <strong>{len(sections)} focused sections</strong>:</p>

<div class="sections-overview">
"""
    
    for i, section in enumerate(sections):
        section_title = section.get('title', f'Section {i+1}')
        section_desc = section.get('description', 'Content section with detailed information')
        section_id = section_title.lower().replace(' ', '-').replace('&', 'and')
        
        overview_content += f"""
<h3 id="section-{section_id}">{i+1}. {section_title}</h3>
<p>{section_desc}</p>"""
    
    overview_content += """
</div>

<div class="callout callout-tip">
<div class="callout-title">💡 Navigation Guide</div>
<div class="callout-content">Use the table of contents above to jump directly to any section that interests you most.</div>
</div>"""
    
    return {
        "id": str(uuid.uuid4()),
        "title": f"{doc_title} - Guide Overview",
        "content": overview_content,
        "status": "published",
        "article_type": "overview",
        "source_document": metadata.get("original_filename", "Unknown"),
        "tags": ["overview", "moderate_split", "navigation"],
        "priority": "high",
        "created_at": datetime.utcnow(),
        "metadata": {
            "granularity_level": "moderate",
            "processing_approach": "moderate_split",
            "article_sequence": 1,
            "has_mini_toc": True,
            "sections_count": len(sections),
            **metadata
        }
    }

async def create_section_article(content: str, section: Dict[str, Any], metadata: Dict[str, Any], analysis: Dict[str, Any], sequence: int) -> Dict[str, Any]:
    """Create enhanced section article using high-quality content generation"""
    section_title = section.get('title', f'Section {sequence}')
    section_description = section.get('description', 'This section covers important information.')
    content_focus = section.get('content_focus', 'General content')
    
    print(f"📝 Creating enhanced section article: {section_title}")
    
    # Create focused content for this section using enhanced generation
    # Add section context to metadata for better content generation
    section_metadata = {
        **metadata,
        'section_title': section_title,
        'section_description': section_description,
        'content_focus': content_focus,
        'article_sequence': sequence
    }
    
    # Generate enhanced content for this specific section
    section_content = await create_high_quality_article_content(
        content, 
        "section_guide",  # Use section-specific content type
        section_metadata
    )
    
    # Fallback if enhanced generation fails
    if not section_content or len(section_content.strip()) < 100:
        # Create basic structured content with mini-TOC
        section_content = f"""<div class="mini-toc">
<h3>📋 Section Contents</h3>
<ul>
<li><a href="#overview">Overview</a></li>
<li><a href="#details">Key Details</a></li>
<li><a href="#implementation">Implementation</a></li>
</ul>
</div>

<h2 id="overview">{section_title}</h2>

<div class="callout callout-info">
<div class="callout-title">ℹ️ Section Overview</div>
<div class="callout-content">{section_description}</div>
</div>

<h3 id="details">Key Details</h3>
<p><strong>Content Focus:</strong> {content_focus}</p>

<h3 id="implementation">Implementation</h3>
<p>This section provides detailed information and practical guidance on {section_title.lower()}.</p>

<div class="callout callout-tip">
<div class="callout-title">💡 Tip</div>
<div class="callout-content">Review related sections for comprehensive understanding of the complete process.</div>
</div>"""
    
    return {
        "id": str(uuid.uuid4()),
        "title": section_title,
        "content": section_content,
        "status": "published",
        "article_type": "section_guide",
        "source_document": metadata.get("original_filename", "Unknown"),
        "tags": ["section_guide", "moderate_split", analysis.get('content_classification', {}).get('content_type', 'guide')],
        "priority": "medium",
        "created_at": datetime.utcnow(),
        "metadata": {
            "granularity_level": "moderate",
            "processing_approach": "moderate_split",
            "article_sequence": sequence,
            "section_description": section_description,
            "content_focus": content_focus,
            "has_mini_toc": True,
            **metadata
        }
    }

async def create_simple_moderate_split(content: str, metadata: Dict[str, Any], analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Enhanced simple moderate split with actual content generation"""
    try:
        print(f"🔧 CREATING SIMPLE MODERATE SPLIT with actual content")
        
        articles = []
        doc_title = metadata.get('original_filename', 'Guide').replace('.docx', '').replace('.pdf', '').replace('_', ' ')
        content_type = analysis['content_classification']['content_type']
        
        # Check if content already has introduction/overview sections to avoid duplication
        has_intro = any(
            keyword in content.lower() 
            for keyword in ['## introduction', '## overview', '## getting started', '## intro', 
                           '<h2>introduction', '<h2>overview', '<h2>getting started', '<h2>intro']
        )
        
        if has_intro:
            print("🚫 SKIPPING overview creation in simple moderate split - Introduction section already exists in content")
        
        # 1. Overview article using enhanced content generation
        print(f"📖 Creating enhanced overview article using high-quality content generation")
        overview_content = await create_high_quality_article_content(content, "overview", metadata)
        
        if not overview_content or len(overview_content.strip()) < 100:
            # Emergency fallback with some actual content
            overview_content = f"<h2>{doc_title} - Overview</h2>\n<p>This comprehensive guide provides detailed information about {doc_title.lower()}.</p>\n<p>The content includes step-by-step instructions, technical details, and practical examples.</p>"
            
            overview = {
                "id": str(uuid.uuid4()),
                "title": f"{doc_title} - Overview",
                "content": overview_content,
                "status": "published",
                "article_type": "overview",
                "source_document": metadata.get("original_filename", "Unknown"),
                "tags": ["overview", "moderate_split", content_type],
                "priority": "high",
                "created_at": datetime.utcnow(),
                "metadata": {
                    "article_sequence": 1, 
                    "content_type": content_type,
                    "processing_approach": "simple_moderate_split",
                    **metadata
                }
            }
            articles.append(overview)
        
        # 2. Main content article using enhanced content generation
        print(f"📚 Creating enhanced main content article using high-quality content generation")
        main_content = await create_high_quality_article_content(content, "complete_guide", metadata)
        
        if not main_content or len(main_content.strip()) < 100:
            # Emergency fallback - at least include some actual content
            content_preview = content[:2000] if len(content) > 2000 else content
            main_content = f"<h2>{doc_title} - Complete Guide</h2>\n<div class=\"content-section\">{content_preview}</div>"
        
        main = {
            "id": str(uuid.uuid4()),
            "title": f"{doc_title} - Complete Guide", 
            "content": main_content,
            "status": "published",
            "article_type": "main_content",
            "source_document": metadata.get("original_filename", "Unknown"),
            "tags": ["main_content", "moderate_split", content_type, "complete_guide"],
            "priority": "high",
            "created_at": datetime.utcnow(),
            "metadata": {
                "article_sequence": 2,
                "content_type": content_type, 
                "processing_approach": "simple_moderate_split",
                **metadata
            }
        }
        articles.append(main)
        
        print(f"✅ SIMPLE MODERATE SPLIT: Created {len(articles)} articles with actual content")
        return articles
        
    except Exception as e:
        print(f"❌ Error in simple moderate split: {e}")
        import traceback
        traceback.print_exc()
        return []

async def create_enhanced_overview_article(outline_topics: List[Dict[str, Any]], metadata: Dict[str, Any], analysis: Dict[str, Any]) -> Dict[str, Any]:
    """Create enhanced overview article - only if no Introduction article exists"""
    
    # Check if there's already an Introduction/Overview topic to avoid duplication
    has_intro = any(
        topic.get('topic_title', '').lower() in ['introduction', 'overview', 'getting started']
        for topic in outline_topics
    )
    
    if has_intro:
        print(f"🚫 SKIPPING overview creation - Introduction article already exists")
        return None
    
    doc_title = clean_document_title(metadata.get('original_filename', 'Guide'))
    
    # Create comprehensive overview with mini-TOC and enhanced content
    toc_items = []
    for i, topic in enumerate(outline_topics[:8]):  # Show first 8 topics in TOC
        topic_id = topic.get('topic_title', f'topic-{i}').lower().replace(' ', '-').replace('&', 'and')
        toc_items.append(f'<li><a href="#section-{topic_id}">{topic.get("topic_title", f"Topic {i+1}")}</a></li>')
    
    mini_toc = f"""<div class="mini-toc">
<h3>📋 Contents</h3>
<ul>
{chr(10).join(toc_items)}
</ul>
</div>"""
    
    # Create sections for each topic with enhanced content
    sections_content = []
    for i, topic in enumerate(outline_topics):
        topic_title = topic.get('topic_title', f'Section {i+1}')
        topic_focus = topic.get('content_focus', 'Key information and procedures')
        topic_id = topic_title.lower().replace(' ', '-').replace('&', 'and')
        
        section_html = f"""<h2 id="section-{topic_id}">{topic_title}</h2>
<p><strong>Focus:</strong> {topic_focus}</p>"""
        
        # Add key points if available
        key_points = topic.get('key_points', [])
        if key_points:
            section_html += "\n<h3>Key Topics Covered:</h3>\n<ul>\n"
            for point in key_points[:5]:  # Limit to 5 key points
                section_html += f"<li>{point}</li>\n"
            section_html += "</ul>"
        
        sections_content.append(section_html)
    
    overview_content = f"""{mini_toc}

<h2 id="introduction">About This Guide</h2>
<div class="callout callout-info">
<div class="callout-title">ℹ️ Overview</div>
<div class="callout-content">This comprehensive guide covers all aspects of <strong>{doc_title}</strong> with detailed explanations, step-by-step instructions, and practical examples.</div>
</div>

<p>The content is organized into <strong>{len(outline_topics)} focused sections</strong> for easy navigation and reference. Each section provides complete information on its topic with practical examples and technical details.</p>

<h2 id="structure">Guide Structure</h2>

{chr(10).join(sections_content)}

<div class="callout callout-tip">
<div class="callout-title">💡 Navigation Tip</div>
<div class="callout-content">Use the table of contents above to jump directly to any section, or read sequentially for complete coverage.</div>
</div>"""
    
    return {
        "id": str(uuid.uuid4()),
        "title": f"{doc_title} - Complete Guide Overview",
        "content": overview_content,
        "status": "published",
        "article_type": "enhanced_overview",
        "source_document": metadata.get("original_filename", "Unknown"),
        "tags": ["overview", "deep_split", "enhanced", "navigation"],
        "priority": "high",
        "created_at": datetime.utcnow(),
        "metadata": {
            "granularity_level": "deep",
            "processing_approach": "deep_split",
            "article_sequence": 1,
            "has_mini_toc": True,
            "sections_count": len(outline_topics),
            **metadata
        }
    }

async def generate_enhanced_topic_article(content: str, topic: Dict[str, Any], metadata: Dict[str, Any], analysis: Dict[str, Any], sequence: int, total_topics: int) -> Dict[str, Any]:
    """Placeholder for enhanced topic article generation"""
    # Use existing topic article generation for now
    return await generate_topic_article(content, topic, metadata, sequence, total_topics)

async def adaptive_granularity_processor(content: str, metadata: Dict[str, Any], analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    """PHASE 6: Adaptive granularity processing engine"""
    try:
        granularity_level = analysis['granularity_decision']['level']
        processing_approach = analysis['processing_strategy']['approach']
        
        print(f"🎯 ADAPTIVE GRANULARITY PROCESSOR: {granularity_level} | {processing_approach}")
        
        generated_articles = []
        
        if processing_approach == "unified":
            # Use enhanced unified approach with high-quality content generation 
            print(f"📄 UNIFIED PROCESSING: Single comprehensive article")
            
            # Generate high-quality content with proper formatting
            article_content = await create_high_quality_article_content(content, "complete_guide", metadata)
            
            # Create the complete article object with proper title structure
            doc_title = clean_document_title(metadata.get('original_filename', 'Guide'))
            
            unified_article = {
                "id": str(uuid.uuid4()),
                "title": f"{doc_title}",  # FIXED: Remove redundant "Complete Guide" suffix 
                "content": article_content,
                "status": "published",
                "article_type": "complete_guide",
                "source_document": metadata.get("original_filename", "Unknown"),
                "tags": ["complete_guide", "unified", analysis.get('content_classification', {}).get('content_type', 'guide')],
                "priority": "high",
                "created_at": datetime.utcnow(),
                "metadata": {
                    "unified_article": True,
                    "processing_approach": "unified",
                    "granularity_level": analysis.get('granularity_decision', {}).get('level', 'shallow'),
                    "content_type": analysis.get('content_classification', {}).get('content_type', 'guide'),
                    **metadata
                }
            }
            generated_articles.append(unified_article)
            
            # FIXED: Do NOT create redundant overview for unified processing - complete guide is self-contained
                
        elif processing_approach == "shallow_split":
            print(f"📖 SHALLOW SPLIT PROCESSING: 2-3 articles")
            articles = await create_shallow_split_articles(content, metadata, analysis)
            generated_articles.extend(articles)
            
        elif processing_approach == "moderate_split":
            print(f"📚 MODERATE SPLIT PROCESSING: 4-6 articles")
            articles = await create_moderate_split_articles(content, metadata, analysis)
            generated_articles.extend(articles)
            
        elif processing_approach == "deep_split":
            print(f"📊 DEEP SPLIT PROCESSING: 7+ articles")
            articles = await create_deep_split_articles(content, metadata, analysis)
            generated_articles.extend(articles)
        
        # Always create FAQ for substantial content with proper standardization
        if len(content) > 2000:
            print(f"📚 CREATING STANDARDIZED FAQ ARTICLE")
            
            # Extract subject from filename for standardized title
            doc_title = metadata.get('original_filename', 'Guide').replace('.docx', '').replace('.pdf', '').replace('_', ' ').replace('-', ' ')
            subject = doc_title if doc_title != 'Guide' else analysis.get('content_classification', {}).get('content_type', 'Content')
            faq_title = f"Frequently Asked Questions & Troubleshooting – {subject}"
            
            # Generate high-quality FAQ content with cross-references
            faq_system_message = f"""You are an expert technical writer creating a comprehensive FAQ & Troubleshooting article.

CRITICAL FAQ STANDARDIZATION REQUIREMENTS:
1. Use proper HTML formatting (NO Markdown)
2. Include comprehensive cross-references with clickable links
3. Add organized "Related Links" block at the end
4. Use proper FAQ structure with categories
5. Include WYSIWYG editor features (callouts, mini-TOC)

ENHANCED FAQ STRUCTURE:
- Mini-TOC with anchor links to categories
- FAQ Categories with proper HTML headings
- Cross-references within answers: <a href="#section-id" class="cross-ref">Link text</a>
- Related Links block with grid layout

WYSIWYG FEATURES TO INCLUDE:
- Mini-TOC: <div class="mini-toc"><h3>FAQ Categories</h3><ul><li><a href="#getting-started">Getting Started</a></li></ul></div>
- Callouts: <div class="callout callout-tip"><div class="callout-title">💡 Tip</div><div class="callout-content">Helpful information</div></div>
- Cross-references: <a href="#implementation" class="cross-ref">Implementation Guide</a>

REQUIRED SECTIONS:
1. Mini-TOC of FAQ categories
2. Getting Started & Setup (5-7 questions)
3. Implementation & Configuration (5-7 questions)
4. Common Issues & Troubleshooting (5-7 questions)
5. Advanced Topics & Best Practices (3-5 questions)
6. Related Links block with organized categories

RELATED LINKS STRUCTURE:
<div class="related-articles">
  <h3>🔗 Related Links</h3>
  <div class="links-grid">
    <div class="links-category">
      <h4>Setup & Configuration</h4>
      <ul>
        <li><a href="/content-library/setup-guide">Setup Guide</a></li>
        <li><a href="/content-library/configuration">Configuration Options</a></li>
      </ul>
    </div>
    <div class="links-category">  
      <h4>Implementation</h4>
      <ul>
        <li><a href="/content-library/complete-guide">Complete Implementation Guide</a></li>
        <li><a href="/content-library/best-practices">Best Practices</a></li>
      </ul>
    </div>
  </div>
</div>

Create comprehensive FAQ content with proper HTML formatting and cross-references."""

            faq_response = await call_llm_with_fallback(
                system_message=faq_system_message,
                user_message=f"Create a comprehensive FAQ & Troubleshooting article for {subject}. Use proper HTML formatting and include cross-references.\n\nContent to base FAQs on:\n\n{content[:15000]}"
            )
            
            if faq_response:
                faq_content = await apply_quality_fixes(faq_response)
                
                faq_article = {
                    "id": str(uuid.uuid4()),
                    "title": faq_title,
                    "content": faq_content,
                    "status": "published",
                    "article_type": "faq",
                    "source_document": metadata.get("original_filename", "Unknown"),
                    "tags": ["faq", "troubleshooting", subject.lower().replace(' ', '-'), "standardized"],
                    "priority": "high",
                    "created_at": datetime.utcnow(),
                    "metadata": {
                        "standardized_faq": True,
                        "subject": subject,
                        "cross_references": True,
                        "related_links": True,
                        **metadata
                    }
                }
                generated_articles.append(faq_article)
        
        # Add cross-references and related links
        if len(generated_articles) > 1:
            enhanced_articles = await add_enhanced_cross_references(generated_articles, analysis)
            generated_articles = enhanced_articles
        
        print(f"✅ ADAPTIVE PROCESSING COMPLETE: {len(generated_articles)} articles generated")
        return generated_articles
        
    except Exception as e:
        print(f"❌ Error in adaptive granularity processor: {e}")
        import traceback
        traceback.print_exc()
        return []

async def create_shallow_split_articles(content: str, metadata: Dict[str, Any], analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Create 2-3 articles with shallow splitting"""
    try:
        print(f"📖 CREATING SHALLOW SPLIT ARTICLES")
        
        articles = []
        content_type = analysis['content_classification']['content_type']
        doc_title = metadata.get('original_filename', 'Guide').replace('.docx', '').replace('.pdf', '').replace('_', ' ')
        
        # 1. Overview Article using enhanced content generation
        print(f"📖 Creating enhanced overview article for shallow split")
        overview_content_html = await create_high_quality_article_content(content, "overview", metadata)
        
        if overview_content_html and len(overview_content_html.strip()) > 100:
            overview_article = {
                "id": str(uuid.uuid4()),
                "title": f"{doc_title} - Overview",
                "content": overview_content_html,
                "status": "published",
                "article_type": "overview",
                "source_document": metadata.get("original_filename", "Unknown"),
                "tags": ["overview", content_type, "shallow_split"],
                "priority": "high",
                "created_at": datetime.utcnow(),
                "metadata": {
                    "granularity_level": "shallow",
                    "processing_approach": "shallow_split",
                    "article_sequence": 1,
                    "content_type": content_type,
                    **metadata
                }
            }
            articles.append(overview_article)
        
        # 2. Main Content Article using enhanced content generation
        print(f"📚 Creating enhanced main content article for shallow split")
        main_content_html = await create_high_quality_article_content(content, "complete_guide", metadata)
        
        if main_content_html and len(main_content_html.strip()) > 100:
            main_article = {
                "id": str(uuid.uuid4()),
                "title": f"{doc_title} - Complete Guide",
                "content": main_content_html,
                "status": "published",
                "article_type": "main_content",
                "source_document": metadata.get("original_filename", "Unknown"),
                "tags": ["main_content", content_type, "shallow_split"],
                "priority": "high",
                "created_at": datetime.utcnow(),
                "metadata": {
                    "granularity_level": "shallow",
                    "processing_approach": "shallow_split",
                    "article_sequence": 2,
                    "content_type": content_type,
                    **metadata
                }
            }
            articles.append(main_article)
        
        print(f"✅ SHALLOW SPLIT: Created {len(articles)} articles")
        return articles
        
    except Exception as e:
        print(f"❌ Error creating shallow split articles: {e}")
        return []

async def create_moderate_split_articles(content: str, metadata: Dict[str, Any], analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Create 4-6 articles with moderate splitting"""
    try:
        print(f"📚 CREATING MODERATE SPLIT ARTICLES")
        
        articles = []
        content_type = analysis['content_classification']['content_type']
        doc_title = metadata.get('original_filename', 'Guide').replace('.docx', '').replace('.pdf', '').replace('_', ' ')
        
        # Generate content outline for moderate split
        outline_system = f"""Create a moderate split outline for this content (4-6 main sections).

CONTENT TYPE: {content_type}
APPROACH: Moderate split (4-6 articles)

Analyze the content and create 4-6 logical sections that:
1. Can stand alone but connect logically
2. Have balanced content distribution
3. Maintain context for technical elements
4. Provide comprehensive coverage

Return JSON format:
{{
  "sections": [
    {{
      "title": "Section Title",
      "description": "What this section covers",
      "content_focus": "Key topics and elements"
    }}
  ]
}}"""

        outline_response = await call_llm_with_fallback(
            system_message=outline_system,
            user_message=f"Create moderate split outline for:\n\n{content[:15000]}"
        )
        
        if outline_response:
            try:
                outline_data = json.loads(outline_response.strip())
                sections = outline_data.get('sections', [])
                
                # Create overview article
                overview_article = await create_overview_article_with_sections(content, sections, metadata, analysis)
                if overview_article:
                    articles.append(overview_article)
                
                # Create articles for each section
                for i, section in enumerate(sections[:6]):  # Limit to 6 sections
                    section_article = await create_section_article(content, section, metadata, analysis, i + 2)
                    if section_article:
                        articles.append(section_article)
                        
            except json.JSONDecodeError:
                # Fallback to simple splitting
                print(f"⚠️ Outline parsing failed, using simple moderate split")
                articles = await create_simple_moderate_split(content, metadata, analysis)
        
        print(f"✅ MODERATE SPLIT: Created {len(articles)} articles")
        return articles
        
    except Exception as e:
        print(f"❌ Error creating moderate split articles: {e}")
        return []

async def create_deep_split_articles(content: str, metadata: Dict[str, Any], analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Create 7+ articles with deep splitting"""
    try:
        print(f"📊 CREATING DEEP SPLIT ARTICLES")
        
        # Use existing outline-first approach for deep splitting
        outline_data = await generate_comprehensive_outline(content, metadata)
        if not outline_data:
            print(f"❌ Could not generate outline for deep split")
            return []
        
        articles = []
        outline_topics = outline_data.get('comprehensive_outline', [])
        content_type = analysis['content_classification']['content_type']
        
        print(f"📋 DEEP SPLIT: Processing {len(outline_topics)} topics")
        
        # Create overview article
        overview_article = await create_enhanced_overview_article(outline_topics, metadata, analysis)
        if overview_article:
            articles.append(overview_article)
        
        # Create article for each topic
        for i, topic in enumerate(outline_topics):
            article = await generate_enhanced_topic_article(content, topic, metadata, analysis, i + 2, len(outline_topics))
            if article:
                articles.append(article)
                print(f"💾 DEEP SPLIT ARTICLE {i+2}/{len(outline_topics)+1}: {article['title']}")
        
        print(f"✅ DEEP SPLIT: Created {len(articles)} articles")
        return articles
        
    except Exception as e:
        print(f"❌ Error creating deep split articles: {e}")
        return []

async def intelligent_content_processing_pipeline(content: str, metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
    """PHASE 6: ENHANCED INTELLIGENT PIPELINE with adaptive granularity processing"""
    try:
        print(f"🧠 PHASE 6: ENHANCED INTELLIGENT CONTENT PROCESSING PIPELINE STARTED")
        print(f"📄 Content: {len(content)} characters from {metadata.get('original_filename', 'Unknown')}")
        
        # STEP 1: Enhanced multi-dimensional content analysis
        print(f"🔍 STEP 1: Enhanced multi-dimensional content analysis")
        analysis = await enhanced_multi_dimensional_analysis(content, metadata)
        
        # STEP 2: Adaptive granularity processing
        print(f"🎯 STEP 2: Adaptive granularity processing")
        generated_articles = await adaptive_granularity_processor(content, metadata, analysis)
        
        # STEP 3: Save all articles to database
        print(f"💾 STEP 3: Saving {len(generated_articles)} articles to database")
        for i, article in enumerate(generated_articles, 1):
            try:
                # KE-PR9.3: Use repository pattern for article insertion
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(unified_article)
                generated_articles.append(unified_article)
                print(f"💾 SAVED UNIFIED ARTICLE: {unified_article['title']}")
            
        else:
            # SPLIT APPROACH: Use the original outline-first approach
            print(f"🔀 SPLIT PROCESSING: Using outline-first approach for multiple articles")
            
            # Generate comprehensive outline
            outline_data = await generate_comprehensive_outline(content, metadata)
            if not outline_data:
                print(f"❌ PIPELINE FAILED: Could not generate outline")
                return []
            
            # Generate articles for each topic
            outline_topics = outline_data.get('comprehensive_outline', [])
            print(f"📋 PROCESSING {len(outline_topics)} TOPICS individually")
            
            for i, topic in enumerate(outline_topics, 1):
                article = await generate_topic_article(content, topic, metadata, i, len(outline_topics))
                if article:
                    # KE-PR9.3: Use repository pattern for content_library operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(overview_article)
                    generated_articles.insert(0, overview_article)
                    print(f"💾 SAVED OVERVIEW ARTICLE with mini-TOC")
        
        # ALWAYS: Create FAQ article for substantial content
        if len(content) > 2000:
            faq_article = await create_faq_article(content, generated_articles, metadata)
            if faq_article:
                # KE-PR9.3: Use assets repository
                from engine.stores.mongo import RepositoryFactory
                assets_repo = RepositoryFactory.get_assets()
                result = await assets_repo.insert_assets(self.pending_assets)
                            print(f"📚 Successfully added {len(self.pending_assets)} assets to Asset Library")
                            
                        # Clear pending assets
                        self.pending_assets = []
                        
                    except Exception as db_error:
                        print(f"⚠️ Failed to insert assets into Asset Library: {db_error}")
                        # Continue processing even if Asset Library insertion fails
                        
            elif file_type.lower() == 'pdf':
                html_content, images = await self._convert_pdf_to_html(file_path)
                
                # FIXED: Add Asset Library insertion for PDF images (same as DOCX)
                print(f"🔍 DEBUG: Checking for PDF pending_assets: hasattr={hasattr(self, 'pending_assets')}")
                if hasattr(self, 'pending_assets'):
                    print(f"🔍 DEBUG: Found {len(self.pending_assets)} pending PDF assets to insert")
                    
                if hasattr(self, 'pending_assets') and self.pending_assets:
                    try:
                        # Batch insert PDF images into Asset Library
                        result = # KE-PR9.3: Use repository pattern for find operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                result = await content_repo.collection.find_one({"id": article_id})
        
        if not article:
            raise HTTPException(status_code=404, detail="Article not found")
        
        # Get article content and title
        title = article.get("title", "Generated Article")
        content = article.get("content", "")
        
        # Debug logging for content validation
        print(f"📄 Article title: '{title}'")
        print(f"📏 Content length: {len(content)} characters")
        print(f"📝 Content preview: {content[:200]}...")
        
        if not content or len(content.strip()) < 50:
            # Provide fallback content if article is empty or too short
            print("⚠️ Article content is empty or too short, providing fallback content")
            content = f"""
            <h1>{title}</h1>
            <div class="article-metadata">
                <p><strong>Article ID:</strong> {article_id}</p>
                <p><strong>Created:</strong> {article.get('created_at', 'Unknown')}</p>
                <p><strong>Source:</strong> {article.get('source_type', 'Unknown')}</p>
            </div>
            <div class="no-content-notice">
                <h2>Content Not Available</h2>
                <p>This article appears to have empty or insufficient content for PDF generation.</p>
                <p>Please ensure the article has been properly processed and contains meaningful content.</p>
            </div>
            """
        
        # Clean title for filename
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = f"{safe_title[:50]}.pdf"  # Limit filename length
        
        # Generate PDF
        print(f"🎨 Starting PDF generation with {len(content)} characters of content")
        pdf_bytes = generate_pdf_from_html(content, title)
        
        # Validate PDF was actually generated
        if not pdf_bytes or len(pdf_bytes) < 1000:
            print(f"❌ Generated PDF is too small: {len(pdf_bytes) if pdf_bytes else 0} bytes")
            raise HTTPException(status_code=500, detail="PDF generation produced invalid output")
        
        print(f"✅ PDF generated successfully: {len(pdf_bytes)} bytes")
        
        # Create PDF stream response
        def generate():
            yield pdf_bytes
        
        # Return streaming response with proper headers
        headers = {
            'Content-Disposition': f'attachment; filename="{filename}"',
            'Content-Type': 'application/pdf',
            'Content-Length': str(len(pdf_bytes))
        }
        
        return StreamingResponse(
            generate(),
            media_type='application/pdf',
            headers=headers
        )
        
    except HTTPException:
        # Re-raise HTTPExceptions (like 404) without modification
        raise
    except Exception as e:
        print(f"❌ Content Library PDF download error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/training/article/{session_id}/{article_index}/download-pdf")
async def download_training_article_pdf(session_id: str, article_index: int):
    """Download a Training Interface article as PDF"""
    try:
        print(f"🔍 Generating PDF for Training article: {session_id}, index: {article_index}")
        
        # Find the training session
        training_session = await db.training_sessions.find_one({"session_id": session_id})
        
        if not training_session:
            raise HTTPException(status_code=404, detail="Training session not found")
        
        # Get the specific article
        articles = training_session.get("articles", [])
        
        if article_index >= len(articles) or article_index < 0:
            raise HTTPException(status_code=404, detail="Article not found in session")
        
        article = articles[article_index]
        
        # Get article content and title
        title = article.get("title", f"Training Article {article_index + 1}")
        content = article.get("content", "")
        
        # Debug logging for content validation
        print(f"📄 Training article title: '{title}'")
        print(f"📏 Content length: {len(content)} characters")
        print(f"📝 Content preview: {content[:200]}...")
        
        if not content or len(content.strip()) < 50:
            # Provide fallback content if article is empty or too short
            print("⚠️ Training article content is empty or too short, providing fallback content")
            content = f"""
            <h1>{title}</h1>
            <div class="article-metadata">
                <p><strong>Training Session:</strong> {session_id}</p>
                <p><strong>Article Index:</strong> {article_index}</p>
                <p><strong>Template:</strong> {training_session.get('template', 'Unknown')}</p>
                <p><strong>Filename:</strong> {training_session.get('filename', 'Unknown')}</p>
            </div>
            <div class="no-content-notice">
                <h2>Content Not Available</h2>
                <p>This training article appears to have empty or insufficient content for PDF generation.</p>
                <p>Please ensure the document was properly processed and contains meaningful content.</p>
                <p>Try re-processing the document or check the original file for content.</p>
            </div>
            """
        
        # Clean title for filename
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = f"Lab_{safe_title[:40]}.pdf"  # Limit filename length
        
        # Generate PDF
        print(f"🎨 Starting PDF generation with {len(content)} characters of content")
        pdf_bytes = generate_pdf_from_html(content, title)
        
        # Validate PDF was actually generated
        if not pdf_bytes or len(pdf_bytes) < 1000:
            print(f"❌ Generated PDF is too small: {len(pdf_bytes) if pdf_bytes else 0} bytes")
            raise HTTPException(status_code=500, detail="PDF generation produced invalid output")
        
        print(f"✅ PDF generated successfully: {len(pdf_bytes)} bytes")
        
        # Create PDF stream response
        def generate():
            yield pdf_bytes
        
        # Return streaming response with proper headers
        headers = {
            'Content-Disposition': f'attachment; filename="{filename}"',
            'Content-Type': 'application/pdf',
            'Content-Length': str(len(pdf_bytes))
        }
        
        return StreamingResponse(
            generate(),
            media_type='application/pdf',
            headers=headers
        )
        
    except HTTPException:
        # Re-raise HTTPExceptions (like 404) without modification
        raise
    except Exception as e:
        print(f"❌ Training PDF download error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


    # PDF TIMEOUT FIX: Add processing limits and timeout handling
    PDF_PROCESSING_TIMEOUT = 120  # 2 minutes maximum processing time
    MAX_PDF_SIZE = 10 * 1024 * 1024  # 10MB maximum file size
    
    import signal
    import asyncio
    from contextlib import asynccontextmanager
    
    @asynccontextmanager
    async def pdf_processing_timeout(timeout_seconds=120):
        """Context manager for PDF processing with timeout"""
        try:
            yield
        except asyncio.TimeoutError:
            print(f"❌ PDF processing timed out after {timeout_seconds} seconds")
            raise HTTPException(status_code=408, detail="PDF processing timeout - file too large or complex")
        except Exception as e:
            print(f"❌ PDF processing error: {e}")
            raise HTTPException(status_code=500, detail=f"PDF processing failed: {str(e)}")
    
async def process_pdf_with_template(file_path: str, template_data: dict, training_session: dict) -> list:
    """Process PDF file with comprehensive text and image extraction"""
    try:
        print(f"🔍 Starting comprehensive PDF processing: {file_path}")
        
        # Import PDF processing libraries
        try:
            import PyPDF2
            import fitz  # PyMuPDF for image extraction
        except ImportError:
            print("⚠️ PDF processing libraries not fully available, using fallback")
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return await process_text_with_template(content, template_data, training_session)
            except:
                return []
        
        # Try to process as actual PDF file
        try:
            # Open PDF with PyMuPDF for comprehensive extraction
            doc = fitz.open(file_path)
            print(f"✅ Successfully loaded PDF file with {len(doc)} pages")
        except Exception as pdf_error:
            print(f"⚠️ Failed to load as PDF file: {pdf_error}")
            print(f"🔄 Falling back to text processing")
            # If it's not a valid PDF file, treat it as text
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return await process_text_with_template(content, template_data, training_session)
            except Exception as text_error:
                print(f"❌ Text fallback also failed: {text_error}")
                return []
        
        # Extract text and images with contextual positioning
        full_text = ""
        all_images = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Extract text from page
            page_text = page.get_text()
            if page_text.strip():
                full_text += f"\n\n=== Page {page_num + 1} ===\n{page_text}\n"
            
            # Extract images from page with filtering
            image_list = page.get_images(full=True)
            page_rect = page.rect
            
            for img_index, img in enumerate(image_list):
                try:
                    # Get image data and properties
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    # Get image position on page
                    img_rects = page.get_image_rects(xref)
                    
                    # Filter out header/footer images and small decorative elements
                    should_include = True
                    
                    if img_rects:
                        for rect in img_rects:
                            # Check if image is in header (top 15% of page)
                            if rect.y1 < page_rect.height * 0.15:
                                print(f"🚫 Skipping header image on page {page_num + 1}")
                                should_include = False
                                break
                            # Check if image is in footer (bottom 15% of page)  
                            elif rect.y0 > page_rect.height * 0.85:
                                print(f"🚫 Skipping footer image on page {page_num + 1}")
                                should_include = False
                                break
                            # Check if image is very small (likely decorative)
                            elif rect.width < 100 or rect.height < 100:
                                print(f"🚫 Skipping small decorative image on page {page_num + 1}")
                                should_include = False
                                break
                    
                    # Skip if image is too small in actual pixels or filtered out
                    if not should_include or pix.width < 100 or pix.height < 100:
                        pix = None
                        continue
                    
                    # Skip if image appears to be a logo (very wide and short, or very narrow and tall)
                    aspect_ratio = pix.width / pix.height
                    if aspect_ratio > 5 or aspect_ratio < 0.2:
                        print(f"🚫 Skipping logo-like image (aspect ratio: {aspect_ratio:.2f}) on page {page_num + 1}")
                        pix = None
                        continue
                    
                    # Convert to RGB if CMYK
                    if pix.n - pix.alpha < 4:
                        img_data = pix.tobytes("png")
                    else:
                        pix1 = fitz.Pixmap(fitz.csRGB, pix)
                        img_data = pix1.tobytes("png")
                        pix1 = None
                    
                    # Generate unique filename
                    safe_prefix = "".join(c for c in training_session.get('filename', 'pdf') if c.isalnum())[:10]
                    unique_filename = f"{safe_prefix}_page{page_num + 1}_img{img_index + 1}_{str(uuid.uuid4())[:8]}.png"
                    file_path_static = f"static/uploads/{unique_filename}"
                    
                    # Ensure upload directory exists
                    os.makedirs("static/uploads", exist_ok=True)
                    
                    # Save image to disk
                    with open(file_path_static, "wb") as f:
                        f.write(img_data)
                    
                    # Generate URL
                    image_url = f"/api/static/uploads/{unique_filename}"
                    
                    # Store image info with contextual position
                    all_images.append({
                        "filename": unique_filename,
                        "url": image_url,
                        "page": page_num + 1,
                        "position": img_index + 1,
                        "width": pix.width,
                        "height": pix.height,
                        "size": len(img_data),
                        "is_svg": False
                    })
                    
                    print(f"✅ Extracted PDF image: Page {page_num + 1}, Image {img_index + 1} -> {image_url}")
                    
                    pix = None
                    
                except Exception as img_error:
                    print(f"⚠️ Error extracting image {img_index + 1} from page {page_num + 1}: {img_error}")
                    continue
        
        doc.close()
        
        # Basic PDF metadata
        try:
            pdf_file = open(file_path, 'rb')
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            pdf_info = pdf_reader.metadata
            if pdf_info:
                metadata_text = "\n\n=== Document Information ===\n"
                if pdf_info.get('/Title'):
                    metadata_text += f"Title: {pdf_info['/Title']}\n"
                if pdf_info.get('/Author'):
                    metadata_text += f"Author: {pdf_info['/Author']}\n"
                if pdf_info.get('/Subject'):
                    metadata_text += f"Subject: {pdf_info['/Subject']}\n"
                full_text = metadata_text + full_text
            pdf_file.close()
        except Exception as e:
            print(f"⚠️ Error extracting PDF metadata: {e}")
        
        print(f"✅ PDF extraction complete: {len(full_text)} characters, {len(all_images)} images")
        
        # Check if we have meaningful content
        if not full_text.strip():
            print("⚠️ No text content extracted from PDF")
            return []
        
        # INTELLIGENT APPROACH: Use the intelligent content processing pipeline
        print(f"🧠 Using INTELLIGENT CONTENT PROCESSING PIPELINE for PDF content")
        articles = await intelligent_content_processing_pipeline(full_text, {
            "source": "pdf",
            "original_filename": template_data.get("filename", "document.pdf"),
            "images": all_images,
            "template_data": template_data,
            "training_session": training_session
        })
        
        print(f"✅ PDF processing generated {len(articles)} articles")
        
        return articles
        
    except Exception as e:
        print(f"❌ PDF processing error: {e}")
        import traceback
        traceback.print_exc()
        return []

async def process_ppt_with_template(file_path: str, template_data: dict, training_session: dict) -> list:
    """Process PowerPoint file with comprehensive text and image extraction"""
    try:
        print(f"🔍 Starting comprehensive PowerPoint processing: {file_path}")
        
        # Import PowerPoint processing library
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            print("⚠️ python-pptx not installed, using fallback processing")
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return await process_text_with_template(content, template_data, training_session)
            except:
                return []
        
        # Try to process as actual PowerPoint file
        try:
            prs = Presentation(file_path)
            print(f"✅ Successfully loaded PowerPoint file with {len(prs.slides)} slides")
        except Exception as ppt_error:
            print(f"⚠️ Failed to load as PowerPoint file: {ppt_error}")
            print(f"🔄 Falling back to text processing")
            # If it's not a valid PowerPoint file, treat it as text
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return await process_text_with_template(content, template_data, training_session)
            except Exception as text_error:
                print(f"❌ Text fallback also failed: {text_error}")
                return []
        
        # Extract text content and images from all slides
        full_text = ""
        all_images = []
        slide_count = 0
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"\n\n=== Slide {slide_num + 1} ===\n"
            slide_count += 1
            
            # Extract text and images from all shapes in the slide
            for shape_index, shape in enumerate(slide.shapes):
                # Extract text from shape
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                
                # Extract images from shape
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Get image data
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Skip if image is too small
                        if len(image_bytes) < 1000:
                            continue
                        
                        # Determine image format
                        image_ext = image.ext
                        if not image_ext:
                            image_ext = "png"  # Default to PNG
                        
                        # Generate unique filename
                        safe_prefix = "".join(c for c in training_session.get('filename', 'ppt') if c.isalnum())[:10]
                        unique_filename = f"{safe_prefix}_slide{slide_num + 1}_shape{shape_index + 1}_{str(uuid.uuid4())[:8]}.{image_ext}"
                        file_path_static = f"static/uploads/{unique_filename}"
                        
                        # Ensure upload directory exists
                        os.makedirs("static/uploads", exist_ok=True)
                        
                        # Save image to disk
                        with open(file_path_static, "wb") as f:
                            f.write(image_bytes)
                        
                        # Generate URL
                        image_url = f"/api/static/uploads/{unique_filename}"
                        
                        # Store image info with contextual position
                        all_images.append({
                            "filename": unique_filename,
                            "url": image_url,
                            "slide": slide_num + 1,
                            "shape_index": shape_index + 1,
                            "format": image_ext,
                            "size": len(image_bytes),
                            "is_svg": False
                        })
                        
                        print(f"✅ Extracted PPT image: Slide {slide_num + 1}, Shape {shape_index + 1} -> {image_url}")
                        
                    except Exception as img_error:
                        print(f"⚠️ Error extracting image from slide {slide_num + 1}, shape {shape_index + 1}: {img_error}")
                        continue
                
                # Extract images from grouped shapes
                elif hasattr(shape, "shapes"):
                    for sub_shape in shape.shapes:
                        if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                # Get image data
                                image = sub_shape.image
                                image_bytes = image.blob
                                
                                # Skip if image is too small
                                if len(image_bytes) < 1000:
                                    continue
                                
                                # Determine image format
                                image_ext = image.ext
                                if not image_ext:
                                    image_ext = "png"
                                
                                # Generate unique filename
                                safe_prefix = "".join(c for c in training_session.get('filename', 'ppt') if c.isalnum())[:10]
                                unique_filename = f"{safe_prefix}_slide{slide_num + 1}_subshape_{str(uuid.uuid4())[:8]}.{image_ext}"
                                file_path_static = f"static/uploads/{unique_filename}"
                                
                                # Ensure upload directory exists
                                os.makedirs("static/uploads", exist_ok=True)
                                
                                # Save image to disk
                                with open(file_path_static, "wb") as f:
                                    f.write(image_bytes)
                                
                                # Generate URL
                                image_url = f"/api/static/uploads/{unique_filename}"
                                
                                # Store image info
                                all_images.append({
                                    "filename": unique_filename,
                                    "url": image_url,
                                    "slide": slide_num + 1,
                                    "format": image_ext,
                                    "size": len(image_bytes),
                                    "is_svg": False
                                })
                                
                                print(f"✅ Extracted PPT grouped image: Slide {slide_num + 1} -> {image_url}")
                                
                            except Exception as img_error:
                                print(f"⚠️ Error extracting grouped image from slide {slide_num + 1}: {img_error}")
                                continue
            
            # Add slide text if it has content
            if slide_text.strip() != f"=== Slide {slide_num + 1} ===":
                full_text += slide_text
        
        print(f"✅ PowerPoint extraction complete: {len(full_text)} characters, {len(all_images)} images from {slide_count} slides")
        
        # Check if we have meaningful content
        if not full_text.strip():
            print("⚠️ No text content extracted from PowerPoint")
            return []
        
        # Process with template including images
        articles = await create_articles_with_template(full_text, all_images, template_data, training_session)
        
        print(f"✅ PowerPoint processing generated {len(articles)} articles")
        
        return articles
        
    except Exception as e:
        print(f"❌ PowerPoint processing error: {e}")
        import traceback
        traceback.print_exc()
        return []

async def process_text_with_template(content: str, template_data: dict, training_session: dict) -> list:
    """Process text content with training template"""
    try:
        # Apply template processing instructions
        processing_instructions = template_data.get("processing_instructions", [])
        
        # For plain text files, let's add some structure
        if not content.strip():
            content = f"Text content from {training_session.get('filename', 'unknown file')}"
        
        # Generate articles based on template
        articles = await create_articles_with_template(content, [], template_data, training_session)
        
        return articles
        
    except Exception as e:
        print(f"Text processing error: {e}")
        return []

async def process_doc_with_template(file_path: str, template_data: dict, training_session: dict) -> list:
    """Process DOC file with comprehensive text extraction"""
    try:
        print(f"🔍 Starting DOC processing: {file_path}")
        
        # Import doc processing libraries
        try:
            import subprocess
            import os
        except ImportError:
            print("⚠️ Required libraries not available, using fallback processing")
            return await process_text_with_template("", template_data, training_session)
        
        # Try to extract text using antiword (if available)
        try:
            result = subprocess.run(['antiword', file_path], 
                                  capture_output=True, text=True, timeout=30)
            if result.returncode == 0:
                full_text = result.stdout
                print(f"✅ DOC text extracted using antiword: {len(full_text)} characters")
            else:
                raise Exception("antiword failed")
        except (subprocess.TimeoutExpired, FileNotFoundError, Exception):
            # Fallback: read as binary and extract readable text
            try:
                with open(file_path, 'rb') as f:
                    content = f.read()
                
                # Basic text extraction from DOC binary
                # DOC files contain text mixed with binary data
                text_parts = []
                current_text = ""
                
                for byte in content:
                    if 32 <= byte <= 126:  # Printable ASCII
                        current_text += chr(byte)
                    else:
                        if len(current_text) > 3:  # Only keep meaningful text chunks
                            text_parts.append(current_text)
                        current_text = ""
                
                if current_text:
                    text_parts.append(current_text)
                
                full_text = " ".join(text_parts)
                print(f"✅ DOC text extracted using binary parsing: {len(full_text)} characters")
                
            except Exception as e:
                print(f"⚠️ Error extracting DOC content: {e}")
                full_text = f"Error processing DOC file: {training_session.get('filename', 'unknown')}"
        
        # Process with template (DOC files typically don't have easily accessible embedded images)
        articles = await create_articles_with_template(full_text, [], template_data, training_session)
        
        print(f"✅ DOC processing generated {len(articles)} articles")
        
        return articles
        
    except Exception as e:
        print(f"❌ DOC processing error: {e}")
        return []

async def process_excel_with_template(file_path: str, template_data: dict, training_session: dict) -> list:
    """Process Excel file with comprehensive text and image extraction"""
    try:
        print(f"🔍 Starting comprehensive Excel processing: {file_path}")
        
        # Import Excel processing libraries
        try:
            from openpyxl import load_workbook
            from openpyxl.drawing.image import Image as OpenpyxlImage
            import pandas as pd
        except ImportError:
            print("⚠️ openpyxl/pandas not installed, using fallback processing")
            return await process_text_with_template("", template_data, training_session)
        
        # Read Excel file with openpyxl for comprehensive extraction
        try:
            workbook = load_workbook(file_path)
            full_text = ""
            all_images = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                full_text += f"\n\n=== Sheet: {sheet_name} ===\n"
                
                # Extract text content from cells
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        full_text += row_text + "\n"
                
                # Extract images from sheet
                if hasattr(sheet, '_images'):
                    for img_index, img in enumerate(sheet._images):
                        try:
                            # Get image data
                            image_data = img.ref
                            
                            # Skip if no image data
                            if not hasattr(image_data, 'read'):
                                continue
                            
                            image_bytes = image_data.read()
                            
                            # Skip if image is too small
                            if len(image_bytes) < 1000:
                                continue
                            
                            # Generate unique filename
                            safe_prefix = "".join(c for c in training_session.get('filename', 'excel') if c.isalnum())[:10]
                            unique_filename = f"{safe_prefix}_{sheet_name}_img{img_index + 1}_{str(uuid.uuid4())[:8]}.png"
                            file_path_static = f"static/uploads/{unique_filename}"
                            
                            # Ensure upload directory exists
                            os.makedirs("static/uploads", exist_ok=True)
                            
                            # Save image to disk
                            with open(file_path_static, "wb") as f:
                                f.write(image_bytes)
                            
                            # Generate URL
                            image_url = f"/api/static/uploads/{unique_filename}"
                            
                            # Store image info
                            all_images.append({
                                "filename": unique_filename,
                                "url": image_url,
                                "sheet": sheet_name,
                                "size": len(image_bytes),
                                "is_svg": False
                            })
                            
                            print(f"✅ Extracted Excel image: {sheet_name} -> {image_url}")
                            
                        except Exception as img_error:
                            print(f"⚠️ Error extracting image from sheet {sheet_name}: {img_error}")
                            continue
                
            workbook.close()
            
        except Exception as e:
            print(f"⚠️ Error reading Excel file with openpyxl: {e}")
            
            # Fallback with pandas
            try:
                excel_file = pd.ExcelFile(file_path)
                full_text = ""
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    full_text += f"\n\n=== Sheet: {sheet_name} ===\n"
                    
                    if not df.empty:
                        # Add column headers
                        full_text += "Columns: " + ", ".join(df.columns.astype(str)) + "\n\n"
                        
                        # Add data rows (limit to first 100 rows)
                        for idx, row in df.head(100).iterrows():
                            row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                            if row_text.strip():
                                full_text += row_text + "\n"
                        
                        if len(df) > 100:
                            full_text += f"\n... and {len(df) - 100} more rows\n"
                    else:
                        full_text += "Empty sheet\n"
                
                all_images = []  # No image extraction in pandas fallback
                
            except Exception as pandas_error:
                print(f"⚠️ Error with pandas fallback: {pandas_error}")
                full_text = f"Error processing Excel file: {training_session.get('filename', 'unknown')}"
                all_images = []
        
        print(f"✅ Excel extraction complete: {len(full_text)} characters, {len(all_images)} images")
        
        # Process with template
        articles = await create_articles_with_template(full_text, all_images, template_data, training_session)
        
        print(f"✅ Excel processing generated {len(articles)} articles")
        
        return articles
        
    except Exception as e:
        print(f"❌ Excel processing error: {e}")
        return []

async def process_html_with_template(file_path: str, template_data: dict, training_session: dict) -> list:
    """Process HTML file with comprehensive text and image extraction"""
    try:
        print(f"🔍 Starting comprehensive HTML processing: {file_path}")
        
        # Read HTML content
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Parse HTML and extract text content and images
        try:
            from bs4 import BeautifulSoup
            import requests
            import base64
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extract text content
            full_text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in full_text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            full_text = ' '.join(chunk for chunk in chunks if chunk)
            
            # Extract images
            all_images = []
            img_tags = soup.find_all('img')
            
            for i, img in enumerate(img_tags):
                src = img.get('src', '')
                alt = img.get('alt', f'Image {i+1}')
                
                if src.startswith('data:image'):
                    # Base64 embedded image
                    try:
                        # Extract base64 data
                        header, data = src.split(',', 1)
                        image_data = base64.b64decode(data)
                        
                        # Determine file extension
                        if 'svg' in header:
                            ext = 'svg'
                            is_svg = True
                        elif 'png' in header:
                            ext = 'png'
                            is_svg = False
                        elif 'jpg' in header or 'jpeg' in header:
                            ext = 'jpg'
                            is_svg = False
                        else:
                            ext = 'png'
                            is_svg = False
                        
                        if is_svg:
                            # SVG images keep base64 format
                            all_images.append({
                                "filename": f"html_image_{i+1}.{ext}",
                                "data": src,
                                "size": len(src),
                                "is_svg": True,
                                "alt": alt
                            })
                        else:
                            # Non-SVG images save to disk
                            safe_prefix = "".join(c for c in training_session.get('filename', 'html') if c.isalnum())[:10]
                            unique_filename = f"{safe_prefix}_img{i+1}_{str(uuid.uuid4())[:8]}.{ext}"
                            file_path_static = f"static/uploads/{unique_filename}"
                            
                            # Ensure upload directory exists
                            os.makedirs("static/uploads", exist_ok=True)
                            
                            # Save image to disk
                            with open(file_path_static, "wb") as f:
                                f.write(image_data)
                            
                            # Generate URL
                            image_url = f"/api/static/uploads/{unique_filename}"
                            
                            all_images.append({
                                "filename": unique_filename,
                                "url": image_url,
                                "size": len(image_data),
                                "is_svg": False,
                                "alt": alt
                            })
                            
                            print(f"✅ Extracted HTML embedded image: {unique_filename}")
                    
                    except Exception as img_error:
                        print(f"⚠️ Error processing embedded image {i+1}: {img_error}")
                        continue
                
                elif src.startswith('http') or src.startswith('/'):
                    # External or relative URL - add reference to text
                    full_text += f"\n[Referenced image: {src} - Alt: {alt}]\n"
                    
        except ImportError:
            print("⚠️ BeautifulSoup not available, using basic text extraction")
            # Fallback: basic HTML tag removal
            import re
            full_text = re.sub(r'<[^>]+>', '', html_content)
            all_images = []
        except Exception as e:
            print(f"⚠️ Error parsing HTML: {e}")
            full_text = html_content
            all_images = []
        
        print(f"✅ HTML extraction complete: {len(full_text)} characters, {len(all_images)} images")
        
        # Process with template
        articles = await create_articles_with_template(full_text, all_images, template_data, training_session)
        
        print(f"✅ HTML processing generated {len(articles)} articles")
        
        return articles
        
    except Exception as e:
        print(f"❌ HTML processing error: {e}")
        return []

async def create_recovery_articles(content: str, images: list, template_data: dict, training_session: dict) -> list:
    """Recovery function for enhanced processing failures - creates simplified articles"""
    try:
        print(f"🔧 Recovery processing: {len(content)} chars, {len(images)} images")
        
        # Create a single simplified article with all content
        article_id = str(uuid.uuid4())
        
        # Basic HTML structure
        html_content = f"<h1>Document Content from {training_session.get('filename', 'Document')}</h1>\n"
        html_content += content
        
        # Add images at the end
        for i, img in enumerate(images):
            html_content += f'\n<figure class="embedded-image">'
            html_content += f'<img src="{img["url"]}" alt="{img["alt_text"]}" style="max-width: 100%; height: auto; margin: 1rem 0;">'
            html_content += f'<figcaption>{img.get("caption", f"Figure {i + 1}")}</figcaption>'
            html_content += f'</figure>\n'
        
        # Create media array
        media_array = []
        for img in images:
            media_array.append({
                "url": img['url'],
                "alt": img['alt_text'],
                "caption": img.get('caption', ''),
                "placement": img.get('placement', 'inline'),
                "filename": img['filename']
            })
        
        article = {
            "id": article_id,
            "title": f"Recovery Article from {training_session['filename']}",
            "html": html_content,
            "markdown": html_content.replace('<h1>', '# ').replace('</h1>', '').replace('<p>', '').replace('</p>', '\n'),
            "content": html_content,
            "media": media_array,
            "tags": ["extracted", "recovery", "simplified"],
            "status": "training",
            "template_id": training_session['template_id'],
            "session_id": training_session['session_id'],
            "word_count": len(html_content.split()),
            "image_count": len(media_array),
            "format": "html",
            "created_at": datetime.utcnow().isoformat(),
            "ai_processed": False,
            "ai_model": "recovery_mode",
            "training_mode": True,
            "metadata": {
                "article_number": 1,
                "source_filename": training_session['filename'],
                "template_applied": training_session['template_id'],
                "phase": "recovery_processing"
            }
        }
        
        print(f"✅ Recovery article created: {len(html_content)} chars, {len(media_array)} images")
        return [article]
        
    except Exception as e:
        print(f"❌ Recovery processing also failed: {e}")
        return []

async def create_recovery_articles(content: str, images: list, template_data: dict, training_session: dict) -> list:
    """Recovery function for enhanced processing failures - creates simplified articles"""
    try:
        print(f"🔧 Recovery processing: {len(content)} chars, {len(images)} images")
        
        # Create a single simplified article with all content
        article_id = str(uuid.uuid4())
        
        # Basic HTML structure
        html_content = f"<h1>Document Content from {training_session.get('filename', 'Document')}</h1>\n"
        html_content += content
        
        # Add images at the end
        for i, img in enumerate(images):
            html_content += f'\n<figure class="embedded-image">'
            html_content += f'<img src="{img["url"]}" alt="{img["alt_text"]}" style="max-width: 100%; height: auto; margin: 1rem 0;">'
            html_content += f'<figcaption>{img.get("caption", f"Figure {i + 1}")}</figcaption>'
            html_content += f'</figure>\n'
        
        # Create media array
        media_array = []
        for img in images:
            media_array.append({
                "url": img['url'],
                "alt": img['alt_text'],
                "caption": img.get('caption', ''),
                "placement": img.get('placement', 'inline'),
                "filename": img['filename']
            })
        
        article = {
            "id": article_id,
            "title": f"Recovery Article from {training_session['filename']}",
            "html": html_content,
            "markdown": html_content.replace('<h1>', '# ').replace('</h1>', '').replace('<p>', '').replace('</p>', '\n'),
            "content": html_content,
            "media": media_array,
            "tags": ["extracted", "recovery", "simplified"],
            "status": "training",
            "template_id": training_session['template_id'],
            "session_id": training_session['session_id'],
            "word_count": len(html_content.split()),
            "image_count": len(media_array),
            "format": "html",
            "created_at": datetime.utcnow().isoformat(),
            "ai_processed": False,
            "ai_model": "recovery_mode",
            "training_mode": True,
            "metadata": {
                "article_number": 1,
                "source_filename": training_session['filename'],
                "template_applied": training_session['template_id'],
                "phase": "recovery_processing"
            }
        }
        
        print(f"✅ Recovery article created: {len(html_content)} chars, {len(media_array)} images")
        return [article]
        
    except Exception as e:
        print(f"❌ Recovery processing also failed: {e}")
        return []

async def create_articles_with_template(content: str, images: list, template_data: dict, training_session: dict) -> list:
    """Create articles using template specifications - processing full content without artificial limits"""
    try:
        content_length = len(content)
        image_count = len(images)
        
        print(f"📊 Content analysis: {content_length} chars, {image_count} images")
        print(f"🔍 DEBUG - Template data keys: {list(template_data.keys())}")
        print(f"🔍 DEBUG - Training session keys: {list(training_session.keys())}")
        print(f"🔍 DEBUG - Content preview: {content[:200]}...")
        
        # Check if content is empty
        if not content or not content.strip():
            print("❌ DEBUG - Content is empty or only whitespace")
            return []
            
        # Remove artificial limits and process full content
        # Use natural content structure to determine article splitting
        articles = []
        
        # Analyze content for natural breaking points
        natural_sections = []
        
        # CRITICAL FIX: Add Markdown to HTML conversion before content analysis
        # Check if content contains Markdown H1 headers and convert to HTML for proper chunking
        if '#' in content and re.search(r'^#+\s', content, re.MULTILINE):
            print(f"📝 Markdown H1 headers detected - converting to HTML for proper chunking")
            original_content = content
            content = convert_markdown_to_html_for_text_processing(content)
            h1_count = len(re.findall(r'<h1>', content))
            print(f"✅ Markdown conversion complete - {h1_count} H1 tags found")
            
            # If we found H1 tags after conversion, we should create multiple articles
            if h1_count > 1:
                print(f"🎯 Multiple H1 sections detected - will create {h1_count} separate articles")
        
        # Look for major headings and section breaks
        if '<h1>' in content or '<h2>' in content or '\n\n' in content:
            # Content has structure, split intelligently
            
            # ENHANCED FIX: Create separate articles for each H1 section
            if '<h1>' in content:
                h1_sections = content.split('<h1>')
                print(f"🔍 Splitting content on H1 tags - found {len(h1_sections)} sections")
                
                for i, section in enumerate(h1_sections):
                    if section.strip():
                        if i == 0 and not section.startswith('<h1>'):
                            # First section without H1 prefix - could be intro content
                            if len(section.strip()) > 200:  # Only include substantial intro content
                                natural_sections.append(section)
                                print(f"✅ Added intro section: {len(section)} chars")
                        else:
                            # Restore H1 tag and treat as separate article
                            if not section.startswith('<h1>'):
                                section = '<h1>' + section
                            natural_sections.append(section)
                            print(f"✅ Added H1 section {i}: {len(section)} chars")
                            
                print(f"🎯 Created {len(natural_sections)} separate sections from H1 content")
            elif '<h2>' in content:
                sections = content.split('<h2>')
                for i, section in enumerate(sections):
                    if section.strip():
                        if i > 0:  # Add back the h2 tag
                            section = '<h2>' + section
                        natural_sections.append(section)
            else:
                # Split on double line breaks for paragraph-based content
                sections = content.split('\n\n')
                current_section = ""
                
                for section in sections:
                    # Aim for sections of reasonable length (2000-8000 chars)
                    if len(current_section + section) > 8000 and current_section:
                        natural_sections.append(current_section.strip())
                        current_section = section
                    else:
                        current_section += "\n\n" + section if current_section else section
                
                if current_section.strip():
                    natural_sections.append(current_section.strip())
        else:
            # No clear structure, treat as single section
            natural_sections = [content]
        
        # Filter out very small sections and merge with adjacent ones
        filtered_sections = []
        for i, section in enumerate(natural_sections):
            if len(section.strip()) < 500 and i < len(natural_sections) - 1:
                # Merge small section with next one
                natural_sections[i + 1] = section + "\n\n" + natural_sections[i + 1]
            elif len(section.strip()) >= 100:  # Only include sections with substantial content
                filtered_sections.append(section)
        
        natural_sections = filtered_sections if filtered_sections else [content]
        
        print(f"📝 Identified {len(natural_sections)} natural content sections")
        
        # Enhanced image distribution with contextual matching
        distributed_images = []
        
        # ENHANCED: For text files, check for image references in content
        if not images and any(indicator in content.lower() for indicator in ['[image:', '![', '<img', 'image.png', 'image.jpg', 'figure']):
            print("🖼️ Text content contains image references - attempting to extract image information")
            # Extract potential image references from text content
            image_refs = extract_image_references_from_text(content)
            if image_refs:
                distributed_images = image_refs
                print(f"📸 Found {len(image_refs)} image references in text content")
        elif images:
            # Existing image distribution logic for actual images
            if len(natural_sections) > 1 and len(images) > 0:
                images_per_section = max(1, len(images) // len(natural_sections))
                for i, section in enumerate(natural_sections):
                    start_idx = i * images_per_section
                    end_idx = min(start_idx + images_per_section, len(images))
                    section_images = images[start_idx:end_idx]
                    distributed_images.extend(section_images)
            else:
                distributed_images = images
        section_images = distribute_images_contextually(natural_sections, images)
        
        # Create articles from natural sections with enhanced processing
        for i, section in enumerate(natural_sections):
            if section.strip():
                assigned_images = section_images[i] if i < len(section_images) else []
                
                print(f"📄 Creating comprehensive article {i+1} with {len(assigned_images)} images")
                
                article = await create_single_article_with_template(
                    section, 
                    assigned_images,
                    template_data, 
                    training_session,
                    i + 1,
                    len(natural_sections)
                )
                
                if article:
                    articles.append(article)
        
        print(f"✅ Created {len(articles)} comprehensive articles from {content_length} chars and {image_count} images")
        return articles
        
    except Exception as e:
        print(f"❌ Enhanced article creation error: {e}")
        import traceback
        traceback.print_exc()
        return []

def analyze_document_structure(content: str) -> list:
    """
    Analyze document structure to determine natural breaking points
    """
    sections = []
    
    # Enhanced structure detection
    if '<h1>' in content or '<h2>' in content:
        # Split on major headings with intelligent merging
        if '<h1>' in content:
            sections = intelligent_split_on_headings(content, 'h1')
        else:
            sections = intelligent_split_on_headings(content, 'h2')
    elif content.count('\n\n') > 10:
        # Handle paragraph-based documents
        sections = intelligent_paragraph_grouping(content)
    else:
        # Single comprehensive section for shorter documents
        sections = [content]
    
    # Filter and validate sections
    validated_sections = []
    for section in sections:
        if len(section.strip()) > 200:  # Minimum meaningful content
            validated_sections.append(section.strip())
        elif validated_sections:
            # Merge small sections with previous one
            validated_sections[-1] += "\n\n" + section.strip()
    
    return validated_sections if validated_sections else [content]

def intelligent_split_on_headings(content: str, heading_tag: str) -> list:
    """
    Intelligently split content on headings while preserving context
    """
    import re
    
    # Pattern to match heading tags
    pattern = f'<{heading_tag}[^>]*>(.*?)</{heading_tag}>'
    headings = re.findall(pattern, content, re.IGNORECASE | re.DOTALL)
    
    if not headings:
        return [content]
    
    # Split and reconstruct sections
    sections = []
    parts = re.split(f'<{heading_tag}[^>]*>.*?</{heading_tag}>', content, flags=re.IGNORECASE | re.DOTALL)
    heading_matches = re.finditer(pattern, content, re.IGNORECASE | re.DOTALL)
    
    # Reconstruct sections with their headings
    for i, match in enumerate(heading_matches):
        section_content = match.group(0)  # Include the heading
        if i + 1 < len(parts):
            section_content += parts[i + 1]  # Add content after heading
        
        if section_content.strip():
            sections.append(section_content.strip())
    
    return sections if sections else [content]

def intelligent_paragraph_grouping(content: str) -> list:
    """
    Group paragraphs intelligently to form coherent sections
    """
    paragraphs = content.split('\n\n')
    sections = []
    current_section = ""
    
    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
            
        # Check if this paragraph starts a new logical section
        if is_section_boundary(paragraph, current_section):
            if current_section.strip():
                sections.append(current_section.strip())
            current_section = paragraph
        else:
            current_section += "\n\n" + paragraph if current_section else paragraph
        
        # Ensure sections don't get too long (but allow substantial content)
        if len(current_section) > 15000:  # Increased from 8000
            sections.append(current_section.strip())
            current_section = ""
    
    if current_section.strip():
        sections.append(current_section.strip())
    
    return sections

def is_section_boundary(paragraph: str, current_section: str) -> bool:
    """
    Determine if a paragraph represents a natural section boundary
    """
    # Check for topic transition indicators
    transition_indicators = [
        'next step', 'following section', 'in addition', 'furthermore',
        'alternatively', 'however', 'on the other hand', 'similarly',
        'step ', 'phase ', 'part ', 'section ', 'chapter '
    ]
    
    paragraph_lower = paragraph.lower()
    
    # Strong indicators of new section
    if any(indicator in paragraph_lower for indicator in transition_indicators):
        return True
    
    # Check if starting a numbered list or procedure
    if re.match(r'^\d+\.', paragraph.strip()):
        return len(current_section) > 1000  # Only if current section has substance
    
    return False

def distribute_images_contextually(sections: list, images: list) -> list:
    """
    Distribute images across sections based on contextual relevance
    """
    if not images:
        return [[] for _ in sections]
    
    section_images = [[] for _ in sections]
    
    for image in images:
        best_section_idx = find_best_section_for_image(image, sections)
        if best_section_idx is not None:
            section_images[best_section_idx].append(image)
        else:
            # Distribute evenly if no clear match
            min_images_idx = min(range(len(section_images)), key=lambda i: len(section_images[i]))
            section_images[min_images_idx].append(image)
    
    return section_images

def find_best_section_for_image(image: dict, sections: list) -> int:
    """
    Find the best section for an image based on contextual matching
    """
    image_chapter = image.get('chapter', '').lower()
    image_context = image.get('context_text', '').lower()
    
    best_score = 0
    best_section = None
    
    for i, section in enumerate(sections):
        section_lower = section.lower()
        score = 0
        
        # Match based on chapter name
        if image_chapter and image_chapter in section_lower:
            score += 50
        
        # Match based on context keywords
        context_words = image_context.split()[:10]  # First 10 words of context
        for word in context_words:
            if len(word) > 3 and word in section_lower:
                score += 5
        
        # Prefer sections with visual references
        if any(keyword in section_lower for keyword in ['figure', 'diagram', 'image', 'shows', 'illustrates']):
            score += 20
        
        if score > best_score:
            best_score = score
            best_section = i
    
    return best_section if best_score > 10 else None

async def create_single_article_with_template(content: str, images: list, template_data: dict, training_session: dict, article_number: int, total_articles: int = 1) -> dict:
    """Create a single comprehensive article with enhanced structure and quality using segmented generation"""
    try:
        print(f"🔍 Creating comprehensive article {article_number}/{total_articles} with {len(images)} images")
        
        # Generate intelligent title based on content, not filename
        title = extract_h1_title_from_content(content) or generate_contextual_title(content, article_number, training_session)
        
        # CRITICAL FIX: Use segmented generation for comprehensive coverage - OPTIMIZED
        if len(content) > 5000:  # Increased threshold from 2000 to 5000 chars
            print("📝 Using segmented generation for comprehensive coverage")
            final_content = await generate_comprehensive_article_segmented(content, images, template_data, title)
        else:
            print("📝 Using single-pass generation for shorter content")
            final_content = await generate_single_pass_article(content, images, template_data, title)
        
        if not final_content:
            print("⚠️ No AI content generated, creating enhanced fallback")
            final_content = create_structured_fallback_content(content, images, title)
        
        # Post-process the AI content for quality assurance
        final_content = enhance_content_quality(final_content, images)
        
        # Extract or generate final title from content
        final_title = extract_content_title(final_content) or title
        
        # Create comprehensive article with enhanced metadata
        article = {
            "id": str(uuid.uuid4()),
            "title": final_title,
            "content": final_content,
            "word_count": len(final_content.split()),
            "image_count": len(images),
            "status": "training",
            "template_id": template_data.get("template_id", "enhanced_processing"),
            "session_id": training_session.get("session_id"),
            "training_mode": True,
            "ai_processed": True,
            "ai_model": "enhanced_knowledge_engine",
            "has_images": len(images) > 0,
            "has_structure": check_content_structure(final_content),
            "processing_metadata": {
                "article_number": article_number,
                "total_articles": total_articles,
                "original_length": len(content),
                "enhanced_length": len(final_content),
                "images_embedded": count_embedded_images(final_content),
                "generation_method": "segmented" if len(content) > 2000 else "single_pass"
            }
        }
        
        print(f"✅ Comprehensive article created: '{final_title}' ({article['word_count']} words)")
        return article
        
    except Exception as e:
        print(f"❌ Enhanced article creation error: {e}")
        import traceback
        traceback.print_exc()
        return None

async def generate_comprehensive_article_segmented(content: str, images: list, template_data: dict, title: str) -> str:
    """Generate comprehensive article using segmented approach for full coverage"""
    try:
        print("🔄 Starting segmented generation for comprehensive coverage")
        
        # Step 1: Generate article outline and structure
        outline = await generate_article_outline(content, title, template_data)
        if not outline:
            print("⚠️ Could not generate outline, falling back to single-pass")
            return await generate_single_pass_article(content, images, template_data, title)
        
        print(f"📋 Generated article outline with sections")
        
        # Step 2: Split content into logical segments based on outline
        content_segments = split_content_into_segments(content, outline)
        print(f"📄 Split content into {len(content_segments)} segments")
        
        # Step 3: Distribute images across segments
        segment_images = distribute_images_to_segments(images, content_segments)
        
        # Step 4: Generate each segment comprehensively
        generated_sections = []
        generated_sections.append(f"<h1>{title}</h1>\n")
        
        for i, (segment, segment_imgs) in enumerate(zip(content_segments, segment_images)):
            print(f"🔄 Generating comprehensive section {i+1}/{len(content_segments)}")
            
            section_content = await generate_content_segment(
                segment, 
                segment_imgs, 
                template_data, 
                i + 1, 
                len(content_segments),
                outline
            )
            
            if section_content:
                generated_sections.append(section_content)
            else:
                print(f"⚠️ Failed to generate section {i+1}, using fallback")
                generated_sections.append(create_fallback_segment(segment, segment_imgs))
        
        # Step 5: Combine all sections
        final_content = "\n\n".join(generated_sections)
        
        print(f"✅ Segmented generation completed: {len(final_content.split())} words")
        return final_content
        
    except Exception as e:
        print(f"❌ Segmented generation error: {e}")
        return await generate_single_pass_article(content, images, template_data, title)

async def generate_article_outline(content: str, title: str, template_data: dict) -> dict:
    """Generate a structured outline for the article"""
    try:
        system_message = """You are an expert content strategist creating comprehensive article outlines.

Generate a detailed outline for a comprehensive knowledge base article. Return ONLY a JSON structure with no additional text.

Required JSON format:
{
    "title": "Article Title",
    "sections": [
        {
            "heading": "Section Title",
            "level": 2,
            "key_points": ["point 1", "point 2", "point 3"],
            "estimated_words": 500
        }
    ],
    "total_estimated_words": 2500
}

Create 2-3 major sections with detailed key points. Aim for balanced coverage with 1500-2500 total words."""

        user_message = f"""Create a comprehensive outline for this content:

TITLE: {title}

CONTENT TO OUTLINE:
{content[:2000]}...

Create a detailed outline with major sections, key points, and estimated word counts for comprehensive coverage."""

        outline_response = await call_llm_with_fallback(system_message, user_message)
        
        if outline_response:
            import json
            try:
                outline_data = json.loads(outline_response)
                if "sections" in outline_data and len(outline_data["sections"]) > 0:
                    return outline_data
            except json.JSONDecodeError:
                print("⚠️ Could not parse outline JSON")
        
        return None
        
    except Exception as e:
        print(f"❌ Outline generation error: {e}")
        return None

def split_content_into_segments(content: str, outline: dict) -> list:
    """Split content into segments based on outline structure"""
    if not outline or "sections" not in outline:
        # Fallback: split into 2 balanced segments for faster processing
        words = content.split()
        segments_count = 2  # Reduced from 4 to 2 segments
        segment_size = max(200, len(words) // segments_count)  # At least 200 words per segment
        
        segments = []
        for i in range(0, len(words), segment_size):
            segment_words = words[i:i + segment_size]
            if segment_words:  # Only add non-empty segments
                segments.append(" ".join(segment_words))
        
        return segments
    
    # Try to split based on outline sections - OPTIMIZED to 2 segments
    sections = outline["sections"]
    # Limit to maximum 2 segments for faster processing
    target_segments = min(2, len(sections))  # Reduced from 4-6 to 2 segments
    total_words = len(content.split())
    words_per_section = total_words // target_segments
    
    segments = []
    words = content.split()
    
    for i in range(target_segments):
        start_idx = i * words_per_section
        end_idx = min((i + 1) * words_per_section, len(words)) if i < target_segments - 1 else len(words)
        
        segment_words = words[start_idx:end_idx]
        if segment_words:
            segments.append(" ".join(segment_words))
    
    return segments

def distribute_images_to_segments(images: list, segments: list) -> list:
    """Distribute images across content segments"""
    if not images:
        return [[] for _ in segments]
    
    # Distribute images roughly evenly across segments
    segment_images = [[] for _ in segments]
    
    for i, image in enumerate(images):
        segment_idx = i % len(segments)
        segment_images[segment_idx].append(image)
    
    return segment_images

async def generate_content_segment(segment_content: str, segment_images: list, template_data: dict, segment_num: int, total_segments: int, outline: dict) -> str:
    """Generate a comprehensive content segment"""
    try:
        section_info = ""
        if outline and "sections" in outline and segment_num <= len(outline["sections"]):
            section_data = outline["sections"][segment_num - 1]
            section_info = f"""
SECTION FOCUS: {section_data.get('heading', f'Section {segment_num}')}
KEY POINTS TO COVER: {', '.join(section_data.get('key_points', []))}
TARGET LENGTH: {section_data.get('estimated_words', 600)} words"""

        system_message = f"""You are an expert technical writer creating comprehensive section content.

CRITICAL REQUIREMENTS:
1. Generate detailed, comprehensive content for this section
2. Write 400-800 words for thorough coverage - BALANCED LENGTH
3. Use professional technical documentation style with good detail
4. Include detailed explanations, comprehensive step-by-step procedures, and thorough information
5. Use proper HTML structure: <h2>, <h3>, <p>, <ul>, <ol>, <li>, <strong>, <em>
6. Embed provided images with proper figure elements
7. NO meta-commentary - only detailed article content
8. NO truncation or summarization - provide COMPLETE detailed content
9. Focus on quality over quantity - comprehensive but concise

{section_info}

QUALITY STANDARDS:
- Comprehensive, detailed explanations with good depth
- Professional enterprise technical writing with thorough coverage
- Complete step-by-step procedures with detailed instructions
- Thorough coverage of all aspects with focused information
- Proper HTML semantic structure with rich formatting
- Target 400-800 words - comprehensive yet efficient content"""

        user_message = f"""Create comprehensive content for section {segment_num} of {total_segments}:

CONTENT TO EXPAND:
{segment_content}

AVAILABLE IMAGES: {len(segment_images)}
{format_available_images(segment_images)}

CRITICAL REQUIREMENTS:
- Write 400-800 words for comprehensive coverage - BALANCED LENGTH
- Include detailed explanations and complete comprehensive procedures
- Use proper HTML structure with headings and rich formatting
- Embed images contextually with provided HTML code
- Focus on thorough, professional technical documentation with good detail
- NO truncation or summarization - provide complete detailed comprehensive content
- Balance depth with efficiency - comprehensive but focused content
- Provide complete step-by-step instructions where applicable
- Include thorough background information and context

Generate comprehensive section content with good detail, proper HTML structure, and target 400-800 words."""

        segment_response = await call_llm_with_fallback(system_message, user_message)
        
        if segment_response:
            return segment_response.strip()
        
        return None
        
    except Exception as e:
        print(f"❌ Segment generation error: {e}")
        return None

def extract_h1_title_from_content(content: str) -> str:
    """Extract the first H1 heading from HTML content as the title"""
    import re
    from bs4 import BeautifulSoup
    
    try:
        # Try with regex first (faster)
        h1_match = re.search(r'<h1[^>]*>(.*?)</h1>', content, re.IGNORECASE | re.DOTALL)
        if h1_match:
            # Clean up the title text by removing any HTML tags
            title_text = re.sub(r'<[^>]+>', '', h1_match.group(1)).strip()
            # Remove anchor links and IDs
            title_text = re.sub(r'<a[^>]*>.*?</a>', '', title_text)
            if 5 < len(title_text) < 100:
                print(f"📋 Extracted H1 title: '{title_text}'")
                return title_text
        
        # Fallback to BeautifulSoup for complex HTML
        soup = BeautifulSoup(content, 'html.parser')
        h1_tag = soup.find('h1')
        if h1_tag:
            title_text = h1_tag.get_text().strip()
            if 5 < len(title_text) < 100:
                print(f"📋 Extracted H1 title (BeautifulSoup): '{title_text}'")
                return title_text
        
        print("⚠️ No valid H1 title found in content")
        return ""
        
    except Exception as e:
        print(f"❌ H1 title extraction failed: {e}")
        return ""

def generate_contextual_title(content: str, article_number: int, training_session: dict) -> str:
    """Generate intelligent title based on content analysis"""
    import re
    
    # Extract from existing headings first
    h1_match = re.search(r'<h1[^>]*>(.*?)</h1>', content, re.IGNORECASE | re.DOTALL)
    if h1_match:
        title_text = re.sub(r'<[^>]+>', '', h1_match.group(1)).strip()
        if 5 < len(title_text) < 100:
            return title_text
    
    h2_match = re.search(r'<h2[^>]*>(.*?)</h2>', content, re.IGNORECASE | re.DOTALL)
    if h2_match:
        title_text = re.sub(r'<[^>]+>', '', h2_match.group(1)).strip()
        if 5 < len(title_text) < 100:
            return title_text
    
    # Analyze content for key topics
    clean_content = re.sub(r'<[^>]+>', '', content)
    words = clean_content.split()[:100]  # First 100 words
    
    # Look for topic keywords
    topic_keywords = ['guide', 'process', 'procedure', 'step', 'method', 'system', 'overview', 'management']
    found_topics = []
    
    for i, word in enumerate(words):
        if word.lower() in topic_keywords and i > 0:
            # Get context around the keyword
            start_idx = max(0, i-3)
            end_idx = min(len(words), i+4)
            context = ' '.join(words[start_idx:end_idx])
            found_topics.append(context)
    
    if found_topics:
        # Use the first relevant topic
        topic = found_topics[0]
        topic = re.sub(r'[^\w\s]', '', topic).strip()
        if len(topic) > 10:
            return topic.title()
    
    # Fallback to source-based naming
    source_name = training_session.get('filename', 'Document').replace('.docx', '').replace('.pdf', '')
    source_name = source_name.replace('_', ' ').replace('-', ' ')
    
    if article_number == 1:
        return f"Understanding {source_name}"
    else:
        return f"{source_name} - Section {article_number}"

def format_available_images(images: list) -> str:
    """Format available images for LLM reference with URLs for embedding"""
    if not images:
        return "No images available - focus on comprehensive text content"
    
    image_refs = []
    for i, img in enumerate(images, 1):
        chapter = img.get('chapter', 'Document')
        filename = img.get('filename', f'image{i}')
        caption = img.get('caption', f'Figure {i}')
        url = img.get('url', '')
        alt_text = img.get('alt_text', f'Figure {i}')
        
        image_refs.append(f"""IMAGE_{i}: 
- Filename: {filename}
- URL: {url}
- Caption: {caption}
- Alt Text: {alt_text}
- Chapter: {chapter}
- HTML to embed: <figure class="embedded-image"><img src="{url}" alt="{alt_text}" style="max-width: 100%; height: auto;"><figcaption>{caption}</figcaption></figure>""")
    
    return "\n\n".join(image_refs)

def create_structured_fallback_content(content: str, images: list, title: str) -> str:
    """Create structured fallback when AI is unavailable"""
    html = f"<h1>{title}</h1>\n\n"
    
    # Process content into structured HTML
    paragraphs = content.split('\n\n')
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
            
        # Detect if paragraph looks like a heading
        if len(para) < 80 and (para.isupper() or para.endswith(':') or para.count(' ') < 8):
            html += f"<h2>{para.rstrip(':')}</h2>\n\n"
        elif para.startswith('•') or para.startswith('-'):
            # Convert to list
            html += f"<ul>\n<li>{para[1:].strip()}</li>\n</ul>\n\n"
        else:
            html += f"<p>{para}</p>\n\n"
    
    # Add images if available
    if images:
        html += "<h2>Related Figures</h2>\n\n"
        for img in images[:2]:  # Limit for fallback
            html += f"""<figure class="embedded-image">
<img src="{img.get('url', '')}" alt="{img.get('alt_text', 'Figure')}" style="max-width: 100%; height: auto;">
<figcaption>{img.get('caption', 'Figure')}</figcaption>
</figure>\n\n"""
    
    return html

def enhance_content_quality(content: str, images: list) -> str:
    """Post-process content for quality and image integration"""
    if not content:
        return content
    
    # CRITICAL FIX: Remove HTML document wrappers that may come from LLM
    content = clean_html_wrappers(content)
    
    # Add missing images if AI didn't embed them
    if images and not any(img.get('url', '') in content for img in images):
        content = add_missing_images(content, images)
    
    # Clean HTML formatting
    import re
    content = re.sub(r'\n\s*\n\s*\n', '\n\n', content)  # Remove excessive newlines
    content = re.sub(r'</p>\s*<p>', '</p>\n\n<p>', content)  # Proper paragraph spacing
    
    return content.strip()

def clean_html_wrappers(content: str) -> str:
    """Remove HTML document wrappers and extract only body content"""
    import re
    
    # Remove HTML document structure wrappers
    content = re.sub(r'<!DOCTYPE[^>]*>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'<html[^>]*>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'</html>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'<head>.*?</head>', '', content, flags=re.IGNORECASE | re.DOTALL)
    content = re.sub(r'<body[^>]*>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'</body>', '', content, flags=re.IGNORECASE)
    
    # ENHANCED FIX: Convert markdown code blocks to proper HTML <pre><code> tags instead of removing them
    # This preserves code content while ensuring proper HTML formatting
    
    # Convert markdown code blocks with language specifiers to HTML
    content = re.sub(r'```html\s*(.*?)```', r'<pre><code class="language-html">\1</code></pre>', content, flags=re.IGNORECASE | re.DOTALL)
    content = re.sub(r'```javascript\s*(.*?)```', r'<pre><code class="language-javascript">\1</code></pre>', content, flags=re.IGNORECASE | re.DOTALL)
    content = re.sub(r'```css\s*(.*?)```', r'<pre><code class="language-css">\1</code></pre>', content, flags=re.IGNORECASE | re.DOTALL)
    content = re.sub(r'```json\s*(.*?)```', r'<pre><code class="language-json">\1</code></pre>', content, flags=re.IGNORECASE | re.DOTALL)
    
    # Convert generic markdown code blocks to HTML
    content = re.sub(r'```\s*(.*?)```', r'<pre><code>\1</code></pre>', content, flags=re.IGNORECASE | re.DOTALL)
    
    # FIXED: Only remove truly empty code blocks (no content at all), preserve formatted code
    content = re.sub(r'<pre><code[^>]*></code></pre>', '', content, flags=re.IGNORECASE)
    
    # CRITICAL FIX: Remove duplicate h1 tags that repeat the title
    h1_matches = re.findall(r'<h1[^>]*>(.*?)</h1>', content, flags=re.IGNORECASE | re.DOTALL)
    if len(h1_matches) > 1:
        # Keep only the first h1, remove the rest
        first_h1_found = False
        def replace_h1(match):
            nonlocal first_h1_found
            if not first_h1_found:
                first_h1_found = True
                return match.group(0)  # Keep the first one
            else:
                # Convert subsequent h1s to h2s
                inner_text = match.group(1)
                return f'<h2>{inner_text}</h2>'
        
        content = re.sub(r'<h1[^>]*>(.*?)</h1>', replace_h1, content, flags=re.IGNORECASE | re.DOTALL)
    
    # Remove any remaining div wrappers that might have been added
    content = re.sub(r'^<div[^>]*>(.*)</div>$', r'\1', content.strip(), flags=re.IGNORECASE | re.DOTALL)
    
    # Ensure we still have content
    if not content.strip():
        return content
    
    # Ensure content starts with proper HTML tags (but not document structure)
    content = content.strip()
    if content and not content.startswith('<'):
        # Only wrap in paragraph if it's plain text
        if '\n' not in content[:100]:  # Single line text
            content = f"<p>{content}</p>"
        else:
            # Multi-line content, process paragraphs
            paragraphs = content.split('\n\n')
            html_parts = []
            for para in paragraphs:
                para = para.strip()
                if para:
                    html_parts.append(f"<p>{para}</p>")
            content = '\n\n'.join(html_parts)
    
    return content

def clean_article_html_content(content: str) -> str:
    """Clean HTML content to remove document structure while preserving rich formatting - WYSIWYG OPTIMIZED"""
    import re
    
    # CRITICAL FIX 1: Remove markdown HTML code block wrappers that break WYSIWYG editors
    if content.strip().startswith('```html'):
        print(f"🚨 CRITICAL FIX: Removing markdown HTML code block wrapper")
        content = content.strip()
        # Remove opening ```html
        content = re.sub(r'^```html\s*', '', content)
        # Remove closing ```
        content = re.sub(r'\s*```$', '', content)
    
    # CRITICAL FIX 2: Remove any full-article code block wrapping that breaks WYSIWYG editors
    if content.strip().startswith('<pre><code class="language-html">') and content.strip().endswith('</code></pre>'):
        print(f"🚨 CRITICAL FIX: Removing WYSIWYG-breaking code block wrapper")
        content = content.strip()
        content = content.replace('<pre><code class="language-html">', '', 1)
        content = content.rsplit('</code></pre>', 1)[0]
    
    # Remove any other variations of full-article wrapping
    content = re.sub(r'^<pre><code[^>]*>(.*)</code></pre>$', r'\1', content, flags=re.DOTALL)
    
    # ENHANCED FIX 3: Remove HTML document structure elements
    content = re.sub(r'<!DOCTYPE[^>]*>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'</?html[^>]*>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'</?head[^>]*>', '', content, flags=re.IGNORECASE) 
    content = re.sub(r'</?body[^>]*>', '', content, flags=re.IGNORECASE)
    
    # First apply the existing HTML wrapper cleaning
    content = clean_html_wrappers(content)
    
    # Additional cleaning specific to article content
    # Remove only document structure elements, preserve content elements
    content = re.sub(r'<meta[^>]*>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'<link[^>]*>', '', content, flags=re.IGNORECASE)
    content = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.IGNORECASE | re.DOTALL)
    content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.IGNORECASE | re.DOTALL)
    
    # PRESERVE rich formatting elements - DO NOT remove:
    # - Lists (ul, ol, li)
    # - Code blocks (pre, code) - but only legitimate ones, not wrappers
    # - Emphasis (strong, em, b, i)
    # - Callouts and notes (div with classes)
    # - Tables (table, tr, td, th)
    # - Blockquotes
    # - Definition lists (dl, dt, dd)
    
    # Clean up excessive whitespace but preserve code block formatting
    content = re.sub(r'\n\s*\n\s*\n', '\n\n', content)
    content = re.sub(r'^\s+|\s+$', '', content)
    
    # Preserve spacing in code blocks - don't collapse whitespace inside <pre> tags
    def preserve_pre_spacing(match):
        return match.group(0)  # Return unchanged
    
    # Apply whitespace normalization only outside of <pre> tags
    parts = re.split(r'(<pre[^>]*>.*?</pre>)', content, flags=re.IGNORECASE | re.DOTALL)
    cleaned_parts = []
    
    for i, part in enumerate(parts):
        if i % 2 == 0:  # Not inside <pre> tags
            # Normal whitespace cleaning for non-code content
            part = re.sub(r'[ \t]+', ' ', part)  # Normalize spaces but preserve line breaks
        else:  # Inside <pre> tags
            # Preserve all whitespace in code blocks
            pass
        cleaned_parts.append(part)
    
    content = ''.join(cleaned_parts)
    
    return content

def fix_fragmented_ordered_lists(content: str) -> str:
    """Fix fragmented ordered lists by consolidating them into continuous numbering"""
    import re
    from bs4 import BeautifulSoup
    
    try:
        # Parse the HTML content
        soup = BeautifulSoup(content, 'html.parser')
        
        # Find all ordered lists
        ol_tags = soup.find_all('ol')
        
        if len(ol_tags) <= 1:
            return content  # No fragmentation to fix
        
        # Group consecutive ordered lists
        ol_groups = []
        current_group = []
        
        for ol in ol_tags:
            # Check if this ol is immediately after the previous one (allowing for whitespace)
            if current_group:
                prev_ol = current_group[-1]
                # Get all siblings between previous and current
                siblings_between = []
                current = prev_ol.next_sibling
                while current and current != ol:
                    if hasattr(current, 'name') and current.name:  # It's a tag
                        siblings_between.append(current)
                    current = current.next_sibling
                
                # If there are significant elements between, start new group
                significant_elements = [s for s in siblings_between if hasattr(s, 'name') and s.name in ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'div', 'section']]
                
                if significant_elements:
                    # Start new group
                    ol_groups.append(current_group)
                    current_group = [ol]
                else:
                    # Continue current group
                    current_group.append(ol)
            else:
                current_group = [ol]
        
        # Add the last group
        if current_group:
            ol_groups.append(current_group)
        
        # Process each group to consolidate numbering
        for group in ol_groups:
            if len(group) > 1:
                # Consolidate the group
                first_ol = group[0]
                
                # Collect all list items from the group
                all_items = []
                for ol in group:
                    items = ol.find_all('li', recursive=False)
                    all_items.extend(items)
                
                # Clear the first ol and add all items
                first_ol.clear()
                for item in all_items:
                    # Clone the item to avoid moving issues
                    new_item = soup.new_tag('li')
                    new_item.string = item.get_text()
                    # Preserve any nested content
                    for child in item.children:
                        if hasattr(child, 'name'):
                            new_item.append(child.extract())
                        else:
                            new_item.append(str(child))
                    first_ol.append(new_item)
                
                # Remove the other ols in the group
                for ol in group[1:]:
                    ol.decompose()
        
        return str(soup)
        
    except Exception as e:
        print(f"⚠️ Error in fix_fragmented_ordered_lists: {e}")
        # Return original content if there's an error
        return content

def add_missing_images(content: str, images: list) -> str:
    """Add missing images at appropriate locations"""
    import re
    
    # Find insertion points after headings
    h2_positions = [m.end() for m in re.finditer(r'</h2>', content, re.IGNORECASE)]
    
    if not h2_positions:
        # Fall back to paragraph positions
        h2_positions = [m.end() for m in re.finditer(r'</p>', content, re.IGNORECASE)][:2]
    
    # Insert images at found positions
    offset = 0
    images_to_add = images[:min(len(h2_positions), 2)]  # Limit insertions
    
    for i, image in enumerate(images_to_add):
        if i < len(h2_positions):
            insert_pos = h2_positions[i] + offset
            
            img_html = f"""
<figure class="embedded-image">
<img src="{image.get('url', '')}" alt="{image.get('alt_text', 'Figure')}" style="max-width: 100%; height: auto;">
<figcaption>{image.get('caption', 'Figure')}</figcaption>
</figure>
"""
            content = content[:insert_pos] + img_html + content[insert_pos:]
            offset += len(img_html)
    
    return content

def extract_content_title(content: str) -> str:
    """Extract title from generated content"""
    import re
    
    h1_match = re.search(r'<h1[^>]*>(.*?)</h1>', content, re.IGNORECASE | re.DOTALL)
    if h1_match:
        title_text = re.sub(r'<[^>]+>', '', h1_match.group(1)).strip()
        if 5 < len(title_text) < 120:
            return title_text
    return None

def _extract_title_from_html(html_content: str, fallback_title: str = 'Generated Article') -> str:
    """Extract title from HTML content, with fallback to provided title"""
    import re
    from bs4 import BeautifulSoup
    
    try:
        # Try to extract from h1 tag first
        h1_match = re.search(r'<h1[^>]*>(.*?)</h1>', html_content, re.IGNORECASE | re.DOTALL)
        if h1_match:
            title_text = re.sub(r'<[^>]+>', '', h1_match.group(1)).strip()
            if 5 < len(title_text) < 120:
                return title_text
        
        # Try to extract from h2 tag as fallback
        h2_match = re.search(r'<h2[^>]*>(.*?)</h2>', html_content, re.IGNORECASE | re.DOTALL)
        if h2_match:
            title_text = re.sub(r'<[^>]+>', '', h2_match.group(1)).strip()
            if 5 < len(title_text) < 120:
                return title_text
        
        # Try to extract from title tag
        title_match = re.search(r'<title[^>]*>(.*?)</title>', html_content, re.IGNORECASE | re.DOTALL)
        if title_match:
            title_text = re.sub(r'<[^>]+>', '', title_match.group(1)).strip()
            if 5 < len(title_text) < 120:
                return title_text
        
        # If no suitable title found, return fallback
        return fallback_title
        
    except Exception as e:
        print(f"⚠️ Error extracting title from HTML: {e}")
        return fallback_title

def check_content_structure(content: str) -> bool:
    """Check if content has proper structure"""
    import re
    has_headings = bool(re.search(r'<h[1-6][^>]*>', content, re.IGNORECASE))
    has_lists = bool(re.search(r'<[uo]l[^>]*>', content, re.IGNORECASE))
    return has_headings or has_lists

def count_embedded_images(content: str) -> int:
    """Count embedded images in content"""
    import re
    return len(re.findall(r'<img[^>]*>', content, re.IGNORECASE))

# PHASE 2: ADVANCED REFINED ENGINE v2.1 - Process content endpoint
@app.post("/api/content/process")
async def process_content(request: ContentProcessRequest):
    """V2 ENGINE: Process text content with AI"""
    print(f"🚀 V2 ENGINE: Processing text content - engine=v2")
    
    # KE-PR1: Add structured logging with job_id
    start_time = time.time()
    job_id = str(uuid.uuid4())
    if logger:
        logger.info({
            "event": "content_process_start", 
            "job_id": job_id, 
            "content_type": request.content_type,
            "stage": "content_process"
        })
    
    try:
        job = ProcessingJob(
            job_id=job_id,  # Use our job_id
            input_type=request.content_type,
            status="processing"
        )
        
        # Store job in database
        await db.processing_jobs.insert_one(job.dict())
        print(f"📊 V2 ENGINE: Job created - {job.job_id} - engine=v2")
        
        # V2 PROCESSING: Route to V2 text processor with timeout protection
        if request.content_type in ["text", "markdown"]:
            # Add timeout wrapper similar to file upload endpoint
            async def process_with_timeout():
                return await process_text_content_v2_pipeline(request.content, request.metadata)
            
            try:
                # 10-minute timeout for V2 processing pipeline (same as upload endpoint)
                articles = await asyncio.wait_for(process_with_timeout(), timeout=600)
                
                # TICKET 1 FIX: Ensure V2 articles are saved to content_library with proper format
                if articles:
                    print(f"💾 V2 ENGINE: Saving {len(articles)} V2-processed articles to content library")
                    saved_articles = []
                    for article in articles:
                        try:
                            # Ensure V2 metadata is present
                            if not article.get('metadata'):
                                article['metadata'] = {}
                            article['metadata']['engine'] = 'v2'
                            article['metadata']['processing_version'] = '2.0'
                            article['metadata']['format'] = article.get('format', 'html_canonical')
                            
                            # Save to content library
                            # KE-PR9.3: Use V2 validation repository
                from engine.stores.mongo import RepositoryFactory
                validation_repo = RepositoryFactory.get_v2_validation()
                await validation_repo.store_validation(validation_result)
            print(f"💾 V2 ENGINE: Stored validation result for diagnostics - validation_id: {validation_result.get('validation_id')} - engine=v2")
        except Exception as validation_storage_error:
            print(f"❌ V2 ENGINE: Error storing validation result - {validation_storage_error} - engine=v2")
        
        print(f"✅ V2 ENGINE: Step 8 complete - Validation status: {validation_status} - engine=v2")
        
        # V2 STEP 9: Cross-Article QA (dedupe, link validation, FAQ consolidation, terminology)
        print(f"🔍 V2 ENGINE: Starting Step 9 - Cross-Article QA for coherence and consistency - engine=v2")
        
        # Perform cross-article quality assurance
        qa_result = await v2_cross_article_qa_system.perform_cross_article_qa(
            generated_articles_result, run_id
        )
        
        # Check QA status and update articles accordingly
        qa_status = qa_result.get('qa_status', 'unknown')
        issues_found = qa_result.get('summary', {}).get('issues_found', 0)
        
        if qa_status == 'error':
            print(f"❌ V2 ENGINE: Step 9 QA failed with error - run {run_id} - engine=v2")
        elif issues_found == 0:
            print(f"✅ V2 ENGINE: Step 9 QA passed - No coherence issues found - engine=v2")
            # Add QA metadata to articles
            for article in articles:
                article.setdefault('metadata', {})['qa_result'] = qa_result
                article['qa_status'] = 'passed'
        else:
            print(f"⚠️ V2 ENGINE: Step 9 QA found {issues_found} issues - Articles require consolidation - engine=v2")
            # Mark articles with QA issues
            for article in articles:
                article.setdefault('metadata', {})['qa_result'] = qa_result
                article['qa_status'] = 'issues_found'
                article['qa_issues_count'] = issues_found
        
        # Store QA result separately for analysis
        try:
            # KE-PR9.3: Use repository pattern for find operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                result = await content_repo.collection.find_one({"id": article["id"]}):
                        # KE-PR9.3: Use repository pattern for content_library operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(article)
            print(f"✅ Created article with related links: '{article['title']}'")
        return articles
    else:
        # Create single comprehensive article
        article = await create_single_article_from_content(full_content, metadata)
        if article:
            print(f"🔍 DEBUG: Before DB insert - article keys: {list(article.keys())}")
            print(f"🔍 DEBUG: Before DB insert - has content: {'content' in article}")
            print(f"🔍 DEBUG: Before DB insert - content preview: {article.get('content', 'NO CONTENT')[:100]}...")
            # KE-PR9.3: Use repository pattern for content_library operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(article_record)
    print(f"✅ Created basic Content Library article: {article_record['title']}")
    return [article_record]


async def inject_real_images_into_articles():
    """Post-processing function to inject real extracted images into Content Library articles"""
    print("🎯 STARTING: Real image injection into existing articles")
    
    try:
        # Get all articles from Content Library that DON'T already have real images
        articles_cursor = db.content_library.find({
            "$or": [
                {"has_images": {"$ne": True}},
                {"has_images": {"$exists": False}}
            ]
        })
        articles = await articles_cursor.to_list(length=None)
        
        print(f"📚 Found {len(articles)} articles to process for image injection")
        
        # Get all available real images from static/uploads directory
        uploads_dir = "/app/backend/static/uploads"
        available_images = []
        
        if os.path.exists(uploads_dir):
            for filename in os.listdir(uploads_dir):
                if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp')):
                    # Include ALL images - both direct uploads and session images
                    image_url = f"/api/static/uploads/{filename}"
                    # Create a more generic source name for matching
                    source_parts = filename.replace('_img_', '_').split('_')
                    source_name = source_parts[0] if source_parts else filename.split('.')[0]
                    
                    available_images.append({
                        'filename': filename,
                        'url': image_url,
                        'source_name': source_name.lower()
                    })
        
        print(f"🖼️ Found {len(available_images)} real images available for injection")
        
        injection_count = 0
        
        for article in articles:
            article_id = article.get('_id')
            title = article.get('title', '')
            content = article.get('content', '')
            
            # Skip articles that already have contextual images or are too short
            if 'contextual-image' in content or len(content) < 500:
                continue
                
            # Find matching images - be more generous with matching
            matching_images = []
            
            # Extract key terms for matching
            title_lower = title.lower()
            content_lower = content.lower()
            
            # Try to match images by source name
            for image in available_images:
                source_name = image['source_name']
                
                # Match by keywords
                if (any(word in source_name for word in ['billing', 'management', 'google', 'map', 'javascript', 'api', 'test']) or
                    any(word in title_lower for word in [source_name]) or
                    any(part in content_lower for part in source_name.split('-')[:2] if len(part) > 3)):
                    matching_images.append(image)
                    if len(matching_images) >= 3:  # Max 3 per article
                        break
            
            # If no specific matches, use some available images anyway
            if len(matching_images) == 0:
                matching_images = available_images[:2]  # Use first 2 available images
            
            if matching_images:
                print(f"📝 Processing article: '{title[:50]}...' with {len(matching_images)} images")
                
                # Inject images contextually into the article
                enhanced_content = inject_images_contextually(content, matching_images)
                
                if enhanced_content != content and len(enhanced_content) > len(content):
                    # Update article in database
                    await db.content_library.update_one(
                        {"_id": article_id},
                        {"$set": {
                            "content": enhanced_content, 
                            "has_images": True, 
                            "image_injection_timestamp": datetime.utcnow(),
                            "injected_image_count": len(matching_images)
                        }}
                    )
                    
                    injection_count += 1
                    print(f"✅ Injected {len(matching_images)} images into: '{title[:50]}...'")
        
        print(f"🎉 COMPLETED: Successfully injected images into {injection_count} articles")
        return injection_count
        
    except Exception as e:
        print(f"❌ Image injection failed: {e}")
        import traceback
        print(traceback.format_exc())
        return 0

def inject_images_contextually(content: str, images: list) -> str:
    """Inject images contextually throughout article content"""
    if not images or not content:
        return content
    
    print(f"🔧 Injecting {len(images)} images contextually into content")
    
    # Split content into paragraphs
    paragraphs = content.split('</p>')
    if len(paragraphs) < 2:
        # Try splitting by line breaks if no HTML paragraphs
        paragraphs = content.split('\n\n')
    
    # Calculate optimal image placement positions
    total_paragraphs = len(paragraphs)
    images_to_place = len(images)
    
    if total_paragraphs < 3:
        # Short content - add images at the end
        enhanced_content = content
        for i, image in enumerate(images):
            figure_html = f"""

<figure class="contextual-image" style="margin: 20px 0; text-align: center;">
    <img src="{image['url']}" alt="Figure {i+1}: Illustration for {image['source_name']}" style="max-width: 100%; height: auto; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);" />
    <figcaption style="margin-top: 10px; font-style: italic; color: #666; font-size: 14px;">
        Figure {i+1}: {image['source_name'].replace('_', ' ').replace('-', ' ').title()}
    </figcaption>
</figure>
"""
            enhanced_content += figure_html
        return enhanced_content
    
    # Long content - distribute images throughout
    enhanced_paragraphs = []
    image_index = 0
    
    # Calculate spacing - place images every N paragraphs
    spacing = max(2, total_paragraphs // (images_to_place + 1))
    
    for i, paragraph in enumerate(paragraphs):
        # Add the paragraph
        if paragraph.strip():
            enhanced_paragraphs.append(paragraph + ('</p>' if not paragraph.endswith('</p>') else ''))
        
        # Insert image at calculated intervals
        if (image_index < len(images) and 
            i > 0 and 
            (i + 1) % spacing == 0 and 
            i < total_paragraphs - 1):  # Don't add at the very end
            
            image = images[image_index]
            figure_html = f"""

<figure class="contextual-image" style="margin: 20px 0; text-align: center;">
    <img src="{image['url']}" alt="Figure {image_index+1}: Illustration for {image['source_name']}" style="max-width: 100%; height: auto; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);" />
    <figcaption style="margin-top: 10px; font-style: italic; color: #666; font-size: 14px;">
        Figure {image_index+1}: {image['source_name'].replace('_', ' ').replace('-', ' ').title()}
    </figcaption>
</figure>
"""
            enhanced_paragraphs.append(figure_html)
            image_index += 1
    
    # Add any remaining images at the end
    while image_index < len(images):
        image = images[image_index]
        figure_html = f"""

<figure class="contextual-image" style="margin: 20px 0; text-align: center;">
    <img src="{image['url']}" alt="Figure {image_index+1}: Additional illustration for {image['source_name']}" style="max-width: 100%; height: auto; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);" />
    <figcaption style="margin-top: 10px; font-style: italic; color: #666; font-size: 14px;">
        Figure {image_index+1}: {image['source_name'].replace('_', ' ').replace('-', ' ').title()}
    </figcaption>
</figure>
"""
        enhanced_paragraphs.append(figure_html)
        image_index += 1
    
    enhanced_content = ''.join(enhanced_paragraphs)
    print(f"✅ Successfully injected {image_index} images into content")
    return enhanced_content

# Add endpoint to trigger image injection manually
@app.post("/api/inject-images")
async def inject_images_endpoint():
    """Manually trigger image injection into existing articles"""
    try:
        injected_count = await inject_real_images_into_articles()
        return {
            "success": True,
            "message": f"Successfully injected images into {injected_count} articles",
            "articles_updated": injected_count
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# NEW REFINED ENGINE - File upload endpoint  
@app.post("/api/content/upload")
async def upload_file(
    file: UploadFile = File(...),
    metadata: str = Form("{}")
):
    """V2 ENGINE: Upload and process files (text, audio, video, images)"""
    print(f"🚀 V2 ENGINE: Processing file upload - {file.filename} - engine=v2")
    
    # KE-PR1: Add structured logging with job_id
    start_time = time.time()
    job_id = str(uuid.uuid4())
    if logger:
        logger.info({
            "event": "content_upload_start", 
            "job_id": job_id, 
            "filename": file.filename,
            "content_type": file.content_type,
            "stage": "content_upload"
        })
    
    try:
        # Parse metadata
        file_metadata = json.loads(metadata)
        
        # Create processing job
        job = ProcessingJob(
            input_type="file",
            original_filename=file.filename,
            status="processing"
        )
        
        await db.processing_jobs.insert_one(job.dict())
        
        # OPTIMIZED: Add progress tracking for better UI feedback
        async def update_job_progress(stage: str, details: str = ""):
            """Update job progress to prevent UI timeout"""
            await db.processing_jobs.update_one(
                {"job_id": job.job_id},
                {"$set": {
                    "status": "processing", 
                    "current_stage": stage,
                    "stage_details": details,
                    "last_updated": datetime.utcnow().isoformat()
                }}
            )
            print(f"📊 PROGRESS: {stage} - {details}")
        
        await update_job_progress("initializing", "Reading file content...")
        
        # Read file content
        file_content = await file.read()
        
        # Get file extension for proper handling
        file_extension = file.filename.split('.')[-1].lower() if '.' in file.filename else ''
        
        await update_job_progress("analyzing", f"Processing {file_extension.upper()} file ({len(file_content)} bytes)")
        print(f"Processing file: {file.filename}, Extension: {file_extension}, Size: {len(file_content)} bytes")
        
        extracted_content = ""
        
        # Extract content based on file type
        await update_job_progress("extracting", f"Extracting content from {file_extension.upper()} file...")
        
        if file_extension in ['txt', 'md', 'csv']:
            try:
                extracted_content = file_content.decode('utf-8')
                print(f"✅ Extracted {len(extracted_content)} characters from text file")
            except UnicodeDecodeError:
                extracted_content = file_content.decode('latin-1', errors='ignore')
                print(f"⚠️ Used latin-1 fallback, extracted {len(extracted_content)} characters")
                
        elif file_extension == 'pdf':
            await update_job_progress("extracting", "Processing PDF with comprehensive image extraction...")
            try:
                # FIXED: Use DocumentPreprocessor for comprehensive PDF processing with image extraction
                # Create temporary file for DocumentPreprocessor processing
                temp_pdf_path = f"/app/backend/temp_uploads/temp_{file.filename}"
                os.makedirs(os.path.dirname(temp_pdf_path), exist_ok=True)
                
                with open(temp_pdf_path, 'wb') as temp_file:
                    temp_file.write(file_content)
                
                # Use DocumentPreprocessor for comprehensive PDF processing
                doc_processor = DocumentPreprocessor(session_id=job.job_id[:8])
                html_content, pdf_images = await doc_processor._convert_pdf_to_html(temp_pdf_path)
                
                # Convert HTML back to text for extracted_content
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                extracted_content = soup.get_text()
                
                await update_job_progress("extracting", f"Extracted content and {len(pdf_images)} images from PDF")
                
                # FIXED: Save PDF images to Asset Library
                if hasattr(doc_processor, 'pending_assets') and doc_processor.pending_assets:
                    try:
                        result = # KE-PR9.3: Use V2 validation repository
                from engine.stores.mongo import RepositoryFactory
                validation_repo = RepositoryFactory.get_v2_validation()
                await validation_repo.store_validation(validation_result)
                print(f"💾 V2 ENGINE: Stored file validation result for diagnostics - validation_id: {validation_result.get('validation_id')} - engine=v2")
            except Exception as validation_storage_error:
                print(f"❌ V2 ENGINE: Error storing file validation result - {validation_storage_error} - engine=v2")
            
            print(f"✅ V2 ENGINE: Step 8 complete for file upload - Validation status: {validation_status} - engine=v2")
            
            # V2 STEP 9: Cross-Article QA (dedupe, link validation, FAQ consolidation, terminology)
            print(f"🔍 V2 ENGINE: Starting Step 9 - Cross-Article QA for file upload - engine=v2")
            
            # Perform cross-article quality assurance
            qa_result = await v2_cross_article_qa_system.perform_cross_article_qa(
                generated_articles_result, run_id
            )
            
            # Check QA status and update chunks accordingly
            qa_status = qa_result.get('qa_status', 'unknown')
            issues_found = qa_result.get('summary', {}).get('issues_found', 0)
            
            if qa_status == 'error':
                print(f"❌ V2 ENGINE: Step 9 QA failed with error for file upload - run {run_id} - engine=v2")
            elif issues_found == 0:
                print(f"✅ V2 ENGINE: Step 9 QA passed for file upload - No coherence issues found - engine=v2")
                # Add QA metadata to chunks
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['qa_result'] = qa_result
                    chunk['qa_status'] = 'passed'
            else:
                print(f"⚠️ V2 ENGINE: Step 9 QA found {issues_found} issues for file upload - Articles require consolidation - engine=v2")
                # Mark chunks with QA issues
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['qa_result'] = qa_result
                    chunk['qa_status'] = 'issues_found'
                    chunk['qa_issues_count'] = issues_found
            
            # Store QA result separately for analysis
            try:
                # KE-PR9.3: Use repository pattern for content_library operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(article_record)
                    
                    return {
                        "success": True,
                        "article_id": article_record["id"],
                        "title": article_record["title"],
                        "message": "Article created successfully in Content Library"
                    }
                    
                except json.JSONDecodeError:
                    # Fallback if AI doesn't return proper JSON
                    article_record = {
                        "id": str(uuid.uuid4()),
                        "title": title,
                        "content": content,
                        "summary": f"Extracted from {source_type}",
                        "tags": [source_type],
                        "source_job_id": source_job_id,
                        "source_type": source_type,
                        "status": "draft",
                        "metadata": article_metadata,
                        "created_at": datetime.utcnow(),
                        "updated_at": datetime.utcnow()
                    }
                    
                    # KE-PR9.3: Use repository pattern for content_library operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(article_record)
        
        return {
            "success": True,
            "article_id": article_record["id"],
            "title": article_record["title"],
            "message": "Article created successfully in Content Library"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Get Content Library articles
@app.get("/api/content-library")
async def get_content_library_articles():
    """Get all Content Library articles"""
    try:
        articles = []
        async for article in db.content_library.find().sort("created_at", -1):
            # Convert ObjectId to string for JSON serialization
            clean_article = objectid_to_str(article)
            
            articles.append({
                "id": clean_article["id"],
                "title": clean_article["title"],
                "content": clean_article.get("content", ""),  # Added content field!
                "summary": clean_article.get("summary", ""),
                "tags": clean_article.get("tags", []),
                "status": clean_article.get("status", "draft"),
                "source_type": clean_article.get("source_type", ""),
                "takeaways": clean_article.get("takeaways", []),  # Added takeaways too
                "metadata": clean_article.get("metadata", {}),    # Added metadata - now ObjectId safe
                "processing_metadata": clean_article.get("processing_metadata", {}),  # Added processing metadata
                "engine": clean_article.get("engine"),  # Added engine field for V2 articles
                # V2 Publishing comprehensive structure fields
                "html": clean_article.get("html"),
                "markdown": clean_article.get("markdown"),
                "toc": clean_article.get("toc"),
                "faq": clean_article.get("faq"),
                "related_links": clean_article.get("related_links"),
                "provenance_map": clean_article.get("provenance_map"),
                "metrics": clean_article.get("metrics"),
                "media_references": clean_article.get("media_references"),
                "created_at": clean_article.get("created_at"),
                "updated_at": clean_article.get("updated_at")
            })
        
        return {
            "articles": articles,
            "total": len(articles)
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# URL processing endpoint  
@app.post("/api/content/process-url")
async def process_url_content(
    url: str = Form(...),
    metadata: str = Form("{}")
):
    """V2 ENGINE: Process URL content by scraping and generating articles"""
    print(f"🚀 V2 ENGINE: Processing URL content - {url} - engine=v2")
    try:
        # Parse metadata
        url_metadata = json.loads(metadata)
        
        # Create processing job
        job = ProcessingJob(
            input_type="url",
            url=url,
            status="processing"
        )
        
        await db.processing_jobs.insert_one(job.dict())
        
        print(f"🌐 Processing URL: {url}")
        
        # Scrape website content
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        
        # Parse HTML content
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Extract title
        page_title = soup.find('title')
        title = page_title.get_text().strip() if page_title else url
        
        # Extract meta description
        meta_desc = soup.find('meta', attrs={'name': 'description'})
        description = meta_desc.get('content', '').strip() if meta_desc else ""
        
        # Extract main content
        main_content = ""
        
        # Try to find main content areas
        content_selectors = [
            'main', 'article', '.content', '#content', 
            '.post', '.entry', '.article-body', '.main-content'
        ]
        
        content_found = False
        for selector in content_selectors:
            content_area = soup.select_one(selector)
            if content_area:
                main_content = content_area.get_text(separator='\n', strip=True)
                content_found = True
                break
        
        # If no main content area found, extract from body
        if not content_found:
            body = soup.find('body')
            if body:
                # Remove navigation, footer, sidebar elements
                for element in body.find_all(['nav', 'footer', 'aside', '.sidebar', '.navigation']):
                    element.decompose()
                main_content = body.get_text(separator='\n', strip=True)
        
        # Clean up content
        lines = main_content.split('\n')
        cleaned_lines = [line.strip() for line in lines if line.strip() and len(line.strip()) > 10]
        extracted_content = '\n'.join(cleaned_lines)
        
        # Create enriched content
        enriched_content = f"""Website: {title}
URL: {url}

{f"Description: {description}" if description else ""}

=== Main Content ===

{extracted_content}

---
Source Information:
- Original URL: {url}
- Page Title: {title}
- Scraped Date: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}
- Content Length: {len(extracted_content)} characters
- Source: Knowledge Engine URL Processing"""
        
        print(f"✅ Extracted {len(extracted_content)} characters from URL")
        
        # V2 ENGINE: Use direct URL extraction with V2 system
        try:
            print(f"🌐 V2 ENGINE: Starting URL processing with V2 extractor - engine=v2")
            
            # V2 STEP 2: Use V2 Content Extractor for URL processing
            v2_extractor = V2ContentExtractor()
            
            # Extract content using V2 URL extractor
            normalized_doc = await v2_extractor.extract_url_content(url, response.text)
            
            print(f"📋 V2 ENGINE: Extracted {len(normalized_doc.blocks)} blocks, {len(normalized_doc.media)} media from URL - engine=v2")
            
            # Store normalized document in database
            await store_normalized_document(normalized_doc)
            
            # V2 STEP 4: Perform multi-dimensional analysis
            run_id = f"run_{int(datetime.utcnow().timestamp())}_{uuid.uuid4().hex[:8]}"
            analysis_result = await get_v2_analyzer().run(normalized_doc, run_id)
            
            # Extract analysis for use in processing
            analysis = analysis_result.get('analysis', {}) if analysis_result else {}
            audience = analysis.get('audience', 'end_user')
            granularity = analysis.get('granularity', 'shallow')
            
            print(f"🎯 V2 ENGINE: Analysis complete - {analysis.get('content_type', 'unknown')} content for {audience} audience with {granularity} granularity - engine=v2")
            
            # V2 STEP 5: Create global outline with 100% block assignment
            global_outline = await v2_global_planner.create_global_outline(normalized_doc, analysis, run_id)
            
            # Extract outline for use in article generation
            outline = global_outline.get('outline', {}) if global_outline else {}
            article_outlines = outline.get('articles', [])
            discarded_blocks = outline.get('discarded_blocks', [])
            
            print(f"📋 V2 ENGINE: Global outline created - {len(article_outlines)} articles planned, {len(discarded_blocks)} blocks discarded - engine=v2")
            
            # V2 STEP 6: Create detailed per-article outlines
            per_article_outlines_result = await v2_per_article_outline_planner.create_per_article_outlines(
                normalized_doc, outline, analysis, run_id
            )
            
            # Extract per-article outlines
            per_article_outlines = per_article_outlines_result.get('per_article_outlines', []) if per_article_outlines_result else []
            
            print(f"✅ V2 ENGINE: Step 6 complete - Per-article outlines created for URL processing - engine=v2")
            
            # V2 STEP 6.5: Section-Grounded Prewrite Pass (facts → claims)
            print(f"🔄 V2 ENGINE: Starting Step 6.5 - Section-grounded prewrite pass for URL processing - engine=v2")
            
            prewrite_result = await v2_prewrite_system.execute_prewrite_pass(
                enriched_content, 'url', article_outlines, 
                per_article_outlines_result.get('per_article_outlines', {}), 
                analysis, run_id
            )
            
            prewrite_status = prewrite_result.get('prewrite_status', 'unknown')
            successful_prewrites = prewrite_result.get('successful_prewrites', 0)
            
            if prewrite_status == 'success':
                print(f"✅ V2 ENGINE: Step 6.5 prewrite successful for URL - {successful_prewrites} articles prepared - engine=v2")
            elif prewrite_status == 'partial':
                print(f"⚠️ V2 ENGINE: Step 6.5 prewrite partial for URL - {successful_prewrites} articles prepared, some failed - engine=v2")
            else:
                print(f"❌ V2 ENGINE: Step 6.5 prewrite failed for URL - {prewrite_status} - engine=v2")
            
            # Store prewrite result for diagnostics
            try:
                await db.v2_prewrite_results.insert_one(prewrite_result)
                print(f"💾 V2 ENGINE: Stored URL prewrite result for diagnostics - prewrite_id: {prewrite_result.get('prewrite_id')} - engine=v2")
            except Exception as prewrite_storage_error:
                print(f"❌ V2 ENGINE: Error storing URL prewrite result - {prewrite_storage_error} - engine=v2")
            
            print(f"✅ V2 ENGINE: Step 6.5 complete - Section-grounded prewrite pass complete for URL - engine=v2")
            
            # V2 STEP 7: Generate final articles with strict format and audience-aware styling
            print(f"🎯 V2 ENGINE: Starting Step 7 - Final article generation with V2ArticleGenerator - engine=v2")
            
            # Use V2ArticleGenerator for final article generation
            generated_articles_result = await v2_article_generator.generate_final_articles(
                normalized_doc, 
                per_article_outlines, 
                analysis, 
                run_id
            )
            
            # Convert V2ArticleGenerator output to expected format
            chunks = []
            if generated_articles_result and 'generated_articles' in generated_articles_result:
                for generated_article in generated_articles_result['generated_articles']:
                    article_data = generated_article.get('article_data', {})
                    if article_data:
                        # Extract title from HTML content
                        article_title = v2_article_generator._extract_title_from_html(article_data.get('html', ''), generated_article.get('article_id', 'Generated Article'))
                        
                        # Create article in expected format for content library storage
                        chunk = {
                            "id": str(uuid.uuid4()),
                            "title": article_title,
                            "content": article_data.get('html', ''),
                            "summary": article_data.get('summary', ''),
                            "status": "draft",
                            "created_at": datetime.utcnow().isoformat(),
                            "updated_at": datetime.utcnow().isoformat(),
                            "source_content": f"V2 Engine processed content from URL: {url}",
                            "source_type": "v2_generated",
                            "markdown": article_data.get('markdown', ''),
                            "takeaways": [],
                            "metadata": {
                                "engine": "v2",
                                "processing_version": "2.0",
                                "normalized_doc_id": normalized_doc.doc_id,
                                "run_id": run_id,
                                "analysis": analysis,
                                "audience": audience,
                                "granularity": granularity,
                                "article_id": generated_article.get('article_id', 'unknown'),
                                "validation_metadata": article_data.get('validation_metadata', {}),
                                "generated_by": "v2_article_generator",
                                "extraction_method": "v2_url_extraction",
                                "source_url": url
                            }
                        }
                        chunks.append(chunk)
            
            print(f"✅ V2 ENGINE: Step 7 complete - Generated {len(chunks)} final articles from URL {url} - engine=v2")
            
            # V2 STEP 7.6: Evidence Tagging for URL processing
            print(f"🏷️ V2 ENGINE: Starting Step 7.6 - Evidence tagging for URL processing - engine=v2")
            
            # Tag paragraphs with evidence block IDs using prewrite data
            evidence_tagging_result = await v2_evidence_tagging_system.tag_content_with_evidence(
                chunks, normalized_doc.blocks, prewrite_result, run_id
            )
            
            evidence_tagging_status = evidence_tagging_result.get('evidence_tagging_status', 'unknown')
            if evidence_tagging_status == 'success':
                total_paragraphs = evidence_tagging_result.get('total_paragraphs', 0)
                tagged_paragraphs = evidence_tagging_result.get('tagged_paragraphs', 0)
                tagging_rate = evidence_tagging_result.get('overall_tagging_rate', 0)
                target_achieved = evidence_tagging_result.get('target_achieved', False)
                
                print(f"✅ V2 ENGINE: Step 7.6 evidence tagging successful for URL processing - {tagged_paragraphs}/{total_paragraphs} paragraphs tagged ({tagging_rate:.1f}%) - Target ≥95%: {'✅' if target_achieved else '⚠️'} - engine=v2")
            else:
                print(f"⚠️ V2 ENGINE: Step 7.6 evidence tagging failed for URL processing - Status: {evidence_tagging_status} - engine=v2")
            
            # Store evidence tagging result for diagnostics
            try:
                await db.v2_evidence_tagging_results.insert_one(evidence_tagging_result)
                print(f"💾 V2 ENGINE: Stored evidence tagging result for URL processing diagnostics - evidence_tagging_id: {evidence_tagging_result.get('evidence_tagging_id')} - engine=v2")
            except Exception as evidence_storage_error:
                print(f"❌ V2 ENGINE: Error storing evidence tagging result for URL processing - {evidence_storage_error} - engine=v2")
            
            # Add evidence tagging metadata to chunks
            for i, chunk in enumerate(chunks):
                evidence_result = next((r for r in evidence_tagging_result.get('evidence_tagging_results', []) if r.get('article_index') == i), None)
                if evidence_result:
                    chunk.setdefault('metadata', {})['evidence_tagging_result'] = evidence_result
                    chunk['tagged_paragraphs'] = evidence_result.get('tagged_paragraphs', 0)
                    chunk['tagging_rate'] = evidence_result.get('tagging_rate', 0)
            
            print(f"✅ V2 ENGINE: Step 7.6 complete for URL processing - Evidence tagging complete - engine=v2")
            
            # V2 STEP 7.5: Woolf-aligned Technical Writing Style + Structural Lint
            print(f"🔄 V2 ENGINE: Starting Step 7.5 - Woolf-aligned style formatting - engine=v2")
            
            style_result = await v2_style_processor.apply_style_formatting(
                enriched_content, 'url', chunks, 
                generated_articles_result, analysis, run_id
            )
            
            style_status = style_result.get('style_status', 'unknown')
            successful_formatting = style_result.get('successful_formatting', 0)
            style_compliance = style_result.get('style_compliance', {})
            
            if style_status == 'success':
                print(f"✅ V2 ENGINE: Step 7.5 style formatting successful - {successful_formatting} articles formatted, {style_compliance.get('overall_compliance', 0):.1f}% compliance - engine=v2")
            elif style_status == 'partial':
                print(f"⚠️ V2 ENGINE: Step 7.5 style formatting partial - {successful_formatting} articles formatted, some failed - engine=v2")
            else:
                print(f"❌ V2 ENGINE: Step 7.5 style formatting failed - {style_status} - engine=v2")
            
            # Store style result for diagnostics
            try:
                await db.v2_style_results.insert_one(style_result)
                print(f"💾 V2 ENGINE: Stored style result for diagnostics - style_id: {style_result.get('style_id')} - engine=v2")
            except Exception as style_storage_error:
                print(f"❌ V2 ENGINE: Error storing style result - {style_storage_error} - engine=v2")
            
            print(f"✅ V2 ENGINE: Step 7.5 complete - Woolf-aligned style formatting complete - engine=v2")
            
            # V2 STEP 7.7: Related Links Generation (internal + external from source) for URL processing
            print(f"🔗 V2 ENGINE: Starting Step 7.7 - Related links generation for URL processing - engine=v2")
            
            # Generate related links for each article using content library and source links
            related_links_results = []
            for i, chunk in enumerate(chunks):
                try:
                    related_links_result = await v2_related_links_system.generate_related_links(
                        chunk, enriched_content, normalized_doc.blocks, run_id
                    )
                    
                    # Add related links to chunk
                    if related_links_result.get('related_links_status') == 'success':
                        related_links = related_links_result.get('related_links', [])
                        
                        # Update chunk with related links
                        chunk['related_links'] = related_links
                        chunk.setdefault('metadata', {})['related_links_result'] = related_links_result
                        chunk['related_links_count'] = len(related_links)
                        
                        internal_count = related_links_result.get('internal_links_count', 0)
                        external_count = related_links_result.get('external_links_count', 0)
                        print(f"✅ V2 ENGINE: Added {len(related_links)} related links to '{chunk['title'][:50]}...' ({internal_count} internal, {external_count} external) - engine=v2")
                    else:
                        print(f"⚠️ V2 ENGINE: Failed to generate related links for '{chunk['title'][:50]}...' - engine=v2")
                    
                    related_links_results.append(related_links_result)
                    
                except Exception as chunk_links_error:
                    print(f"❌ V2 ENGINE: Error generating related links for chunk {i+1} - {chunk_links_error} - engine=v2")
                    related_links_results.append({
                        "related_links_status": "error",
                        "error": str(chunk_links_error),
                        "run_id": run_id,
                        "engine": "v2"
                    })
            
            # Store related links results for diagnostics
            try:
                for result in related_links_results:
                    await db.v2_related_links_results.insert_one(result)
                print(f"💾 V2 ENGINE: Stored {len(related_links_results)} related links results for URL processing diagnostics - engine=v2")
            except Exception as related_links_storage_error:
                print(f"❌ V2 ENGINE: Error storing related links results for URL processing - {related_links_storage_error} - engine=v2")
            
            successful_related_links = len([r for r in related_links_results if r.get('related_links_status') == 'success'])
            total_related_links = sum([r.get('total_links_count', 0) for r in related_links_results])
            
            print(f"✅ V2 ENGINE: Step 7.7 complete for URL processing - Related links generation complete - {successful_related_links}/{len(chunks)} articles, {total_related_links} total links - engine=v2")
            
            # V2 STEP 7.8: Intelligent Gap Filling for URL processing
            print(f"🔍 V2 ENGINE: Starting Step 7.8 - Intelligent gap filling for URL processing - engine=v2")
            
            # Fill gaps in chunks using in-corpus retrieval and pattern synthesis
            gap_filling_result = await v2_gap_filling_system.fill_content_gaps(
                chunks, enriched_content, normalized_doc.blocks, run_id, enrich_mode="internal"
            )
            
            gap_filling_status = gap_filling_result.get('gap_filling_status', 'unknown')
            if gap_filling_status == 'success':
                total_gaps_found = gap_filling_result.get('total_gaps_found', 0)
                total_gaps_filled = gap_filling_result.get('total_gaps_filled', 0)
                gap_fill_rate = gap_filling_result.get('gap_fill_rate', 0)
                
                print(f"✅ V2 ENGINE: Step 7.8 gap filling successful for URL processing - {total_gaps_filled}/{total_gaps_found} gaps filled ({gap_fill_rate:.1f}% success rate) - engine=v2")
            else:
                print(f"⚠️ V2 ENGINE: Step 7.8 gap filling failed for URL processing - Status: {gap_filling_status} - engine=v2")
            
            # Store gap filling result for diagnostics
            try:
                await db.v2_gap_filling_results.insert_one(gap_filling_result)
                print(f"💾 V2 ENGINE: Stored gap filling result for URL processing diagnostics - gap_filling_id: {gap_filling_result.get('gap_filling_id')} - engine=v2")
            except Exception as gap_storage_error:
                print(f"❌ V2 ENGINE: Error storing gap filling result for URL processing - {gap_storage_error} - engine=v2")
            
            # Add gap filling metadata to chunks
            for i, chunk in enumerate(chunks):
                gap_result = next((r for r in gap_filling_result.get('gap_filling_results', []) if r.get('article_index') == i), None)
                if gap_result:
                    chunk.setdefault('metadata', {})['gap_filling_result'] = gap_result
                    chunk['gaps_filled'] = gap_result.get('gaps_filled', 0)
            
            print(f"✅ V2 ENGINE: Step 7.8 complete for URL processing - Intelligent gap filling complete - engine=v2")
            
            # V2 STEP 7.9: Code Block Normalization & Beautification for URL processing
            print(f"🎨 V2 ENGINE: Starting Step 7.9 - Code block normalization for URL processing - engine=v2")
            
            # Normalize and beautify code blocks for Prism rendering
            code_normalization_result = await v2_code_normalization_system.normalize_code_blocks(
                chunks, normalized_doc.blocks, prewrite_result, run_id
            )
            
            code_normalization_status = code_normalization_result.get('code_normalization_status', 'unknown')
            if code_normalization_status == 'success':
                total_code_blocks = code_normalization_result.get('total_code_blocks', 0)
                normalized_blocks = code_normalization_result.get('normalized_blocks', 0)
                normalization_rate = code_normalization_result.get('overall_normalization_rate', 0)
                
                print(f"✅ V2 ENGINE: Step 7.9 code normalization successful for URL processing - {normalized_blocks}/{total_code_blocks} blocks normalized ({normalization_rate:.1f}%) - engine=v2")
            else:
                print(f"⚠️ V2 ENGINE: Step 7.9 code normalization failed for URL processing - Status: {code_normalization_status} - engine=v2")
            
            # Store code normalization result for diagnostics
            try:
                await db.v2_code_normalization_results.insert_one(code_normalization_result)
                print(f"💾 V2 ENGINE: Stored code normalization result for URL processing diagnostics - code_normalization_id: {code_normalization_result.get('code_normalization_id')} - engine=v2")
            except Exception as code_storage_error:
                print(f"❌ V2 ENGINE: Error storing code normalization result for URL processing - {code_storage_error} - engine=v2")
            
            # Add code normalization metadata to chunks
            for i, chunk in enumerate(chunks):
                code_result = next((r for r in code_normalization_result.get('code_normalization_results', []) if r.get('article_index') == i), None)
                if code_result:
                    chunk.setdefault('metadata', {})['code_normalization_result'] = code_result
                    chunk['normalized_code_blocks'] = code_result.get('normalized_blocks', 0)
                    chunk['code_normalization_rate'] = code_result.get('normalization_rate', 0)
            
            print(f"✅ V2 ENGINE: Step 7.9 complete for URL processing - Code block normalization complete - engine=v2")
            
            # V2 STEP 8: Implement Validators (fidelity, 100% coverage, placeholders, style)
            print(f"🔍 V2 ENGINE: Starting Step 8 - Comprehensive validation for URL processing - engine=v2")
            
            # Perform comprehensive validation of generated articles
            validation_result = await v2_validation_system.validate_generated_articles(
                normalized_doc, generated_articles_result, analysis, run_id
            )
            
            # Check validation status and update chunks accordingly
            validation_status = validation_result.get('validation_status', 'unknown')
            if validation_status == 'passed':
                print(f"✅ V2 ENGINE: Step 8 validation passed for URL processing - Articles meet all quality thresholds - engine=v2")
                # Add validation metadata to chunks
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['validation_result'] = validation_result
                    chunk['validation_status'] = 'passed'
            else:
                print(f"⚠️ V2 ENGINE: Step 8 validation failed for URL processing - Status: {validation_status} - engine=v2")
                # Mark chunks as partial with diagnostics
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['validation_result'] = validation_result
                    chunk['validation_status'] = 'partial'
                    chunk['validation_diagnostics'] = validation_result.get('diagnostics', [])
            
            # Store validation result separately for diagnostics endpoint
            try:
                # KE-PR9.3: Use QA results repository
                from engine.stores.mongo import RepositoryFactory
                qa_repo = RepositoryFactory.get_qa_results()
                await qa_repo.insert_qa_report(qa_result)
                print(f"💾 V2 ENGINE: Stored URL QA result for analysis - qa_id: {qa_result.get('qa_id')} - engine=v2")
            except Exception as qa_storage_error:
                print(f"❌ V2 ENGINE: Error storing URL QA result - {qa_storage_error} - engine=v2")
            
            print(f"✅ V2 ENGINE: Step 9 complete for URL processing - QA status: {qa_status}, Issues found: {issues_found} - engine=v2")
            
            # V2 STEP 10: Adaptive Adjustment (balance splits/length)
            print(f"⚖️ V2 ENGINE: Starting Step 10 - Adaptive adjustment for URL processing - engine=v2")
            
            # Perform adaptive adjustment for article length and split optimization
            adjustment_result = await v2_adaptive_adjustment_system.perform_adaptive_adjustment(
                generated_articles_result, analysis, run_id
            )
            
            # Check adjustment status and update chunks accordingly
            adjustment_status = adjustment_result.get('adjustment_status', 'unknown')
            total_adjustments = adjustment_result.get('adjustment_summary', {}).get('total_adjustments', 0)
            readability_score = adjustment_result.get('readability_score', 0.5)
            
            if adjustment_status == 'error':
                print(f"❌ V2 ENGINE: Step 10 adjustment failed with error for URL processing - run {run_id} - engine=v2")
            elif total_adjustments == 0:
                print(f"✅ V2 ENGINE: Step 10 adjustment complete for URL processing - No adjustments needed, optimal balance achieved - engine=v2")
                # Add adjustment metadata to chunks
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['adjustment_result'] = adjustment_result
                    chunk['adjustment_status'] = 'optimal'
                    chunk['readability_score'] = readability_score
            else:
                print(f"⚖️ V2 ENGINE: Step 10 found {total_adjustments} adjustments for URL processing - Articles balanced for optimal readability - engine=v2")
                # Mark chunks with adjustment recommendations
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['adjustment_result'] = adjustment_result
                    chunk['adjustment_status'] = 'adjusted'
                    chunk['adjustments_applied'] = total_adjustments
                    chunk['readability_score'] = readability_score
            
            # Store adjustment result separately for analysis
            try:
                await db.v2_adjustment_results.insert_one(adjustment_result)
                print(f"💾 V2 ENGINE: Stored URL adjustment result for analysis - adjustment_id: {adjustment_result.get('adjustment_id')} - engine=v2")
            except Exception as adjustment_storage_error:
                print(f"❌ V2 ENGINE: Error storing URL adjustment result - {adjustment_storage_error} - engine=v2")
            
            print(f"✅ V2 ENGINE: Step 10 complete for URL processing - Adjustment status: {adjustment_status}, Readability score: {readability_score:.2f} - engine=v2")
            
            # V2 STEP 11: Publishing Flow (V2 only)
            print(f"📚 V2 ENGINE: Starting Step 11 - V2-only publishing flow for URL processing - engine=v2")
            
            # Publish finalized V2 content as single source of truth
            publishing_result = await v2_publishing_system.publish_v2_content(
                chunks, generated_articles_result, validation_result, qa_result, adjustment_result, run_id
            )
            
            # Check publishing status and update chunks accordingly
            publishing_status = publishing_result.get('publishing_status', 'unknown')
            published_count = publishing_result.get('published_articles', 0)
            coverage_achieved = publishing_result.get('coverage_achieved', 0)
            
            if publishing_status == 'success':
                print(f"✅ V2 ENGINE: Step 11 publishing successful for URL processing - {published_count} articles published with {coverage_achieved}% coverage - engine=v2")
                # Add publishing metadata to chunks
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['publishing_result'] = publishing_result
                    chunk['publishing_status'] = 'published'
                    chunk['published_at'] = datetime.utcnow().isoformat()
            else:
                print(f"⚠️ V2 ENGINE: Step 11 publishing issue for URL processing - Status: {publishing_status} - Coverage: {coverage_achieved}% - engine=v2")
                # Mark chunks with publishing status
                for chunk in chunks:
                    chunk.setdefault('metadata', {})['publishing_result'] = publishing_result
                    chunk['publishing_status'] = publishing_status
            
            # Store publishing result separately for analysis
            try:
                await db.v2_publishing_results.insert_one(publishing_result)
                print(f"💾 V2 ENGINE: Stored URL publishing result for analysis - publishing_id: {publishing_result.get('publishing_id')} - engine=v2")
            except Exception as publishing_storage_error:
                print(f"❌ V2 ENGINE: Error storing URL publishing result - {publishing_storage_error} - engine=v2")
            
            print(f"✅ V2 ENGINE: Step 11 complete for URL processing - Publishing status: {publishing_status}, Articles published: {published_count} - engine=v2")
            
            # V2 STEP 12: Versioning & Diff (reprocessing support) for URL processing
            print(f"🔄 V2 ENGINE: Starting Step 12 - Versioning and diff management for URL processing - engine=v2")
            
            # Determine content type and get original content for hashing
            content_for_versioning = enriched_content[:5000] if enriched_content else ""  # Use first 5000 chars for hash  
            content_type = f"url:{url}"
            
            # Perform versioning management
            versioning_result = await v2_versioning_system.manage_versioning(
                content_for_versioning, content_type, chunks, generated_articles_result, publishing_result, run_id
            )
            
            # Check versioning status and update articles accordingly
            versioning_status = versioning_result.get('versioning_status', 'unknown')
            version_number = versioning_result.get('version_metadata', {}).get('version', 1)
            is_update = versioning_result.get('version_metadata', {}).get('supersedes') is not None
            
            if versioning_status == 'success':
                update_text = "(update)" if is_update else "(new)"
                print(f"✅ V2 ENGINE: Step 12 versioning successful for URL processing - Version {version_number} {update_text} - engine=v2")
                # Add versioning metadata to articles
                for article in chunks:
                    if isinstance(article, dict):
                        article.setdefault('metadata', {})['versioning_result'] = versioning_result
                        article['version_metadata'] = versioning_result.get('version_metadata', {})
                        article['version_number'] = version_number
                        article['is_version_update'] = is_update
                        if versioning_result.get('diff_result'):
                            article['version_diff'] = versioning_result['diff_result']
            else:
                print(f"❌ V2 ENGINE: Step 12 versioning failed for URL processing - Status: {versioning_status} - engine=v2")
                # Mark articles with versioning error
                for article in chunks:
                    if isinstance(article, dict):
                        article.setdefault('metadata', {})['versioning_result'] = versioning_result
                        article['versioning_status'] = 'failed'
            
            # Store versioning result separately for diagnostics
            try:
                await db.v2_versioning_results.insert_one(versioning_result)
                print(f"💾 V2 ENGINE: Stored URL versioning result for diagnostics - versioning_id: {versioning_result.get('versioning_id')} - engine=v2")
            except Exception as versioning_storage_error:
                print(f"❌ V2 ENGINE: Error storing URL versioning result - {versioning_storage_error} - engine=v2")
            
            print(f"✅ V2 ENGINE: Step 12 complete for URL processing - Versioning status: {versioning_status}, Version: {version_number} - engine=v2")
            
        except Exception as v2_error:
            print(f"⚠️ V2 ENGINE: URL extraction failed, falling back to legacy processing - {v2_error} - engine=v2")
            
            # Fallback to legacy processing
            enhanced_metadata = {
                **url_metadata,
                "url": url,
                "page_title": title,
                "page_description": description,
                "content_length": len(extracted_content),
                "type": "url_processing"
            }
            
            chunks = await process_text_content_v2_pipeline(enriched_content, enhanced_metadata)
        
        # Update job
        job.chunks = chunks
        job.status = "completed"
        job.completed_at = datetime.utcnow()
        
        await db.processing_jobs.update_one(
            {"job_id": job.job_id},
            {"$set": job.dict()}
        )
        
        print(f"✅ V2 ENGINE: URL processing complete - {len(chunks)} chunks created - engine=v2")
        return {
            "job_id": job.job_id,
            "status": job.status,
            "url": url,
            "page_title": title,
            "extracted_content_length": len(extracted_content),
            "chunks_created": len(chunks),
            "message": "V2 Engine: URL processed successfully",
            "engine": "v2"
        }
        
    except requests.RequestException as e:
        # Update job with error
        if 'job' in locals():
            await db.processing_jobs.update_one(
                {"job_id": job.job_id},
                {"$set": {"status": "failed", "error_message": f"Failed to fetch URL: {str(e)}"}}
            )
        raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {str(e)}")
    except Exception as e:
        # Update job with error
        if 'job' in locals():
            await db.processing_jobs.update_one(
                {"job_id": job.job_id},
                {"$set": {"status": "failed", "error_message": str(e)}}
            )
        raise HTTPException(status_code=500, detail=str(e))

# Recording processing endpoint
@app.post("/api/content/process-recording")
async def process_recording(
    recording_type: str = Form(...),  # 'screen', 'audio', 'video', 'screenshot'
    duration: int = Form(0),
    title: str = Form(""),
    metadata: str = Form("{}")
):
    """V2 ENGINE: Process recorded content (screen, audio, video, screenshots)"""
    print(f"🚀 V2 ENGINE: Processing {recording_type} recording - {title} - engine=v2")
    try:
        recording_metadata = json.loads(metadata)
        
        job = ProcessingJob(
            input_type="recording",
            original_filename=f"{recording_type}_recording_{title}",
            status="processing"
        )
        
        await db.processing_jobs.insert_one(job.dict())
        
        # Simulate recording processing
        if recording_type == "screenshot":
            content = f"Screenshot captured: {title}\n\nThis screenshot shows important visual information that has been processed and can be referenced in conversations."
        elif recording_type == "screen":
            content = f"Screen recording captured: {title} (Duration: {duration}s)\n\nThis screen recording demonstrates key processes and workflows that have been analyzed and processed."
        elif recording_type == "audio":
            content = f"Audio recording processed: {title} (Duration: {duration}s)\n\nThis audio content has been transcribed and processed for searchability."
        else:
            content = f"Video recording processed: {title} (Duration: {duration}s)\n\nThis video content has been analyzed and key information extracted."
        
        # V2 PROCESSING: Process the content with V2 engine
        chunks = await process_text_content_v2_pipeline(content, {**recording_metadata, "recording_type": recording_type, "duration": duration})
        
        # Update job
        job.chunks = chunks
        job.status = "completed" 
        job.completed_at = datetime.utcnow()
        
        await db.processing_jobs.update_one(
            {"job_id": job.job_id},
            {"$set": job.dict()}
        )
        
        print(f"✅ V2 ENGINE: Recording processing complete - {len(chunks)} chunks created - engine=v2")
        return {
            "job_id": job.job_id,
            "status": job.status,
            "recording_type": recording_type,
            "duration": duration,
            "chunks_created": len(chunks),
            "message": "V2 Engine: Recording processed successfully",
            "engine": "v2"
        }
        
    except Exception as e:
        if 'job' in locals():
            await db.processing_jobs.update_one(
                {"job_id": job.job_id},
                {"$set": {"status": "failed", "error_message": str(e)}}
            )
        raise HTTPException(status_code=500, detail=str(e))

# Enhanced document listing with article links
@app.get("/api/documents")
async def list_documents():
    """List all processed documents with article information"""
    try:
        documents = []
        
        async for chunk in db.document_chunks.find().sort("created_at", -1):
            # Find related articles for this document
            related_articles = []
            if 'original_filename' in chunk.get('metadata', {}):
                filename = chunk['metadata']['original_filename']
                async for article in db.content_library.find({
                    "metadata.original_filename": filename
                }):
                    related_articles.append({
                        "id": article["id"],
                        "title": article["title"],
                        "summary": article["summary"]
                    })
            
            documents.append({
                "id": chunk["id"],
                "content": chunk["content"][:200] + "..." if len(chunk["content"]) > 200 else chunk["content"],
                "metadata": chunk.get("metadata", {}),
                "created_at": chunk.get("created_at"),
                "related_articles": related_articles,
                "articles_count": len(related_articles)
            })
        
        return {
            "documents": documents,
            "total": len(documents)
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Update Content Library article
@app.put("/api/content-library/{article_id}")
async def update_content_library_article(
    article_id: str,
    title: str = Form(...),
    content: str = Form(...),
    status: str = Form("draft"),
    tags: str = Form("[]"),
    metadata: str = Form("{}")
):
    """Update an existing Content Library article with version history"""
    try:
        # Parse JSON fields
        tags_list = json.loads(tags) if tags else []
        metadata_dict = json.loads(metadata) if metadata else {}
        
        # Get existing article for version history
        existing_article = # KE-PR9.3: Use repository pattern for content_library operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(new_article)
        
        return {
            "success": True,
            "article_id": article_id,
            "message": "Article created successfully"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Get article version history
@app.get("/api/content-library/{article_id}/versions")
async def get_article_version_history(article_id: str):
    """Get version history for an article"""
    try:
        article = # KE-PR9.3: Use repository pattern for find operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                result = await content_repo.collection.find_one({"id": article_id})
        if not article:
            raise HTTPException(status_code=404, detail="Article not found")
        
        version_history = article.get("version_history", [])
        target_version = None
        
        for v in version_history:
            if v.get("version") == version:
                target_version = v
                break
        
        if not target_version:
            raise HTTPException(status_code=404, detail="Version not found")
        
        # Save current version to history before restoring
        current_version_entry = {
            "version": article.get("version", 1),
            "title": article.get("title", ""),
            "content": article.get("content", ""),
            "status": article.get("status", "draft"),
            "tags": article.get("tags", []),
            "updated_at": article.get("updated_at"),
            "updated_by": article.get("updated_by", "system")
        }
        version_history.append(current_version_entry)
        
        # Restore to target version
        new_version = article.get("version", 1) + 1
        restored_article = {
            "title": target_version.get("title", ""),
            "content": target_version.get("content", ""),
            "status": target_version.get("status", "draft"),
            "tags": target_version.get("tags", []),
            "version": new_version,
            "version_history": version_history,
            "updated_at": datetime.utcnow().isoformat(),
            "updated_by": "user",
            "restored_from_version": version
        }
        
        await db.content_library.update_one(
            {"id": article_id},
            {"$set": restored_article}
        )
        
        return {
            "success": True,
            "article_id": article_id,
            "restored_from_version": version,
            "new_version": new_version,
            "message": f"Article restored to version {version}"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/content-library")
async def create_content_library_article(request: Request):
    """Create a new article in the Content Library"""
    try:
        data = await request.json()
        
        # Generate ID if not provided
        if 'id' not in data:
            data['id'] = f"article_{int(datetime.now().timestamp())}_{uuid.uuid4().hex[:8]}"
        
        # Set timestamps
        data['created_at'] = datetime.now().isoformat()
        data['updated_at'] = datetime.now().isoformat()
        
        # Set defaults
        data.setdefault('version', 1)
        data.setdefault('views', 0)
        data.setdefault('status', 'draft')
        data.setdefault('source', 'manual')
        data.setdefault('tags', [])
        data.setdefault('metadata', {})
        
        # Calculate word count
        content = data.get('content', '')
        if content:
            # Remove HTML tags and count words
            import re
            text_content = re.sub(r'<[^>]+>', '', content)
            data['wordCount'] = len(text_content.split())
        else:
            data['wordCount'] = 0
        
        # Insert into database
        await content_library_collection.insert_one(data)
        
        return {
            "success": True,
            "message": "Article created successfully",
            "article": data
        }
        
    except Exception as e:
        print(f"❌ Error creating article: {str(e)}")
        return {"success": False, "error": str(e)}


@app.delete("/api/content-library/{article_id}")
async def delete_content_library_article(article_id: str):
    """Delete an article from the Content Library"""
    try:
        result = await content_library_collection.delete_one({"id": article_id})
        
        if result.deleted_count == 0:
            raise HTTPException(status_code=404, detail="Article not found")
        
        return {
            "success": True,
            "message": "Article deleted successfully"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error deleting article: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

# V2 ENGINE HEALTH CHECK
@app.get("/api/engine")
async def get_engine_status():
    """V2 ENGINE: Get current engine status and configuration with QA summaries"""
    try:
        print(f"📊 V2 ENGINE: Health check requested - engine=v2")
        
        # Get recent QA summaries for KE-PR7
        qa_summaries = await get_recent_qa_summaries(limit=5)
        
        return {
            "engine": "v2",
            "legacy": "disabled",
            "status": "active",
            "version": "2.0",
            "endpoints": {
                "text_processing": "/api/content/process",
                "file_upload": "/api/content/upload", 
                "url_processing": "/api/content/process-url",
                "recording_processing": "/api/content/process-recording",
                "validation_diagnostics": "/api/validation/diagnostics",
                "qa_diagnostics": "/api/qa/diagnostics",
                "adjustment_diagnostics": "/api/adjustment/diagnostics",
                "publishing_diagnostics": "/api/publishing/diagnostics",
                "versioning_diagnostics": "/api/versioning/diagnostics",
                "prewrite_diagnostics": "/api/prewrite/diagnostics",
                "style_diagnostics": "/api/style/diagnostics",
                "related_links_diagnostics": "/api/related-links/diagnostics",
                "gap_filling_diagnostics": "/api/gap-filling/diagnostics",
                "evidence_tagging_diagnostics": "/api/evidence-tagging/diagnostics",
                "code_normalization_diagnostics": "/api/code-normalization/diagnostics",
                "review_runs": "/api/review/runs",
                "review_approve": "/api/review/approve",
                "review_reject": "/api/review/reject",
                "review_rerun": "/api/review/rerun"
            },
            "features": [
                "multi_dimensional_analysis",
                "adaptive_granularity",
                "intelligent_chunking",
                "cross_referencing",
                "comprehensive_file_support",
                "image_extraction",
                "progress_tracking",
                "comprehensive_validation",
                "fidelity_checking",
                "coverage_validation",
                "placeholder_detection",
                "style_guard",
                "cross_article_qa",
                "content_deduplication",
                "link_validation",
                "faq_consolidation", 
                "terminology_consistency",
                "woolf_style_processing",
                "structural_linting",
                "microsoft_style_guide",
                "technical_writing_standards",
                "adaptive_adjustment",
                "length_balancing",
                "split_optimization",
                "merge_suggestions",
                "readability_optimization",
                "v2_publishing_flow",
                "content_library_persistence",
                "comprehensive_metadata",
                "provenance_mapping",
                "media_reference_handling",
                "version_management",
                "content_versioning",
                "diff_analysis",
                "reprocessing_support",
                "version_chain_tracking",
                "human_in_the_loop_review",
                "quality_badges",
                "approval_workflow",
                "rejection_tracking",
                "step_rerun_capability",
                "section_grounded_prewrite",
                "fact_extraction",
                "evidence_based_writing",
                "concrete_examples_extraction",
                "gap_analysis",
                "terminology_extraction",
                "content_library_indexing",
                "internal_links_discovery",
                "external_links_extraction",
                "similarity_matching",
                "related_links_generation",
                "intelligent_gap_filling",
                "in_corpus_retrieval",
                "pattern_synthesis",
                "missing_content_detection",
                "safe_content_patching",
                "evidence_paragraph_tagging",
                "block_id_attribution",
                "fidelity_enforcement",
                "source_traceability",
                "code_block_normalization",
                "prism_integration",
                "syntax_highlighting_ready",
                "language_detection",
                "code_beautification",
                "copy_to_clipboard_ready",
                # KE-PR7: QA Report Features
                "qa_reports_machine_readable",
                "coverage_analysis_advanced",
                "unsupported_claims_detection",
                "placeholder_detection_comprehensive",
                "duplicate_content_detection",
                "broken_links_detection",
                "missing_media_detection",
                "content_quality_checks",
                "technical_accuracy_validation",
                "publish_gates_p0_blocking",
                "qa_database_persistence",
                "qa_summaries_api_exposure",
                "first_class_validation_module",
                "reusable_qa_validation"
            ],
            "message": "V2 Engine is active with comprehensive validation, cross-article QA, adaptive adjustment, V2-only publishing, version management, section-grounded prewrite pass, Woolf-aligned style processing, content library related links generation, intelligent gap filling with in-corpus retrieval, evidence tagging for fidelity enforcement, code normalization for Prism rendering, human-in-the-loop review systems, and first-class QA report module",
            # KE-PR7: QA Summaries
            "qa_summaries": qa_summaries,
            "qa_summary_count": len(qa_summaries),
            "qa_features": {
                "coverage_analysis": True,
                "unsupported_claims_detection": True,
                "placeholder_detection": True,
                "duplicate_content_detection": True,
                "broken_links_detection": True,
                "missing_media_detection": True,
                "publish_gates": True,
                "machine_readable_reports": True,
                "database_persistence": True,
                "api_exposure": True
            }
        }
        
    except Exception as e:
        print(f"❌ V2 ENGINE: Error getting engine status - {e}")
        return {
            "engine": "v2",
            "status": "error", 
            "error": str(e),
            "qa_summaries": [],
            "qa_summary_count": 0
        }

@app.get("/api/validation/diagnostics")
async def get_validation_diagnostics(run_id: str = None, validation_id: str = None):
    """V2 ENGINE: Get validation diagnostics for runs"""
    try:
        print(f"🔍 V2 VALIDATION: Diagnostics requested - run_id: {run_id}, validation_id: {validation_id} - engine=v2")
        
        # Build query filter
        query_filter = {}
        if run_id:
            query_filter["run_id"] = run_id
        if validation_id:
            query_filter["validation_id"] = validation_id
        
        # If no specific filters, get recent validation results
        if not query_filter:
            validation_results = await db.v2_validation_results.find().sort("timestamp", -1).limit(10).to_list(10)
        else:
            validation_results = await db.v2_validation_results.find(query_filter).sort("timestamp", -1).to_list(100)
        
        # Convert ObjectId to string for JSON serialization
        for result in validation_results:
            result['_id'] = str(result['_id'])
        
        diagnostics_summary = {
            "total_validations": len(validation_results),
            "passed_validations": len([r for r in validation_results if r.get('validation_status') == 'passed']),
            "partial_validations": len([r for r in validation_results if r.get('validation_status') == 'partial']),
            "failed_validations": len([r for r in validation_results if r.get('validation_status') not in ['passed', 'partial']]),
            "validation_results": validation_results
        }
        
        print(f"📊 V2 VALIDATION: Returning {len(validation_results)} validation results - engine=v2")
        return diagnostics_summary
        
    except Exception as e:
        print(f"❌ V2 VALIDATION: Error retrieving diagnostics - {e} - engine=v2")
        raise HTTPException(status_code=500, detail=f"Error retrieving validation diagnostics: {str(e)}")

@app.get("/api/validation/diagnostics/{validation_id}")
async def get_specific_validation_diagnostics(validation_id: str):
    """V2 ENGINE: Get specific validation result with detailed diagnostics"""
    try:
        print(f"🔍 V2 VALIDATION: Specific diagnostics requested - validation_id: {validation_id} - engine=v2")
        
        validation_result = await db.v2_validation_results.find_one({"validation_id": validation_id})
        
        if not validation_result:
            raise HTTPException(status_code=404, detail=f"Validation result not found: {validation_id}")
        
        # Convert ObjectId to string for JSON serialization
        validation_result['_id'] = str(validation_result['_id'])
        
        # Add summary for easy consumption
        summary_scores = validation_result.get('summary_scores', {})
        threshold_compliance = validation_result.get('threshold_compliance', {})
        
        enhanced_result = {
            **validation_result,
            "compliance_summary": {
                "overall_status": validation_result.get('validation_status', 'unknown'),
                "fidelity_score": f"{summary_scores.get('fidelity_score', 0):.2f}",
                "coverage_percent": f"{summary_scores.get('coverage_percent', 0):.1f}%",
                "placeholder_count": summary_scores.get('placeholder_count', 0),
                "style_compliance": f"{summary_scores.get('style_compliance', 0):.2f}",
                "passed_checks": sum(threshold_compliance.values()) if threshold_compliance else 0,
                "total_checks": len(threshold_compliance) if threshold_compliance else 0
            }
        }
        
        print(f"✅ V2 VALIDATION: Returning specific validation result - Status: {validation_result.get('validation_status')} - engine=v2")
        return enhanced_result
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ V2 VALIDATION: Error retrieving specific diagnostics - {e} - engine=v2")
        raise HTTPException(status_code=500, detail=f"Error retrieving validation diagnostics: {str(e)}")

@app.post("/api/validation/rerun")
async def rerun_validation(run_id: str = Form(...)):
    """V2 ENGINE: Rerun validation for a specific processing run"""
    try:
        print(f"🔄 V2 VALIDATION: Rerun validation requested - run_id: {run_id} - engine=v2")
        
        # Find the original processing result
        generated_articles_result = await v2_article_generator.get_generated_articles_for_run(run_id)
        
        if not generated_articles_result:
            raise HTTPException(status_code=404, detail=f"Processing run not found: {run_id}")
        
        # Find the normalized document and analysis (this is a simplified approach)
        # In a real implementation, you'd store these separately or reconstruct them
        print(f"⚠️ V2 VALIDATION: Rerun validation requires original normalized_doc and analysis - run_id: {run_id} - engine=v2")
        
        return {
            "message": "Validation rerun requested but requires original processing context",
            "run_id": run_id,
            "status": "pending_implementation",
            "note": "This endpoint requires storing normalized_doc and analysis for rerun capability"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ V2 VALIDATION: Error in validation rerun - {e} - engine=v2")
        raise HTTPException(status_code=500, detail=f"Error rerunning validation: {str(e)}")

@app.get("/api/qa/diagnostics")
async def get_qa_diagnostics(run_id: str = None, qa_id: str = None):
    """V2 ENGINE: Get cross-article QA diagnostics"""
    try:
        print(f"🔍 V2 CROSS-ARTICLE QA: Diagnostics requested - run_id: {run_id}, qa_id: {qa_id} - engine=v2")
        
        # Build query filter
        query_filter = {}
        if run_id:
            query_filter["run_id"] = run_id
        if qa_id:
            query_filter["qa_id"] = qa_id
        
        # If no specific filters, get recent QA results
        if not query_filter:
            qa_results = await db.v2_qa_results.find().sort("timestamp", -1).limit(10).to_list(10)
        else:
            qa_results = await db.v2_qa_results.find(query_filter).sort("timestamp", -1).to_list(100)
        
        # Convert ObjectId to string for JSON serialization
        for result in qa_results:
            result['_id'] = str(result['_id'])
        
        diagnostics_summary = {
            "total_qa_runs": len(qa_results),
            "passed_qa_runs": len([r for r in qa_results if r.get('summary', {}).get('issues_found', 0) == 0]),
            "qa_runs_with_issues": len([r for r in qa_results if r.get('summary', {}).get('issues_found', 0) > 0]),
            "error_qa_runs": len([r for r in qa_results if r.get('qa_status') == 'error']),
            "qa_results": qa_results
        }
        
        print(f"📊 V2 CROSS-ARTICLE QA: Returning {len(qa_results)} QA results - engine=v2")
        return diagnostics_summary
        
    except Exception as e:
        print(f"❌ V2 CROSS-ARTICLE QA: Error retrieving QA diagnostics - {e} - engine=v2")
        raise HTTPException(status_code=500, detail=f"Error retrieving QA diagnostics: {str(e)}")

@app.get("/api/qa/diagnostics/{qa_id}")
async def get_specific_qa_diagnostics(qa_id: str):
    """V2 ENGINE: Get specific QA result with detailed analysis"""
    try:
        print(f"🔍 V2 CROSS-ARTICLE QA: Specific QA diagnostics requested - qa_id: {qa_id} - engine=v2")
        
        qa_result = await db.v2_qa_results.find_one({"qa_id": qa_id})
        
        if not qa_result:
            raise HTTPException(status_code=404, detail=f"QA result not found: {qa_id}")
        
        # Convert ObjectId to string for JSON serialization
        qa_result['_id'] = str(qa_result['_id'])
        
        # Add summary for easy consumption
        summary = qa_result.get('summary', {})
        consolidation_result = qa_result.get('consolidation_result', {})
        
        enhanced_result = {
            **qa_result,
            "qa_summary": {
                "overall_status": qa_result.get('qa_status', 'unknown'),
                "total_issues": summary.get('issues_found', 0),
                "duplicates_found": summary.get('total_duplicates', 0),
                "invalid_links_found": summary.get('total_invalid_links', 0),
                "duplicate_faqs_found": summary.get('total_duplicate_faqs', 0),
                "terminology_issues_found": summary.get('total_terminology_issues', 0),
                "consolidation_actions": consolidation_result.get('total_actions', 0),
                "successful_consolidations": consolidation_result.get('successful_actions', 0)
            }
        }
        
        print(f"✅ V2 CROSS-ARTICLE QA: Returning specific QA result - Status: {qa_result.get('qa_status')} - engine=v2")
        return enhanced_result
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ V2 CROSS-ARTICLE QA: Error retrieving specific QA diagnostics - {e} - engine=v2")
        raise HTTPException(status_code=500, detail=f"Error retrieving QA diagnostics: {str(e)}")

@app.post("/api/qa/rerun")
async def rerun_qa_analysis(run_id: str = Form(...)):
    """V2 ENGINE: Rerun cross-article QA for a specific processing run"""
    try:
        print(f"🔄 V2 CROSS-ARTICLE QA: Rerun QA requested - run_id: {run_id} - engine=v2")
        
        # Find the original processing result
        generated_articles_result = await v2_article_generator.get_generated_articles_for_run(run_id)
        
        if not generated_articles_result:
            raise HTTPException(status_code=404, detail=f"Processing run not found: {run_id}")
        
        # Perform cross-article QA analysis
        qa_result = await v2_cross_article_qa_system.perform_cross_article_qa(
            generated_articles_result, run_id
        )
        
        # Store new QA result
        try:
            # KE-PR9.3: Use repository pattern for content_library operations
                from engine.stores.mongo import RepositoryFactory
                content_repo = RepositoryFactory.get_content_library()
                await content_repo.insert_article(article)
            
            article['_id'] = str(result.inserted_id)
            created_articles.append(article)
            
            print(f"✅ Created seed article: {article_data['title']}")
        
        print(f"🌱 Successfully created {len(created_articles)} seed articles for related links testing")
        
        return {
            "status": "success",
            "success": True,
            "message": f"Successfully created {len(created_articles)} seed articles",
            "articles_created": len(created_articles),
            "articles": [
                {
                    "id": article["id"],
                    "title": article["title"],
                    "summary": article["summary"]
                }
                for article in created_articles
            ]
        }
        
    except Exception as e:
        print(f"❌ Error creating seed articles: {e}")
        return {
            "success": False,
            "error": str(e)
        }

# ========================================
# TICKET 3: UNIVERSAL BOOKMARKS & DURABLE LINKS API ENDPOINTS
# ========================================

@app.post("/api/ticket3/backfill-bookmarks")
async def backfill_bookmarks(limit: int = None):
    """TICKET 3: Backfill existing v2 articles with bookmark registry data"""
    # KE-PR2: Use extracted linking module
    try:
        result = await backfill_registry(limit)
        
        return {
            "status": result.get("status", "success"),
            "message": f"Backfilled {result['articles_processed']} articles",
            "data": result
        }
    except Exception as e:
        return {
            "status": "error",
            "message": f"Backfill failed: {str(e)}"
        }

@app.get("/api/ticket3/validate-links/{doc_uid}")
async def validate_document_links(doc_uid: str):
    """TICKET 3: Validate cross-document links for a specific document"""
    try:
        # Get document from content library
        doc = await db.content_library.find_one({"doc_uid": doc_uid})
        
        if not doc:
            return {
                "status": "error",
                "message": f"Document not found: {doc_uid}"
            }
        
        v2_style_processor = V2StyleProcessor()
        xrefs = doc.get("xrefs", [])
        related_links = doc.get("related_links", [])
        
        validation_result = await v2_style_processor.validate_cross_document_links(
            doc_uid, xrefs, related_links
        )
        
        return {
            "status": "success",
            "doc_uid": doc_uid,
            "title": doc.get("title", "Unknown"),
            "validation": validation_result
        }
        
    except Exception as e:
        return {
            "status": "error",
            "message": f"Link validation failed: {str(e)}"
        }

@app.get("/api/ticket3/build-link")
async def build_cross_document_link(
    target_doc_uid: str, 
    anchor_id: str = None, 
    environment: str = "content_library"
):
    """TICKET 3: Build environment-aware cross-document link"""
    # KE-PR2: Use extracted linking modules
    try:
        # KE-PR9: Get target document using repository pattern
        if mongo_repo_available:
            content_repo = RepositoryFactory.get_content_library()
            target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
        else:
            # KE-PR9.3: Fallback to repository pattern
            content_repo = RepositoryFactory.get_content_library()
            target_doc = await content_repo.find_by_doc_uid(target_doc_uid)
        
        if not target_doc:
            return {
                "status": "error",
                "message": f"Target document not found: {target_doc_uid}"
            }
        
        route_map = get_default_route_map(environment)
        
        # Build href
        href = build_href(target_doc, anchor_id or "", route_map)
        
        return {
            "status": "success",
            "href": href,
            "target_doc": {
                "doc_uid": target_doc_uid,
                "doc_slug": target_doc.get("doc_slug"),
                "title": target_doc.get("title")
            },
            "anchor_id": anchor_id,
            "environment": environment
        }
        
    except Exception as e:
        return {
            "status": "error", 
            "message": f"Link building failed: {str(e)}"
        }

@app.get("/api/ticket3/document-registry/{doc_uid}")
async def get_document_bookmark_registry(doc_uid: str):
    """TICKET 3: Get bookmark registry for a specific document"""
    # KE-PR2: Use extracted linking modules
    try:
        registry = await get_registry(doc_uid)
        
        if not registry:
            return {
                "status": "error",
                "message": f"Document not found: {doc_uid}"
            }
        
        return {
            "status": "success",
            "doc_uid": doc_uid,
            "doc_slug": registry.get("doc_slug"),
            "title": registry.get("title"),
            "headings": registry.get("headings", []),
            "bookmark_count": len(registry.get("headings", []))
        }
        
    except Exception as e:
        return {
            "status": "error",
            "message": f"Registry retrieval failed: {str(e)}"
        }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("server:app", host="0.0.0.0", port=8001, reload=True)